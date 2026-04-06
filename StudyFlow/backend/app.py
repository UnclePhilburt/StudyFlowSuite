"""StudyFlow Suite Backend - Flask API Server"""

from flask import Flask, request, jsonify, send_from_directory, render_template_string
from PIL import Image
from io import BytesIO
import pytesseract
import subprocess
import os
import json
import openai
import re
import cv2
import numpy as np
import psycopg2
import traceback
import uuid
import hashlib
import requests
import stripe
import google.generativeai as genai


from StudyFlow.backend.image_processing import preprocess_image
from StudyFlow.config import TESSERACT_PATH
from StudyFlow.logging_utils import debug_log
from StudyFlow.backend.submit_button_storage import register_submit_button_upload
from StudyFlow.backend.tasks import process_question_async, update_note_votes, celery_app
from StudyFlow.backend import tasks  # registers the Celery task
import redis as redis_lib
from StudyFlow.backend.supabase_auth import supabase_auth_required, account_not_frozen  # Supabase Auth decorators
from StudyFlow.backend.supabase_client import supabase  # Supabase client for database operations

BACKEND_URL = os.environ.get("BACKEND_URL", "https://studyflowsuite.onrender.com")

stripe.api_key = os.environ['STRIPE_SECRET_KEY']
WEBHOOK_SECRET    = os.environ['STRIPE_WEBHOOK_SECRET']

# Redis cache (same instance as Celery broker)
_redis_url = os.getenv("CELERY_BROKER_URL", os.getenv("REDIS_URL", "redis://localhost:6379/0"))
try:
    redis_cache = redis_lib.from_url(_redis_url, decode_responses=True)
    redis_cache.ping()
    print(f"Redis cache connected: {_redis_url}")
except Exception as _e:
    print(f"Redis cache not available: {_e}")
    redis_cache = None

PROFILE_CACHE_TTL = 3600   # 1 hour
BROWSE_CACHE_TTL = 1800   # 30 minutes

def get_cached_profile(user_id):
    """Get user profile from Redis cache or Supabase. Returns {username, is_public} or None."""
    cache_key = f"profile:{user_id}"

    # Try cache first
    if redis_cache:
        try:
            cached = redis_cache.get(cache_key)
            if cached:
                return json.loads(cached)
        except:
            pass

    # Fetch from Supabase
    try:
        resp = supabase.table("user_profiles").select("username, is_public").eq("id", user_id).single().execute()
        if resp.data:
            profile = {"username": resp.data.get("username") or "Anonymous", "is_public": resp.data.get("is_public", True)}
            # Cache it
            if redis_cache:
                try:
                    redis_cache.setex(cache_key, PROFILE_CACHE_TTL, json.dumps(profile))
                except:
                    pass
            return profile
    except:
        pass
    return None


BREVO_API_KEY = os.environ.get("BREVO_API_KEY")
if not BREVO_API_KEY:
    raise RuntimeError("Missing BREVO_API_KEY environment variable")

# Brevo API helper function using requests
def send_brevo_email(to_email, to_name, subject, html_content, text_content=None, sender_email="info@studyflowsuite.com", sender_name="StudyFlow Suite"):
    """Send email via Brevo REST API"""
    url = "https://api.brevo.com/v3/smtp/email"

    headers = {
        "accept": "application/json",
        "api-key": BREVO_API_KEY,
        "content-type": "application/json"
    }

    payload = {
        "sender": {"email": sender_email, "name": sender_name},
        "to": [{"email": to_email, "name": to_name}] if to_name else [{"email": to_email}],
        "subject": subject,
        "htmlContent": html_content
    }

    if text_content:
        payload["textContent"] = text_content

    response = requests.post(url, json=payload, headers=headers)
    return response

# Import AI clients
from StudyFlow.backend.ai_manager import triple_call_ai_api_json_final
from StudyFlow.backend.deepflow import get_deepflow_question

# Set up Tesseract
pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH

# Set up OpenAI
openai.api_key = os.getenv("OPENAI_API_KEY")

# Log Tesseract version
try:
    version_output = subprocess.check_output([TESSERACT_PATH, "--version"]).decode("utf-8")
    debug_log("✅ Tesseract version output:\n" + version_output)
except Exception as e:
    debug_log("❌ Failed to call Tesseract: " + str(e) + "\n" + traceback.format_exc())

app = Flask(__name__)

# Enable CORS for website
from flask_cors import CORS
CORS(app,
     resources={r"/api/*": {"origins": [
         "https://unclephilburt.github.io",
         "https://studyflowsuite.com",
         "https://www.studyflowsuite.com",
         "http://studyflowsuite.com",
         "http://www.studyflowsuite.com",
         "http://localhost:*",
         "http://127.0.0.1:*"
     ]}, r"/admin/*": {"origins": [
         "https://unclephilburt.github.io",
         "https://studyflowsuite.com",
         "https://www.studyflowsuite.com",
         "http://studyflowsuite.com",
         "http://www.studyflowsuite.com",
         "http://localhost:*",
         "http://127.0.0.1:*"
     ]}},
     allow_headers=["Content-Type", "Authorization"],
     methods=["GET", "POST", "PUT", "PATCH", "DELETE", "OPTIONS"],
     supports_credentials=True,
     expose_headers=["Content-Type", "Authorization"]
)

# Ensure CORS headers on error responses (Flask-CORS misses unhandled exceptions)
@app.errorhandler(500)
def handle_500(e):
    response = jsonify({"error": "Internal server error"})
    response.status_code = 500
    return response

# ============ DMCA EMAIL NOTIFICATIONS ============

def send_dmca_report_notification_to_admin(takedown_id: str, note_filename: str, reporter_email: str, reason: str) -> bool:
    """Notify admin when new DMCA report is submitted"""
    admin_email = os.getenv("ADMIN_EMAIL", "info@studyflowsuite.com")

    html_content = f"""
    <div style="font-family: Arial, sans-serif; max-width: 600px; margin: 0 auto; padding: 20px;">
        <h2 style="color: #dc2626;">New DMCA Takedown Report</h2>
        <p>A copyright infringement report has been submitted and requires your review.</p>

        <div style="background: #f3f4f6; padding: 16px; border-radius: 8px; margin: 20px 0;">
            <p><strong>Note File:</strong> {note_filename}</p>
            <p><strong>Reporter:</strong> {reporter_email}</p>
            <p><strong>Reason:</strong></p>
            <p style="margin-left: 16px; font-style: italic;">{reason}</p>
        </div>

        <p><strong>Next Steps:</strong></p>
        <ol>
            <li>Review the report in Supabase dashboard</li>
            <li>Verify the claim is legitimate</li>
            <li>Process with admin endpoint:</li>
        </ol>

        <pre style="background: #1f2937; color: #10b981; padding: 12px; border-radius: 4px; font-size: 12px; overflow-x: auto;">
POST https://studyflowsuite.onrender.com/api/dmca/process/{takedown_id}
{{
  "admin_key": "YOUR_ADMIN_KEY",
  "action": "approve"  // or "reject"
}}
        </pre>

        <p style="margin-top: 24px; padding-top: 16px; border-top: 1px solid #e5e7eb; color: #6b7280; font-size: 12px;">
            Legal Requirement: DMCA Safe Harbor requires responding to takedown requests within a reasonable timeframe (typically 24-48 hours).
        </p>
    </div>
    """

    plain_text = f"""
New DMCA Takedown Report

Note File: {note_filename}
Reporter: {reporter_email}
Reason: {reason}

Review this report in your Supabase dashboard and process using the admin endpoint.

Legal Requirement: DMCA Safe Harbor requires responding within 24-48 hours.
    """

    try:
        response = send_brevo_email(
            to_email=admin_email,
            to_name=None,
            subject=f"[DMCA] New Takedown Report - {note_filename}",
            html_content=html_content,
            text_content=plain_text,
            sender_name="StudyFlow DMCA System"
        )

        if response.status_code in (200, 201):
            app.logger.info("DMCA admin notification sent successfully")
            return True
        else:
            app.logger.error(f"Failed to send DMCA admin notification: {response.text}")
            return False
    except Exception as e:
        app.logger.error(f"Failed to send DMCA admin notification: {e}")
        return False


def send_dmca_strike_notification_to_user(user_email: str, user_name: str, strike_count: int, note_filename: str, is_permanent_ban: bool) -> bool:
    """Notify user when they receive a DMCA strike"""

    if is_permanent_ban:
        subject = "Account Suspended - 3 DMCA Strikes"
        html_content = f"""
        <div style="font-family: Arial, sans-serif; max-width: 600px; margin: 0 auto; padding: 20px;">
            <h2 style="color: #dc2626;">Account Suspended</h2>
            <p>Hello {user_name},</p>

            <p>Your StudyFlow account has been <strong>permanently suspended</strong> from Good Standing privileges due to repeated copyright violations.</p>

            <div style="background: #fee2e2; border-left: 4px solid #dc2626; padding: 16px; margin: 20px 0;">
                <p style="margin: 0;"><strong>Strike Count: 3/3 (PERMANENT BAN)</strong></p>
                <p style="margin: 8px 0 0 0;">Most recent violation: {note_filename}</p>
            </div>

            <p><strong>What this means:</strong></p>
            <ul>
                <li>You can NO LONGER download notes from StudyFlow</li>
                <li>This suspension is PERMANENT and cannot be appealed</li>
                <li>Uploading new notes will NOT restore your access</li>
            </ul>

            <p><strong>Why this happened:</strong></p>
            <p>You uploaded copyrighted material (lecture slides, textbooks, etc.) owned by professors or publishers. This violates copyright law and our Terms of Service.</p>

            <p style="margin-top: 24px; padding-top: 16px; border-top: 1px solid #e5e7eb; color: #6b7280; font-size: 12px;">
                StudyFlow maintains a Three-Strikes Policy in compliance with DMCA Safe Harbor requirements. Repeat copyright infringers are permanently banned from the platform.
            </p>
        </div>
        """
    else:
        subject = f"DMCA Strike {strike_count}/3 - Warning"
        html_content = f"""
        <div style="font-family: Arial, sans-serif; max-width: 600px; margin: 0 auto; padding: 20px;">
            <h2 style="color: #f59e0b;">Copyright Strike Warning</h2>
            <p>Hello {user_name},</p>

            <p>You have received a copyright strike for uploading infringing content to StudyFlow.</p>

            <div style="background: #fef3c7; border-left: 4px solid #f59e0b; padding: 16px; margin: 20px 0;">
                <p style="margin: 0;"><strong>Strike Count: {strike_count}/3</strong></p>
                <p style="margin: 8px 0 0 0;">File removed: {note_filename}</p>
            </div>

            <p><strong>What happened:</strong></p>
            <p>A copyright owner (likely a professor or publisher) reported that your uploaded file contained their copyrighted material.</p>

            <p><strong>What this means:</strong></p>
            <ul>
                <li>Your file has been removed from public access</li>
                <li>{'One more strike and you will be PERMANENTLY BANNED' if strike_count == 2 else 'You have 2 strikes remaining before permanent ban'}</li>
                <li>You currently maintain Good Standing (if you have recent uploads)</li>
            </ul>

            <p><strong>How to avoid future strikes:</strong></p>
            <ul>
                <li>Only upload YOUR OWN notes that you personally wrote</li>
                <li>Do NOT upload professor's lecture slides or PowerPoints</li>
                <li>Do NOT upload textbook pages or publisher materials</li>
                <li>Do NOT upload copyrighted exam/quiz materials</li>
            </ul>

            <p style="margin-top: 24px; padding-top: 16px; border-top: 1px solid #e5e7eb; color: #6b7280; font-size: 12px;">
                <strong>Three-Strikes Policy:</strong> Strike 1 = Warning. Strike 2 = Final Warning. Strike 3 = Permanent ban from downloading any notes.
            </p>
        </div>
        """

    try:
        response = send_brevo_email(
            to_email=user_email,
            to_name=user_name,
            subject=subject,
            html_content=html_content,
            sender_name="StudyFlow Legal"
        )

        if response.status_code in (200, 201):
            app.logger.info("DMCA strike notification sent to user successfully")
            return True
        else:
            app.logger.error(f"Failed to send DMCA strike notification: {response.text}")
            return False
    except Exception as e:
        app.logger.error(f"Failed to send DMCA strike notification: {e}")
        return False


def send_dmca_report_confirmation_to_reporter(reporter_email: str, reporter_name: str, note_filename: str, action: str) -> bool:
    """Notify reporter when their DMCA report is processed"""

    if action == "approve":
        subject = "DMCA Report Processed - Content Removed"
        status_html = """
        <div style="background: #d1fae5; border-left: 4px solid #10b981; padding: 16px; margin: 20px 0;">
            <p style="margin: 0;"><strong>Status: APPROVED</strong></p>
            <p style="margin: 8px 0 0 0;">The infringing content has been removed and the uploader has received a copyright strike.</p>
        </div>
        """
    else:
        subject = "DMCA Report Processed - No Action Taken"
        status_html = """
        <div style="background: #fef3c7; border-left: 4px solid #f59e0b; padding: 16px; margin: 20px 0;">
            <p style="margin: 0;"><strong>Status: REJECTED</strong></p>
            <p style="margin: 8px 0 0 0;">After review, we determined the content does not constitute copyright infringement.</p>
        </div>
        """

    html_content = f"""
    <div style="font-family: Arial, sans-serif; max-width: 600px; margin: 0 auto; padding: 20px;">
        <h2 style="color: #1f2937;">DMCA Report Update</h2>
        <p>Hello {reporter_name or 'there'},</p>

        <p>Your copyright infringement report has been processed by our team.</p>

        {status_html}

        <p><strong>Reported File:</strong> {note_filename}</p>

        <p>Thank you for helping us maintain a legal and respectful learning community.</p>

        <p style="margin-top: 24px; padding-top: 16px; border-top: 1px solid #e5e7eb; color: #6b7280; font-size: 12px;">
            If you have questions or wish to submit additional reports, please contact us at info@studyflowsuite.com
        </p>
    </div>
    """

    try:
        response = send_brevo_email(
            to_email=reporter_email,
            to_name=reporter_name,
            subject=subject,
            html_content=html_content,
            sender_name="StudyFlow Legal"
        )

        if response.status_code in (200, 201):
            app.logger.info("DMCA confirmation sent to reporter successfully")
            return True
        else:
            app.logger.error(f"Failed to send DMCA confirmation: {response.text}")
            return False
    except Exception as e:
        app.logger.error(f"Failed to send DMCA confirmation: {e}")
        return False


def send_access_key_email(to_email: str, stripe_id: str) -> bool:
    # Build your message bodies
    html_content = (
        "<p>Welcome to StudyFlow!</p>"
        "<p>Your access key is:</p>"
        f"<pre style='background:#f4f4f4;padding:8px;border-radius:4px;'>{stripe_id}</pre>"
        "<p>Keep it safe—enter it when launching the app.</p>"
    )
    plain_text_content = (
        f"Welcome to StudyFlow!\n\n"
        f"Your access key is: {stripe_id}\n\n"
        "Keep it safe—enter it when launching the app."
    )

    try:
        app.logger.debug(f"Sending access key email to {to_email}")

        response = send_brevo_email(
            to_email=to_email,
            to_name=None,
            subject="Your StudyFlow Access Key",
            html_content=html_content,
            text_content=plain_text_content
        )

        if response.status_code in (200, 201):
            app.logger.info("Access key email sent successfully")
            return True
        else:
            app.logger.error(f"Brevo error sending access key email: {response.text}")
            return False
    except Exception as e:
        app.logger.error("Brevo error sending access key email", exc_info=e)
        return False

import logging
logging.basicConfig(level=logging.INFO)
app.logger.setLevel(logging.INFO)
register_submit_button_upload(app)


def init_postgres_db():
    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()
        cur.execute("""
            CREATE TABLE IF NOT EXISTS qa_pairs (
                id SERIAL PRIMARY KEY,
                question TEXT UNIQUE,
                answer TEXT,
                count INTEGER DEFAULT 1,
                timestamp TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
            CREATE TABLE IF NOT EXISTS app_config (
                key   TEXT PRIMARY KEY,
                value TEXT NOT NULL
            );
        """)
        conn.commit()
        conn.close()
        print("✅ PostgreSQL: qa_pairs and app_config tables ready.")
    except Exception as e:
        print(f"❌ DB init error: {e}")

# Initialize DB on startup
init_postgres_db()


# ============================================================================
# USER AUTHENTICATION & SUBSCRIPTION ROUTES
# ============================================================================
import bcrypt
import jwt
from datetime import datetime, timedelta
from functools import wraps

JWT_SECRET = os.environ.get("JWT_SECRET", "a7f3e9d2c4b8a1f6e3d9c2b7a5f8e1d4c9b6a3f7e2d8c5b1a6f9e4d3c8b2a7f5e1d8c4b9a6f3")

def hash_password(password: str) -> str:
    """Hash a password using bcrypt"""
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(password: str, hashed: str) -> bool:
    """Verify a password against its hash"""
    if not hashed:
        return False
    return bcrypt.checkpw(password.encode('utf-8'), hashed.encode('utf-8'))

def create_token(user_id: int, email: str) -> str:
    """Create JWT token for user"""
    payload = {
        'user_id': user_id,
        'email': email,
        'exp': datetime.utcnow() + timedelta(days=30)
    }
    return jwt.encode(payload, JWT_SECRET, algorithm='HS256')

def check_question_limit(user_id):
    """Check if user has exceeded their daily question limit. Returns (can_proceed, remaining_questions)"""
    # Everything is free for now - unlimited questions for everyone
    return True, -1  # -1 means unlimited

    # Legacy code (disabled):
    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()

        # Check user's subscription status and beta status
        cur.execute("SELECT subscription_tier, created_at FROM user_profiles WHERE id = %s", (user_id,))
        user = cur.fetchone()
        if not user:
            conn.close()
            return False, 0

        subscription_status, is_beta = user

        # Beta testers (grandfathered) have unlimited questions
        if is_beta:
            conn.close()
            return True, -1  # -1 means unlimited

        # Pro/trialing users have unlimited questions
        if subscription_status in ['active', 'trialing']:
            conn.close()
            return True, -1  # -1 means unlimited

        # Free users have 10 questions per day limit
        today = datetime.utcnow().date()

        # Get or create today's usage record
        cur.execute("""
            INSERT INTO question_usage (user_id, question_date, question_count)
            VALUES (%s, %s, 0)
            ON CONFLICT (user_id, question_date) DO NOTHING
        """, (user_id, today))

        cur.execute("""
            SELECT question_count FROM question_usage
            WHERE user_id = %s AND question_date = %s
        """, (user_id, today))

        result = cur.fetchone()
        current_count = result[0] if result else 0

        FREE_LIMIT = 10
        remaining = FREE_LIMIT - current_count

        if current_count >= FREE_LIMIT:
            conn.close()
            return False, 0

        # Increment count
        cur.execute("""
            UPDATE question_usage
            SET question_count = question_count + 1
            WHERE user_id = %s AND question_date = %s
        """, (user_id, today))

        conn.commit()
        conn.close()

        return True, remaining - 1

    except Exception as e:
        app.logger.error(f"Error checking question limit: {e}")
        return True, -1  # Allow on error to not block users

def token_required(f):
    """Decorator to protect routes with JWT authentication"""
    @wraps(f)
    def decorated(*args, **kwargs):
        token = request.headers.get('Authorization')
        if not token:
            return jsonify({'error': 'Token is missing'}), 401
        try:
            if token.startswith('Bearer '):
                token = token[7:]
            data = jwt.decode(token, JWT_SECRET, algorithms=['HS256'])
            request.user_id = data['user_id']
            request.user_email = data['email']
        except:
            return jsonify({'error': 'Token is invalid'}), 401
        return f(*args, **kwargs)
    return decorated

def init_users_table():
    """Create users table if it doesn't exist, and add missing columns if needed"""
    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()

        print("🔧 Starting users table initialization...")

        # Create table if it doesn't exist
        cur.execute("""
            CREATE TABLE IF NOT EXISTS users (
                id SERIAL PRIMARY KEY,
                email VARCHAR(255) UNIQUE NOT NULL,
                name VARCHAR(255),
                password_hash VARCHAR(255),
                stripe_customer_id VARCHAR(255) UNIQUE,
                stripe_subscription_id VARCHAR(255),
                subscription_status VARCHAR(50) DEFAULT 'free',
                is_beta BOOLEAN DEFAULT FALSE,
                trial_ends_at TIMESTAMP,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                updated_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
        """)

        # Create question usage tracking table
        cur.execute("""
            CREATE TABLE IF NOT EXISTS question_usage (
                id SERIAL PRIMARY KEY,
                user_id INTEGER REFERENCES users(id) ON DELETE CASCADE,
                question_date DATE NOT NULL DEFAULT CURRENT_DATE,
                question_count INTEGER DEFAULT 0,
                UNIQUE(user_id, question_date)
            );
        """)

        # Create questions table for storing all questions and answers
        cur.execute("""
            CREATE TABLE IF NOT EXISTS questions (
                id SERIAL PRIMARY KEY,
                user_id INTEGER REFERENCES users(id) ON DELETE CASCADE,
                question_text TEXT NOT NULL,
                question_type VARCHAR(50) NOT NULL,
                answers_json TEXT,
                ai_answer TEXT,
                ai_reasoning TEXT,
                correct BOOLEAN,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
        """)

        # Create notes table for NoteFlow
        cur.execute("""
            CREATE TABLE IF NOT EXISTS notes (
                id SERIAL PRIMARY KEY,
                user_id INTEGER REFERENCES users(id) ON DELETE CASCADE,
                filename VARCHAR(255) NOT NULL,
                original_filename VARCHAR(255) NOT NULL,
                file_path TEXT,
                file_type VARCHAR(50),
                file_size INTEGER,
                ocr_text TEXT,
                page_count INTEGER,
                uploaded_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            );
        """)

        # Create indexes for fast lookups
        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_question_hash
            ON questions(MD5(question_text));
        """)

        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_user_questions
            ON questions(user_id, created_at);
        """)

        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_question_type
            ON questions(question_type);
        """)

        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_user_notes
            ON notes(user_id, uploaded_at DESC);
        """)

        cur.execute("""
            CREATE INDEX IF NOT EXISTS idx_notes_ocr_text
            ON notes USING gin(to_tsvector('english', ocr_text));
        """)

        conn.commit()
        print("✅ Table creation/check complete")

        # Add missing columns if table already exists (compatible with older PostgreSQL)
        columns_to_add = [
            ("name", "VARCHAR(255)"),
            ("password_hash", "VARCHAR(255)"),
            ("stripe_customer_id", "VARCHAR(255)"),
            ("stripe_subscription_id", "VARCHAR(255)"),
            ("is_beta", "BOOLEAN DEFAULT FALSE"),
            ("trial_ends_at", "TIMESTAMP"),
            ("created_at", "TIMESTAMP DEFAULT CURRENT_TIMESTAMP"),
            ("updated_at", "TIMESTAMP DEFAULT CURRENT_TIMESTAMP")
        ]

        print(f"🔍 Checking {len(columns_to_add)} columns...")

        for column_name, column_type in columns_to_add:
            try:
                # Check if column exists first
                cur.execute("""
                    SELECT column_name
                    FROM information_schema.columns
                    WHERE table_name='users' AND column_name=%s;
                """, (column_name,))

                result = cur.fetchone()
                if result is None:
                    # Column doesn't exist, add it
                    print(f"➕ Adding missing column: {column_name}")
                    cur.execute(f"ALTER TABLE users ADD COLUMN {column_name} {column_type};")
                    conn.commit()
                    print(f"✅ Added column: {column_name}")
                else:
                    print(f"✓ Column exists: {column_name}")
            except Exception as col_error:
                conn.rollback()
                print(f"⚠️ Column {column_name} error: {col_error}")
                import traceback
                print(traceback.format_exc())

        conn.close()
        print("✅ PostgreSQL: users table ready.")
    except Exception as e:
        print(f"❌ Users table init error: {e}")
        import traceback
        print(traceback.format_exc())

# Initialize users table
init_users_table()

@app.route("/health", methods=["GET"])
def health_check():
    return jsonify({"status": "ok"}), 200


@app.route("/admin/backfill-votes", methods=["POST"])
def admin_backfill_votes():
    """ADMIN: Trigger a vote count backfill for all notes."""
    admin_key = request.args.get("key", "") or (request.get_json() or {}).get("key", "")
    if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
        return jsonify({"error": "Unauthorized"}), 403

    from StudyFlow.backend.tasks import backfill_all_vote_counts
    task = backfill_all_vote_counts.delay()
    return jsonify({"success": True, "task_id": str(task.id), "message": "Backfill queued"}), 200


@app.route("/api/signup", methods=["POST"])
def signup():
    """Create a new user account with Supabase Auth"""
    try:
        from StudyFlow.backend.supabase_auth import create_user_with_metadata

        data = request.json
        email = data.get('email')
        name = data.get('name')
        password = data.get('password')
        collective_brain_opt_in = data.get('collective_brain_opt_in', False)

        if not email or not password:
            return jsonify({'error': 'Email and password are required'}), 400

        # Create user with Supabase Auth + profile
        result, error = create_user_with_metadata(
            email=email,
            password=password,
            full_name=name,
            collective_brain_opt_in=collective_brain_opt_in
        )

        if error:
            # Check if it's a duplicate email error
            if 'already registered' in error.lower() or 'already exists' in error.lower():
                return jsonify({'error': 'Email already exists'}), 409
            return jsonify({'error': error}), 400

        debug_log(f"✅ New user created: {email}")

        # Auto-follow the official StudyFlow account
        official_id = os.getenv("STUDYFLOW_OFFICIAL_USER_ID")
        if official_id and result['user']:
            try:
                new_user_id = result['user'].id
                if new_user_id != official_id:
                    supabase.table("user_followers").insert({
                        "follower_id": new_user_id,
                        "following_id": official_id
                    }).execute()
                    debug_log(f"Auto-followed official account for {email}")
            except Exception as follow_err:
                debug_log(f"Auto-follow failed (non-fatal): {follow_err}")

        return jsonify({
            'success': True,
            'message': 'Account created successfully',
            'access_token': result['session'].access_token if result['session'] else None,
            'refresh_token': result['session'].refresh_token if result['session'] else None,
            'user': {
                'id': result['user'].id,
                'email': result['user'].email,
                'collective_brain_opt_in': collective_brain_opt_in
            }
        }), 201

    except Exception as e:
        debug_log(f"❌ Signup error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Signup failed'}), 500


@app.route("/api/login", methods=["POST"])
def login():
    """Authenticate user with Supabase Auth and return access token"""
    try:
        from StudyFlow.backend.supabase_auth import sign_in_user
        from StudyFlow.backend.supabase_client import get_user_profile

        data = request.json
        email = data.get('email')
        password = data.get('password')

        if not email or not password:
            return jsonify({'error': 'Email and password are required'}), 400

        # Sign in with Supabase Auth
        result, error = sign_in_user(email, password)

        if error:
            return jsonify({'error': error}), 401

        # Get user profile for additional info
        user_profile = get_user_profile(result['user'].id)

        debug_log(f"✅ User logged in: {email}")

        return jsonify({
            'success': True,
            'access_token': result['access_token'],
            'refresh_token': result['refresh_token'],
            'expires_in': result['expires_in'],
            'user': {
                'id': result['user'].id,
                'email': result['user'].email,
                'full_name': user_profile.get('full_name') if user_profile else None,
                'subscription_tier': user_profile.get('subscription_tier') if user_profile else 'free',
                'stripe_customer_id': user_profile.get('stripe_customer_id') if user_profile else None,
                'collective_brain_opt_in': user_profile.get('collective_brain_opt_in') if user_profile else False
            }
        }), 200

    except Exception as e:
        debug_log(f"❌ Login error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Login failed'}), 500


@app.route("/api/refresh", methods=["POST"])
def refresh_token():
    """Refresh an expired access token"""
    try:
        from StudyFlow.backend.supabase_auth import refresh_access_token

        data = request.json
        refresh_token_str = data.get('refresh_token')

        if not refresh_token_str:
            return jsonify({'error': 'Refresh token is required'}), 400

        result, error = refresh_access_token(refresh_token_str)

        if error:
            return jsonify({'error': error}), 401

        return jsonify({
            'success': True,
            'access_token': result['access_token'],
            'refresh_token': result['refresh_token'],
            'expires_in': result['expires_in']
        }), 200

    except Exception as e:
        debug_log(f"❌ Refresh token error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Failed to refresh token'}), 500


@app.route("/api/logout", methods=["POST"])
def logout():
    """Sign out user (invalidate session)"""
    try:
        from StudyFlow.backend.supabase_auth import sign_out_user

        # Get token from Authorization header
        auth_header = request.headers.get('Authorization')
        if not auth_header or not auth_header.startswith('Bearer '):
            return jsonify({'error': 'Missing or invalid Authorization header'}), 401

        token = auth_header[7:]  # Remove 'Bearer ' prefix

        success, error = sign_out_user(token)

        if error:
            return jsonify({'error': error}), 400

        return jsonify({'success': True, 'message': 'Logged out successfully'}), 200

    except Exception as e:
        debug_log(f"❌ Logout error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Logout failed'}), 500


@app.route("/api/reset-password", methods=["POST"])
def reset_password():
    """Send password reset email"""
    try:
        from StudyFlow.backend.supabase_auth import reset_password_email

        data = request.json
        email = data.get('email')

        if not email:
            return jsonify({'error': 'Email is required'}), 400

        success, error = reset_password_email(email)

        if error:
            return jsonify({'error': error}), 400

        return jsonify({
            'success': True,
            'message': 'Password reset email sent. Check your inbox.'
        }), 200

    except Exception as e:
        debug_log(f"❌ Reset password error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Failed to send reset email'}), 500

@app.route("/api/create-subscription", methods=["POST"])
def create_subscription():
    """Create Stripe subscription for user (with Supabase)"""
    from StudyFlow.backend.supabase_client import supabase
    from datetime import datetime

    try:
        data = request.json
        email = data.get('email')
        name = data.get('name')
        payment_method_id = data.get('paymentMethodId')
        plan = data.get('plan', 'pro')

        if not email or not payment_method_id:
            return jsonify({'error': 'Missing required fields'}), 400

        # Check if user already exists in Supabase
        result = supabase.table("user_profiles").select("id, stripe_customer_id").eq("email", email).execute()

        user_id = None
        customer_id = None

        if result.data and len(result.data) > 0:
            # User exists
            user_id = result.data[0]['id']
            customer_id = result.data[0].get('stripe_customer_id')
            debug_log(f"Existing user found: {user_id}")
        else:
            # Create new Supabase auth user (auto-creates user_profile via trigger)
            import secrets
            temp_password = secrets.token_urlsafe(32)

            auth_result = supabase.auth.sign_up({
                "email": email,
                "password": temp_password,
                "options": {
                    "data": {
                        "full_name": name
                    }
                }
            })

            if auth_result.user:
                user_id = auth_result.user.id
                debug_log(f"New user created: {user_id}")
            else:
                raise Exception("Failed to create user account")

        # Create or retrieve Stripe customer
        if customer_id:
            customer = stripe.Customer.retrieve(customer_id)
        else:
            customer = stripe.Customer.create(
                email=email,
                name=name,
                payment_method=payment_method_id,
                invoice_settings={'default_payment_method': payment_method_id},
                metadata={'user_id': user_id}
            )
            customer_id = customer.id

        # Create subscription with 7-day trial - $4.99/month
        subscription = stripe.Subscription.create(
            customer=customer_id,
            items=[{'price': 'price_1TCOih9LWKaKRffVWuR3bQin'}],  # $4.99/month
            trial_period_days=7,
            expand=['latest_invoice.payment_intent'],
            metadata={'user_id': user_id}
        )

        # Calculate trial end date
        trial_ends_at = datetime.fromtimestamp(subscription.trial_end).isoformat()

        # Update user_profile with subscription info
        supabase.table("user_profiles").update({
            "stripe_customer_id": customer_id,
            "stripe_subscription_id": subscription.id,
            "subscription_tier": "pro",
            "subscription_status": "trialing",
            "full_name": name
        }).eq("id", user_id).execute()

        # Send welcome email
        send_access_key_email(email, customer_id)

        # Create JWT token for website auth
        token = create_token(user_id, email)

        debug_log(f"✅ Subscription created for {email}: {subscription.id}")

        return jsonify({
            'success': True,
            'token': token,
            'user': {
                'id': user_id,
                'email': email,
                'name': name,
                'subscription_status': 'trialing',
                'subscription_tier': 'pro',
                'stripe_customer_id': customer_id,
                'trial_ends_at': trial_ends_at
            }
        }), 200

    except stripe.error.CardError as e:
        debug_log(f"❌ Card error: {e}")
        return jsonify({'error': 'Card was declined'}), 400
    except Exception as e:
        debug_log(f"❌ Subscription creation error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500

@app.route("/api/me", methods=["GET"])
@supabase_auth_required
def get_me():
    """Get current user info (protected route)"""
    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()
        cur.execute("""
            SELECT id, email, full_name, subscription_tier, stripe_customer_id, NULL as trial_ends_at, created_at, FALSE as is_beta
            FROM user_profiles WHERE id = %s
        """, (request.user_id,))
        user = cur.fetchone()
        conn.close()

        if not user:
            return jsonify({'error': 'User not found'}), 404

        user_id, email, name, subscription_status, stripe_customer_id, trial_ends_at, created_at, is_beta = user

        return jsonify({
            'id': user_id,
            'email': email,
            'name': name,
            'subscription_status': subscription_status,
            'stripe_customer_id': stripe_customer_id,
            'trial_ends_at': trial_ends_at.isoformat() if trial_ends_at else None,
            'is_beta': is_beta if is_beta is not None else False,
            'created_at': created_at.isoformat()
        }), 200

    except Exception as e:
        app.logger.error(f"❌ Get user error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Failed to get user info'}), 500

@app.route("/api/create-portal-session", methods=["POST"])
@supabase_auth_required
def create_portal_session():
    """Create a Stripe Customer Portal session for managing subscription"""
    try:
        # Get user's stripe_customer_id
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()
        cur.execute("""
            SELECT stripe_customer_id FROM user_profiles WHERE id = %s
        """, (request.user_id,))
        result = cur.fetchone()
        conn.close()

        if not result or not result[0]:
            return jsonify({'error': 'No active subscription found'}), 404

        stripe_customer_id = result[0]

        # Create Stripe Customer Portal session
        session = stripe.billing_portal.Session.create(
            customer=stripe_customer_id,
            return_url='https://unclephilburt.github.io/studyflowwebsite/dashboard.html'
        )

        return jsonify({'url': session.url}), 200

    except Exception as e:
        app.logger.error(f"❌ Create portal session error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Failed to create portal session'}), 500

@app.route("/api/mobile-dashboard", methods=["GET"])
@supabase_auth_required
def mobile_dashboard():
    """Single endpoint for mobile dashboard -- returns everything in one call."""
    try:
        user_id = request.user_id
        cache_key = f"mobile_dash:{user_id}"

        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        # Parallel-ish fetches from Supabase
        profile_resp = supabase.table("user_profiles").select(
            "username, university, subscription_tier, subscription_status"
        ).eq("id", user_id).single().execute()

        notes_resp = supabase.table("notes").select("id").eq("user_id", user_id).execute()

        convos_resp = supabase.table("conversations").select(
            "id, title, updated_at"
        ).eq("user_id", user_id).eq("source", "chat").is_("deleted_at", "null").order("updated_at", desc=True).limit(5).execute()

        groups_resp = supabase.table("study_group_members").select(
            "group_id"
        ).eq("user_id", user_id).execute()

        profile = profile_resp.data or {}
        result = {
            "username": profile.get("username") or "Anonymous",
            "university": profile.get("university"),
            "tier": profile.get("subscription_tier", "free"),
            "notes_count": len(notes_resp.data or []),
            "chats_count": len(convos_resp.data or []),
            "groups_count": len(groups_resp.data or []),
            "recent_chats": convos_resp.data or []
        }

        if redis_cache:
            try: redis_cache.setex(cache_key, 120, json.dumps(result))
            except: pass

        return jsonify(result), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/folders", methods=["GET"])
@supabase_auth_required
def get_user_folders():
    """Get user's folders for dropdown selection"""
    try:
        user_id = request.user_id
        response = supabase.table("folders").select("id, name").eq(
            "user_id", user_id
        ).order("created_at", desc=False).execute()

        folders = response.data or []
        return jsonify({"folders": folders}), 200

    except Exception as e:
        debug_log(f"Get folders error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes", methods=["GET"])
@supabase_auth_required
def get_user_notes_list():
    """Get list of user's notes for dropdowns/selection"""
    try:
        user_id = request.user_id
        from StudyFlow.backend.supabase_client import get_user_notes
        notes = get_user_notes(user_id)

        formatted_notes = []
        for note in notes:
            formatted_notes.append({
                "id": note['id'],
                "original_filename": note.get('original_filename', ''),
                "file_type": note.get('file_type'),
                "thumbnail_url": note.get('thumbnail_url'),
                "uploaded_at": note.get('uploaded_at'),
                "processed": note.get('processed')
            })

        return jsonify({"notes": formatted_notes}), 200
    except Exception as e:
        debug_log(f"Get notes list error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/page-init", methods=["GET"])
@supabase_auth_required
def notes_page_init():
    """Single endpoint for notes page -- returns profile, standing, favorites, folders, notes, and usage stats."""
    try:
        user_id = request.user_id
        cache_key = f"notes_init:{user_id}"

        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        # Profile
        profile_resp = supabase.table("user_profiles").select(
            "username, university, subscription_tier, subscription_status, good_standing, permanent_bad_standing, last_verified_upload_date"
        ).eq("id", user_id).single().execute()
        profile = profile_resp.data or {}

        # Folders
        folders_resp = supabase.table("folders").select("*").eq("user_id", user_id).order("created_at", desc=False).execute()

        # Notes
        from StudyFlow.backend.supabase_client import get_user_notes
        notes = get_user_notes(user_id)
        formatted_notes = []
        for note in notes:
            formatted_notes.append({
                "id": note['id'],
                "filename": note.get('original_filename', ''),
                "file_type": note.get('file_type'),
                "file_size": note.get('file_size'),
                "pages": note.get('page_count'),
                "uploaded_at": note.get('uploaded_at'),
                "processed": note.get('processed'),
                "is_public": note.get('is_public'),
                "university": note.get('university'),
                "course_code": note.get('course_code'),
                "professor": note.get('professor'),
                "folder_id": note.get('folder_id')
            })

        # Favorites
        favs_resp = supabase.table("note_favorites").select("note_id, created_at").eq("user_id", user_id).order("created_at", desc=True).execute()

        # Usage stats
        usage_stats = {"totalUsage": 0, "weeklyUsage": 0, "activeNotes": 0, "noteUsage": {}}
        try:
            note_ids = [n['id'] for n in notes]
            if note_ids:
                logs_resp = supabase.table("ai_response_logs").select("sources_used").execute()
                note_usage = {}
                for log in (logs_resp.data or []):
                    sources = log.get('sources_used') or []
                    for src in sources:
                        nid = src.get('note_id') if isinstance(src, dict) else src
                        if nid in note_ids:
                            note_usage[nid] = note_usage.get(nid, 0) + 1
                usage_stats["totalUsage"] = sum(note_usage.values())
                usage_stats["activeNotes"] = len(note_usage)
                usage_stats["noteUsage"] = note_usage
        except:
            pass

        result = {
            "profile": profile,
            "folders": folders_resp.data or [],
            "notes": formatted_notes,
            "favorites": favs_resp.data or [],
            "usage_stats": usage_stats
        }

        if redis_cache:
            try: redis_cache.setex(cache_key, 120, json.dumps(result))
            except: pass

        return jsonify(result), 200
    except Exception as e:
        debug_log(f"Notes page init error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/warmup", methods=["POST"])
@supabase_auth_required
def user_warmup():
    """Pre-warm all Redis caches for a user on login. Called once, everything loads fast after."""
    try:
        user_id = request.user_id

        # 1. Cache user profile
        try:
            profile = supabase.table("user_profiles").select("*").eq("id", user_id).single().execute()
            if profile.data and redis_cache:
                redis_cache.setex(f"profile:{user_id}", PROFILE_CACHE_TTL, json.dumps({
                    "username": profile.data.get("username") or "Anonymous",
                    "is_public": profile.data.get("is_public", True)
                }))
        except:
            pass

        # 2. Cache dashboard layout
        try:
            layout = supabase.table("user_profiles").select("dashboard_layout").eq("id", user_id).single().execute()
            if layout.data and redis_cache:
                redis_cache.setex(f"dashboard:{user_id}", 600, json.dumps(layout.data.get("dashboard_layout", {})))
        except:
            pass

        # 3. Cache canvas preload (all folders, conversations, messages, sticky notes)
        try:
            folders_resp = supabase.table("chat_folders").select("*").eq("user_id", user_id).order("created_at").execute()
            convos_resp = supabase.table("conversations").select(
                "id, title, folder_id, updated_at, canvas_layout"
            ).eq("user_id", user_id).eq("source", "chat").is_("deleted_at", "null").order("updated_at", desc=True).execute()
            stickies_resp = supabase.table("canvas_sticky_notes").select("*").eq("user_id", user_id).execute()

            convos = convos_resp.data or []
            conversations_with_messages = []
            for conv in convos:
                msgs_resp = supabase.table("conversation_messages").select(
                    "content, role, sources, created_at"
                ).eq("conversation_id", conv["id"]).order("created_at").execute()
                conversations_with_messages.append({
                    "id": conv["id"],
                    "title": conv.get("title"),
                    "folder_id": conv.get("folder_id"),
                    "updated_at": conv.get("updated_at"),
                    "canvas_layout": conv.get("canvas_layout"),
                    "messages": msgs_resp.data or []
                })

            canvas_data = {
                "folders": folders_resp.data or [],
                "conversations": conversations_with_messages,
                "sticky_notes": stickies_resp.data or []
            }

            if redis_cache:
                redis_cache.setex(f"canvas_preload:{user_id}", 600, json.dumps(canvas_data))
        except:
            pass

        # 4. Cache user's notes list (same format as /api/notes/list)
        try:
            notes_resp = supabase.table("notes").select("*").eq("user_id", user_id).order("uploaded_at", desc=True).execute()
            if notes_resp.data and redis_cache:
                formatted_notes = []
                for note in notes_resp.data:
                    formatted_notes.append({
                        "id": note['id'],
                        "filename": note.get('original_filename', ''),
                        "file_type": note.get('file_type'),
                        "file_size": note.get('file_size'),
                        "pages": note.get('page_count'),
                        "uploaded_at": note.get('uploaded_at'),
                        "processed": note.get('processed'),
                        "is_public": note.get('is_public'),
                        "university": note.get('university'),
                        "course_code": note.get('course_code'),
                        "professor": note.get('professor'),
                        "folder_id": note.get('folder_id')
                    })
                redis_cache.setex(f"user_notes:{user_id}", 300, json.dumps(formatted_notes))
        except:
            pass

        # 5. Cache notifications
        try:
            notif_resp = supabase.table("notifications").select("*").eq("user_id", user_id).order("created_at", desc=True).limit(50).execute()
            if redis_cache:
                redis_cache.setex(f"notifications:{user_id}", 120, json.dumps({"notifications": notif_resp.data or []}))
        except:
            pass

        # 6. Cache study groups
        try:
            memberships = supabase.table("study_group_members").select("group_id, role").eq("user_id", user_id).execute()
            groups = []
            for m in (memberships.data or []):
                g = supabase.table("study_groups").select("*").eq("id", m["group_id"]).execute()
                if g.data:
                    group = g.data[0]
                    members = supabase.table("study_group_members").select("id", count="exact").eq("group_id", m["group_id"]).execute()
                    group["member_count"] = members.count if members.count else 0
                    group["role"] = m["role"]
                    groups.append(group)
            if redis_cache:
                redis_cache.setex(f"user_groups:{user_id}", 300, json.dumps({"groups": groups}))
        except:
            pass

        # 7. Cache calendar events
        try:
            events_resp = supabase.table("calendar_events").select("*").eq("user_id", user_id).order("due_date").execute()
            if redis_cache:
                redis_cache.setex(f"calendar:{user_id}:::", 300, json.dumps({"success": True, "events": events_resp.data or []}))
        except:
            pass

        # 8. Cache leaderboard (shared across users)
        try:
            if not redis_cache or not redis_cache.get("leaderboard:main"):
                from collections import defaultdict
                notes_lb = supabase.table("notes").select("user_id, page_count").eq("is_public", True).execute()
                user_pages = defaultdict(int)
                for note in (notes_lb.data or []):
                    user_pages[note['user_id']] += note.get('page_count', 0)
                top_users = sorted(user_pages.items(), key=lambda x: x[1], reverse=True)[:20]
                leaderboard = []
                for uid, pages in top_users:
                    try:
                        p = supabase.table("user_profiles").select("username, university").eq("id", uid).single().execute()
                        if p.data:
                            leaderboard.append({"user_id": uid, "username": p.data.get('username'), "university": p.data.get('university'), "pages_contributed": pages})
                    except: pass
                if redis_cache:
                    redis_cache.setex("leaderboard:main", 1800, json.dumps(leaderboard))
        except:
            pass

        debug_log(f"[Warmup] Pre-warmed all caches for user {user_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Warmup error: {e}")
        return jsonify({"success": True}), 200  # Don't fail login over warmup


@app.route("/api/current-user", methods=["GET"])
@supabase_auth_required
def get_current_user():
    """Get current authenticated user info"""
    try:
        from StudyFlow.backend.supabase_client import get_user_profile, create_user_profile

        # Get user profile from Supabase
        profile = get_user_profile(request.user_id)

        # If profile doesn't exist, create it (handles edge case where auth exists but profile is missing)
        if not profile:
            debug_log(f"⚠️ User profile missing for {request.user_email}, creating now...")
            profile = create_user_profile(
                user_id=request.user_id,
                email=request.user_email,
                full_name=None,
                collective_brain_opt_in=False
            )

            if not profile:
                return jsonify({"error": "Failed to create user profile"}), 500

            # Auto-follow official account for new profiles
            official_id = os.getenv("STUDYFLOW_OFFICIAL_USER_ID")
            if official_id and request.user_id != official_id:
                try:
                    supabase.table("user_followers").insert({
                        "follower_id": request.user_id,
                        "following_id": official_id
                    }).execute()
                except:
                    pass

        return jsonify({
            "id": profile.get("id"),
            "email": profile.get("email"),
            "name": profile.get("full_name"),
            "university": profile.get("university"),
            "username": profile.get("username"),
            "subscription_tier": profile.get("subscription_tier", "free"),
            "subscription_status": profile.get("subscription_tier", "free"),
            "edu_verified": profile.get("edu_email_verified", False),
            "is_beta": False,
            "created_at": profile.get("created_at")
        }), 200

    except Exception as e:
        debug_log(f"❌ Get current user error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/stats", methods=["GET"])
@supabase_auth_required
def get_user_stats():
    """Get user statistics"""
    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()

        # Get user info
        cur.execute("""
            SELECT email, full_name, subscription_tier, FALSE as is_beta, created_at
            FROM user_profiles WHERE id = %s
        """, (request.user_id,))
        user = cur.fetchone()

        if not user:
            conn.close()
            return jsonify({'error': 'User not found'}), 404

        email, name, subscription_status, is_beta, created_at = user

        # Get total questions
        cur.execute("""
            SELECT COUNT(*) FROM questions WHERE user_id = %s
        """, (request.user_id,))
        total_questions = cur.fetchone()[0]

        # Get questions today
        cur.execute("""
            SELECT COUNT(*) FROM questions
            WHERE user_id = %s AND created_at >= CURRENT_DATE
        """, (request.user_id,))
        questions_today = cur.fetchone()[0]

        # Get questions this week
        cur.execute("""
            SELECT COUNT(*) FROM questions
            WHERE user_id = %s AND created_at >= CURRENT_DATE - INTERVAL '7 days'
        """, (request.user_id,))
        questions_this_week = cur.fetchone()[0]

        # Get question type breakdown
        cur.execute("""
            SELECT question_type, COUNT(*)
            FROM questions
            WHERE user_id = %s
            GROUP BY question_type
        """, (request.user_id,))
        question_types = dict(cur.fetchall())

        # Get most recent question date
        cur.execute("""
            SELECT MAX(created_at) FROM questions WHERE user_id = %s
        """, (request.user_id,))
        last_activity = cur.fetchone()[0]

        conn.close()

        # Determine tier
        tier = 'Free'
        if is_beta:
            tier = 'Beta (Pro)'
        elif subscription_status in ['active', 'trialing']:
            tier = 'Pro'

        return jsonify({
            'user': {
                'name': name,
                'email': email,
                'tier': tier,
                'created_at': created_at.isoformat() if created_at else None
            },
            'stats': {
                'total_questions': total_questions,
                'questions_today': questions_today,
                'questions_this_week': questions_this_week,
                'last_activity': last_activity.isoformat() if last_activity else None
            },
            'question_types': question_types
        }), 200

    except Exception as e:
        app.logger.error(f"❌ Get stats error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': 'Failed to get stats'}), 500


@app.route("/api/user/dashboard-layout", methods=["GET"])
@supabase_auth_required
def get_dashboard_layout():
    """Get user's dashboard widget layout and canvas pan position"""
    try:
        # Use Supabase to get dashboard_layout
        result = supabase.table("user_profiles").select("dashboard_layout").eq("id", request.user_id).execute()

        if result.data and len(result.data) > 0 and result.data[0].get('dashboard_layout'):
            data = result.data[0]['dashboard_layout']
            # Handle both old format (array) and new format (object with layout + canvasPan)
            if isinstance(data, list):
                # Old format - just array of widgets
                return jsonify({'layout': data, 'canvasPan': {'x': 0, 'y': 0}}), 200
            elif isinstance(data, dict):
                # New format - object with layout and canvasPan
                return jsonify(data), 200
            else:
                # Return default
                return jsonify({
                    'layout': ['recentConversations', 'upcomingEvents', 'studyGroups'],
                    'canvasPan': {'x': 0, 'y': 0}
                }), 200
        else:
            # Return default layout
            return jsonify({
                'layout': ['recentConversations', 'upcomingEvents', 'studyGroups'],
                'canvasPan': {'x': 0, 'y': 0}
            }), 200

    except Exception as e:
        debug_log(f"❌ Get dashboard layout error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500


@app.route("/api/user/dashboard-layout", methods=["POST"])
@supabase_auth_required
def save_dashboard_layout():
    """Save user's dashboard widget layout and canvas pan position"""
    try:
        data = request.get_json()
        layout = data.get('layout', [])
        canvas_pan = data.get('canvasPan', {'x': 0, 'y': 0})

        if not isinstance(layout, list):
            return jsonify({'error': 'Layout must be an array'}), 400

        # Store as a single JSON object with both layout and canvasPan
        dashboard_data = {
            'layout': layout,
            'canvasPan': canvas_pan
        }

        # Use Supabase to update dashboard_layout
        supabase.table("user_profiles").update({
            "dashboard_layout": dashboard_data
        }).eq("id", request.user_id).execute()

        return jsonify({'success': True}), 200

    except Exception as e:
        debug_log(f"❌ Save dashboard layout error: {e}\n{traceback.format_exc()}")
        return jsonify({'error': str(e)}), 500


@app.route("/api/user/preferences", methods=["GET"])
@supabase_auth_required
def get_user_preferences():
    """Get user's preferences (theme, etc.)"""
    try:
        # Try to get preferences column, but fallback gracefully if it doesn't exist
        try:
            result = supabase.table("user_profiles").select("preferences").eq("id", request.user_id).execute()
            if result.data and len(result.data) > 0:
                preferences = result.data[0].get('preferences') or {}
                return jsonify(preferences), 200
        except Exception as column_error:
            # Column might not exist yet - return defaults
            debug_log(f"⚠️ Preferences column not found, using defaults: {column_error}")
            pass

        # Return default preferences
        return jsonify({'theme': 'default'}), 200

    except Exception as e:
        debug_log(f"❌ Get preferences error: {e}\n{traceback.format_exc()}")
        # Don't fail - return defaults
        return jsonify({'theme': 'default'}), 200


@app.route("/api/user/preferences", methods=["PATCH"])
@supabase_auth_required
def update_user_preferences():
    """Update user's preferences (theme, etc.)"""
    try:
        updates = request.get_json()

        # Try to update preferences, but gracefully handle if column doesn't exist
        try:
            # Get current preferences
            result = supabase.table("user_profiles").select("preferences").eq("id", request.user_id).execute()
            current_prefs = {}
            if result.data and len(result.data) > 0:
                current_prefs = result.data[0].get('preferences') or {}

            # Merge updates with current preferences
            current_prefs.update(updates)

            # Save back to database
            supabase.table("user_profiles").update({
                "preferences": current_prefs
            }).eq("id", request.user_id).execute()

            return jsonify({'success': True, 'preferences': current_prefs}), 200

        except Exception as db_error:
            # Column might not exist - just acknowledge the update locally
            debug_log(f"⚠️ Could not save preferences to DB (column may not exist): {db_error}")
            return jsonify({'success': True, 'preferences': updates, 'note': 'Saved locally only'}), 200

    except Exception as e:
        debug_log(f"❌ Update preferences error: {e}\n{traceback.format_exc()}")
        return jsonify({'success': True, 'preferences': {}}), 200


@app.route("/api/leaderboard", methods=["GET"])
@supabase_auth_required
def get_leaderboard():
    """Get top contributors to the Nexus (leaderboard)"""
    try:
        # Check Redis cache
        if redis_cache:
            try:
                cached = redis_cache.get("leaderboard:main")
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        # Try to use the stored procedure first
        try:
            result = supabase.rpc("get_nexus_leaderboard", {
                "limit_count": 20
            }).execute()

            if result.data:
                return jsonify(result.data), 200
        except Exception as rpc_error:
            debug_log(f"RPC call failed (expected if stored procedure doesn't exist): {rpc_error}")
            # Continue to fallback logic below

        # Fallback: manually aggregate if RPC doesn't exist or returns no data
        from collections import defaultdict

        # Get all notes with user info
        notes = supabase.table("notes").select("user_id, page_count").eq("is_public", True).execute()

        # Aggregate pages by user
        user_pages = defaultdict(int)
        for note in notes.data:
            user_pages[note['user_id']] += note.get('page_count', 0)

        # Get user profiles for top contributors
        top_users = sorted(user_pages.items(), key=lambda x: x[1], reverse=True)[:20]

        leaderboard = []
        for user_id, pages in top_users:
            try:
                profile = supabase.table("user_profiles").select("username, university").eq("id", user_id).single().execute()
                if profile.data:
                    leaderboard.append({
                        "user_id": user_id,
                        "username": profile.data.get('username'),
                        "university": profile.data.get('university'),
                        "pages_contributed": pages
                    })
            except:
                pass

        if redis_cache:
            try: redis_cache.setex("leaderboard:main", 1800, json.dumps(leaderboard))
            except: pass
        return jsonify(leaderboard), 200

    except Exception as e:
        debug_log(f"Leaderboard error: {e}\n{traceback.format_exc()}")
        return jsonify([]), 200


@app.route("/api/leaderboard/cited", methods=["GET"])
@supabase_auth_required
def get_most_cited_notes():
    """Get most cited notes from AI chat responses"""
    try:
        from collections import defaultdict

        # Get all AI response logs with sources
        logs = supabase.table("ai_response_logs").select("sources_used").execute()

        # Count citations per note_id
        citation_counts = defaultdict(int)
        note_details = {}  # Store first occurrence details

        for log in logs.data:
            sources = log.get('sources_used') or []
            for source in sources:
                # Each source in the array counts as one citation
                note_id = source.get('note_id')
                if note_id:
                    citation_counts[note_id] += 1
                    # Store details from first occurrence
                    if note_id not in note_details:
                        note_details[note_id] = {
                            'filename': source.get('filename'),
                            'username': source.get('contributor_username'),
                        }

        # Get full note metadata for top cited notes
        top_note_ids = sorted(citation_counts.items(), key=lambda x: x[1], reverse=True)[:50]

        cited_notes = []
        for note_id, count in top_note_ids:
            try:
                # Get note metadata
                note = supabase.table("notes").select(
                    "id, original_filename, user_id, university, course_code"
                ).eq("id", note_id).single().execute()

                if note.data:
                    # Get username if not already in note_details
                    username = note_details[note_id].get('username')
                    if not username:
                        try:
                            profile = supabase.table("user_profiles").select("username").eq("id", note.data['user_id']).single().execute()
                            username = profile.data.get('username') if profile.data else None
                        except:
                            username = None

                    cited_notes.append({
                        "note_id": note_id,
                        "filename": note.data.get('original_filename') or note_details[note_id].get('filename'),
                        "username": username,
                        "university": note.data.get('university'),
                        "course_code": note.data.get('course_code'),
                        "citation_count": count
                    })
            except:
                # Note might have been deleted, skip it
                pass

        return jsonify(cited_notes[:20]), 200

    except Exception as e:
        debug_log(f"❌ Most cited error: {e}\n{traceback.format_exc()}")
        return jsonify([]), 200


@app.route("/api/leaderboard/downloaded", methods=["GET"])
@supabase_auth_required
def get_most_downloaded_notes():
    """Get most downloaded notes"""
    try:
        # Get notes with highest download counts
        # TODO: Implement download tracking (add download_count column to notes table)

        # Get all public notes
        result = supabase.table("notes").select(
            "id, original_filename, user_id, university, course_code, page_count"
        ).eq("is_public", True).limit(50).execute()

        downloaded_notes = []
        for note in result.data:
            # Skip Wikipedia entries
            university = note.get('university', '')
            if university and 'wikipedia' in university.lower():
                continue

            # Get username
            try:
                profile = supabase.table("user_profiles").select("username").eq("id", note['user_id']).single().execute()
                username = profile.data.get('username') if profile.data else None
            except:
                username = None

            downloaded_notes.append({
                "note_id": note['id'],
                "filename": note['original_filename'],
                "username": username,
                "university": note.get('university'),
                "course_code": note.get('course_code'),
                "download_count": note.get('page_count', 0)  # Temporary: simulate downloads
            })

        # Sort by download count
        downloaded_notes.sort(key=lambda x: x['download_count'], reverse=True)

        return jsonify(downloaded_notes[:20]), 200

    except Exception as e:
        debug_log(f"❌ Most downloaded error: {e}\n{traceback.format_exc()}")
        return jsonify([]), 200


@app.route("/api/leaderboard/rising", methods=["GET"])
@supabase_auth_required
def get_rising_stars():
    """Get rising stars (new contributors from last 30 days)"""
    try:
        from datetime import datetime, timedelta

        # Get notes uploaded in last 30 days
        thirty_days_ago = (datetime.now() - timedelta(days=30)).isoformat()

        result = supabase.table("notes").select("user_id, page_count").eq("is_public", True).gte("uploaded_at", thirty_days_ago).execute()

        # Aggregate pages by user
        from collections import defaultdict
        user_pages = defaultdict(int)
        for note in result.data:
            user_pages[note['user_id']] += note.get('page_count', 0)

        # Get top contributors
        top_users = sorted(user_pages.items(), key=lambda x: x[1], reverse=True)[:20]

        rising_stars = []
        for user_id, pages in top_users:
            try:
                profile = supabase.table("user_profiles").select("username, university").eq("id", user_id).single().execute()
                if profile.data:
                    rising_stars.append({
                        "user_id": user_id,
                        "username": profile.data.get('username'),
                        "university": profile.data.get('university'),
                        "pages_contributed": pages
                    })
            except:
                pass

        return jsonify(rising_stars), 200

    except Exception as e:
        debug_log(f"❌ Rising stars error: {e}\n{traceback.format_exc()}")
        return jsonify([]), 200


@app.route("/api/leaderboard/most-helpful", methods=["GET", "OPTIONS"])
@supabase_auth_required
def get_most_helpful_notes():
    """Get notes ranked by helpfulness votes"""
    # Handle OPTIONS preflight
    if request.method == 'OPTIONS':
        return '', 200

    try:
        from collections import defaultdict

        # Get all response ratings
        ratings = supabase.table("ai_response_ratings").select("cited_note_ids, vote").execute()

        # Calculate score for each note (upvotes - downvotes)
        note_scores = defaultdict(lambda: {"upvotes": 0, "downvotes": 0, "score": 0})

        for rating in ratings.data:
            cited_note_ids = rating.get('cited_note_ids', [])
            vote = rating.get('vote', 0)

            # Apply vote to all notes cited in that response
            for note_id in cited_note_ids:
                if vote == 1:
                    note_scores[note_id]["upvotes"] += 1
                elif vote == -1:
                    note_scores[note_id]["downvotes"] += 1
                note_scores[note_id]["score"] = note_scores[note_id]["upvotes"] - note_scores[note_id]["downvotes"]

        # Get top notes by score
        sorted_notes = sorted(note_scores.items(), key=lambda x: x[1]['score'], reverse=True)[:20]

        # Fetch note details
        helpful_notes = []
        for note_id, scores in sorted_notes:
            try:
                note = supabase.table("notes").select(
                    "id, original_filename, user_id, university, course_code"
                ).eq("id", note_id).single().execute()

                if note.data:
                    # Get username
                    profile = supabase.table("user_profiles").select("username").eq("id", note.data['user_id']).single().execute()
                    username = profile.data.get('username') if profile.data else None

                    helpful_notes.append({
                        "note_id": note_id,
                        "filename": note.data.get('original_filename'),
                        "username": username,
                        "university": note.data.get('university'),
                        "course_code": note.data.get('course_code'),
                        "helpfulness_score": scores['score'],
                        "upvotes": scores['upvotes'],
                        "downvotes": scores['downvotes']
                    })
            except:
                pass

        return jsonify(helpful_notes), 200

    except Exception as e:
        debug_log(f"❌ Most helpful notes error: {e}\n{traceback.format_exc()}")
        return jsonify([]), 200


@app.route("/api/leaderboard/universities", methods=["GET"])
@supabase_auth_required
def get_university_leaderboard():
    """Get university leaderboard - Springfield Showdown style"""
    try:
        from collections import defaultdict

        # Get list of valid universities from user_profiles (universities users have selected from dropdown)
        profiles = supabase.table("user_profiles").select("university").execute()
        valid_universities = set()
        for profile in profiles.data:
            uni = profile.get('university')
            if uni and 'wikipedia' not in uni.lower():
                valid_universities.add(uni)

        debug_log(f"[*] Valid universities from user profiles: {len(valid_universities)}")

        # Get all public notes with university info
        result = supabase.table("notes").select("university, page_count, user_id").eq("is_public", True).execute()

        # Aggregate stats by university (only valid ones)
        university_stats = defaultdict(lambda: {"total_pages": 0, "total_notes": 0, "contributors": set()})

        for note in result.data:
            university = note.get('university')
            # Only include universities that exist in user_profiles (from official dropdown)
            if university and university in valid_universities:
                university_stats[university]["total_pages"] += note.get('page_count', 0)
                university_stats[university]["total_notes"] += 1
                university_stats[university]["contributors"].add(note['user_id'])

        # Convert to list and sort by total pages
        leaderboard = []
        for university, stats in university_stats.items():
            leaderboard.append({
                "university": university,
                "total_pages": stats["total_pages"],
                "total_notes": stats["total_notes"],
                "contributor_count": len(stats["contributors"])
            })

        # Sort by total pages contributed
        leaderboard.sort(key=lambda x: x['total_pages'], reverse=True)

        return jsonify(leaderboard), 200

    except Exception as e:
        debug_log(f"❌ University leaderboard error: {e}\n{traceback.format_exc()}")
        return jsonify([]), 200


@app.route("/api/my-stats/citation-rank", methods=["GET"])
@supabase_auth_required
def get_my_citation_rank():
    """Get current user's rank by citation count"""
    try:
        from collections import defaultdict

        # Get all AI response logs with sources
        logs = supabase.table("ai_response_logs").select("sources_used").execute()

        # Count citations per note_id (and track which user owns each note)
        note_citations = defaultdict(int)
        note_owners = {}

        for log in logs.data:
            sources = log.get('sources_used') or []
            for source in sources:
                note_id = source.get('note_id')
                if note_id:
                    note_citations[note_id] += 1

        # Get note owners
        if note_citations:
            note_ids = list(note_citations.keys())
            notes = supabase.table("notes").select("id, user_id").in_("id", note_ids).execute()
            for note in notes.data:
                note_owners[note['id']] = note['user_id']

        # Aggregate citations by user
        user_citations = defaultdict(int)
        for note_id, count in note_citations.items():
            user_id = note_owners.get(note_id)
            if user_id:
                user_citations[user_id] += count

        # Sort users by citation count
        sorted_users = sorted(user_citations.items(), key=lambda x: x[1], reverse=True)

        # Find current user's rank
        my_rank = None
        my_citations = user_citations.get(request.user_id, 0)
        total_ranked_users = len(sorted_users)

        for index, (user_id, count) in enumerate(sorted_users):
            if user_id == request.user_id:
                my_rank = index + 1
                break

        # If user has no citations, they're unranked
        if my_rank is None:
            my_rank = total_ranked_users + 1 if total_ranked_users > 0 else 1

        return jsonify({
            "rank": my_rank,
            "citation_count": my_citations,
            "total_users": max(total_ranked_users, my_rank)
        }), 200

    except Exception as e:
        debug_log(f"❌ Citation rank error: {e}\n{traceback.format_exc()}")
        return jsonify({"rank": 0, "citation_count": 0, "total_users": 0}), 200


@app.route("/api/my-stats/university-rank", methods=["GET"])
@supabase_auth_required
def get_my_university_rank():
    """Get current user's university rank"""
    try:
        from collections import defaultdict

        # Get user's university
        profile = supabase.table("user_profiles").select("university").eq("id", request.user_id).single().execute()
        my_university = profile.data.get('university') if profile.data else None

        if not my_university:
            return jsonify({"rank": 0, "university": None, "total_pages": 0, "total_universities": 0}), 200

        # Get list of valid universities
        profiles = supabase.table("user_profiles").select("university").execute()
        valid_universities = set()
        for p in profiles.data:
            uni = p.get('university')
            if uni and 'wikipedia' not in uni.lower():
                valid_universities.add(uni)

        # Get all public notes with university info
        result = supabase.table("notes").select("university, page_count").eq("is_public", True).execute()

        # Aggregate stats by university
        university_pages = defaultdict(int)
        for note in result.data:
            university = note.get('university')
            if university and university in valid_universities:
                university_pages[university] += note.get('page_count', 0)

        # Sort universities by total pages
        sorted_universities = sorted(university_pages.items(), key=lambda x: x[1], reverse=True)

        # Find my university's rank
        my_rank = None
        my_pages = university_pages.get(my_university, 0)
        total_universities = len(sorted_universities)

        for index, (university, pages) in enumerate(sorted_universities):
            if university == my_university:
                my_rank = index + 1
                break

        # If university has no pages, it's unranked
        if my_rank is None:
            my_rank = total_universities + 1 if total_universities > 0 else 1

        return jsonify({
            "rank": my_rank,
            "university": my_university,
            "total_pages": my_pages,
            "total_universities": max(total_universities, my_rank)
        }), 200

    except Exception as e:
        debug_log(f"❌ University rank error: {e}\n{traceback.format_exc()}")
        return jsonify({"rank": 0, "university": None, "total_pages": 0, "total_universities": 0}), 200


# ============================================================================
# END USER AUTHENTICATION & SUBSCRIPTION ROUTES
# ============================================================================


@app.route("/api/process", methods=["POST"])
@supabase_auth_required
def process_data():
    try:
        # Check rate limit for free users
        can_proceed, remaining = check_question_limit(request.user_id)
        if not can_proceed:
            debug_log(f"❌ Rate limit exceeded for user {request.user_id}")
            return jsonify({
                "error": "Daily question limit exceeded",
                "message": "Free users are limited to 10 questions per day. Upgrade to Pro for unlimited questions!",
                "limit_exceeded": True
            }), 429

        debug_log(f"✅ Rate limit check passed. Remaining questions: {remaining if remaining >= 0 else 'unlimited'}")

        ocr_json = request.get_json()
        if not ocr_json:
            debug_log("❌ No JSON provided")
            return jsonify({"error": "No JSON provided"}), 400

        question_text = ocr_json.get("question", "").strip()
        if not question_text:
            debug_log("❌ No question found in input")
            return jsonify({"error": "No question text provided"}), 400

        # Connect and check cache
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()
        cur.execute("SELECT answer, count FROM qa_pairs WHERE question = %s", (question_text,))
        row = cur.fetchone()

        if row:
            saved_answer, current_count = row
            new_count = current_count + 1
            cur.execute(
                "UPDATE qa_pairs SET count = %s WHERE question = %s",
                (new_count, question_text)
            )
            conn.commit()

            debug_log("📦 Using cached answer from DB")
            debug_log(f"✅ Q: {question_text[:100]}")
            debug_log(f"✅ A: {saved_answer}")
            debug_log(f"📈 Count incremented to {new_count}")

            # Try to find which key in the incoming answers matches the saved text
            for key, val in ocr_json.get("answers", {}).items():
                if val.get("text", "").strip() == saved_answer.strip():
                    conn.close()
                    return jsonify({"result": int(key)})

    
            # Fallback if none matched exactly
            conn.close()
            return jsonify({"result": 1})
    

        # Not cached → queue an async AI task
        debug_log("📨 About to queue async task...")
        task = process_question_async.delay(ocr_json)
        conn.close()
        debug_log(f"✅ Task queued: {task.id}")
        return jsonify({
            "status": "processing",
            "message": "Question sent for async processing",
            "task_id": task.id
        })

    except Exception as e:
        debug_log(f"🔥 Error in /api/process: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

from StudyFlow.logging_utils import debug_log

@app.route("/api/stripe_webhook", methods=["POST"])
def stripe_webhook():
    """
    Handle Stripe webhook events for Scholar's Club subscriptions

    Events handled:
    - checkout.session.completed: User completes payment, upgrade to 'pro' tier
    - customer.subscription.updated: Subscription status changes
    - customer.subscription.deleted: Subscription canceled, downgrade to 'free' tier
    """
    from StudyFlow.backend.supabase_client import supabase

    # 1) Raw body + signature header
    payload = request.get_data()
    sig_header = request.headers.get("Stripe-Signature")

    # 2) Verify & parse
    try:
        event = stripe.Webhook.construct_event(payload, sig_header, WEBHOOK_SECRET)
        debug_log(f"✅ Stripe webhook verified: {event['id']} → {event['type']}")
    except (ValueError, stripe.error.SignatureVerificationError) as e:
        debug_log(f"⚠️ Webhook validation failed: {e}")
        return "", 400

    evt_type = event["type"]
    obj = event["data"]["object"]

    # 3) Handle checkout.session.completed - User just subscribed to Scholar's Club
    if evt_type == "checkout.session.completed":
        try:
            cust_id = obj["customer"]
            subscription_id = obj.get("subscription")
            user_id = obj["metadata"].get("user_id")  # From checkout session metadata

            if not user_id:
                debug_log(f"⚠️ No user_id in checkout session metadata")
                return jsonify({"received": True}), 200

            # Update user_profiles to 'pro' tier
            supabase.table("user_profiles").update({
                "subscription_tier": "pro",
                "subscription_status": "active",
                "stripe_customer_id": cust_id,
                "stripe_subscription_id": subscription_id
            }).eq("id", user_id).execute()

            debug_log(f"🎉 Scholar's Club activated for user {user_id}")

        except Exception as e:
            debug_log(f"❌ checkout.session.completed error: {e}\n{traceback.format_exc()}")
            return "", 500

    # 4) Handle customer.subscription.updated - Subscription status changed
    elif evt_type == "customer.subscription.updated":
        try:
            cust_id = obj["customer"]
            status = obj["status"]
            subscription_id = obj["id"]

            # Update subscription status in user_profiles
            supabase.table("user_profiles").update({
                "subscription_status": status,
                "stripe_subscription_id": subscription_id
            }).eq("stripe_customer_id", cust_id).execute()

            debug_log(f"🔄 Subscription updated: {cust_id} → {status}")

        except Exception as e:
            debug_log(f"❌ customer.subscription.updated error: {e}\n{traceback.format_exc()}")
            return "", 500

    # 5) Handle customer.subscription.deleted - User canceled or subscription expired
    elif evt_type == "customer.subscription.deleted":
        try:
            cust_id = obj["customer"]

            # Downgrade user back to free tier
            supabase.table("user_profiles").update({
                "subscription_tier": "free",
                "subscription_status": "canceled",
                "stripe_subscription_id": None
            }).eq("stripe_customer_id", cust_id).execute()

            debug_log(f"🗑️ Subscription canceled, user downgraded to free tier: {cust_id}")

        except Exception as e:
            debug_log(f"❌ customer.subscription.deleted error: {e}\n{traceback.format_exc()}")
            return "", 500

    return jsonify({"received": True}), 200

@app.route("/api/deepflow_question", methods=["POST"])
def deepflow_question():
    try:
        data = request.get_json()
        topic = data.get("topic", "default topic")
        previous_questions = data.get("previous_questions", [])
        debug_log(f"Received deepflow question request for topic '{topic}' with previous questions: {previous_questions}")

        question_data = get_deepflow_question(topic, previous_questions)
        if question_data is None:
            debug_log("Failed to generate deepflow question.")
            return jsonify({"error": "Failed to generate question"}), 500

        debug_log("Deepflow question generated successfully.")
        return jsonify(question_data), 200

    except Exception as e:
        debug_log(f"🔥 Error in /api/deepflow_question: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/check_subscription")
def check_subscription():
    sid = request.args.get("stripe_id")
    if not sid:
        return jsonify({"error": "Missing stripe_id"}), 400

    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()
        cur.execute(
            "SELECT subscription_tier FROM user_profiles WHERE stripe_customer_id = %s",
            (sid,)
        )
        row = cur.fetchone()
        cur.close()
        conn.close()

        if not row:
            return jsonify({"error": "Unknown customer"}), 404

        return jsonify({"subscription_status": row[0]}), 200

    except Exception as e:
        app.logger.error(f"❌ check_subscription error: {e}")
        return jsonify({"error": "Server error"}), 500


@app.route("/ocr", methods=["POST"])
def ocr_endpoint():
    debug_log("🔍 /ocr endpoint hit")
    if "image" not in request.files:
        debug_log("❌ No image in request")
        return jsonify({"error": "No image provided"}), 400
    try:
        file = request.files["image"]
        image = Image.open(file.stream)
        processed = preprocess_image(image)

        ocr_text = pytesseract.image_to_string(processed)
        data = pytesseract.image_to_data(
            processed,
            output_type=pytesseract.Output.DICT,
            config="--psm 6 --oem 3"
        )
        mapping = {}
        tag_number = 1
        for i, txt in enumerate(data["text"]):
            text = txt.strip()
            try:
                conf = float(data["conf"][i])
            except:
                continue
            if text and conf > 0:
                mapping[str(tag_number)] = {
                    "text": text,
                    "left": data["left"][i],
                    "top": data["top"][i],
                    "width": data["width"][i],
                    "height": data["height"][i],
                    "line_num": data["line_num"][i]
                }
                tag_number += 1

        tagged_text = " ".join(f"[{k}] {v['text']}" for k, v in mapping.items())
        return jsonify({"ocr_text": tagged_text, "mapping": mapping})

    except Exception as e:
        debug_log(f"🔥 OCR processing failed: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/layout", methods=["POST"])
def layout():
    data = request.get_json()
    text = data.get("text", "")
    if not text:
        return jsonify({"error": "No text provided"}), 400

    try:
        # 1) Let the model extract question + answers, preserving tags
        resp = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            temperature=0,
            messages=[
                {
                    "role": "system",
                    "content": (
                        "You are an OCR layout engine. The input text is annotated with "
                        "bracketed numeric tags, for example:\n\n"
                        "[12] A. Clostridium difficile\n"
                        "[15] B. Helicobacter pylori\n"
                        "[18] C. Helicobacter baculiformis\n"
                        "[21] D. Vibrio vulnificus\n\n"
                        "Extract the question and the answer options into JSON using "
                        "exactly those same bracket numbers as both the keys and the "
                        "`tag` values. Do NOT renumber anything.\n\n"
                        "Return JSON in this shape:\n"
                        "{\n"
                        '  "question": "<the question text>",\n'
                        '  "answers": {\n'
                        '    "12": {"text": "A. Clostridium difficile",       "tag": 12},\n'
                        '    "15": {"text": "B. Helicobacter pylori",        "tag": 15},\n'
                        '    "18": {"text": "C. Helicobacter baculiformis", "tag": 18},\n'
                        '    "21": {"text": "D. Vibrio vulnificus",         "tag": 21}\n'
                        "  }\n"
                        "}"
                    )
                },
                {"role": "user", "content": text}
            ]
        )

        # 2) Parse the model’s output
        extracted = json.loads(resp.choices[0].message.content.strip())
        raw_answers = extracted.get("answers", {})

        # 3) Remap into position-based keys "1", "2", … while keeping the original tag
        #    sorted by numeric tag so display order matches the OCR tags in reading order
        sorted_tags = sorted(raw_answers.keys(), key=lambda k: int(k))
        positioned = {}
        for i, tag_str in enumerate(sorted_tags, start=1):
            entry = raw_answers[tag_str]
            positioned[str(i)] = {
                "text": entry["text"],
                "tag": int(tag_str)
            }

        # 4) Overwrite and return
        extracted["answers"] = positioned
        return jsonify({"structured_ai": extracted}), 200

    except Exception as e:
        debug_log(f"🔥 /api/layout error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/fallback", methods=["POST"])
def fallback():
    try:
        data = request.get_json()
        mapping = data.get("ocr_mapping")
        expected = int(data.get("expected_answers", 4))
        if not mapping:
            return jsonify({"error": "Missing ocr_mapping"}), 400

        from StudyFlow.backend.ocr_logic import fallback_structure
        return jsonify(fallback_structure(mapping, expected)), 200

    except Exception as e:
        debug_log(f"🔥 /api/fallback error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/merge", methods=["POST"])
def merge():
    try:
        data = request.get_json()
        from StudyFlow.backend.ocr_logic import merge_ai_and_fallback
        merged = merge_ai_and_fallback(
            data.get("ai_json", {}),
            data.get("fallback_json", {}),
            data.get("ocr_mapping", {})
        )
        return jsonify(merged), 200

    except Exception as e:
        debug_log(f"🔥 /api/merge error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/select-best-ocr", methods=["POST"])
def select_best_ocr():
    data = request.get_json()
    cands = data.get("candidates")
    if not isinstance(cands, list) or not cands:
        return jsonify({"error": "You must provide a list of candidates"}), 400
    if len(cands) == 1:
        return jsonify({"chosen_index": 1}), 200

    prompt = "Below are OCR candidate outputs:\n\n" + "\n\n".join(
        f"Candidate {i+1}:\n{txt}" for i, txt in enumerate(cands)
    ) + f"\n\nWhich is best? Return only the number 1–{len(cands)}."
    try:
        resp = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0
        )
        m = re.search(r"(\d+)", resp.choices[0].message.content)
        return jsonify({"chosen_index": int(m.group(1)) if m else 1}), 200
    except Exception as e:
        debug_log(f"🔥 /api/select-best-ocr error: {e}")
        return jsonify({"chosen_index": 1}), 200


@app.route("/api/log", methods=["POST"])
def receive_log():
    try:
        data = request.get_json()
        msg = data.get("message", "")
        if msg:
            print(msg)
            with open("backend_log.txt", "a", encoding="utf-8") as f:
                f.write(msg + "\n")
        return jsonify({"status": "ok"}), 200
    except Exception as e:
        debug_log(f"🔥 /api/log error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/explanation", methods=["POST"])
def generate_explanation():
    try:
        data = request.get_json()
        ocr_json = data.get("ocr_json")
        idx = data.get("chosen_index")
        if ocr_json is None or idx is None:
            return jsonify({"error": "Missing data"}), 400

        prompt = (
            f"Here is the OCR output in JSON:\n{ocr_json}\n"
            f"Explain why answer option {idx} is correct (max 100 words)."
        )
        resp = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0
        )
        return jsonify({"explanation": resp.choices[0].message.content.strip()}), 200

    except Exception as e:
        debug_log(f"🔥 /api/explanation error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/focusflow", methods=["POST"])
def api_focusflow():
    try:
        from StudyFlow.backend.ocr_logic import fallback_structure, merge_ai_and_fallback

        # 1️⃣ grab the uploaded image
        if "image" not in request.files:
            return jsonify({"error":"No image file provided"}), 400
        file = request.files["image"]
        img = Image.open(file.stream).convert("RGB")

        # 2️⃣ preprocess & OCR
        proc = preprocess_image(img)
        data = pytesseract.image_to_data(proc, output_type=pytesseract.Output.DICT)
        mapping = {}
        tag = 1
        for i, txt in enumerate(data["text"]):
            w = txt.strip()
            try: conf = float(data["conf"][i])
            except: continue
            if w and conf > 0:
                mapping[str(tag)] = {
                    "text": w,
                    "left":   data["left"][i],
                    "top":    data["top"][i],
                    "width":  data["width"][i],
                    "height": data["height"][i],
                    "line_num": data["line_num"][i]
                }
                tag += 1
        tagged = " ".join(f"[{k}] {v['text']}" for k,v in mapping.items())

        # 3️⃣ layout via API
        layout_resp = requests.post(f"{BACKEND_URL}/api/layout", json={"text": tagged})
        if layout_resp.status_code != 200:
            return jsonify({"error":"Layout failed","details":layout_resp.text}), 500
        ai_json = layout_resp.json().get("structured_ai", {})

        # 4️⃣ merge with fallback
        expected = len(ai_json.get("answers", {})) if ai_json else 4
        fb_json  = fallback_structure(mapping, expected)
        merged = merge_ai_and_fallback(ai_json, fb_json, mapping) if ai_json and fb_json.get("answers") else ai_json or fb_json

           # 5️⃣ AI vote
        idx = triple_call_ai_api_json_final(merged)
        full = merged["answers"].get(str(idx), {}).get("text", "Unknown")

    # 6️⃣ Generate explanation inline
        prompt = (
            f"Here is the OCR output in JSON:\n{json.dumps(merged)}\n"
            f"Explain why answer option {idx} is correct (max 100 words)."
        )
        resp = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0
        )
        explanation = resp.choices[0].message.content.strip()

    # 7️⃣ Return everything
        return jsonify({
            "full_answer": full,
            "explanation": explanation,
            "merged_json": merged,
            "tagged_text": tagged
        }), 200

    except Exception as e:
        debug_log(f"🔥 /api/focusflow error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500



@app.route("/api/find_button", methods=["POST"])
def find_button():
    try:
        if "image" not in request.files or "template" not in request.files:
            return jsonify({"error": "Missing image or template"}), 400

        npimg = np.frombuffer(request.files["image"].read(), np.uint8)
        img = cv2.imdecode(npimg, cv2.IMREAD_GRAYSCALE)

        tnp = np.frombuffer(request.files["template"].read(), np.uint8)
        tmpl = cv2.imdecode(tnp, cv2.IMREAD_GRAYSCALE)

        best_val, best_loc, best_shape = 0, None, None
        for scale in np.linspace(0.8, 1.2, 10):
            tpl = cv2.resize(tmpl, None, fx=scale, fy=scale)
            res = cv2.matchTemplate(img, tpl, cv2.TM_CCOEFF_NORMED)
            _, mx, _, ml = cv2.minMaxLoc(res)
            if mx > best_val:
                best_val, best_loc, best_shape = mx, ml, tpl.shape

        if best_val >= 0.7 and best_loc and best_shape:
            h, w = best_shape
            cx, cy = best_loc[0] + w//2, best_loc[1] + h//2
            return jsonify({
                "center_x": int(cx),
                "center_y": int(cy),
                "confidence": float(best_val)
            }), 200

        return jsonify({"error": "Button not found", "confidence": best_val}), 404

    except Exception as e:
        debug_log(f"🔥 /api/find_button error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/button-templates")
def admin_view_button_templates():
    try:
        templates_dir = os.path.join(os.path.dirname(__file__), "static", "button_templates")
        meta = os.path.join(templates_dir, "submit_template_index.json")
        data = {}
        if os.path.exists(meta):
            data = json.load(open(meta, encoding="utf-8"))
        items = sorted(data.items(), key=lambda kv: -kv[1].get("count", 0))

        html = """
        <!DOCTYPE html><html><head><title>Submit Button Templates</title>
        <style>
        body{font-family:Arial;background:#f4f4f4;padding:20px}
        .grid{display:grid;grid-template-columns:repeat(auto-fill,minmax(200px,1fr));gap:20px}
        .item{background:#fff;padding:10px;border-radius:8px;box-shadow:0 2px 5px rgba(0,0,0,0.1);text-align:center}
        .item img{max-width:100%;border-radius:6px;margin-bottom:10px}
        h1{text-align:center}
        </style></head><body>
        <h1>Submit Button Templates</h1><div class="grid">
        {% for filename,data in templates %}
          <div class="item">
            <img src="/admin/button-image/{{ filename }}" alt="{{ filename }}">
            <div><strong>{{ filename }}</strong></div>
            <div>Matches: {{ data.get('count',0) }}</div>
          </div>
        {% endfor %}
        </div></body></html>
        """
        return render_template_string(html, templates=items)

    except Exception as e:
        debug_log(f"🔥 /admin/button-templates error: {e}\n{traceback.format_exc()}")
        return f"<h1>Error:</h1><p>{e}</p>"


@app.route("/admin/button-image/<path:filename>")
def serve_button_template(filename):
    try:
        templates_dir = os.path.join(os.path.dirname(__file__), "static", "button_templates")
        return send_from_directory(templates_dir, filename)
    except Exception as e:
        debug_log(f"🔥 /admin/button-image error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/view-qa")
def view_qa():
    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()
        cur.execute("SELECT question, answer, timestamp, count FROM qa_pairs ORDER BY count DESC")
        rows = cur.fetchall()
        cur.execute("SELECT COUNT(*), SUM(count) FROM qa_pairs")
        total, total_count = cur.fetchone()
        conn.close()

        html = f"<h1>Stored Questions & Answers</h1><p>Total Questions: {total} | Total Attempts: {total_count}</p><ul>"
        for q,a,t,c in rows:
            html += f"<li><b>Q:</b> {q}<br><b>A:</b> {a}<br><small>{t} | Count: {c}</small><br><br></li>"
        html += "</ul>"
        return html

    except Exception as e:
        debug_log(f"🔥 /admin/view-qa error: {e}\n{traceback.format_exc()}")
        return f"<h1>Error:</h1><p>{e}</p>"


@app.route("/admin/view-questions")
def view_questions():
    """Admin page to view all questions from the questions table"""
    try:
        conn = psycopg2.connect(os.environ["DATABASE_URL"])
        cur = conn.cursor()

        # Get filter parameters
        user_id_filter = request.args.get('user_id', None)
        question_type_filter = request.args.get('type', None)
        limit = int(request.args.get('limit', 100))

        # Build query with filters
        query = "SELECT q.id, q.user_id, u.email, q.question_text, q.question_type, q.answers_json, q.ai_answer, q.ai_reasoning, q.created_at FROM questions q LEFT JOIN user_profiles u ON q.user_id = u.id WHERE 1=1"
        params = []

        if user_id_filter:
            query += " AND q.user_id = %s"
            params.append(int(user_id_filter))

        if question_type_filter:
            query += " AND q.question_type = %s"
            params.append(question_type_filter)

        query += " ORDER BY q.created_at DESC LIMIT %s"
        params.append(limit)

        cur.execute(query, params)
        rows = cur.fetchall()

        # Get statistics
        cur.execute("SELECT COUNT(*) FROM questions")
        total_questions = cur.fetchone()[0]

        cur.execute("SELECT question_type, COUNT(*) FROM questions GROUP BY question_type")
        type_counts = cur.fetchall()

        cur.execute("SELECT COUNT(DISTINCT user_id) FROM questions")
        unique_users = cur.fetchone()[0]

        conn.close()

        # Build HTML with stats
        stats_html = ""
        for q_type, count in type_counts:
            stats_html += f'<div class="stat-card"><h3>{count}</h3><p>{q_type.replace("_", " ").title()}</p></div>'

        # Build table rows
        rows_html = ""
        for row in rows:
            q_id, user_id, email, question_text, question_type, answers_json, ai_answer, ai_reasoning, created_at = row
            question_display = (question_text[:100] + '...') if len(question_text) > 100 else question_text
            answer_display = (ai_answer[:80] + '...') if ai_answer and len(ai_answer) > 80 else (ai_answer or 'N/A')
            timestamp_str = created_at.strftime('%Y-%m-%d %H:%M') if created_at else 'N/A'

            rows_html += f'''<tr>
                <td>{q_id}</td>
                <td><div><strong>ID {user_id}</strong></div><div class="email">{email or 'Unknown'}</div></td>
                <td><span class="badge badge-{question_type}">{question_type.replace("_", " ").title()}</span></td>
                <td class="question-text" title="{question_text}">{question_display}</td>
                <td class="answer-text" title="{ai_answer or ''}">{answer_display}</td>
                <td class="timestamp">{timestamp_str}</td>
            </tr>'''

        html = f'''<!DOCTYPE html>
<html>
<head>
    <title>Question Database</title>
    <style>
        body {{ font-family: Arial, sans-serif; background: #f4f4f4; padding: 20px; }}
        .header {{ background: white; padding: 20px; border-radius: 8px; margin-bottom: 20px; box-shadow: 0 2px 5px rgba(0,0,0,0.1); }}
        .stats {{ display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 15px; margin-bottom: 20px; }}
        .stat-card {{ background: #667eea; color: white; padding: 15px; border-radius: 8px; text-align: center; }}
        .stat-card h3 {{ margin: 0; font-size: 32px; }}
        .stat-card p {{ margin: 5px 0 0 0; font-size: 14px; }}
        .filters {{ background: white; padding: 15px; border-radius: 8px; margin-bottom: 20px; box-shadow: 0 2px 5px rgba(0,0,0,0.1); }}
        .filters form {{ display: flex; gap: 10px; align-items: center; flex-wrap: wrap; }}
        .filters input, .filters select, .filters button {{ padding: 8px 12px; border: 1px solid #ddd; border-radius: 4px; }}
        .filters button {{ background: #667eea; color: white; cursor: pointer; border: none; }}
        .filters button:hover {{ background: #5568d3; }}
        .filters a {{ padding: 8px 12px; background: #f5f5f5; border-radius: 4px; text-decoration: none; color: #333; }}
        table {{ width: 100%; background: white; border-collapse: collapse; border-radius: 8px; overflow: hidden; box-shadow: 0 2px 5px rgba(0,0,0,0.1); }}
        th {{ background: #667eea; color: white; padding: 12px; text-align: left; font-weight: 600; }}
        td {{ padding: 12px; border-bottom: 1px solid #eee; }}
        tr:hover {{ background: #f9f9f9; }}
        .question-text {{ max-width: 400px; overflow: hidden; text-overflow: ellipsis; }}
        .answer-text {{ max-width: 300px; overflow: hidden; text-overflow: ellipsis; color: #333; }}
        .badge {{ display: inline-block; padding: 4px 8px; border-radius: 4px; font-size: 11px; font-weight: 600; }}
        .badge-multiple_choice {{ background: #e3f2fd; color: #1976d2; }}
        .badge-multiple_answer {{ background: #f3e5f5; color: #7b1fa2; }}
        .badge-essay {{ background: #fff3e0; color: #f57c00; }}
        .badge-short_answer {{ background: #e8f5e9; color: #388e3c; }}
        .badge-fill_in_blank {{ background: #fce4ec; color: #c2185b; }}
        .email {{ color: #666; font-size: 12px; }}
        .timestamp {{ color: #999; font-size: 12px; }}
    </style>
</head>
<body>
    <div class="header">
        <h1>📊 Question Database</h1>
        <p>All questions answered by StudyFlow users</p>
    </div>
    <div class="stats">
        <div class="stat-card"><h3>{total_questions}</h3><p>Total Questions</p></div>
        <div class="stat-card"><h3>{unique_users}</h3><p>Unique Users</p></div>
        {stats_html}
    </div>
    <div class="filters">
        <form method="get">
            <input type="number" name="user_id" placeholder="User ID" value="{user_id_filter or ''}">
            <select name="type">
                <option value="">All Types</option>
                <option value="multiple_choice" {"selected" if question_type_filter == "multiple_choice" else ""}>Multiple Choice</option>
                <option value="multiple_answer" {"selected" if question_type_filter == "multiple_answer" else ""}>Multiple Answer</option>
                <option value="essay" {"selected" if question_type_filter == "essay" else ""}>Essay</option>
                <option value="short_answer" {"selected" if question_type_filter == "short_answer" else ""}>Short Answer</option>
                <option value="fill_in_blank" {"selected" if question_type_filter == "fill_in_blank" else ""}>Fill in Blank</option>
            </select>
            <select name="limit">
                <option value="50" {"selected" if limit == 50 else ""}>50 results</option>
                <option value="100" {"selected" if limit == 100 else ""}>100 results</option>
                <option value="500" {"selected" if limit == 500 else ""}>500 results</option>
                <option value="1000" {"selected" if limit == 1000 else ""}>1000 results</option>
            </select>
            <button type="submit">Filter</button>
            <a href="/admin/view-questions">Clear</a>
        </form>
    </div>
    <table>
        <thead>
            <tr>
                <th>ID</th>
                <th>User</th>
                <th>Type</th>
                <th>Question</th>
                <th>AI Answer</th>
                <th>Date</th>
            </tr>
        </thead>
        <tbody>
            {rows_html}
        </tbody>
    </table>
</body>
</html>'''

        return html

    except Exception as e:
        debug_log(f"🔥 /admin/view-questions error: {e}\n{traceback.format_exc()}")
        return f"<h1>Error:</h1><p>{e}</p><pre>{traceback.format_exc()}</pre>"


@app.route("/api/status/<task_id>")
def get_task_status(task_id):
    try:
        task = celery_app.AsyncResult(task_id)
        if task.state == "PENDING":
            return jsonify({"status": "pending"}), 202
        if task.state == "SUCCESS":
            return jsonify({"status": "complete", "result": task.result}), 200
        if task.state == "FAILURE":
            return jsonify({"status": "failed"}), 500
        return jsonify({"status": task.state}), 202

    except Exception as e:
        debug_log(f"🔥 /api/status error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

@app.route("/api/home_message", methods=["GET", "POST"])
def home_message():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        # Read the new message from the JSON body
        msg = request.get_json().get("message", "").strip()
        # Insert or update the key/value in app_config
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('home_message', %s)
            ON CONFLICT (key)
            DO UPDATE SET value = EXCLUDED.value
        """, (msg,))
        conn.commit()
        conn.close()
        return jsonify({"message": msg})

    # GET: fetch the stored message or return a default
    cur.execute("SELECT value FROM app_config WHERE key = 'home_message'")
    row = cur.fetchone()
    conn.close()
    return jsonify({
        "message": row[0] if row else "Welcome to StudyFlow!"
    })

@app.route("/api/freeflow_message", methods=["GET", "POST"])
def freeflow_message():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        msg = request.get_json().get("message", "").strip()
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('freeflow_message', %s)
            ON CONFLICT (key)
            DO UPDATE SET value = EXCLUDED.value
        """, (msg,))
        conn.commit()
        conn.close()
        return jsonify({"message": msg})

    cur.execute("SELECT value FROM app_config WHERE key = 'freeflow_message'")
    row = cur.fetchone()
    conn.close()
    return jsonify({
        "message": row[0] if row else "Welcome to FreeFlow!"
    })

@app.route("/api/focusflow_message", methods=["GET", "POST"])
def focusflow_message():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        msg = request.get_json().get("message", "").strip()
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('focusflow_message', %s)
            ON CONFLICT (key)
            DO UPDATE SET value = EXCLUDED.value
        """, (msg,))
        conn.commit()
        conn.close()
        return jsonify({"message": msg})

    cur.execute("SELECT value FROM app_config WHERE key = 'focusflow_message'")
    row = cur.fetchone()
    conn.close()
    return jsonify({
        "message": row[0] if row else "Welcome to FocusFlow!"
    })

@app.route("/api/deepflow_message", methods=["GET", "POST"])
def deepflow_message():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        msg = request.get_json().get("message", "").strip()
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('deepflow_message', %s)
            ON CONFLICT (key)
            DO UPDATE SET value = EXCLUDED.value
        """, (msg,))
        conn.commit()
        conn.close()
        return jsonify({"message": msg})

    cur.execute("SELECT value FROM app_config WHERE key = 'deepflow_message'")
    row = cur.fetchone()
    conn.close()
    return jsonify({
        "message": row[0] if row else "Welcome to DeepFlow!"
    })

@app.route("/admin/freeflow_message", methods=["GET", "POST"])
def admin_freeflow_message():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        new_msg = request.form.get("message", "").strip()
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('freeflow_message', %s)
            ON CONFLICT (key) DO UPDATE SET value = EXCLUDED.value
        """, (new_msg,))
        conn.commit()
        current = new_msg
    else:
        cur.execute("SELECT value FROM app_config WHERE key = 'freeflow_message'")
        row = cur.fetchone()
        current = row[0] if row else ""

    conn.close()

    return render_template_string("""
    <!doctype html>
    <html>
      <head>
        <meta charset="utf-8">
        <title>Admin: FreeFlow Message</title>
      </head>
      <body style="font-family: sans-serif; padding: 2rem;">
        <h1>Update FreeFlow Message</h1>
        <form method="post">
          <textarea name="message" rows="4" cols="60"
            style="font-size:1rem; padding:0.5rem;">{{ message }}</textarea><br><br>
          <button type="submit" style="font-size:1rem; padding:0.5rem 1rem;">
            Save
          </button>
        </form>
      </body>
    </html>
    """, message=current)


@app.route("/admin/focusflow_message", methods=["GET", "POST"])
def admin_focusflow_message():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        new_msg = request.form.get("message", "").strip()
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('focusflow_message', %s)
            ON CONFLICT (key) DO UPDATE SET value = EXCLUDED.value
        """, (new_msg,))
        conn.commit()
        current = new_msg
    else:
        cur.execute("SELECT value FROM app_config WHERE key = 'focusflow_message'")
        row = cur.fetchone()
        current = row[0] if row else ""

    conn.close()

    return render_template_string("""
    <!doctype html>
    <html>
      <head>
        <meta charset="utf-8">
        <title>Admin: FocusFlow Message</title>
      </head>
      <body style="font-family: sans-serif; padding: 2rem;">
        <h1>Update FocusFlow Message</h1>
        <form method="post">
          <textarea name="message" rows="4" cols="60"
            style="font-size:1rem; padding:0.5rem;">{{ message }}</textarea><br><br>
          <button type="submit" style="font-size:1rem; padding:0.5rem 1rem;">
            Save
          </button>
        </form>
      </body>
    </html>
    """, message=current)


@app.route("/admin/deepflow_message", methods=["GET", "POST"])
def admin_deepflow_message():
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        new_msg = request.form.get("message", "").strip()
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('deepflow_message', %s)
            ON CONFLICT (key) DO UPDATE SET value = EXCLUDED.value
        """, (new_msg,))
        conn.commit()
        current = new_msg
    else:
        cur.execute("SELECT value FROM app_config WHERE key = 'deepflow_message'")
        row = cur.fetchone()
        current = row[0] if row else ""

    conn.close()

    return render_template_string("""
    <!doctype html>
    <html>
      <head>
        <meta charset="utf-8">
        <title>Admin: DeepFlow Message</title>
      </head>
      <body style="font-family: sans-serif; padding: 2rem;">
        <h1>Update DeepFlow Message</h1>
        <form method="post">
          <textarea name="message" rows="4" cols="60"
            style="font-size:1rem; padding:0.5rem;">{{ message }}</textarea><br><br>
          <button type="submit" style="font-size:1rem; padding:0.5rem 1rem;">
            Save
          </button>
        </form>
      </body>
    </html>
    """, message=current)



@app.route("/admin/home_message", methods=["GET", "POST"])
def admin_home_message():
    # connect to DB
    conn = psycopg2.connect(os.environ["DATABASE_URL"])
    cur = conn.cursor()

    if request.method == "POST":
        # grab the form value and upsert
        new_msg = request.form.get("message", "").strip()
        cur.execute("""
            INSERT INTO app_config (key, value)
            VALUES ('home_message', %s)
            ON CONFLICT (key) DO UPDATE SET value = EXCLUDED.value
        """, (new_msg,))
        conn.commit()
        current = new_msg
    else:
        # fetch the existing message
        cur.execute("SELECT value FROM app_config WHERE key = 'home_message'")
        row = cur.fetchone()
        current = row[0] if row else ""

    conn.close()

    # render a minimal HTML form
    return render_template_string("""
    <!doctype html>
    <html>
      <head>
        <meta charset="utf-8">
        <title>Admin: Home Message</title>
      </head>
      <body style="font-family: sans-serif; padding: 2rem;">
        <h1>Update Home Message</h1>
        <form method="post">
          <textarea name="message" rows="4" cols="60"
            style="font-size:1rem; padding:0.5rem;">{{ message }}</textarea><br><br>
          <button type="submit" style="font-size:1rem; padding:0.5rem 1rem;">
            Save
          </button>
        </form>
      </body>
    </html>
    """, message=current)




# ============ TEXT-ONLY ANSWER ENDPOINT (CHEAP) ============
@app.route("/api/answer", methods=["POST"])
@supabase_auth_required
def answer_question():
    """
    Answer a quiz question using text only (no image). Uses OpenAI API.
    MUCH cheaper than Vision API - use this when OCR extraction works.

    Expects JSON:
    {
        "question": "What is the capital of France?",
        "answers": ["London", "Paris", "Berlin", "Madrid"]
    }

    Returns:
    {
        "correct_answer_index": 2,
        "correct_answer_text": "Paris",
        "confidence": "high",
        "reasoning": "Paris is the capital..."
    }
    """
    try:
        # Check rate limit for free users
        can_proceed, remaining = check_question_limit(request.user_id)
        if not can_proceed:
            debug_log(f"❌ Rate limit exceeded for user {request.user_id}")
            return jsonify({
                "error": "Daily question limit exceeded",
                "message": "Free users are limited to 10 questions per day. Upgrade to Pro for unlimited questions!",
                "limit_exceeded": True
            }), 429

        debug_log(f"✅ Rate limit check passed. Remaining questions: {remaining if remaining >= 0 else 'unlimited'}")

        data = request.get_json()
        if not data:
            return jsonify({"error": "No JSON provided"}), 400

        question = data.get("question", "")
        answers = data.get("answers", [])
        is_multiple_answer = data.get("isMultipleAnswer", False)  # New field for checkbox questions
        model = data.get("model", "gemini-2.5-flash")  # Get model from request, default to 2.5 Flash

        if not question or not answers:
            return jsonify({"error": "Missing question or answers"}), 400

        # Use OpenAI API (gpt-4o-mini for speed/cost, gpt-4o for accuracy)
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        # Use model name directly (extension sends gpt-4o-mini or gpt-4o)
        openai_model = model if model in ["gpt-4o-mini", "gpt-4o"] else "gpt-4o-mini"

        answers_text = "\n".join(f"{i+1}. {a}" for i, a in enumerate(answers))

        # Different prompts for single-answer vs multiple-answer questions
        if is_multiple_answer:
            prompt = f"""Answer this quiz question that allows MULTIPLE correct answers (checkboxes). Return ONLY a JSON object.

Question: {question}

Options:
{answers_text}

This question allows selecting MULTIPLE correct answers. Identify ALL answers that are correct.

Return this exact JSON format:
{{
    "correct_answer_indices": [<array of numbers, e.g., [1, 3, 4]>],
    "correct_answer_texts": ["<first correct answer>", "<second correct answer>", ...],
    "confidence": "high/medium/low",
    "reasoning": "<brief explanation why these answers are correct>"
}}

RULES:
- correct_answer_indices is 1-based (1, 2, 3, etc.)
- Include ALL correct answers in the array
- If only one answer is correct, return array with one element: [2]
- Return ONLY valid JSON, no markdown, no other text."""
        else:
            prompt = f"""Answer this quiz question and return ONLY a JSON object.

Question: {question}

Options:
{answers_text}

Return this exact JSON format:
{{
    "correct_answer_index": <number 1-{len(answers)}>,
    "correct_answer_text": "<answer text>",
    "confidence": "high",
    "reasoning": "<brief explanation>"
}}

Return ONLY valid JSON, no markdown, no other text."""

        response = client.chat.completions.create(
            model=openai_model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
            response_format={"type": "json_object"}
        )

        result = json.loads(response.choices[0].message.content)

        # Validate result based on question type
        if is_multiple_answer:
            if not result or not result.get('correct_answer_indices'):
                return jsonify({"error": "AI processing failed"}), 500
            debug_log(f"TEXT API (OpenAI {openai_model}) MULTIPLE-ANSWER: Answers={result.get('correct_answer_indices')}, Confidence={result.get('confidence')}")
        else:
            if not result or not result.get('correct_answer_index'):
                return jsonify({"error": "AI processing failed"}), 500
            debug_log(f"TEXT API (OpenAI {openai_model}): Answer={result.get('correct_answer_index')}, Confidence={result.get('confidence')}")

        # Store question in database for analytics/caching
        try:
            conn = psycopg2.connect(os.environ["DATABASE_URL"])
            cur = conn.cursor()

            # Prepare answer text based on question type
            if is_multiple_answer:
                answer_indices = result.get('correct_answer_indices', [])
                ai_answer = ', '.join([answers[idx - 1] for idx in answer_indices if idx <= len(answers)])
            else:
                answer_idx = result.get('correct_answer_index', 1)
                ai_answer = answers[answer_idx - 1] if answer_idx <= len(answers) else 'Unknown'

            cur.execute("""
                INSERT INTO questions (user_id, question_text, question_type, answers_json, ai_answer, ai_reasoning)
                VALUES (%s, %s, %s, %s, %s, %s)
            """, (
                request.user_id,
                question,
                'multiple_answer' if is_multiple_answer else 'multiple_choice',
                json.dumps(answers),
                ai_answer,
                result.get('reasoning', '')
            ))
            conn.commit()
            conn.close()
            debug_log(f"✅ Question stored in database for user {request.user_id}")
        except Exception as db_error:
            debug_log(f"⚠️ Failed to store question in database: {db_error}")
            # Don't fail the request if database storage fails

        return jsonify(result), 200

    except json.JSONDecodeError as e:
        debug_log(f"TEXT API JSON parse error: {e}")
        return jsonify({"error": "JSON parse failed"}), 500
    except Exception as e:
        debug_log(f"TEXT API error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ LEGAL TUTOR EXPLANATION ENDPOINT ============
@app.route("/api/tutor/explain", methods=["POST"])
@supabase_auth_required
def tutor_explain():
    """
    Legal AI tutoring endpoint - Returns step-by-step explanations WITHOUT selecting answers.
    Complies with Missouri HB 2271 by acting as a tutor, not a cheating tool.

    Human-in-the-loop: Student must read explanation and click answer themselves.

    Expects JSON:
    {
        "question": "What is the capital of France? A) London B) Paris C) Berlin D) Madrid"
    }

    Returns:
    {
        "explanation": "1. This question asks about...\n2. Consider the geography...\n3. Paris is located in...",
        "concept": "European Geography - Capital Cities",
        "engagement_time": 0  // Server tracks when explanation was generated
    }
    """
    try:
        # Check rate limit
        can_proceed, remaining = check_question_limit(request.user_id)
        if not can_proceed:
            return jsonify({
                "error": "Daily question limit exceeded",
                "message": "Free users are limited to 10 questions per day. Upgrade to Pro for unlimited access!",
                "limit_exceeded": True
            }), 429

        data = request.get_json()
        if not data:
            return jsonify({"error": "No JSON provided"}), 400

        question = data.get("question", "")
        if not question:
            return jsonify({"error": "Missing question"}), 400

        # Use GPT-4o-mini for cost-effectiveness (tutoring requires explanation, not just accuracy)
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        prompt = f"""You are an AI tutor helping a student understand a quiz question. Your job is to explain the CONCEPT and REASONING, not to directly give the answer.

Question: {question}

Provide a step-by-step explanation that helps the student understand:
1. What concept is being tested
2. How to approach this type of question
3. Key facts or formulas they should consider
4. How to eliminate wrong answers (if multiple choice)
5. A hint toward the correct approach (WITHOUT saying "the answer is X")

Format your response as numbered steps. Be educational and helpful, like a tutor.

IMPORTANT: Do NOT say "The answer is..." or "Click option B". Guide them to figure it out themselves.

Your explanation:"""

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful AI tutor focused on teaching concepts, not giving direct answers."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=500
        )

        explanation = response.choices[0].message.content.strip()

        # Identify concept being tested (simple extraction)
        concept_response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You identify educational concepts in one short phrase."},
                {"role": "user", "content": f"What concept is this question testing? Return ONLY the concept name in 3-5 words.\n\nQuestion: {question}"}
            ],
            temperature=0.3,
            max_tokens=20
        )

        concept = concept_response.choices[0].message.content.strip()

        # Store in database for analytics (as tutoring session, not answer)
        try:
            conn = psycopg2.connect(os.environ.get("DATABASE_URL"))
            cur = conn.cursor()
            cur.execute("""
                INSERT INTO questions (user_id, question_text, question_type, answer_text, ai_reasoning)
                VALUES (%s, %s, %s, %s, %s)
            """, (
                request.user_id,
                question,
                'tutoring_session',  # Different type to distinguish from automated answers
                '',  # No answer provided in tutoring mode
                explanation
            ))
            conn.commit()
            conn.close()
            debug_log(f"✅ Tutoring session stored for user {request.user_id}")
        except Exception as db_error:
            debug_log(f"⚠️ Failed to store tutoring session: {db_error}")

        result = {
            "explanation": explanation,
            "concept": concept,
            "engagement_time": 0,  # Client will track actual read time
            "legal_mode": True  # Flag indicating this is legal tutoring, not automation
        }

        return jsonify(result), 200

    except Exception as e:
        debug_log(f"TUTOR API error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ ESSAY ANSWER ENDPOINT ============
@app.route("/api/essay", methods=["POST"])
@supabase_auth_required
def essay_answer():
    """
    Generate an essay answer for a short-answer or essay question.

    Expects JSON:
    {
        "question": "Explain the process of photosynthesis."
    }

    Returns:
    {
        "essay_answer": "Photosynthesis is the process by which..."
    }
    """
    try:
        # Check rate limit for free users
        can_proceed, remaining = check_question_limit(request.user_id)
        if not can_proceed:
            debug_log(f"❌ Rate limit exceeded for user {request.user_id}")
            return jsonify({
                "error": "Daily question limit exceeded",
                "message": "Free users are limited to 10 questions per day. Upgrade to Pro for unlimited questions!",
                "limit_exceeded": True
            }), 429

        debug_log(f"✅ Rate limit check passed. Remaining questions: {remaining if remaining >= 0 else 'unlimited'}")

        data = request.get_json()
        if not data:
            return jsonify({"error": "No JSON provided"}), 400

        question = data.get("question", "")
        model = data.get("model", "gemini-2.5-flash")  # Get model from request, default to 2.5 Flash

        if not question:
            return jsonify({"error": "Missing question"}), 400

        prompt = f"""You are answering a quiz question that requires a text response (essay, short answer, or fill-in-the-blank).

Question: {question}

Provide an accurate answer. Follow these rules:
- For fill-in-the-blank questions (asking for a single word/phrase/number): Give ONLY the answer, no extra text
- For short answer questions: Provide 2-5 complete sentences
- For essay questions: Provide 1-3 well-structured paragraphs
- Be factually accurate and direct
- Use a natural, student-like tone
- If the question asks for a definition, name, date, or specific term, provide just that without elaboration

Return ONLY the answer text, nothing else."""

        # Use OpenAI API for essay generation
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        # Use model name directly (extension sends gpt-4o-mini or gpt-4o)
        openai_model = model if model in ["gpt-4o-mini", "gpt-4o"] else "gpt-4o-mini"

        response = client.chat.completions.create(
            model=openai_model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=500
        )

        essay_answer = response.choices[0].message.content.strip()

        if not essay_answer:
            return jsonify({"error": "AI processing failed"}), 500

        debug_log(f"ESSAY API (OpenAI {openai_model}): Generated {len(essay_answer)} characters")

        # Store question in database for analytics/caching
        try:
            conn = psycopg2.connect(os.environ["DATABASE_URL"])
            cur = conn.cursor()

            # Determine question type based on answer length
            if len(essay_answer) < 50:
                question_type = 'fill_in_blank'
            elif len(essay_answer) < 200:
                question_type = 'short_answer'
            else:
                question_type = 'essay'

            cur.execute("""
                INSERT INTO questions (user_id, question_text, question_type, answers_json, ai_answer, ai_reasoning)
                VALUES (%s, %s, %s, %s, %s, %s)
            """, (
                request.user_id,
                question,
                question_type,
                None,  # No multiple choice options for essay questions
                essay_answer,
                None  # No reasoning for essay questions
            ))
            conn.commit()
            conn.close()
            debug_log(f"✅ Essay question stored in database for user {request.user_id}")
        except Exception as db_error:
            debug_log(f"⚠️ Failed to store essay question in database: {db_error}")
            # Don't fail the request if database storage fails

        return jsonify({"essay_answer": essay_answer}), 200

    except Exception as e:
        debug_log(f"ESSAY API error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ VISION API ENDPOINT ============
@app.route("/api/send-download-link", methods=["POST"])
def send_download_link():
    """Send download link email to mobile users"""
    try:
        data = request.get_json()
        email = data.get('email')

        if not email:
            return jsonify({'error': 'Email is required'}), 400

        html_content = """
            <div style="font-family: Arial, sans-serif; max-width: 600px; margin: 0 auto; padding: 20px;">
                <h2 style="color: #667eea;">Open StudyFlow on Your Computer</h2>
                <p>You visited StudyFlow on your mobile device, but our Chrome extension only works on desktop computers.</p>
                <p style="margin: 30px 0;">
                    <a href="https://studyflowsuite.com/"
                       style="display: inline-block; background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                              color: white; padding: 15px 30px; text-decoration: none; border-radius: 8px;
                              font-weight: bold;">
                        Open on Desktop
                    </a>
                </p>
                <h3 style="color: #1e293b;">Requirements:</h3>
                <ul style="line-height: 1.8;">
                    <li>✅ Desktop or laptop computer</li>
                    <li>✅ Chrome browser</li>
                    <li>✅ Follow the installation guide on the website</li>
                </ul>
                <p style="margin-top: 30px; color: #64748b; font-size: 14px;">
                    Questions? Visit our <a href="https://studyflowsuite.com/docs.html" style="color: #667eea;">FAQ page</a>.
                </p>
                <hr style="margin: 30px 0; border: none; border-top: 1px solid #e2e8f0;">
                <p style="color: #94a3b8; font-size: 12px; text-align: center;">
                    StudyFlow Suite - AI-Powered Quiz Automation<br>
                    <a href="https://studyflowsuite.com/" style="color: #667eea;">studyflowsuite.com</a>
                </p>
            </div>
        """

        response = send_brevo_email(
            to_email=email,
            to_name=None,
            subject="StudyFlow Download Link - Open on Desktop",
            html_content=html_content
        )

        if response.status_code in (200, 201):
            app.logger.info(f"Download link sent to {email}")
            return jsonify({'success': True, 'message': 'Email sent successfully'}), 200
        else:
            app.logger.error(f"Error sending download link: {response.text}")
            return jsonify({'error': 'Failed to send email'}), 500

    except Exception as e:
        app.logger.error(f"Error sending download link: {e}")
        return jsonify({'error': 'Failed to send email'}), 500


@app.route("/api/vision", methods=["POST"])
def vision_analyze():
    """
    Analyze a screenshot using GPT-4o Vision.
    Expects JSON with 'image' field containing base64-encoded PNG.
    Returns question, answers, correct answer, and button text.
    """
    try:
        data = request.get_json()
        if not data or 'image' not in data:
            return jsonify({"error": "No image provided"}), 400

        image_base64 = data['image']

        # Use OpenAI client (v1.0+ style)
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        prompt = """Analyze this screenshot of a quiz/test interface.

Your task:
1. Find the quiz QUESTION being asked
2. Find ALL answer options (usually 2-4 options)
3. Determine which answer is CORRECT
4. Identify the Submit/Next button
5. PROVIDE CLICK COORDINATES for the correct answer and the button

Return a JSON object with this EXACT structure:
{
    "question": "the full question text",
    "answers": [
        {"index": 1, "text": "first answer option text exactly as shown"},
        {"index": 2, "text": "second answer option text exactly as shown"},
        {"index": 3, "text": "third answer option text exactly as shown"},
        {"index": 4, "text": "fourth answer option text exactly as shown"}
    ],
    "correct_answer_index": 2,
    "correct_answer_text": "the exact text of the correct answer",
    "answer_click_x": 500,
    "answer_click_y": 300,
    "button_text": "Submit or Next or Continue - exact text shown",
    "button_click_x": 800,
    "button_click_y": 450,
    "confidence": "high/medium/low",
    "reasoning": "Brief explanation why this answer is correct"
}

CRITICAL RULES:
- correct_answer_index is 1-based (1, 2, 3, or 4)
- answer_click_x and answer_click_y should be the CENTER of the correct answer option (where to click)
- button_click_x and button_click_y should be the CENTER of the Submit/Next button
- Coordinates are in pixels from top-left (0,0) of the image
- Look for radio buttons or checkboxes next to answers - provide coordinates there
- If no quiz visible, set question to null
- Return ONLY valid JSON, no other text"""

        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/png;base64,{image_base64}",
                                "detail": "high"
                            }
                        }
                    ]
                }
            ],
            max_tokens=1000,
            temperature=0.1
        )

        response_text = response.choices[0].message.content.strip()
        debug_log(f"📨 Vision API response received ({len(response_text)} chars)")

        # Parse JSON - handle markdown code blocks
        if "```" in response_text:
            lines = response_text.split("\n")
            json_lines = []
            in_json = False
            for line in lines:
                if line.startswith("```json") or line.startswith("```"):
                    in_json = not in_json
                    continue
                if in_json:
                    json_lines.append(line)
            response_text = "\n".join(json_lines)

        result = json.loads(response_text)
        return jsonify(result), 200

    except json.JSONDecodeError as e:
        debug_log(f"❌ Vision JSON parse error: {e}")
        return jsonify({"error": "JSON parse failed", "raw": response_text}), 500
    except Exception as e:
        debug_log(f"❌ Vision API error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ NOTEFLOW ENDPOINTS ============

@app.route("/api/notes/upload", methods=["POST"])
@supabase_auth_required
def upload_note():
    """
    Upload a note file (PDF, image, Word doc, TXT) to Supabase.

    NEW: Uses Supabase Storage + pgvector + background processing

    Expects: multipart/form-data with:
        - file: The document file
        - university: (optional) University name
        - course_code: (optional) Course code (e.g., "BIO 101")
        - professor: (optional) Professor name
        - semester: (optional) Semester (e.g., "Fall 2024")

    Returns:
    {
        "success": true,
        "note_id": "uuid",
        "filename": "Biology_Chapter_3.pdf",
        "pages": 5,
        "processing": true
    }
    """
    try:
        from StudyFlow.backend.supabase_client import (
            check_page_limit, upload_file_to_storage, create_note_record, increment_page_count, log_upload, get_user_profile
        )
        from StudyFlow.backend.tasks import process_note_async
        import hashlib

        # Get user profile to extract username
        user_profile = get_user_profile(request.user_id)
        username = user_profile.get('username') if user_profile else None

        # Check if file was uploaded
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400

        # Get course metadata from form
        course_metadata = {
            "university": request.form.get('university'),
            "course_code": request.form.get('course_code'),
            "professor": request.form.get('professor'),
            "semester": request.form.get('semester')
        }

        # Get Nexus sharing preference (Missouri SB 1324 compliance)
        share_with_nexus = request.form.get('share_with_nexus', 'false').lower() == 'true'

        # Get user's IP address (for security logging)
        ip_address = request.headers.get('X-Forwarded-For', request.remote_addr)
        if ip_address and ',' in ip_address:
            # X-Forwarded-For can have multiple IPs, take the first
            ip_address = ip_address.split(',')[0].strip()

        # Get file info
        original_filename = file.filename
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)  # Reset file pointer

        # Determine file type
        file_ext = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else ''
        allowed_extensions = ['pdf', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp', 'txt', 'doc', 'docx']
        if file_ext not in allowed_extensions:
            return jsonify({"error": "Unsupported file type. Allowed: PDF, images, TXT, DOC/DOCX"}), 400

        # Read file content
        file_content = file.read()

        # Calculate SHA-256 hash for tamper verification (Missouri SB 1324 compliance)
        file_hash = hashlib.sha256(file_content).hexdigest()
        debug_log(f"[*] File hash (SHA-256): {file_hash}")

        # Log upload to database (7-year retention for Missouri compliance)
        log_upload(
            user_id=request.user_id,
            file_name=original_filename,
            file_hash=file_hash,
            file_size=file_size,
            shared_with_nexus=share_with_nexus,
            ip_address=ip_address,
            university=course_metadata.get('university'),
            course_code=course_metadata.get('course_code')
        )

        # Extract text based on file type
        import google.generativeai as genai
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        model = genai.GenerativeModel('gemini-3.1-flash-lite-preview')

        if file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
            # Image file
            import PIL.Image
            import io
            image = PIL.Image.open(io.BytesIO(file_content))

            prompt = """Extract ALL text from this image. This is a student's class notes.

            Return ONLY the extracted text, preserving:
            - All text content
            - Line breaks and structure
            - Headings and bullet points

            Do not add any commentary. Just return the extracted text."""

            response = model.generate_content([prompt, image])
            ocr_text = response.text.strip()
            page_count = 1

            try:
                from StudyFlow.backend.cost_tracker import track_ai_call
                track_ai_call("gemini", "flash-lite", "ocr_image")
            except:
                pass

        elif file_ext == 'pdf':
            # PDF file
            import PyPDF2
            import io

            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            page_count = len(pdf_reader.pages)

            # Try extracting text first (faster if PDF has text layer)
            ocr_text = ""
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    ocr_text += page_text + "\n\n"

            # If no text extracted, skip for now
            if not ocr_text.strip():
                ocr_text = f"[PDF document: {page_count} pages - requires OCR]"

        elif file_ext == 'txt':
            # Plain text file
            ocr_text = file_content.decode('utf-8', errors='ignore')
            page_count = 1

        elif file_ext in ['doc', 'docx']:
            # Word documents -- extract text with python-docx
            try:
                import io
                from docx import Document
                doc = Document(io.BytesIO(file_content))
                ocr_text = '\n'.join(p.text for p in doc.paragraphs)
                if not ocr_text.strip():
                    ocr_text = f"[{file_ext.upper()} document - no text extracted]"
            except Exception:
                ocr_text = f"[{file_ext.upper()} document]"
            page_count = 1

        else:
            ocr_text = f"[{file_ext.upper()} document]"
            page_count = 1

        # Check page limit BEFORE uploading
        allowed, message = check_page_limit(request.user_id, page_count)
        if not allowed:
            return jsonify({"error": message, "limit_exceeded": True}), 403

        # Convert all files to PDF at upload time
        # Text extraction already happened above on the original format
        from StudyFlow.backend.pdf_flatten import convert_to_pdf
        pdf_data, pdf_filename = convert_to_pdf(file_content, original_filename)
        if pdf_data is None:
            return jsonify({"error": f"Failed to convert {file_ext.upper()} to PDF"}), 500

        # Use converted PDF for storage
        file_content = pdf_data
        original_filename = pdf_filename
        file_ext = 'pdf'
        file_size = len(file_content)
        debug_log(f"[*] Converted to PDF: {pdf_filename} ({file_size} bytes)")

        # Upload PDF to Supabase Storage
        import uuid
        unique_filename = f"{request.user_id}/{uuid.uuid4()}_{original_filename}"
        content_type = 'application/pdf'

        file_url = upload_file_to_storage(file_content, unique_filename, content_type)
        if not file_url:
            return jsonify({"error": "Failed to upload file to storage"}), 500

        # Create note record in Supabase
        debug_log(f"[*] Creating note for user_id: {request.user_id}")
        note = create_note_record(
            user_id=request.user_id,
            filename=original_filename,
            file_type=file_ext,
            file_size=file_size,
            file_path=unique_filename,
            page_count=page_count,
            course_metadata=course_metadata,
            username=username
        )

        if not note:
            return jsonify({"error": "Failed to create note record"}), 500

        note_id = note['id']
        debug_log(f"[*] Created note {note_id} for user {request.user_id}")

        # Increment user's page count
        increment_page_count(request.user_id, page_count)

        # Trigger background processing (chunking, embedding, anonymization)
        process_note_async.delay(note_id, request.user_id, ocr_text, course_metadata, file_hash, username)

        debug_log(f"Note uploaded: {original_filename} ({file_size} bytes, {page_count} pages)")
        debug_log(f"Background processing started for note {note_id}")

        # GOOD STANDING: Verify upload with AI (Missouri HB 2271 compliance)
        verification_result = None
        try:
            if len(ocr_text.strip()) >= 50:
                print(f"[GOOD STANDING] Verifying upload for note {note_id}...", flush=True)

                # Quick AI verification with Gemini
                import google.generativeai as genai
                genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
                model = genai.GenerativeModel('gemini-3.1-flash-lite-preview')

                prompt = f"""Analyze this document to determine if it's legitimate study material.

DOCUMENT TEXT (first 2000 chars):
{ocr_text[:2000]}

Respond with ONLY a JSON object:
{{
    "is_study_material": true/false,
    "confidence": 0.0-1.0,
    "reason": "brief explanation",
    "category": "notes" | "textbook" | "slides" | "homework" | "spam" | "blank" | "other"
}}

Study material includes: class notes, textbook chapters, lecture slides, study guides, homework, assignments, etc.
NOT study material: grocery lists, personal emails, blank pages, random text, advertisements."""

                response = model.generate_content(prompt)
                ai_text = response.text.strip()

                try:
                    from StudyFlow.backend.cost_tracker import track_ai_call
                    track_ai_call("gemini", "flash-lite", "verification")
                except:
                    pass

                # Parse JSON response
                import json
                import re
                if ai_text.startswith("```"):
                    ai_text = re.sub(r'^```json?\s*', '', ai_text)
                    ai_text = re.sub(r'\s*```$', '', ai_text)

                result = json.loads(ai_text)
                is_verified = result.get('is_study_material', False)
                confidence = result.get('confidence', 0.5)
                category = result.get('category', 'other')

                # Determine rejection reason
                rejection_reason = None
                if not is_verified:
                    if category == 'spam':
                        rejection_reason = 'spam'
                    elif category == 'blank':
                        rejection_reason = 'blank'
                    else:
                        rejection_reason = 'invalid_format'

                # Log verification
                from StudyFlow.backend.supabase_client import supabase
                supabase.table("upload_verifications").insert({
                    "user_id": request.user_id,
                    "note_id": note_id,
                    "verification_status": "verified" if is_verified else "rejected",
                    "rejection_reason": rejection_reason,
                    "ai_confidence": confidence
                }).execute()

                # If verified, update user's Good Standing status
                if is_verified:
                    from datetime import datetime
                    supabase.table("user_profiles").update({
                        "good_standing": True,
                        "last_verified_upload_date": datetime.utcnow().isoformat()
                    }).eq("id", request.user_id).execute()

                    print(f"[GOOD STANDING] User {request.user_id} verified upload, status updated", flush=True)
                    verification_result = {"verified": True, "confidence": confidence}
                else:
                    print(f"[GOOD STANDING] Upload rejected as {category}: {result.get('reason')}", flush=True)
                    verification_result = {"verified": False, "reason": rejection_reason}
            else:
                print(f"[GOOD STANDING] Skipping verification (text too short)", flush=True)
        except Exception as verify_error:
            print(f"[GOOD STANDING] Verification failed: {verify_error}", flush=True)
            # Don't fail the upload if verification fails

        # Notify other users in the same course about the new note
        try:
            uni = course_metadata.get('university')
            ccode = course_metadata.get('course_code')
            if uni and ccode:
                # Find all users who have uploaded to the same university+course_code
                course_users = supabase.table("notes") \
                    .select("user_id") \
                    .eq("university", uni) \
                    .eq("course_code", ccode) \
                    .neq("user_id", request.user_id) \
                    .execute()
                notified_ids = set()
                for row in (course_users.data or []):
                    uid = row.get("user_id")
                    if uid and uid not in notified_ids:
                        notified_ids.add(uid)
                        create_notification(
                            user_id=uid,
                            notif_type="course_new_note",
                            title="New Note in Course",
                            message=f"New note in {ccode}: {original_filename} by @{username or 'someone'}",
                            note_id=note_id,
                            actor_username=username,
                        )
                if notified_ids:
                    debug_log(f"[+] Notified {len(notified_ids)} users about new note in {ccode}")
        except Exception as notif_err:
            debug_log(f"[-] Upload notification failed (non-blocking): {notif_err}")

        # Invalidate notes cache
        if redis_cache:
            try: redis_cache.delete(f"user_notes:{request.user_id}")
            except: pass

        return jsonify({
            "success": True,
            "note_id": note_id,
            "filename": original_filename,
            "pages": page_count,
            "processing": True,
            "verification": verification_result,
            "message": "Note uploaded! Processing in background..."
        }), 200

    except Exception as e:
        debug_log(f"Upload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/list", methods=["GET", "OPTIONS"])
@supabase_auth_required
def list_notes():
    """
    Get list of all notes for current user.

    Returns:
    [
        {
            "id": "uuid",
            "filename": "Biology_Chapter_3.pdf",
            "pages": 5,
            "uploaded_at": "2024-01-15T10:30:00",
            "processed": true,
            "university": "MSU",
            "course_code": "BIO 101"
        },
        ...
    ]
    """
    # Handle OPTIONS preflight
    if request.method == 'OPTIONS':
        return '', 200

    try:
        cache_key = f"user_notes:{request.user_id}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        from StudyFlow.backend.supabase_client import get_user_notes
        notes = get_user_notes(request.user_id)

        # Format response
        formatted_notes = []
        for note in notes:
            formatted_notes.append({
                "id": note['id'],
                "filename": note['original_filename'],
                "file_type": note['file_type'],
                "file_size": note['file_size'],
                "pages": note['page_count'],
                "uploaded_at": note['uploaded_at'],
                "processed": note['processed'],
                "is_public": note['is_public'],
                "university": note.get('university'),
                "course_code": note.get('course_code'),
                "professor": note.get('professor'),
                "folder_id": note.get('folder_id')  # Include folder assignment
            })

        if redis_cache:
            try: redis_cache.setex(cache_key, 300, json.dumps(formatted_notes))
            except: pass
        return jsonify(formatted_notes), 200

    except Exception as e:
        debug_log(f"List notes error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/usage-stats", methods=["GET"])
@supabase_auth_required
def get_notes_usage_stats():
    """
    Get usage statistics for user's notes showing how many times they've been used to help answer questions.

    Returns:
    {
        "totalUsage": 47,
        "weeklyUsage": 12,
        "activeNotes": 5,
        "noteUsage": {
            "note_id_1": 23,
            "note_id_2": 15,
            ...
        },
        "topNotes": [
            {"filename": "Biology_Chapter_3.pdf", "count": 23},
            {"filename": "Chemistry_Notes.pdf", "count": 15},
            ...
        ]
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from datetime import datetime, timedelta

        user_id = request.user_id

        # Get user's notes
        notes_response = supabase.table("notes").select("id, original_filename").eq("user_id", user_id).execute()
        user_notes = {note['id']: note['original_filename'] for note in notes_response.data}

        # Get usage counts from conversation sources
        # This queries how many times each note was used in conversations
        note_usage = {}
        total_usage = 0
        weekly_usage = 0
        week_ago = (datetime.now() - timedelta(days=7)).isoformat()

        # Query all assistant messages with sources
        all_messages_response = supabase.table("conversation_messages").select("sources, created_at").eq("role", "assistant").not_.is_("sources", "null").execute()

        for message in all_messages_response.data:
            sources = message.get('sources', [])
            created_at = message.get('created_at', '')

            # Check each source in the message
            for source in sources:
                source_note_id = source.get('note_id')
                if source_note_id in user_notes:
                    # Count this usage
                    note_usage[source_note_id] = note_usage.get(source_note_id, 0) + 1
                    total_usage += 1

                    # Count weekly usage
                    if created_at >= week_ago:
                        weekly_usage += 1

        # Get top notes
        top_notes = []
        for note_id, count in sorted(note_usage.items(), key=lambda x: x[1], reverse=True)[:5]:
            top_notes.append({
                "filename": user_notes[note_id],
                "count": count
            })

        active_notes = len([c for c in note_usage.values() if c > 0])

        return jsonify({
            "totalUsage": total_usage,
            "weeklyUsage": weekly_usage,
            "activeNotes": active_notes,
            "noteUsage": note_usage,
            "topNotes": top_notes
        }), 200

    except Exception as e:
        debug_log(f"❌ Usage stats error: {e}\n{traceback.format_exc()}")
        # Return empty stats on error
        return jsonify({
            "totalUsage": 0,
            "weeklyUsage": 0,
            "activeNotes": 0,
            "noteUsage": {},
            "topNotes": []
        }), 200


@app.route("/api/notes/<note_id>", methods=["DELETE"])
@supabase_auth_required
def delete_note_endpoint(note_id):
    """
    Delete a note by ID (only if it belongs to current user).
    """
    try:
        from StudyFlow.backend.supabase_client import delete_note

        success = delete_note(note_id, request.user_id)

        if success:
            return jsonify({"success": True}), 200
        else:
            return jsonify({"error": "Note not found or unauthorized"}), 404

    except Exception as e:
        debug_log(f"❌ Delete note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/rename", methods=["PATCH"])
@supabase_auth_required
def rename_note_endpoint(note_id):
    """
    Rename a note by ID (only if it belongs to current user).
    Request body: { "original_filename": "new_name.pdf" }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        new_filename = data.get("original_filename", "").strip()

        if not new_filename:
            return jsonify({"error": "Filename is required"}), 400

        # Verify ownership and update filename
        response = supabase.table("notes").update({
            "original_filename": new_filename
        }).eq("id", note_id).eq("user_id", request.user_id).execute()

        if response.data:
            debug_log(f"Renamed note {note_id} to {new_filename}")
            return jsonify({"success": True, "original_filename": new_filename}), 200
        else:
            return jsonify({"error": "Note not found or unauthorized"}), 404

    except Exception as e:
        debug_log(f"❌ Rename note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/visibility", methods=["POST"])
@supabase_auth_required
def toggle_note_visibility(note_id):
    """
    Toggle note public/private visibility (only if it belongs to current user).
    Request body: { "is_public": true/false }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        is_public = data.get("is_public")

        if is_public is None:
            return jsonify({"error": "is_public field is required"}), 400

        # Verify ownership and update visibility
        response = supabase.table("notes").update({
            "is_public": is_public
        }).eq("id", note_id).eq("user_id", request.user_id).execute()

        if response.data:
            debug_log(f"Updated note {note_id} visibility to {'public' if is_public else 'private'}")
            return jsonify({"success": True, "is_public": is_public}), 200
        else:
            return jsonify({"error": "Note not found or unauthorized"}), 404

    except Exception as e:
        debug_log(f"❌ Toggle visibility error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


## ====== ASSISTANT ======

@app.route("/api/assistant/ask", methods=["POST"])
@supabase_auth_required
def assistant_ask():
    """Lightweight AI assistant endpoint. Sends a message with page context."""
    try:
        data = request.get_json()
        message = data.get("message", "")
        page_context = data.get("page", "unknown")
        action = data.get("action", "ask")  # ask, summarize, quiz, tips

        if not message:
            return jsonify({"error": "Missing message"}), 400

        # Build a context-aware system instruction
        system_context = f"You are Flo, a friendly study assistant on the StudyFlow platform. The user is currently on the {page_context} page. "

        if action == "summarize":
            system_context += "Provide a brief, helpful summary. Keep it under 3 sentences."
        elif action == "quiz":
            system_context += "Generate a quick 3-question quiz based on the topic. Format each as a question with 4 multiple choice options (A-D) and the correct answer."
        elif action == "tips":
            system_context += "Give 2-3 short, practical tips relevant to what the user is doing. Be specific to StudyFlow features."
        else:
            system_context += "Be concise and helpful. Keep responses under 3 sentences unless more detail is needed. You can suggest StudyFlow features that might help."

        # Use the existing chat infrastructure
        from StudyFlow.backend.supabase_client import search_notes_vector, supabase
        from StudyFlow.backend.embedding_client import generate_embedding
        import time

        start_time = time.time()

        # Generate embedding for search
        query_embedding = generate_embedding(message)
        search_results = []
        if query_embedding:
            search_results = search_notes_vector(
                query_embedding=query_embedding,
                user_id=request.user_id,
                limit=3,
                scope="all"
            ) or []

        # Build context from search results
        context_text = ""
        sources = []
        for r in search_results[:3]:
            if r.get("content"):
                context_text += f"\n---\nFrom {r.get('original_filename', 'a note')}:\n{r['content'][:500]}\n"
                sources.append({
                    "filename": r.get("original_filename", "Unknown"),
                    "similarity": round(r.get("similarity", 0), 2)
                })

        # Call AI
        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        messages = [{"role": "system", "content": system_context}]
        if context_text:
            messages.append({"role": "system", "content": f"Relevant notes from the user's library:\n{context_text}"})
        messages.append({"role": "user", "content": message})

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=messages,
            max_tokens=300,
            temperature=0.7
        )

        ai_response = response.choices[0].message.content
        elapsed = int((time.time() - start_time) * 1000)

        # Generate follow-up button suggestions
        followups = []
        if action == "ask" and sources:
            followups.append({"label": "Tell me more", "action": "ask", "message": f"Explain more about: {message}"})
            followups.append({"label": "Quiz me on this", "action": "quiz", "message": message})
        elif action == "quiz":
            followups.append({"label": "Another quiz", "action": "quiz", "message": message})
        elif action == "tips":
            followups.append({"label": "More tips", "action": "tips", "message": f"More tips for {page_context}"})

        return jsonify({
            "response": ai_response,
            "sources": sources,
            "followups": followups,
            "response_time_ms": elapsed
        }), 200

    except Exception as e:
        debug_log(f"Assistant ask error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": "Assistant is temporarily unavailable. Try again in a moment."}), 500


@app.route("/api/assistant/digest", methods=["GET"])
@supabase_auth_required
def assistant_digest():
    """Get a daily digest for the assistant: events today, notification count, note stats."""
    try:
        user_id = request.user_id

        # Today's events
        today_events = []
        try:
            from datetime import datetime
            today = datetime.utcnow().strftime('%Y-%m-%d')
            events_res = supabase.table("calendar_events").select("title, event_date, event_time").eq("user_id", user_id).gte("event_date", today).lte("event_date", today).execute()
            today_events = events_res.data or []
        except:
            pass

        # Unread notifications
        unread_count = 0
        try:
            notif_res = supabase.table("notifications").select("id", count="exact").eq("user_id", user_id).eq("read", False).execute()
            unread_count = notif_res.count or 0
        except:
            pass

        # Note stats
        total_notes = 0
        shared_notes = 0
        try:
            notes_res = supabase.table("notes").select("id", count="exact").eq("user_id", user_id).execute()
            total_notes = notes_res.count or 0
            shared_res = supabase.table("shared_notes").select("id", count="exact").eq("user_id", user_id).execute()
            shared_notes = shared_res.count or 0
        except:
            pass

        # Build digest message
        parts = []
        hour = datetime.utcnow().hour
        if hour < 12:
            greeting = "Good morning"
        elif hour < 17:
            greeting = "Good afternoon"
        else:
            greeting = "Good evening"

        if today_events:
            parts.append(f"You have {len(today_events)} event{'s' if len(today_events) != 1 else ''} today")
        if unread_count > 0:
            parts.append(f"{unread_count} unread notification{'s' if unread_count != 1 else ''}")
        if total_notes > 0:
            parts.append(f"{total_notes} notes in your library")

        summary = f"{greeting}! " + (". ".join(parts) + "." if parts else "Everything looks good!")

        return jsonify({
            "greeting": greeting,
            "summary": summary,
            "events_today": today_events,
            "unread_notifications": unread_count,
            "total_notes": total_notes,
            "shared_notes": shared_notes
        }), 200

    except Exception as e:
        debug_log(f"Assistant digest error: {e}\n{traceback.format_exc()}")
        return jsonify({"summary": "Hi! Ready to help.", "events_today": [], "unread_notifications": 0, "total_notes": 0, "shared_notes": 0}), 200


@app.route("/api/assistant/summarize/<note_id>", methods=["GET"])
@supabase_auth_required
def assistant_summarize(note_id):
    """Summarize a note's content using AI."""
    try:
        import time
        start_time = time.time()

        # Get note chunks
        chunks_res = supabase.table("note_chunks").select("chunk_text, content_summary").eq("note_id", note_id).order("chunk_index").execute()

        if not chunks_res.data:
            return jsonify({"error": "No content found for this note"}), 404

        # Build full text from chunks
        full_text = ""
        for chunk in chunks_res.data:
            text = chunk.get("content_summary") or chunk.get("chunk_text", "")
            if text:
                full_text += text + "\n\n"

        if not full_text.strip():
            return jsonify({"error": "Note has no readable content"}), 404

        # Truncate to ~6000 chars for the AI call
        content = full_text[:6000]

        # Get note title
        note_res = supabase.table("notes").select("original_filename").eq("id", note_id).execute()
        title = "this note"
        if note_res.data:
            title = note_res.data[0].get("original_filename") or "this note"

        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a study assistant. Summarize the following note concisely. Highlight the key concepts, main arguments, and important details. Use bullet points. Keep it under 200 words."},
                {"role": "user", "content": f"Note title: {title}\n\nContent:\n{content}"}
            ],
            max_tokens=400,
            temperature=0.5
        )

        summary = response.choices[0].message.content
        elapsed = int((time.time() - start_time) * 1000)

        return jsonify({
            "summary": summary,
            "note_title": title,
            "chunks_used": len(chunks_res.data),
            "response_time_ms": elapsed
        }), 200

    except Exception as e:
        debug_log(f"Assistant summarize error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": "Could not summarize this note right now"}), 500


@app.route("/api/assistant/quiz/<note_id>", methods=["GET"])
@supabase_auth_required
def assistant_quiz_note(note_id):
    """Generate a quick quiz from a specific note's content."""
    try:
        import time
        start_time = time.time()

        count = request.args.get("count", 5, type=int)
        count = min(max(count, 3), 10)

        # Get note chunks
        chunks_res = supabase.table("note_chunks").select("chunk_text, content_summary").eq("note_id", note_id).order("chunk_index").execute()

        if not chunks_res.data:
            return jsonify({"error": "No content found for this note"}), 404

        full_text = ""
        for chunk in chunks_res.data:
            text = chunk.get("content_summary") or chunk.get("chunk_text", "")
            if text:
                full_text += text + "\n\n"

        if not full_text.strip():
            return jsonify({"error": "Note has no readable content"}), 404

        content = full_text[:6000]

        # Get note title
        note_res = supabase.table("notes").select("original_filename").eq("id", note_id).execute()
        title = "this note"
        if note_res.data:
            title = note_res.data[0].get("original_filename") or "this note"

        from openai import OpenAI
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a quiz generator. Create quiz questions that test understanding of the provided study material. Return ONLY valid JSON, no other text."},
                {"role": "user", "content": f"""Based on this note titled "{title}", generate exactly {count} multiple choice questions.

Note content:
{content}

Return a JSON array:
[{{"question": "What is...?", "options": ["A) ...", "B) ...", "C) ...", "D) ..."], "correct": 0, "explanation": "Brief explanation"}}]

The "correct" field is the zero-based index of the correct option.
Return ONLY the JSON array, no markdown, no explanation."""}
            ],
            max_tokens=800,
            temperature=0.7
        )

        raw = response.choices[0].message.content.strip()
        # Clean markdown if present
        if raw.startswith("```"):
            raw = raw.split("\n", 1)[1] if "\n" in raw else raw[3:]
            if raw.endswith("```"):
                raw = raw[:-3]
            raw = raw.strip()

        questions = json.loads(raw)
        elapsed = int((time.time() - start_time) * 1000)

        return jsonify({
            "questions": questions,
            "note_title": title,
            "count": len(questions),
            "response_time_ms": elapsed
        }), 200

    except json.JSONDecodeError:
        return jsonify({"error": "Failed to generate quiz. Try again."}), 500
    except Exception as e:
        debug_log(f"Assistant quiz error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": "Could not generate quiz right now"}), 500


## ====== POMODORO STATS ======

@app.route("/api/pomodoro/complete", methods=["POST"])
@supabase_auth_required
def save_pomodoro():
    """Save a completed pomodoro round."""
    try:
        data = request.get_json() or {}
        work_minutes = data.get("work_minutes", 25)

        supabase.table("pomodoro_completions").insert({
            "user_id": request.user_id,
            "work_minutes": work_minutes
        }).execute()

        # Get total count for response
        res = supabase.table("pomodoro_completions").select("id", count="exact").eq("user_id", request.user_id).execute()
        total = res.count or 0

        return jsonify({"success": True, "total_pomodoros": total}), 200
    except Exception as e:
        debug_log(f"Save pomodoro error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/pomodoro/stats", methods=["GET"])
@supabase_auth_required
def get_pomodoro_stats():
    """Get pomodoro stats for the user."""
    try:
        res = supabase.table("pomodoro_completions").select("work_minutes, completed_at").eq("user_id", request.user_id).execute()
        completions = res.data or []

        total = len(completions)
        total_minutes = sum(c.get('work_minutes', 25) for c in completions)
        today = datetime.utcnow().strftime('%Y-%m-%d')
        today_count = sum(1 for c in completions if c.get('completed_at', '')[:10] == today)

        return jsonify({
            "total_pomodoros": total,
            "total_minutes": total_minutes,
            "today_pomodoros": today_count
        }), 200
    except Exception as e:
        debug_log(f"Pomodoro stats error: {e}\n{traceback.format_exc()}")
        return jsonify({"total_pomodoros": 0, "total_minutes": 0, "today_pomodoros": 0}), 200


## ====== ACHIEVEMENTS ======

@app.route("/api/achievements", methods=["GET"])
@supabase_auth_required
def get_achievements():
    """Compute achievement progress from existing database tables."""
    try:
        user_id = request.user_id

        # Gather stats from various tables
        stats = {}

        # Note counts
        try:
            res = supabase.table("notes").select("id, is_public, uploaded_at", count="exact").eq("user_id", user_id).execute()
            stats['notes_uploaded'] = res.count or 0
            stats['notes_public'] = len([n for n in (res.data or []) if n.get('is_public')])
        except:
            stats['notes_uploaded'] = 0
            stats['notes_public'] = 0

        # Folders
        try:
            res = supabase.table("folders").select("id", count="exact").eq("user_id", user_id).execute()
            stats['folders_created'] = res.count or 0
        except:
            stats['folders_created'] = 0

        # Shared notes
        try:
            res = supabase.table("shared_notes").select("id", count="exact").eq("user_id", user_id).execute()
            stats['notes_shared'] = res.count or 0
        except:
            stats['notes_shared'] = 0

        # Conversations
        try:
            res = supabase.table("conversations").select("id", count="exact").eq("user_id", user_id).execute()
            stats['conversations'] = res.count or 0
        except:
            stats['conversations'] = 0

        # Study groups
        try:
            res = supabase.table("study_group_members").select("group_id", count="exact").eq("user_id", user_id).execute()
            stats['groups_joined'] = res.count or 0
        except:
            stats['groups_joined'] = 0

        # Groups created (owner)
        try:
            res = supabase.table("study_groups").select("id", count="exact").eq("owner_id", user_id).execute()
            stats['groups_created'] = res.count or 0
        except:
            stats['groups_created'] = 0

        # Group messages sent
        try:
            res = supabase.table("study_group_messages").select("id", count="exact").eq("user_id", user_id).execute()
            stats['group_messages'] = res.count or 0
        except:
            stats['group_messages'] = 0

        # Notes viewed by others (usage/help count)
        try:
            res = supabase.table("note_views").select("id", count="exact").eq("note_owner_id", user_id).execute()
            stats['notes_helped'] = res.count or 0
        except:
            try:
                # Fallback: count views on user's notes
                user_notes = supabase.table("notes").select("id").eq("user_id", user_id).execute()
                note_ids = [n['id'] for n in (user_notes.data or [])]
                if note_ids:
                    res = supabase.table("note_views").select("id", count="exact").in_("note_id", note_ids).execute()
                    stats['notes_helped'] = res.count or 0
                else:
                    stats['notes_helped'] = 0
            except:
                stats['notes_helped'] = 0

        # Edu verified
        try:
            res = supabase.table("user_profiles").select("edu_verified").eq("id", user_id).execute()
            stats['edu_verified'] = bool(res.data and res.data[0].get('edu_verified'))
        except:
            stats['edu_verified'] = False

        # Preferences (for customization check)
        try:
            res = supabase.table("user_profiles").select("preferences").eq("id", user_id).execute()
            prefs = (res.data[0].get('preferences') or {}) if res.data else {}
            stats['has_theme'] = bool(prefs.get('theme') and prefs['theme'] != 'default')
            stats['has_flo_custom'] = bool(prefs.get('flo', {}).get('color') and prefs['flo']['color'] != 'sage')
            stats['has_wallpaper'] = bool(prefs.get('wallpaper'))
            stats['has_flo_hat'] = bool(prefs.get('flo', {}).get('hat') and prefs['flo']['hat'] != 'none')
            stats['has_custom_accent'] = bool(prefs.get('accent'))
            stats['has_custom_font'] = bool(prefs.get('font') and prefs['font'] != 'inter')
            stats['has_custom_cursor'] = bool(prefs.get('cursor') and prefs['cursor'] != 'default')
            stats['flo_shape'] = prefs.get('flo', {}).get('shape', 'sphere')
        except:
            stats['has_theme'] = False
            stats['has_flo_custom'] = False
            stats['has_wallpaper'] = False
            stats['has_flo_hat'] = False
            stats['has_custom_accent'] = False
            stats['has_custom_font'] = False
            stats['has_custom_cursor'] = False
            stats['flo_shape'] = 'sphere'

        # Dashboard widgets
        try:
            res = supabase.table("user_profiles").select("dashboard_layout").eq("id", user_id).execute()
            layout = (res.data[0].get('dashboard_layout') or {}) if res.data else {}
            widgets = layout.get('layout', [])
            stats['widget_count'] = len(widgets) if isinstance(widgets, list) else 0
        except:
            stats['widget_count'] = 0

        # AI chat messages
        try:
            res = supabase.table("conversation_messages").select("id", count="exact").eq("role", "user").execute()
            # This counts all messages, but we want user's. Use conversations table.
            user_convs = supabase.table("conversations").select("id").eq("user_id", user_id).execute()
            conv_ids = [c['id'] for c in (user_convs.data or [])]
            if conv_ids:
                msgs = supabase.table("conversation_messages").select("id", count="exact").in_("conversation_id", conv_ids[:50]).eq("role", "user").execute()
                stats['ai_messages'] = msgs.count or 0
            else:
                stats['ai_messages'] = 0
        except:
            stats['ai_messages'] = 0

        # Pomodoros
        try:
            res = supabase.table("pomodoro_completions").select("work_minutes", count="exact").eq("user_id", user_id).execute()
            stats['pomodoros'] = res.count or 0
            stats['pomodoro_minutes'] = sum(r.get('work_minutes', 25) for r in (res.data or []))
        except:
            stats['pomodoros'] = 0
            stats['pomodoro_minutes'] = 0

        # Calendar events
        try:
            res = supabase.table("calendar_events").select("id", count="exact").eq("user_id", user_id).execute()
            stats['calendar_events'] = res.count or 0
        except:
            stats['calendar_events'] = 0

        # Annotations
        try:
            res = supabase.table("note_annotations").select("id", count="exact").eq("user_id", user_id).execute()
            stats['annotations'] = res.count or 0
        except:
            stats['annotations'] = 0

        # Downloads
        try:
            res = supabase.table("download_transactions").select("id", count="exact").eq("user_id", user_id).execute()
            stats['downloads'] = res.count or 0
        except:
            stats['downloads'] = 0

        # Favorites
        try:
            res = supabase.table("note_favorites").select("id", count="exact").eq("user_id", user_id).execute()
            stats['favorites'] = res.count or 0
        except:
            stats['favorites'] = 0

        # Current hour (for time-based achievements)
        stats['current_hour'] = datetime.utcnow().hour

        # Account age in days
        try:
            res = supabase.table("user_profiles").select("created_at").eq("id", user_id).execute()
            if res.data and res.data[0].get('created_at'):
                created = datetime.fromisoformat(res.data[0]['created_at'].replace('Z', '+00:00'))
                stats['account_age_days'] = (datetime.now(created.tzinfo) - created).days
            else:
                stats['account_age_days'] = 0
        except:
            stats['account_age_days'] = 0

        # Define achievements
        achievements = [
            # Getting Started
            {"id": "first_steps", "name": "First Steps", "desc": "Create your account", "icon": "rocket", "category": "Getting Started", "unlocked": True, "progress": 1, "goal": 1},
            {"id": "scholar", "name": "Scholar", "desc": "Verify your .edu email", "icon": "graduation", "category": "Getting Started", "unlocked": stats['edu_verified'], "progress": 1 if stats['edu_verified'] else 0, "goal": 1},
            {"id": "librarian", "name": "Librarian", "desc": "Upload your first note", "icon": "book", "category": "Getting Started", "unlocked": stats['notes_uploaded'] >= 1, "progress": min(stats['notes_uploaded'], 1), "goal": 1},
            {"id": "social_butterfly", "name": "Social Butterfly", "desc": "Join your first study group", "icon": "users", "category": "Getting Started", "unlocked": stats['groups_joined'] >= 1, "progress": min(stats['groups_joined'], 1), "goal": 1},
            {"id": "conversationalist", "name": "Conversationalist", "desc": "Start your first AI chat", "icon": "chat", "category": "Getting Started", "unlocked": stats['conversations'] >= 1, "progress": min(stats['conversations'], 1), "goal": 1},

            # Contributing
            {"id": "sharer", "name": "Sharer", "desc": "Share a note via link", "icon": "link", "category": "Contributing", "unlocked": stats['notes_shared'] >= 1, "progress": min(stats['notes_shared'], 1), "goal": 1},
            {"id": "open_book", "name": "Open Book", "desc": "Make 5 notes public on the Nexus", "icon": "globe", "category": "Contributing", "unlocked": stats['notes_public'] >= 5, "progress": min(stats['notes_public'], 5), "goal": 5},
            {"id": "community_pillar", "name": "Community Pillar", "desc": "Make 25 notes public", "icon": "star", "category": "Contributing", "unlocked": stats['notes_public'] >= 25, "progress": min(stats['notes_public'], 25), "goal": 25},
            {"id": "helping_hand", "name": "Helping Hand", "desc": "Your notes help 10 students", "icon": "heart", "category": "Contributing", "unlocked": stats['notes_helped'] >= 10, "progress": min(stats['notes_helped'], 10), "goal": 10},
            {"id": "mentor", "name": "Mentor", "desc": "Your notes help 100 students", "icon": "award", "category": "Contributing", "unlocked": stats['notes_helped'] >= 100, "progress": min(stats['notes_helped'], 100), "goal": 100},
            {"id": "legend", "name": "Legend", "desc": "Your notes help 1000 students", "icon": "crown", "category": "Contributing", "unlocked": stats['notes_helped'] >= 1000, "progress": min(stats['notes_helped'], 1000), "goal": 1000},

            # Collecting
            {"id": "collector", "name": "Collector", "desc": "Upload 10 notes", "icon": "folder", "category": "Collecting", "unlocked": stats['notes_uploaded'] >= 10, "progress": min(stats['notes_uploaded'], 10), "goal": 10},
            {"id": "archivist", "name": "Archivist", "desc": "Upload 50 notes", "icon": "archive", "category": "Collecting", "unlocked": stats['notes_uploaded'] >= 50, "progress": min(stats['notes_uploaded'], 50), "goal": 50},
            {"id": "hoarder", "name": "Hoarder", "desc": "Upload 100 notes", "icon": "box", "category": "Collecting", "unlocked": stats['notes_uploaded'] >= 100, "progress": min(stats['notes_uploaded'], 100), "goal": 100},
            {"id": "organized", "name": "Organized", "desc": "Create 5 folders", "icon": "layers", "category": "Collecting", "unlocked": stats['folders_created'] >= 5, "progress": min(stats['folders_created'], 5), "goal": 5},
            {"id": "annotator", "name": "Annotator", "desc": "Annotate a note", "icon": "pen", "category": "Collecting", "unlocked": stats['annotations'] >= 1, "progress": min(stats['annotations'], 1), "goal": 1},
            {"id": "bookmark_lover", "name": "Bookmark Lover", "desc": "Favorite 10 notes", "icon": "bookmark", "category": "Collecting", "unlocked": stats['favorites'] >= 10, "progress": min(stats['favorites'], 10), "goal": 10},
            {"id": "scholar_reader", "name": "Scholar Reader", "desc": "Download 20 notes", "icon": "download", "category": "Collecting", "unlocked": stats['downloads'] >= 20, "progress": min(stats['downloads'], 20), "goal": 20},

            # Social
            {"id": "team_player", "name": "Team Player", "desc": "Send a message in a study group", "icon": "message", "category": "Social", "unlocked": stats['group_messages'] >= 1, "progress": min(stats['group_messages'], 1), "goal": 1},
            {"id": "leader", "name": "Leader", "desc": "Create a study group", "icon": "flag", "category": "Social", "unlocked": stats['groups_created'] >= 1, "progress": min(stats['groups_created'], 1), "goal": 1},
            {"id": "chatterbox", "name": "Chatterbox", "desc": "Send 50 group messages", "icon": "megaphone", "category": "Social", "unlocked": stats['group_messages'] >= 50, "progress": min(stats['group_messages'], 50), "goal": 50},
            {"id": "social_network", "name": "Social Network", "desc": "Join 5 study groups", "icon": "network", "category": "Social", "unlocked": stats['groups_joined'] >= 5, "progress": min(stats['groups_joined'], 5), "goal": 5},

            # AI & Learning
            {"id": "curious_mind", "name": "Curious Mind", "desc": "Ask AI 10 questions", "icon": "brain", "category": "AI & Learning", "unlocked": stats['ai_messages'] >= 10, "progress": min(stats['ai_messages'], 10), "goal": 10},
            {"id": "deep_thinker", "name": "Deep Thinker", "desc": "Ask AI 100 questions", "icon": "lightbulb", "category": "AI & Learning", "unlocked": stats['ai_messages'] >= 100, "progress": min(stats['ai_messages'], 100), "goal": 100},
            {"id": "knowledge_seeker", "name": "Knowledge Seeker", "desc": "Have 20 AI conversations", "icon": "telescope", "category": "AI & Learning", "unlocked": stats['conversations'] >= 20, "progress": min(stats['conversations'], 20), "goal": 20},
            {"id": "planner", "name": "Planner", "desc": "Add 10 calendar events", "icon": "calendar", "category": "AI & Learning", "unlocked": stats['calendar_events'] >= 10, "progress": min(stats['calendar_events'], 10), "goal": 10},

            # Focus
            {"id": "focused", "name": "Focused", "desc": "Complete your first pomodoro", "icon": "timer", "category": "Focus", "unlocked": stats['pomodoros'] >= 1, "progress": min(stats['pomodoros'], 1), "goal": 1},
            {"id": "in_the_zone", "name": "In the Zone", "desc": "Complete 10 pomodoros", "icon": "fire", "category": "Focus", "unlocked": stats['pomodoros'] >= 10, "progress": min(stats['pomodoros'], 10), "goal": 10},
            {"id": "deep_work", "name": "Deep Work", "desc": "Complete 50 pomodoros", "icon": "mountain", "category": "Focus", "unlocked": stats['pomodoros'] >= 50, "progress": min(stats['pomodoros'], 50), "goal": 50},
            {"id": "flow_state", "name": "Flow State", "desc": "Complete 100 pomodoros", "icon": "wave", "category": "Focus", "unlocked": stats['pomodoros'] >= 100, "progress": min(stats['pomodoros'], 100), "goal": 100},
            {"id": "time_lord", "name": "Time Lord", "desc": "Accumulate 1000 minutes of focus", "icon": "clock", "category": "Focus", "unlocked": stats['pomodoro_minutes'] >= 1000, "progress": min(stats['pomodoro_minutes'], 1000), "goal": 1000},

            # Customization
            {"id": "decorator", "name": "Decorator", "desc": "Change your dashboard theme", "icon": "palette", "category": "Customization", "unlocked": stats['has_theme'], "progress": 1 if stats['has_theme'] else 0, "goal": 1},
            {"id": "wallpaper_fan", "name": "Wallpaper Fan", "desc": "Set a wallpaper", "icon": "image", "category": "Customization", "unlocked": stats['has_wallpaper'], "progress": 1 if stats['has_wallpaper'] else 0, "goal": 1},
            {"id": "interior_designer", "name": "Interior Designer", "desc": "Add 10 widgets to your dashboard", "icon": "grid", "category": "Customization", "unlocked": stats['widget_count'] >= 10, "progress": min(stats['widget_count'], 10), "goal": 10},
            {"id": "flos_friend", "name": "Flo's Friend", "desc": "Customize Flo's appearance", "icon": "sparkle", "category": "Customization", "unlocked": stats['has_flo_custom'], "progress": 1 if stats['has_flo_custom'] else 0, "goal": 1},
            {"id": "hat_collector", "name": "Hat Collector", "desc": "Give Flo a hat", "icon": "tophat", "category": "Customization", "unlocked": stats['has_flo_hat'], "progress": 1 if stats['has_flo_hat'] else 0, "goal": 1},
            {"id": "color_master", "name": "Color Master", "desc": "Set a custom accent color", "icon": "rainbow", "category": "Customization", "unlocked": stats['has_custom_accent'], "progress": 1 if stats['has_custom_accent'] else 0, "goal": 1},
            {"id": "typographer", "name": "Typographer", "desc": "Change the font", "icon": "type", "category": "Customization", "unlocked": stats['has_custom_font'], "progress": 1 if stats['has_custom_font'] else 0, "goal": 1},
            {"id": "cursor_crafter", "name": "Cursor Crafter", "desc": "Use a custom cursor", "icon": "cursor", "category": "Customization", "unlocked": stats['has_custom_cursor'], "progress": 1 if stats['has_custom_cursor'] else 0, "goal": 1},

            # Hidden / Secret (shown as ??? until unlocked)
            {"id": "night_owl", "name": "Night Owl", "desc": "Use StudyFlow between midnight and 4am", "icon": "moon", "category": "Secret", "hidden": True, "unlocked": 0 <= stats['current_hour'] <= 3, "progress": 1 if 0 <= stats['current_hour'] <= 3 else 0, "goal": 1},
            {"id": "early_bird", "name": "Early Bird", "desc": "Use StudyFlow between 4am and 6am", "icon": "sunrise", "category": "Secret", "hidden": True, "unlocked": 4 <= stats['current_hour'] <= 5, "progress": 1 if 4 <= stats['current_hour'] <= 5 else 0, "goal": 1},
            {"id": "veteran", "name": "Veteran", "desc": "Have an account for 30 days", "icon": "medal", "category": "Secret", "hidden": True, "unlocked": stats['account_age_days'] >= 30, "progress": min(stats['account_age_days'], 30), "goal": 30},
            {"id": "old_timer", "name": "Old Timer", "desc": "Have an account for 365 days", "icon": "hourglass", "category": "Secret", "hidden": True, "unlocked": stats['account_age_days'] >= 365, "progress": min(stats['account_age_days'], 365), "goal": 365},
            {"id": "widget_hoarder", "name": "Widget Hoarder", "desc": "Add 25 widgets to your dashboard", "icon": "infinity", "category": "Secret", "hidden": True, "unlocked": stats['widget_count'] >= 25, "progress": min(stats['widget_count'], 25), "goal": 25},
            {"id": "sharing_spree", "name": "Sharing Spree", "desc": "Share 10 notes via link", "icon": "share", "category": "Secret", "hidden": True, "unlocked": stats['notes_shared'] >= 10, "progress": min(stats['notes_shared'], 10), "goal": 10},
            {"id": "ai_addict", "name": "AI Addict", "desc": "Ask AI 500 questions", "icon": "robot", "category": "Secret", "hidden": True, "unlocked": stats['ai_messages'] >= 500, "progress": min(stats['ai_messages'], 500), "goal": 500},
            {"id": "completionist", "name": "Completionist", "desc": "Unlock every other achievement", "icon": "trophy", "category": "Secret", "hidden": True, "unlocked": False, "progress": 0, "goal": 1},
        ]

        # Check completionist (unlocked if all non-completionist achievements are unlocked)
        non_completionist = [a for a in achievements if a['id'] != 'completionist']
        all_unlocked = all(a['unlocked'] for a in non_completionist)
        for a in achievements:
            if a['id'] == 'completionist':
                a['unlocked'] = all_unlocked
                a['progress'] = 1 if all_unlocked else 0

        unlocked_count = sum(1 for a in achievements if a['unlocked'])

        return jsonify({
            "achievements": achievements,
            "unlocked": unlocked_count,
            "total": len(achievements),
            "stats": stats
        }), 200

    except Exception as e:
        debug_log(f"Achievements error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": "Could not load achievements"}), 500


@app.route("/api/study-streak", methods=["GET"])
@supabase_auth_required
def get_study_streak():
    """Get study activity heatmap data for the past year."""
    try:
        user_id = request.user_id
        one_year_ago = (datetime.utcnow() - timedelta(days=365)).isoformat()
        activity = {}

        # Notes uploaded
        try:
            res = supabase.table("notes").select("uploaded_at").eq("user_id", user_id).gte("uploaded_at", one_year_ago).execute()
            for r in (res.data or []):
                day = r['uploaded_at'][:10]
                activity[day] = activity.get(day, 0) + 2
        except: pass

        # Conversations started
        try:
            res = supabase.table("conversations").select("created_at").eq("user_id", user_id).gte("created_at", one_year_ago).execute()
            for r in (res.data or []):
                day = r['created_at'][:10]
                activity[day] = activity.get(day, 0) + 1
        except: pass

        # Group messages
        try:
            res = supabase.table("study_group_messages").select("created_at").eq("user_id", user_id).gte("created_at", one_year_ago).execute()
            for r in (res.data or []):
                day = r['created_at'][:10]
                activity[day] = activity.get(day, 0) + 1
        except: pass

        # Note views
        try:
            res = supabase.table("note_views").select("viewed_at").eq("user_id", user_id).gte("viewed_at", one_year_ago).execute()
            for r in (res.data or []):
                day = r['viewed_at'][:10]
                activity[day] = activity.get(day, 0) + 1
        except: pass

        # Downloads
        try:
            res = supabase.table("download_transactions").select("downloaded_at").eq("user_id", user_id).gte("downloaded_at", one_year_ago).execute()
            for r in (res.data or []):
                day = r['downloaded_at'][:10]
                activity[day] = activity.get(day, 0) + 1
        except: pass

        # Pomodoros
        try:
            res = supabase.table("pomodoro_completions").select("completed_at").eq("user_id", user_id).gte("completed_at", one_year_ago).execute()
            for r in (res.data or []):
                day = r['completed_at'][:10]
                activity[day] = activity.get(day, 0) + 2
        except: pass

        # Calculate streaks
        today = datetime.utcnow().strftime('%Y-%m-%d')
        sorted_days = sorted(activity.keys(), reverse=True)

        current_streak = 0
        longest_streak = 0
        temp_streak = 0
        check_date = datetime.utcnow().date()

        for i in range(365):
            day_str = check_date.strftime('%Y-%m-%d')
            if day_str in activity:
                if i == 0 or (i > 0 and temp_streak > 0):
                    temp_streak += 1
                else:
                    temp_streak = 1
                if current_streak == 0 and i <= 1:
                    current_streak = temp_streak
            else:
                if current_streak == 0 and i == 0:
                    pass  # today not active yet, check yesterday
                elif current_streak == 0 and temp_streak > 0:
                    current_streak = temp_streak
                longest_streak = max(longest_streak, temp_streak)
                temp_streak = 0
            check_date -= timedelta(days=1)

        longest_streak = max(longest_streak, temp_streak)
        if current_streak == 0:
            current_streak = temp_streak

        # Recalculate current streak properly
        current_streak = 0
        check = datetime.utcnow().date()
        # Allow today to not have activity yet
        if check.strftime('%Y-%m-%d') not in activity:
            check -= timedelta(days=1)
        while check.strftime('%Y-%m-%d') in activity:
            current_streak += 1
            check -= timedelta(days=1)

        return jsonify({
            "activity": activity,
            "current_streak": current_streak,
            "longest_streak": longest_streak,
            "total_active_days": len(activity),
            "today": today
        }), 200

    except Exception as e:
        debug_log(f"Study streak error: {e}\n{traceback.format_exc()}")
        return jsonify({"activity": {}, "current_streak": 0, "longest_streak": 0, "total_active_days": 0}), 200


## ====== NOTE SHARING ======

@app.route("/api/notes/<note_id>/share", methods=["POST"])
@supabase_auth_required
def create_share_link(note_id):
    """Create a share link for a note the user owns. Body: { "expires_in": "24h"|"7d"|"30d"|null }"""
    try:
        import secrets as _secrets

        # Verify ownership
        note = supabase.table("notes").select("id, user_id").eq("id", note_id).eq("user_id", request.user_id).execute()
        if not note.data:
            return jsonify({"error": "Note not found or not yours"}), 404

        # Check for existing share link
        existing = supabase.table("shared_notes").select("*").eq("note_id", note_id).eq("user_id", request.user_id).execute()
        if existing.data:
            share = existing.data[0]
            return jsonify({"share_token": share["share_token"], "expires_at": share["expires_at"], "already_existed": True}), 200

        # Parse expiry
        data = request.get_json() or {}
        expires_in = data.get("expires_in")
        expires_at = None
        if expires_in == "24h":
            expires_at = (datetime.utcnow() + timedelta(hours=24)).isoformat()
        elif expires_in == "7d":
            expires_at = (datetime.utcnow() + timedelta(days=7)).isoformat()
        elif expires_in == "30d":
            expires_at = (datetime.utcnow() + timedelta(days=30)).isoformat()
        # else: no expiry

        token = _secrets.token_urlsafe(16)

        result = supabase.table("shared_notes").insert({
            "note_id": note_id,
            "user_id": request.user_id,
            "share_token": token,
            "expires_at": expires_at
        }).execute()

        return jsonify({"share_token": token, "expires_at": expires_at}), 201

    except Exception as e:
        debug_log(f"Create share link error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/share", methods=["DELETE"])
@supabase_auth_required
def revoke_share_link(note_id):
    """Revoke a share link for a note the user owns."""
    try:
        result = supabase.table("shared_notes").delete().eq("note_id", note_id).eq("user_id", request.user_id).execute()
        return jsonify({"success": True}), 200
    except Exception as e:
        debug_log(f"Revoke share link error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/share", methods=["GET"])
@supabase_auth_required
def get_share_link(note_id):
    """Get the share link info for a note the user owns."""
    try:
        result = supabase.table("shared_notes").select("*").eq("note_id", note_id).eq("user_id", request.user_id).execute()
        if not result.data:
            return jsonify({"shared": False}), 200
        share = result.data[0]
        return jsonify({"shared": True, "share_token": share["share_token"], "expires_at": share["expires_at"], "created_at": share["created_at"]}), 200
    except Exception as e:
        debug_log(f"Get share link error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/shared/<token>", methods=["GET"])
def get_shared_note(token):
    """Public endpoint - get a shared note by token. No auth required."""
    try:
        # Look up the share link
        result = supabase.table("shared_notes").select("*, notes(id, original_filename, file_path, file_type, page_count, uploaded_at, user_id, university, course_code, professor, semester)").eq("share_token", token).execute()

        if not result.data:
            return jsonify({"error": "Share link not found"}), 404

        share = result.data[0]

        # Check expiry
        if share.get("expires_at"):
            expires = datetime.fromisoformat(share["expires_at"].replace("Z", "+00:00"))
            if datetime.now(expires.tzinfo) > expires:
                return jsonify({"error": "This share link has expired"}), 410

        note = share.get("notes")
        if not note:
            return jsonify({"error": "Note no longer exists"}), 404

        # Get sharer/uploader info
        uploader = supabase.table("user_profiles").select("full_name, email").eq("id", share["user_id"]).execute()
        uploader_name = "Anonymous"
        uploader_username = "anonymous"
        if uploader.data:
            uploader_name = uploader.data[0].get("full_name") or uploader.data[0].get("email", "").split("@")[0]
            uploader_username = uploader.data[0].get("email", "anonymous").split("@")[0]

        return jsonify({
            "note_id": note["id"],
            "filename": note.get("original_filename"),
            "file_type": note.get("file_type"),
            "page_count": note.get("page_count"),
            "uploaded_at": note.get("uploaded_at"),
            "shared_by": uploader_name,
            "shared_by_username": uploader_username,
            "university": note.get("university"),
            "course_code": note.get("course_code"),
            "professor": note.get("professor"),
            "semester": note.get("semester"),
            "expires_at": share.get("expires_at"),
            "share_id": share.get("id", "")[:8]
        }), 200

    except Exception as e:
        debug_log(f"Get shared note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/shared/<token>/view-file", methods=["GET"])
def view_shared_file(token):
    """Public endpoint - get the actual file for a shared note. No auth required."""
    try:
        from flask import redirect

        # Look up the share link
        result = supabase.table("shared_notes").select("*, notes(id, file_path, original_filename)").eq("share_token", token).execute()

        if not result.data:
            return jsonify({"error": "Share link not found"}), 404

        share = result.data[0]

        # Check expiry
        if share.get("expires_at"):
            expires = datetime.fromisoformat(share["expires_at"].replace("Z", "+00:00"))
            if datetime.now(expires.tzinfo) > expires:
                return jsonify({"error": "This share link has expired"}), 410

        note = share.get("notes")
        if not note or not note.get("file_path"):
            return jsonify({"error": "File not found"}), 404

        # Generate signed URL
        signed_url_response = supabase.storage.from_("note-files").create_signed_url(
            path=note["file_path"],
            expires_in=3600
        )
        file_url = signed_url_response["signedURL"]
        return redirect(file_url)

    except Exception as e:
        debug_log(f"View shared file error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/shared/<token>/view-as-images", methods=["GET"])
def view_shared_as_images(token):
    """Public endpoint - get shared note pages as images. No auth required."""
    try:
        import fitz
        import base64

        # Look up the share link
        result = supabase.table("shared_notes").select("*, notes(id, file_path, file_type, page_count)").eq("share_token", token).execute()

        if not result.data:
            return jsonify({"error": "Share link not found"}), 404

        share = result.data[0]

        # Check expiry
        if share.get("expires_at"):
            expires = datetime.fromisoformat(share["expires_at"].replace("Z", "+00:00"))
            if datetime.now(expires.tzinfo) > expires:
                return jsonify({"error": "This share link has expired"}), 410

        note = share.get("notes")
        if not note or not note.get("file_path"):
            return jsonify({"error": "File not found"}), 404

        # Download the file
        file_data = supabase.storage.from_("note-files").download(note["file_path"])

        # Try to open as PDF first (most files are already PDF)
        # If that fails, convert to PDF on-the-fly
        try:
            doc = fitz.open(stream=file_data, filetype="pdf")
        except:
            # File isn't a PDF - convert it now
            from StudyFlow.backend.pdf_flatten import convert_to_pdf
            original_filename = note.get("file_path", "file.txt")
            pdf_data, _ = convert_to_pdf(file_data, original_filename)
            if not pdf_data:
                return jsonify({"error": "Could not convert file to PDF"}), 400
            doc = fitz.open(stream=pdf_data, filetype="pdf")

        images = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            mat = fitz.Matrix(2, 2)
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            images.append(base64.b64encode(img_bytes).decode("utf-8"))

        doc.close()
        return jsonify({"images": images, "page_count": len(images)}), 200

    except Exception as e:
        debug_log(f"View shared as images error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/user/shared-notes", methods=["GET"])
@supabase_auth_required
def list_user_shared_notes():
    """Get all share links created by the current user."""
    try:
        result = supabase.table("shared_notes").select("*, notes(id, original_filename)").eq("user_id", request.user_id).execute()
        return jsonify(result.data or []), 200
    except Exception as e:
        debug_log(f"List shared notes error: {e}\n{traceback.format_exc()}")
        return jsonify([]), 200


@app.route("/api/user/send-edu-verification", methods=["POST"])
@supabase_auth_required
def send_edu_verification():
    """
    Send verification email to .edu email address.
    Request body: { "edu_email": "user@university.edu" }
    """
    try:
        import uuid
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        edu_email = data.get("edu_email", "").strip().lower()

        if not edu_email:
            return jsonify({"error": "edu_email is required"}), 400

        if not edu_email.endswith('.edu'):
            return jsonify({"error": "Email must be a .edu address"}), 400

        # Generate verification token
        verification_token = str(uuid.uuid4())

        # Store pending verification in user_profiles
        response = supabase.table("user_profiles").update({
            "pending_edu_email": edu_email,
            "edu_verification_token": verification_token
        }).eq("id", request.user_id).execute()

        if not response.data:
            return jsonify({"error": "Failed to update profile"}), 500

        # Send verification email
        verification_url = f"https://studyflowsuite.com/verify-edu.html?token={verification_token}"

        html_content = f"""
        <div style="font-family: Arial, sans-serif; max-width: 600px; margin: 0 auto; padding: 20px;">
            <h2 style="color: #7c9885;">Verify Your University Email</h2>
            <p>Click the button below to verify your .edu email address and unlock access to the StudyFlow Nexus note library.</p>
            <div style="margin: 30px 0; text-align: center;">
                <a href="{verification_url}"
                   style="background: linear-gradient(135deg, #7c9885 0%, #6b8575 100%);
                          color: white;
                          padding: 14px 32px;
                          text-decoration: none;
                          border-radius: 12px;
                          display: inline-block;
                          font-weight: 600;">
                    Verify Email Address
                </a>
            </div>
            <p style="color: #5d5d5d; font-size: 14px;">Or copy and paste this link into your browser:</p>
            <p style="color: #7c9885; word-break: break-all; font-size: 12px;">{verification_url}</p>
            <p style="color: #999; font-size: 12px; margin-top: 30px;">If you didn't request this verification, you can safely ignore this email.</p>
        </div>
        """

        response = send_brevo_email(
            to_email=edu_email,
            to_name=None,
            subject="Verify your university email - StudyFlow Suite",
            html_content=html_content
        )

        if response.status_code in (200, 201):
            debug_log(f"Verification email sent to {edu_email}")
            return jsonify({"success": True, "message": "Verification email sent"}), 200
        else:
            debug_log(f"Failed to send verification email: {response.text}")
            return jsonify({"error": "Failed to send email"}), 500

    except Exception as e:
        debug_log(f"Send .edu verification error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/user/verify-edu-email", methods=["POST"])
def verify_edu_email():
    """
    Verify .edu email using token from verification link.
    Request body: { "token": "uuid-token" }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        token = data.get("token")

        if not token:
            return jsonify({"error": "token is required"}), 400

        # Find user with matching token
        response = supabase.table("user_profiles").select("*").eq("edu_verification_token", token).execute()

        if not response.data or len(response.data) == 0:
            return jsonify({"error": "Invalid or expired verification token"}), 404

        user_profile = response.data[0]
        pending_email = user_profile.get("pending_edu_email")

        if not pending_email:
            return jsonify({"error": "No pending verification found"}), 404

        # Update user profile to mark as verified
        update_response = supabase.table("user_profiles").update({
            "edu_email_verified": True,
            "edu_email": pending_email,
            "pending_edu_email": None,
            "edu_verification_token": None
        }).eq("id", user_profile["id"]).execute()

        if update_response.data:
            debug_log(f"✅ User {user_profile['id']} verified .edu email: {pending_email}")
            return jsonify({
                "success": True,
                "message": "Email verified successfully",
                "edu_email": pending_email
            }), 200
        else:
            return jsonify({"error": "Failed to update verification status"}), 500

    except Exception as e:
        debug_log(f"[-] Verify .edu email error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/user/verify-age", methods=["POST"])
@supabase_auth_required
def verify_age():
    """
    Verify user's age using transactional data (Missouri SB 1455 / GUARD Act).

    Request body: {
        "legal_name": "Cody Williams",
        "zip_code": "65802",
        "birthdate": "1998-03-15"
    }

    Returns verification result and unfreezes account if successful.
    """
    try:
        from StudyFlow.backend.age_verification import verify_age_transactional
        from StudyFlow.backend.supabase_client import update_age_verification

        data = request.get_json()
        if not data:
            return jsonify({"error": "Missing request body"}), 400

        legal_name = data.get('legal_name', '').strip()
        zip_code = data.get('zip_code', '').strip()
        birthdate = data.get('birthdate', '').strip()
        address = data.get('address', '').strip()
        city = data.get('city', '').strip()
        state = data.get('state', '').strip()

        if not legal_name or not zip_code or not birthdate:
            return jsonify({"error": "legal_name, zip_code, and birthdate are required"}), 400

        # Validate zip code format
        if not zip_code.replace('-', '').isdigit() or len(zip_code.replace('-', '')) not in (5, 9):
            return jsonify({"error": "Invalid zip code format"}), 400

        # Call Veratad AgeMatch5.0 verification
        result = verify_age_transactional(
            legal_name, zip_code, birthdate,
            address=address, city=city, state=state,
            reference=request.user_id
        )

        if result["verified"]:
            # Update profile with verification status
            update_age_verification(
                user_id=request.user_id,
                verified=True,
                method=result["method"],
                legal_name=legal_name,
                zip_code=zip_code,
                birthdate=birthdate
            )

            debug_log(f"[+] Age verified for {request.user_id}: {legal_name}, age={result['age']}")

            return jsonify({
                "verified": True,
                "message": "Age verification successful. Your account is now fully active.",
                "method": result["method"],
                "age": result["age"]
            }), 200
        else:
            debug_log(f"[-] Age verification failed for {request.user_id}: {result['reason']}")

            return jsonify({
                "verified": False,
                "message": result["reason"],
                "method": result["method"]
            }), 200

    except Exception as e:
        debug_log(f"[-] Age verification error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": "Age verification failed"}), 500


@app.route("/api/user/verification-status", methods=["GET"])
@supabase_auth_required
def verification_status():
    """
    Get user's age verification and account freeze status.
    Returns current verification state for frontend display.
    """
    try:
        from StudyFlow.backend.supabase_client import get_verification_status

        status = get_verification_status(request.user_id)

        if not status:
            return jsonify({
                "age_verified": False,
                "account_frozen": False,
                "needs_verification": True
            }), 200

        # Check if re-verification is needed (12-month cycle)
        needs_reverification = False
        if status.get('last_verified_at'):
            from datetime import datetime, timezone
            last_verified = datetime.fromisoformat(status['last_verified_at'].replace('Z', '+00:00'))
            months_since = (datetime.now(timezone.utc) - last_verified).days / 30
            needs_reverification = months_since >= 12

        return jsonify({
            "age_verified": status.get("age_verified", False),
            "age_verified_at": status.get("age_verified_at"),
            "age_verification_method": status.get("age_verification_method"),
            "last_verified_at": status.get("last_verified_at"),
            "account_frozen": status.get("account_frozen", False),
            "frozen_reason": status.get("frozen_reason"),
            "needs_reverification": needs_reverification
        }), 200

    except Exception as e:
        debug_log(f"[-] Verification status error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": "Failed to get verification status"}), 500


@app.route("/api/notes/<note_id>/download", methods=["GET"])
@supabase_auth_required
def download_note_endpoint(note_id):
    """
    Download a note file by ID (only if it belongs to current user).
    PDFs and images are flattened (rasterized) with watermarks.
    Every download is logged as a transaction.
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from StudyFlow.backend.pdf_flatten import flatten_pdf, flatten_image, can_flatten, convert_to_pdf
        from flask import send_file
        from datetime import date
        import io
        import uuid

        # DAILY DOWNLOAD LIMIT CHECK (3 downloads/day for free users)
        # Fetch user profile to check subscription and download count
        profile_result = supabase.table("user_profiles").select(
            "subscription_tier, subscription_status, daily_downloads_count, daily_downloads_reset_date, good_standing, permanent_bad_standing, last_verified_upload_date"
        ).eq("id", request.user_id).single().execute()
        profile = profile_result.data

        # CHECK: Good standing expires after 30 days without a quality upload (free users only)
        from datetime import datetime, timedelta
        is_scholar_check = (
            profile.get('subscription_tier') == 'pro' and
            profile.get('subscription_status') == 'active'
        )

        if not is_scholar_check and profile.get('good_standing', False):
            last_upload = profile.get('last_verified_upload_date')
            if last_upload:
                try:
                    last_upload_date = datetime.fromisoformat(last_upload.replace('Z', '+00:00'))
                    days_since_upload = (datetime.utcnow().replace(tzinfo=last_upload_date.tzinfo) - last_upload_date).days

                    if days_since_upload > 30:
                        # Revoke good standing - hasn't uploaded in 30+ days
                        supabase.table("user_profiles").update({
                            "good_standing": False
                        }).eq("id", request.user_id).execute()
                        profile['good_standing'] = False
                        debug_log(f"[!] User {request.user_id} lost good standing (no upload for {days_since_upload} days)")
                except Exception as e:
                    debug_log(f"[-] Error checking last upload date: {e}")

        # Check if Scholar's Club member (unlimited downloads)
        is_scholar = (
            profile.get('subscription_tier') == 'pro' and
            profile.get('subscription_status') == 'active'
        )

        if not is_scholar:
            # Free user - enforce 3 downloads per day limit
            today = date.today()
            reset_date = profile.get('daily_downloads_reset_date')
            download_count = profile.get('daily_downloads_count', 0)

            # Reset counter if new day
            if not reset_date or str(reset_date) != str(today):
                supabase.table("user_profiles").update({
                    "daily_downloads_count": 0,
                    "daily_downloads_reset_date": str(today)
                }).eq("id", request.user_id).execute()
                download_count = 0
                debug_log(f"[+] Reset daily download count for user {request.user_id}")

            # Check if at limit
            if download_count >= 3:
                debug_log(f"[!] User {request.user_id} hit daily download limit ({download_count}/3)")
                return jsonify({
                    "error": "Daily download limit reached",
                    "message": "You've used all 3 downloads today. Join Scholar's Club for unlimited downloads!",
                    "limit_reached": True,
                    "downloads_used": download_count,
                    "downloads_limit": 3,
                    "subscribe_url": "https://unclephilburt.github.io/studyflowwebsite/account.html"
                }), 403

        # Get note metadata -- check own notes first, then public notes
        result = supabase.table("notes").select("*").eq("id", note_id).eq("user_id", request.user_id).execute()
        note_data = result.data[0] if result.data else None

        if not note_data:
            # Not the user's own note -- check if it's a public note
            result = supabase.table("notes").select("*").eq("id", note_id).eq("is_public", True).execute()
            note_data = result.data[0] if result.data else None

            # GOOD STANDING CHECK: Required for downloading public notes (not your own)
            # (Reuse profile fetched earlier for download limit check)
            if note_data:
                # Block if permanently banned
                if profile.get('permanent_bad_standing', False):
                    return jsonify({"error": "Account suspended due to DMCA violations"}), 403

                # Check good standing (Scholar's Club members are always in good standing)
                if not is_scholar and not profile.get('good_standing', False):
                    return jsonify({
                        "error": "Good Standing required to download notes",
                        "message": "Upload a verified note or subscribe to regain access",
                        "good_standing_required": True
                    }), 403

        if not note_data:
            return jsonify({"error": "Note not found or unauthorized"}), 404

        file_path = note_data.get('file_path')
        original_filename = note_data.get('original_filename')

        if not file_path:
            return jsonify({"error": "File not found in storage"}), 404

        # Generate transaction code
        transaction_code = "DL-" + uuid.uuid4().hex[:8]

        # Get UPLOADER's username for watermark (not downloader's)
        username = "Anonymous"
        try:
            uploader_id = note_data.get('user_id')
            if uploader_id:
                profile_result = supabase.table("user_profiles").select("username").eq("id", uploader_id).execute()
                if profile_result.data and profile_result.data[0].get("username"):
                    username = profile_result.data[0]["username"]
        except Exception:
            pass

        # Download file from Supabase Storage
        file_data = supabase.storage.from_('note-files').download(file_path)

        # Log download transaction
        ip_address = request.headers.get('X-Forwarded-For', request.remote_addr)
        if ip_address and ',' in ip_address:
            ip_address = ip_address.split(',')[0].strip()

        try:
            supabase.table("download_transactions").insert({
                "user_id": request.user_id,
                "note_id": note_id,
                "original_filename": original_filename,
                "transaction_code": transaction_code,
                "ip_address": ip_address,
                "user_agent": request.headers.get('User-Agent', '')
            }).execute()
            debug_log(f"[+] Download transaction logged: {transaction_code}")
        except Exception as tx_err:
            debug_log(f"[-] Failed to log download transaction: {tx_err}")

        # Flatten and watermark -- everything becomes a PDF
        flattenable, file_type = can_flatten(original_filename)
        name_base = original_filename.rsplit('.', 1)[0] if '.' in original_filename else original_filename

        if flattenable and file_type == 'pdf':
            file_data = flatten_pdf(file_data, username, transaction_code)
        elif flattenable and file_type == 'image':
            file_data = flatten_image(file_data, username, transaction_code)
        else:
            # Legacy file (txt, docx, etc.) -- convert to PDF first, then flatten
            pdf_data, _ = convert_to_pdf(file_data, original_filename)
            if pdf_data:
                file_data = flatten_pdf(pdf_data, username, transaction_code)

        download_name = f"{name_base}_studyflow.pdf"
        mimetype = 'application/pdf'

        # Increment download counter for free users
        if not is_scholar:
            try:
                new_count = (profile.get('daily_downloads_count', 0) or 0) + 1
                supabase.table("user_profiles").update({
                    "daily_downloads_count": new_count
                }).eq("id", request.user_id).execute()
                debug_log(f"[+] User {request.user_id} download count: {new_count}/3")
            except Exception as count_err:
                debug_log(f"[-] Failed to increment download count: {count_err}")

        return send_file(
            io.BytesIO(file_data),
            mimetype=mimetype,
            as_attachment=True,
            download_name=download_name
        )

    except Exception as e:
        debug_log(f"[-] Download note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/search", methods=["POST"])
@supabase_auth_required
def search_notes():
    """
    Search through ALL notes using vector similarity (pgvector + embeddings).

    COLLECTIVE BRAIN: Searches user's own notes + ALL public notes from ALL students.
    No course filtering - truly shared knowledge base across all subjects.

    Expects JSON:
    {
        "question": "What is photosynthesis?"
    }

    Returns:
    {
        "results": [
            {
                "source": "Biology_Chapter_3.pdf",
                "text": "Photosynthesis is the process by which...",
                "hint": "Look at the section about cellular processes...",
                "from_collective_brain": true,
                "similarity": 0.89
            }
        ]
    }
    """
    try:
        from StudyFlow.backend.supabase_client import search_notes_vector, get_user_profile, supabase
        from StudyFlow.backend.embedding_client import generate_embedding

        data = request.get_json()
        if not data or not data.get('question'):
            return jsonify({"error": "Missing question"}), 400

        question = data.get('question')

        # Get user profile (just to verify user exists)
        user_profile = get_user_profile(request.user_id)
        if not user_profile:
            return jsonify({"error": "User profile not found", "results": []}), 404

        # Generate embedding for the question
        query_embedding = generate_embedding(question)
        if not query_embedding:
            return jsonify({"error": "Failed to generate query embedding", "results": []}), 500

        debug_log(f"🔍 Searching ENTIRE collective brain for: '{question}'")

        # Search using pgvector - searches ALL public notes from ALL students
        search_results = search_notes_vector(
            query_embedding=query_embedding,
            user_id=request.user_id,
            university=None,  # No course filtering - true collective brain
            course_code=None,
            match_threshold=0.4,  # Only return results with >40% similarity
            match_count=5  # Top 5 results
        )

        print(f"🔎 Search returned {len(search_results) if search_results else 0} results")

        if not search_results:
            print("⚠️ No search results found, returning empty array")
            return jsonify({"results": []}), 200

        # Format results with hints
        formatted_results = []
        print(f"📝 Processing {len(search_results)} results...")

        for result in search_results:
            # Get note filename
            note_result = supabase.table("notes").select("original_filename").eq("id", result['note_id']).execute()
            filename = note_result.data[0]['original_filename'] if note_result.data else "Unknown"

            # Check if this is a Wikipedia source
            is_wikipedia = result.get('university') == 'Wikipedia'
            wikipedia_url = None

            if is_wikipedia and filename.endswith('.txt'):
                # Extract article title and create Wikipedia URL
                article_title = filename.replace('.txt', '')
                url_title = article_title.replace(' ', '_')
                wikipedia_url = f"https://en.wikipedia.org/wiki/{url_title}"

            # Decide which text to show
            text_to_show = result['content_summary'] if result['content_summary'] else result['chunk_text']

            # Generate a detailed answer using the text
            print(f"🔍 About to generate detailed answer for: {question[:50]}...")
            hint = generate_hint_from_text(question, text_to_show)
            print(f"💡 Generated answer: {hint[:100]}...")

            # Build source attribution
            if is_wikipedia:
                source_text = f"Wikipedia: {filename.replace('.txt', '')}"
            elif result['university']:
                source_text = f"{filename} ({result['university']} - {result['course_code']})"
            else:
                source_text = filename

            formatted_results.append({
                "source": source_text,
                "text": text_to_show[:300] + "..." if len(text_to_show) > 300 else text_to_show,
                "hint": hint,
                "from_collective_brain": not result['is_own_note'],
                "from_wikipedia": is_wikipedia,
                "wikipedia_url": wikipedia_url,
                "similarity": round(result['similarity'], 2)
            })

        debug_log(f"✅ Found {len(formatted_results)} results (similarity > 0.4)")

        return jsonify({"results": formatted_results}), 200

    except Exception as e:
        debug_log(f"❌ Search notes error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e), "results": []}), 500


def generate_hint_from_text(question, text):
    """
    Generate a detailed answer from the found text using Gemini.
    """
    try:
        import google.generativeai as genai

        # Check if API key exists
        api_key = os.getenv("GEMINI_API_KEY")
        if not api_key:
            debug_log("❌ GEMINI_API_KEY not found in environment! Answer generation disabled.")
            return "Review this section carefully to find the answer."

        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-3-flash-preview')

        prompt = f"""You are a knowledgeable study tutor. A student asked: "{question}"

I found this relevant excerpt from their notes:
"{text[:1000]}"

Generate a DETAILED, COMPREHENSIVE answer to their question using the information from the notes.

Your job:
- Answer the question completely and thoroughly
- Include all relevant facts, definitions, explanations, and examples
- Explain concepts clearly with context
- Be specific and detailed - don't hold back information
- Make it educational and easy to understand
- Use 2-4 paragraphs if needed to fully explain the topic

DO NOT say things like:
- "According to the notes..."
- "The excerpt mentions..."
- "Review the section..."

Just answer the question directly as if you're a tutor who knows this information.

Return your detailed answer:"""

        debug_log(f"Generating detailed answer for question: '{question[:50]}...'")
        response = model.generate_content(prompt)
        answer = response.text.strip()

        try:
            from StudyFlow.backend.cost_tracker import track_ai_call
            track_ai_call("gemini", "flash", "hint_gen")
        except:
            pass

        debug_log(f"✅ Generated answer: '{answer[:80]}...'")

        return answer if answer else "Review this section carefully to find the answer."

    except Exception as e:
        debug_log(f"❌ Error generating answer: {type(e).__name__}: {str(e)}")
        return "Review this section to find the answer."


@app.route("/api/notes/chat", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def chat_with_notes():
    """
    Conversational AI chat with memory - students can have back-and-forth discussions

    Expects JSON:
    {
        "message": "How does DNA replicate?",
        "conversation_id": "optional-uuid" // omit to start new conversation
    }

    Returns:
    {
        "conversation_id": "uuid",
        "response": "DNA replicates through semiconservative replication...",
        "sources": [
            {
                "filename": "DNA.txt (Wikipedia - Biology)",
                "similarity": 0.89
            }
        ]
    }
    """
    try:
        from StudyFlow.backend.supabase_client import search_notes_vector, get_user_profile, supabase, log_ai_response
        from StudyFlow.backend.embedding_client import generate_embedding
        from StudyFlow.backend.conversational_noteflow import (
            create_conversation,
            get_conversation,
            get_conversation_messages,
            add_message,
            generate_conversational_response,
            generate_conversation_title,
            update_conversation_title
        )

        data = request.get_json()
        if not data or not data.get('message'):
            return jsonify({"error": "Missing message"}), 400

        message = data.get('message')
        conv_id = data.get('conversation_id')
        search_scope = data.get('search_scope', 'all')  # 'personal' or 'all'
        source = data.get('source', 'chat')  # 'chat' or 'plugin'

        # Get or create conversation
        if conv_id:
            conversation = get_conversation(conv_id, request.user_id)
            if not conversation:
                return jsonify({"error": "Conversation not found or access denied"}), 404
        else:
            conv_id = create_conversation(request.user_id, source=source)
            conversation = get_conversation(conv_id, request.user_id)
            # Invalidate canvas cache so new conversation appears
            if redis_cache:
                try: redis_cache.delete(f"canvas_preload:{request.user_id}")
                except: pass

        debug_log(f"💬 Chat message in conversation {conv_id}: '{message[:50]}...'")

        # Get user's university from their account settings (MSU Lane prioritization)
        user_university = None
        try:
            # First, try to get university from user profile (account settings)
            from StudyFlow.backend.supabase_client import get_user_profile
            user_profile = get_user_profile(request.user_id)

            if user_profile and user_profile.get('university'):
                user_university = user_profile.get('university')
                debug_log(f"🎓 User's university from profile: {user_university}")
            else:
                # Fallback: Infer from their most common note university
                user_notes = supabase.table("notes").select("university").eq("user_id", request.user_id).execute()
                if user_notes.data:
                    from collections import Counter
                    universities = [n.get('university') for n in user_notes.data if n.get('university')]
                    if universities:
                        user_university = Counter(universities).most_common(1)[0][0]
                        debug_log(f"🎓 User's university inferred from notes: {user_university}")
        except Exception as e:
            debug_log(f"⚠️ Could not determine user university: {e}")

        # Generate embedding for the question
        query_embedding = generate_embedding(message)
        if not query_embedding:
            return jsonify({"error": "AI service is temporarily busy. Please try again in a moment.", "retry": True}), 503

        # Semantic cache check -- if a very similar question was asked before, return cached response
        from StudyFlow.backend.semantic_cache import check_semantic_cache, store_semantic_cache
        sem_cached = check_semantic_cache(request.user_id, query_embedding)
        if sem_cached:
            debug_log(f"[+] Semantic cache HIT ({sem_cached['similarity']}) for: '{message[:50]}...'")

            # Still need a conversation ID
            if not conv_id:
                conv_id = create_conversation(request.user_id, source=source)

            # Add messages to conversation
            add_message(conv_id, 'user', message)
            add_message(conv_id, 'assistant', sem_cached['response'], sem_cached['sources'])

            # SB 1324 Provenance Logging -- even for cached responses
            try:
                source_log_entries = []
                for s in (sem_cached['sources'] or []):
                    source_log_entries.append({
                        "note_id": s.get("note_id"),
                        "filename": s.get("filename"),
                        "similarity": s.get("similarity"),
                        "contributor_username": s.get("username"),
                        "source_type": "nexus_note"
                    })
                client_ip = request.headers.get('X-Forwarded-For', request.remote_addr)
                if client_ip and ',' in client_ip:
                    client_ip = client_ip.split(',')[0].strip()
                log_ai_response(
                    user_id=request.user_id,
                    conversation_id=conv_id,
                    prompt_text=message,
                    response_text=sem_cached['response'],
                    sources_used=source_log_entries,
                    model_used="semantic_cache",
                    response_time_ms=0,
                    ip_address=client_ip
                )
            except Exception as log_error:
                debug_log(f"[-] Provenance logging (cached) failed: {log_error}")

            # Track search
            try:
                from StudyFlow.backend.search_tracker import track_search
                track_search(message, len(sem_cached['sources']), source="chat")
            except:
                pass

            return jsonify({
                "conversation_id": conv_id,
                "response": sem_cached['response'],
                "sources": sem_cached['sources'],
                "followup_suggestions": []
            }), 200

        # Search database for relevant content (cached for 2 min for similar queries)
        rag_cache_key = f"rag:{hashlib.md5(message.encode()).hexdigest()}:{request.user_id}:{search_scope}"
        search_results = None

        if redis_cache:
            try:
                cached_rag = redis_cache.get(rag_cache_key)
                if cached_rag:
                    search_results = json.loads(cached_rag)
                    debug_log(f"RAG cache HIT: {len(search_results)} chunks")
            except:
                pass

        if search_results is None:
            search_results = search_notes_vector(
                query_embedding=query_embedding,
                user_id=request.user_id,
                university=None,  # Not filtering, just using for prioritization
                course_code=None,
                match_threshold=0.7,  # Raised threshold for relevant results only
                match_count=15
            )
            # Cache RAG results for 2 minutes
            if redis_cache and search_results:
                try:
                    redis_cache.setex(rag_cache_key, 600, json.dumps(search_results))
                except:
                    pass

        debug_log(f"Found {len(search_results) if search_results else 0} relevant chunks")

        # DEBUG: Show detailed search results
        if search_results:
            debug_log(f"🔍 SEARCH RESULTS DEBUG:")
            for i, r in enumerate(search_results[:5]):
                debug_log(f"  Result #{i+1}: similarity={r.get('similarity', 0):.3f}, note_id={r.get('note_id')}, content_preview={str(r.get('chunk_text', r.get('content_summary', '')))[:100]}...")
        else:
            debug_log(f"⚠️ NO SEARCH RESULTS FOUND - will fall back to general knowledge")

        # Prioritize results from same university (MSU Lane)
        if user_university and search_results:
            same_uni = [r for r in search_results if r.get('university') == user_university]
            diff_uni = [r for r in search_results if r.get('university') != user_university]
            search_results = same_uni + diff_uni
            debug_log(f"🎯 Prioritized {len(same_uni)} results from {user_university}")

        # Filter to personal notes only if requested
        if search_scope == 'personal' and search_results:
            note_ids = list(set(r['note_id'] for r in search_results))
            ownership = supabase.table("notes").select("id, user_id").in_("id", note_ids).execute()
            personal_note_ids = set(n['id'] for n in ownership.data if n['user_id'] == request.user_id)
            search_results = [r for r in search_results if r['note_id'] in personal_note_ids]
            debug_log(f"Filtered to {len(search_results)} personal chunks")

        # Apply voting-based ranking using pre-computed net_votes
        if search_results:
            # Batch fetch net_votes for all notes in results
            note_ids_in_results = list(set(r['note_id'] for r in search_results))
            net_votes_map = {}
            try:
                nv_resp = supabase.table("notes").select("id, net_votes").in_("id", note_ids_in_results).execute()
                for n in (nv_resp.data or []):
                    net_votes_map[n["id"]] = n.get("net_votes", 0)
            except:
                pass

            # Apply weighted ranking: Similarity (70%) + Helpfulness (30%)
            for result in search_results:
                note_id = result['note_id']
                net = net_votes_map.get(note_id, 0)

                # Normalize net votes to -1 to 1 range (capped at +/-20)
                capped = max(-20, min(20, net))
                helpfulness_score = capped / 20.0

                # Normalize to 0-1 range
                normalized_helpfulness = (helpfulness_score + 1) / 2

                # Combined score: 70% similarity + 30% helpfulness
                similarity = result.get('similarity', 0)
                combined_score = (similarity * 0.7) + (normalized_helpfulness * 0.3)

                result['combined_score'] = combined_score
                result['helpfulness_score'] = helpfulness_score

            # Re-sort by combined score
            search_results.sort(key=lambda x: x.get('combined_score', 0), reverse=True)

            debug_log(f"🎯 Applied voting-based ranking to {len(search_results)} results")
            for i, r in enumerate(search_results[:5]):
                debug_log(f"  #{i+1}: similarity={r.get('similarity', 0):.2f}, helpfulness={r.get('helpfulness_score', 0):.2f}, combined={r.get('combined_score', 0):.2f}, votes={r.get('vote_data')}")

        # Wikipedia priority: only use Wikipedia as backfill when student notes are thin
        if search_results:
            # Separate student notes from Wikipedia
            note_ids_all = list(set(r['note_id'] for r in search_results))
            wiki_check = supabase.table("notes").select("id, university").in_("id", note_ids_all).execute()
            wiki_note_ids = set(n['id'] for n in (wiki_check.data or []) if (n.get('university') or '').lower() == 'wikipedia')

            student_results = [r for r in search_results if r['note_id'] not in wiki_note_ids]
            wiki_results = [r for r in search_results if r['note_id'] in wiki_note_ids]

            # Only include Wikipedia if fewer than 2 quality student results
            quality_student = [r for r in student_results if r.get('similarity', 0) > 0.6]

            if len(quality_student) >= 2:
                search_results = student_results  # Drop Wikipedia entirely
                debug_log(f"Enough student notes ({len(quality_student)}), skipping Wikipedia")
            else:
                # Backfill with max 3 Wikipedia chunks
                search_results = student_results + wiki_results[:3]
                debug_log(f"Low student notes ({len(quality_student)}), backfilling with {min(len(wiki_results), 3)} Wikipedia chunks")

        # Add search result metadata and enrich with missing fields
        # FILTER: Only show sources with similarity >= search threshold to avoid irrelevant citations
        sources = []
        seen_note_ids = set()
        DISPLAY_THRESHOLD = 0.70  # Match search threshold - show all sources AI used
        if search_results:
            for result in search_results:  # Deduplicate by note_id
                if result['note_id'] in seen_note_ids:
                    continue

                # Skip sources below display threshold
                if result.get('similarity', 0) < DISPLAY_THRESHOLD:
                    debug_log(f"Skipping low-similarity source: {result.get('similarity', 0):.2f} < {DISPLAY_THRESHOLD}")
                    continue

                seen_note_ids.add(result['note_id'])
                # REMOVED 3-source limit to show consensus from multiple students
                note_result = supabase.table("notes").select("original_filename, user_id, university, course_code").eq("id", result['note_id']).execute()
                note_data = note_result.data[0] if note_result.data else {}
                filename = note_data.get('original_filename', 'Unknown')

                # Enrich result with username if missing from chunk
                if not result.get('username') and note_data.get('user_id'):
                    try:
                        profile_result = supabase.table("user_profiles").select("username").eq("id", note_data['user_id']).execute()
                        if profile_result.data and profile_result.data[0].get("username"):
                            result['username'] = profile_result.data[0]['username']
                    except Exception:
                        pass

                # Enrich university/course_code from notes table if missing from chunk
                if not result.get('university') and note_data.get('university'):
                    result['university'] = note_data['university']
                if not result.get('course_code') and note_data.get('course_code'):
                    result['course_code'] = note_data['course_code']

                # Detect Wikipedia sources and build URL
                is_wikipedia = result.get('university') == 'Wikipedia' or (note_data.get('university') == 'Wikipedia')
                wikipedia_url = None
                if is_wikipedia and filename.endswith('.txt'):
                    article_title = filename.replace('.txt', '')
                    wikipedia_url = f"https://en.wikipedia.org/wiki/{article_title.replace(' ', '_')}"

                sources.append({
                    "note_id": result['note_id'],
                    "filename": f"{filename} ({result['university']} - {result['course_code']})" if result.get('university') and not is_wikipedia else filename,
                    "similarity": round(result['similarity'], 2),
                    "username": result.get('username'),
                    "university": result.get('university'),
                    "course_code": result.get('course_code'),
                    "from_wikipedia": is_wikipedia,
                    "wikipedia_url": wikipedia_url
                })
                # Add original_filename to result for context
                result['original_filename'] = filename

        # Add user message to conversation
        add_message(conv_id, 'user', message)

        # Get conversation history from database
        conversation_history = get_conversation_messages(conv_id, request.user_id)

        # Generate conversational AI response
        ai_result = generate_conversational_response(
            question=message,
            search_results=search_results or [],
            conversation_history=conversation_history
        )

        ai_response = ai_result["response"]
        model_used = ai_result["model_used"]
        response_time_ms = ai_result["response_time_ms"]
        followup_suggestions = ai_result.get("followup_suggestions", [])

        # Add AI response to conversation
        add_message(conv_id, 'assistant', ai_response, sources)

        # SB 1324 Provenance Logging - log every AI response for 7-year retention
        try:
            source_log_entries = []
            for s in sources:
                source_log_entries.append({
                    "note_id": s.get("note_id"),
                    "filename": s.get("filename"),
                    "similarity": s.get("similarity"),
                    "contributor_username": s.get("username"),
                    "source_type": "nexus_note"
                })

            client_ip = request.headers.get('X-Forwarded-For', request.remote_addr)
            if client_ip and ',' in client_ip:
                client_ip = client_ip.split(',')[0].strip()

            log_ai_response(
                user_id=request.user_id,
                conversation_id=conv_id,
                prompt_text=message,
                response_text=ai_response,
                sources_used=source_log_entries,
                model_used=model_used,
                response_time_ms=response_time_ms,
                ip_address=client_ip
            )
        except Exception as log_error:
            debug_log(f"[-] Provenance logging failed (non-blocking): {log_error}")

        # Notify note owners in background (Celery) -- don't block the response
        if sources:
            try:
                from StudyFlow.backend.tasks import send_citation_notifications
                send_citation_notifications.delay(sources, request.user_id)
            except Exception as celery_err:
                debug_log(f"[-] Failed to queue citation notifications: {celery_err}")

        # Generate title for new conversations (if title is still None)
        if conversation and not conversation.get('title'):
            try:
                title = generate_conversation_title(message, ai_response)
                update_conversation_title(conv_id, title)
                debug_log(f"[+] Auto-generated title: {title}")
            except Exception as title_error:
                debug_log(f"[!] Failed to generate title: {title_error}")

        debug_log(f"[+] Generated conversational response ({len(ai_response)} chars, {response_time_ms}ms)")

        # Store in semantic cache for future similar questions
        try:
            store_semantic_cache(request.user_id, message, query_embedding, ai_response, sources)
        except:
            pass

        # Track chat search
        try:
            from StudyFlow.backend.search_tracker import track_search
            track_search(message, len(sources), source="chat")
        except:
            pass

        return jsonify({
            "conversation_id": conv_id,
            "response": ai_response,
            "sources": sources,
            "followup_suggestions": followup_suggestions
        }), 200

    except Exception as e:
        debug_log(f"Chat error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/conversations", methods=["GET"])
@supabase_auth_required
def list_conversations():
    """
    List all conversations for the authenticated user

    Returns:
    {
        "conversations": [
            {
                "id": "uuid",
                "title": "First question snippet...",
                "created_at": "2026-03-21T10:00:00Z",
                "updated_at": "2026-03-21T10:05:00Z"
            }
        ]
    }
    """
    try:
        cache_key = f"convos_list:{request.user_id}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        from StudyFlow.backend.conversational_noteflow import list_user_conversations
        conversations = list_user_conversations(request.user_id, limit=50)
        resp = {"conversations": conversations}
        if redis_cache:
            try: redis_cache.setex(cache_key, 120, json.dumps(resp))
            except: pass
        return jsonify(resp), 200

    except Exception as e:
        debug_log(f"❌ Error listing conversations: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/canvas/preload", methods=["GET"])
@supabase_auth_required
def canvas_preload():
    """Preload entire canvas state: folders, conversations with messages, sticky notes. Cached in Redis per user."""
    try:
        user_id = request.user_id
        cache_key = f"canvas_preload:{user_id}"

        # Check Redis cache (2 min TTL)
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached:
                    return jsonify(json.loads(cached)), 200
            except:
                pass

        # Fetch everything in parallel-ish (exclude plugin conversations)
        folders_resp = supabase.table("chat_folders").select("*").eq("user_id", user_id).order("created_at").execute()
        convos_resp = supabase.table("conversations").select(
            "id, title, folder_id, updated_at, canvas_layout"
        ).eq("user_id", user_id).eq("source", "chat").is_("deleted_at", "null").order("updated_at", desc=True).execute()

        stickies_resp = supabase.table("canvas_sticky_notes").select("*").eq("user_id", user_id).execute()

        # Fetch messages for all conversations
        convos = convos_resp.data or []
        conversations_with_messages = []

        for conv in convos:
            msgs_resp = supabase.table("conversation_messages").select(
                "content, role, sources, created_at"
            ).eq("conversation_id", conv["id"]).order("created_at").execute()

            conversations_with_messages.append({
                "id": conv["id"],
                "title": conv.get("title"),
                "folder_id": conv.get("folder_id"),
                "updated_at": conv.get("updated_at"),
                "canvas_layout": conv.get("canvas_layout"),
                "messages": msgs_resp.data or []
            })

        result = {
            "folders": folders_resp.data or [],
            "conversations": conversations_with_messages,
            "sticky_notes": stickies_resp.data or []
        }

        # Cache for 10 minutes
        if redis_cache:
            try:
                redis_cache.setex(cache_key, 600, json.dumps(result))
            except:
                pass

        return jsonify(result), 200

    except Exception as e:
        debug_log(f"Canvas preload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/canvas/invalidate", methods=["POST"])
@supabase_auth_required
def canvas_invalidate():
    """Invalidate the canvas preload cache for current user."""
    try:
        if redis_cache:
            redis_cache.delete(f"canvas_preload:{request.user_id}")
        return jsonify({"success": True}), 200
    except:
        return jsonify({"success": True}), 200


@app.route("/api/flashcards/generate", methods=["POST"])
@supabase_auth_required
def generate_flashcards():
    """Generate quiz questions. Searches Nexus first, falls back to general knowledge."""
    try:
        import google.generativeai as genai
        from StudyFlow.backend.embedding_client import generate_embedding
        from StudyFlow.backend.supabase_client import search_notes_vector

        data = request.get_json()
        topic = data.get("topic", "").strip()
        count = data.get("count", 15)

        if not topic:
            return jsonify({"error": "Topic required"}), 400

        # Search Nexus for relevant notes
        note_context = ""
        sources_used = []
        try:
            query_embedding = generate_embedding(topic)
            if query_embedding:
                results = search_notes_vector(
                    query_embedding=query_embedding,
                    user_id=request.user_id,
                    match_threshold=0.5,
                    match_count=10
                )
                if results:
                    chunks = []
                    for r in results[:8]:
                        text = r.get('content_summary') or r.get('chunk_text', '')
                        if text:
                            chunks.append(text[:500])
                            sources_used.append(r.get('note_id'))
                    if chunks:
                        note_context = "\n\n".join(chunks)
                        debug_log(f"Flashcards: found {len(chunks)} relevant note chunks for '{topic}'")
        except Exception as search_err:
            debug_log(f"Flashcard Nexus search failed (using general knowledge): {search_err}")

        # Check Redis cache
        has_notes = len(note_context) > 0
        cache_key = f"flashcards:{hashlib.md5((topic + str(has_notes)).encode()).hexdigest()}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached:
                    debug_log(f"Flashcard cache HIT for '{topic}' (notes={has_notes})")
                    return jsonify({"questions": cached, "used_notes": has_notes, "source_count": len(sources_used), "cached": True}), 200
            except:
                pass

        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        model = genai.GenerativeModel('gemini-3.1-flash-lite-preview')

        if note_context:
            prompt = f"""Based on the following study notes, generate exactly {count} quiz questions about "{topic}".

STUDY NOTES:
{note_context}

Generate questions that test understanding of the material in these notes.
Each question has exactly two answer choices where only one is correct.

Return ONLY a valid JSON array with no other text:
[{{"question":"What is...?","left":"Answer A","right":"Answer B","correct":"left"}}]

The "correct" field must be either "left" or "right".
No explanation, no markdown, just the JSON array."""
        else:
            prompt = f"""Generate exactly {count} quiz questions about "{topic}".
Each question has exactly two answer choices where only one is correct.
Mix multiple choice and true/false style questions.

Return ONLY a valid JSON array with no other text:
[{{"question":"What is...?","left":"Answer A","right":"Answer B","correct":"left"}}]

The "correct" field must be either "left" or "right".
No explanation, no markdown, just the JSON array."""

        response = model.generate_content(
            prompt,
            generation_config=genai.GenerationConfig(temperature=0.7, max_output_tokens=2000)
        )

        text = response.text.strip()

        try:
            from StudyFlow.backend.cost_tracker import track_ai_call
            track_ai_call("gemini", "flash-lite", "flashcard_gen")
        except:
            pass

        cleaned = text.replace('```json', '').replace('```', '').strip()

        # Cache for 1 hour
        if redis_cache:
            try:
                redis_cache.setex(cache_key, 3600, cleaned)
            except:
                pass

        return jsonify({
            "questions": cleaned,
            "used_notes": has_notes,
            "source_count": len(sources_used)
        }), 200

    except Exception as e:
        debug_log(f"Flashcard generation error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat-folders", methods=["GET"])
@supabase_auth_required
def list_chat_folders():
    """List user's chat folders."""
    try:
        resp = supabase.table("chat_folders").select("*").eq("user_id", request.user_id).order("created_at").execute()
        return jsonify({"folders": resp.data or []}), 200
    except Exception as e:
        return jsonify({"folders": []}), 200


@app.route("/api/chat-folders", methods=["POST"])
@supabase_auth_required
def create_chat_folder():
    """Create a new chat folder (supports nesting via parent_id)."""
    try:
        data = request.get_json()
        row = {
            "user_id": request.user_id,
            "name": data.get("name", "New Folder"),
            "color": data.get("color", "#7c9885"),
            "position_x": data.get("position_x", 0),
            "position_y": data.get("position_y", 0)
        }
        if data.get("parent_id"):
            row["parent_id"] = data["parent_id"]
        folder = supabase.table("chat_folders").insert(row).execute()
        if redis_cache:
            try: redis_cache.delete(f"canvas_preload:{request.user_id}")
            except: pass
        return jsonify(folder.data[0] if folder.data else {}), 201
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat-folders/<folder_id>", methods=["PUT"])
@supabase_auth_required
def update_chat_folder(folder_id):
    """Update a chat folder (name, color, position)."""
    try:
        data = request.get_json()
        update = {}
        if "name" in data: update["name"] = data["name"]
        if "color" in data: update["color"] = data["color"]
        if "icon" in data: update["icon"] = data["icon"]
        if "position_x" in data: update["position_x"] = data["position_x"]
        if "position_y" in data: update["position_y"] = data["position_y"]
        if "parent_id" in data: update["parent_id"] = data["parent_id"]
        supabase.table("chat_folders").update(update).eq("id", folder_id).eq("user_id", request.user_id).execute()
        if redis_cache:
            try: redis_cache.delete(f"canvas_preload:{request.user_id}")
            except: pass
        return jsonify({"success": True}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat-folders/<folder_id>", methods=["DELETE"])
@supabase_auth_required
def delete_chat_folder(folder_id):
    """Delete a chat folder (conversations become unfiled)."""
    try:
        # Unfiled all conversations in this folder
        supabase.table("conversations").update({"folder_id": None}).eq("folder_id", folder_id).eq("user_id", request.user_id).execute()
        supabase.table("chat_folders").delete().eq("id", folder_id).eq("user_id", request.user_id).execute()
        if redis_cache:
            try: redis_cache.delete(f"canvas_preload:{request.user_id}")
            except: pass
        return jsonify({"success": True}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/conversations/<conversation_id>/move", methods=["PUT"])
@supabase_auth_required
def move_conversation(conversation_id):
    """Move a conversation to a folder (or null to unfiled)."""
    try:
        data = request.get_json()
        folder_id = data.get("folder_id")  # null = unfiled
        supabase.table("conversations").update({"folder_id": folder_id}).eq("id", conversation_id).eq("user_id", request.user_id).execute()
        if redis_cache:
            try: redis_cache.delete(f"canvas_preload:{request.user_id}")
            except: pass
        return jsonify({"success": True}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/sticky-notes", methods=["GET"])
@supabase_auth_required
def list_sticky_notes():
    """List sticky notes for a folder (or home if no folder_id)."""
    try:
        folder_id = request.args.get("folder_id")
        query = supabase.table("canvas_sticky_notes").select("*").eq("user_id", request.user_id)
        if folder_id:
            query = query.eq("folder_id", folder_id)
        else:
            query = query.is_("folder_id", "null")
        resp = query.execute()
        return jsonify({"notes": resp.data or []}), 200
    except Exception as e:
        return jsonify({"notes": []}), 200


@app.route("/api/sticky-notes", methods=["POST"])
@supabase_auth_required
def create_sticky_note():
    """Create a sticky note."""
    try:
        data = request.get_json()
        row = {
            "user_id": request.user_id,
            "content": data.get("content", ""),
            "position_x": data.get("position_x", 0),
            "position_y": data.get("position_y", 0)
        }
        if data.get("folder_id"):
            row["folder_id"] = data["folder_id"]
        resp = supabase.table("canvas_sticky_notes").insert(row).execute()
        return jsonify(resp.data[0] if resp.data else {}), 201
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/sticky-notes/<note_id>", methods=["PUT"])
@supabase_auth_required
def update_sticky_note(note_id):
    """Update a sticky note (content, position)."""
    try:
        data = request.get_json()
        update = {}
        if "content" in data: update["content"] = data["content"]
        if "position_x" in data: update["position_x"] = data["position_x"]
        if "position_y" in data: update["position_y"] = data["position_y"]
        supabase.table("canvas_sticky_notes").update(update).eq("id", note_id).eq("user_id", request.user_id).execute()
        return jsonify({"success": True}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/sticky-notes/<note_id>", methods=["DELETE"])
@supabase_auth_required
def delete_sticky_note(note_id):
    """Delete a sticky note."""
    try:
        supabase.table("canvas_sticky_notes").delete().eq("id", note_id).eq("user_id", request.user_id).execute()
        return jsonify({"success": True}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat-folders/<folder_id>/conversations", methods=["GET"])
@supabase_auth_required
def list_chat_folder_conversations(folder_id):
    """List conversations in a folder."""
    try:
        resp = supabase.table("conversations").select(
            "id, title, updated_at, canvas_layout"
        ).eq("user_id", request.user_id).eq("folder_id", folder_id).is_("deleted_at", "null").order("updated_at", desc=True).execute()
        return jsonify({"conversations": resp.data or []}), 200
    except Exception as e:
        return jsonify({"conversations": []}), 200


@app.route("/api/conversations/unfiled", methods=["GET"])
@supabase_auth_required
def list_unfiled_conversations():
    """List conversations not in any folder."""
    try:
        resp = supabase.table("conversations").select(
            "id, title, updated_at, canvas_layout"
        ).eq("user_id", request.user_id).eq("source", "chat").is_("folder_id", "null").is_("deleted_at", "null").order("updated_at", desc=True).execute()
        return jsonify({"conversations": resp.data or []}), 200
    except Exception as e:
        return jsonify({"conversations": []}), 200


@app.route("/api/conversations/<conversation_id>/layout", methods=["GET"])
@supabase_auth_required
def get_canvas_layout(conversation_id):
    """Get saved canvas layout for a conversation."""
    try:
        resp = supabase.table("conversations").select("canvas_layout").eq("id", conversation_id).eq("user_id", request.user_id).single().execute()
        if resp.data and resp.data.get('canvas_layout'):
            return jsonify(resp.data['canvas_layout']), 200
        return jsonify(None), 200
    except Exception as e:
        return jsonify(None), 200


@app.route("/api/conversations/<conversation_id>/layout", methods=["PUT"])
@supabase_auth_required
def save_canvas_layout(conversation_id):
    """Save canvas layout for a conversation."""
    try:
        data = request.get_json()
        supabase.table("conversations").update({
            "canvas_layout": data
        }).eq("id", conversation_id).eq("user_id", request.user_id).execute()

        # Invalidate Redis cache so next preload gets fresh data
        if redis_cache:
            try:
                cache_key = f"canvas:{request.user_id}"
                redis_cache.delete(cache_key)
                debug_log(f"Invalidated canvas cache for user {request.user_id}")
            except Exception as cache_err:
                debug_log(f"Cache invalidation warning: {cache_err}")

        return jsonify({"success": True}), 200
    except Exception as e:
        debug_log(f"Save layout error: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/conversations/<conversation_id>/messages", methods=["GET"])
@supabase_auth_required
def get_conversation_history(conversation_id):
    """
    Get all messages in a specific conversation

    Returns:
    {
        "messages": [
            {
                "id": "uuid",
                "role": "user",
                "content": "What is DNA?",
                "created_at": "2026-03-21T10:00:00Z"
            },
            {
                "id": "uuid",
                "role": "assistant",
                "content": "DNA is...",
                "sources": [...],
                "created_at": "2026-03-21T10:00:05Z"
            }
        ]
    }
    """
    try:
        from StudyFlow.backend.conversational_noteflow import get_conversation_messages

        # get_conversation_messages already verifies ownership
        messages = get_conversation_messages(conversation_id, request.user_id)

        if messages is None:
            return jsonify({"error": "Conversation not found or access denied"}), 404

        return jsonify({"messages": messages}), 200

    except Exception as e:
        debug_log(f"❌ Error getting conversation history: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/conversations/<conversation_id>", methods=["DELETE"])
@supabase_auth_required
def delete_conversation_endpoint(conversation_id):
    """
    Delete a conversation and all its messages

    Returns:
    {
        "success": true
    }
    """
    try:
        from StudyFlow.backend.conversational_noteflow import delete_conversation

        success = delete_conversation(conversation_id, request.user_id)

        # Invalidate canvas preload cache
        if redis_cache:
            try:
                redis_cache.delete(f"canvas_preload:{request.user_id}")
            except:
                pass

        # Return 200 even if already deleted (idempotent)
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"❌ Error deleting conversation: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat/rate-response", methods=["POST", "OPTIONS"])
@supabase_auth_required
def rate_chat_response():
    """
    Rate the helpfulness of sources cited in an AI response

    Expects:
    {
        "message_id": "uuid",
        "conversation_id": "uuid",
        "vote": 1 or -1,  # 1 = helpful (upvote), -1 = not helpful (downvote)
        "cited_note_ids": ["note_id1", "note_id2", ...]  # Notes that were cited
    }

    Returns:
    {
        "success": true,
        "vote": 1
    }
    """
    # Handle OPTIONS preflight
    if request.method == 'OPTIONS':
        return '', 200

    try:
        data = request.get_json()
        message_id = data.get('message_id')
        conversation_id = data.get('conversation_id')
        vote = data.get('vote')  # 1 or -1
        cited_note_ids = data.get('cited_note_ids', [])

        if not all([message_id, conversation_id, vote in [1, -1]]):
            return jsonify({"error": "Missing required fields"}), 400

        # Check if user already voted on this response
        existing = supabase.table("ai_response_ratings").select("id, vote").eq("message_id", message_id).eq("user_id", request.user_id).execute()

        if existing.data:
            # Update existing vote
            supabase.table("ai_response_ratings").update({
                "vote": vote
            }).eq("id", existing.data[0]['id']).execute()

            debug_log(f"[+] Updated response rating: message={message_id}, user={request.user_id}, vote={vote}")
        else:
            # Create new vote
            supabase.table("ai_response_ratings").insert({
                "message_id": message_id,
                "conversation_id": conversation_id,
                "user_id": request.user_id,
                "vote": vote,
                "cited_note_ids": cited_note_ids
            }).execute()

            debug_log(f"[+] Created response rating: message={message_id}, user={request.user_id}, vote={vote}")

        # Trigger async vote count update for cited notes
        if cited_note_ids:
            try:
                update_note_votes.delay(cited_note_ids)
            except Exception as celery_err:
                debug_log(f"[!] Failed to queue vote update: {celery_err}")

        return jsonify({"success": True, "vote": vote}), 200

    except Exception as e:
        debug_log(f"❌ Error rating response: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/chat/get-my-vote", methods=["GET", "OPTIONS"])
@supabase_auth_required
def get_my_response_vote():
    """
    Get user's existing vote for a message

    Query params: message_id

    Returns:
    {
        "vote": 1 or -1 or null
    }
    """
    # Handle OPTIONS preflight
    if request.method == 'OPTIONS':
        return '', 200

    try:
        message_id = request.args.get('message_id')
        if not message_id:
            return jsonify({"error": "Missing message_id"}), 400

        result = supabase.table("ai_response_ratings").select("vote").eq("message_id", message_id).eq("user_id", request.user_id).execute()

        if result.data:
            return jsonify({"vote": result.data[0]['vote']}), 200
        else:
            return jsonify({"vote": None}), 200

    except Exception as e:
        debug_log(f"❌ Error getting vote: {e}\n{traceback.format_exc()}")
        return jsonify({"vote": None}), 200


@app.route("/api/settings/collective-brain", methods=["GET"])
@supabase_auth_required
def get_collective_brain_setting():
    """
    Get the user's collective brain opt-in setting.
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        user = supabase.table("user_profiles").select("collective_brain_opt_in").eq("id", request.user_id).single().execute()

        if not user.data:
            return jsonify({"error": "User not found"}), 404

        return jsonify({"enabled": user.data.get('collective_brain_opt_in', True)}), 200

    except Exception as e:
        debug_log(f"Error getting collective brain setting: {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/settings/collective-brain", methods=["POST"])
@supabase_auth_required
def update_collective_brain_setting():
    """
    Update the user's collective brain opt-in setting.

    If disabled (opt-out):
    - User can ONLY see their own notes
    - User's notes are NOT shared with others

    If enabled (opt-in):
    - User can see their own notes + everyone else's notes
    - User's notes are shared with everyone
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        enabled = data.get('enabled', True)

        # Update user setting
        supabase.table("user_profiles").update({
            "collective_brain_opt_in": enabled
        }).eq("id", request.user_id).execute()

        debug_log(f"✅ User {request.user_id} set collective_brain_opt_in to {enabled}")

        return jsonify({"success": True, "enabled": enabled}), 200

    except Exception as e:
        debug_log(f"❌ Error updating collective brain setting: {e}")
        return jsonify({"error": str(e)}), 500


# ============ FOLDER ENDPOINTS ============

@app.route("/api/folders/list", methods=["GET"])
@supabase_auth_required
def list_note_folders():
    """Get all folders for current user"""
    try:
        from StudyFlow.backend.supabase_client import get_user_folders

        folders = get_user_folders(request.user_id)
        return jsonify(folders), 200

    except Exception as e:
        debug_log(f"Error listing folders: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/folders/create", methods=["POST"])
@supabase_auth_required
def create_folder_endpoint():
    """
    Create a new folder.
    Request body: {
        "name": "Folder Name",
        "parent_id": "uuid" (optional),
        "color": "#7c9885" (optional),
        "position_data": {"x": 20, "y": 20} (optional)
    }
    """
    try:
        from StudyFlow.backend.supabase_client import create_folder

        data = request.get_json()
        name = data.get('name')

        if not name:
            return jsonify({"error": "Folder name is required"}), 400

        folder = create_folder(
            user_id=request.user_id,
            name=name,
            parent_id=data.get('parent_id'),
            color=data.get('color', '#7c9885'),
            position_data=data.get('position_data')
        )

        if folder:
            return jsonify(folder), 201
        else:
            return jsonify({"error": "Failed to create folder"}), 500

    except Exception as e:
        debug_log(f"Error creating folder: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/folders/<folder_id>", methods=["PUT"])
@supabase_auth_required
def update_folder_endpoint(folder_id):
    """
    Update a folder (rename, move, change color, update position).
    Request body: {
        "name": "New Name" (optional),
        "parent_id": "uuid" (optional),
        "color": "#7c9885" (optional),
        "position_data": {"x": 100, "y": 200} (optional)
    }
    """
    try:
        from StudyFlow.backend.supabase_client import update_folder

        data = request.get_json()

        # Build updates dict from provided fields
        updates = {}
        if 'name' in data:
            updates['name'] = data['name']
        if 'parent_id' in data:
            updates['parent_id'] = data['parent_id']
        if 'color' in data:
            updates['color'] = data['color']
        if 'position_data' in data:
            updates['position_data'] = data['position_data']

        if not updates:
            return jsonify({"error": "No updates provided"}), 400

        folder = update_folder(folder_id, request.user_id, updates)

        if folder:
            return jsonify(folder), 200
        else:
            return jsonify({"error": "Folder not found or unauthorized"}), 404

    except Exception as e:
        debug_log(f"Error updating folder: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/folders/<folder_id>", methods=["DELETE"])
@supabase_auth_required
def delete_folder_endpoint(folder_id):
    """Delete a folder (CASCADE moves children to parent)"""
    try:
        from StudyFlow.backend.supabase_client import delete_folder

        success = delete_folder(folder_id, request.user_id)

        if success:
            return jsonify({"success": True}), 200
        else:
            return jsonify({"error": "Folder not found or unauthorized"}), 404

    except Exception as e:
        debug_log(f"Error deleting folder: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/folder", methods=["PUT"])
@supabase_auth_required
def move_note_to_folder_endpoint(note_id):
    """
    Assign a note to a folder (or remove from folder).
    Request body: {
        "folder_id": "uuid" (or null to remove from folder)
    }
    """
    try:
        from StudyFlow.backend.supabase_client import update_note_folder

        data = request.get_json()
        folder_id = data.get('folder_id')  # Can be None

        success = update_note_folder(note_id, request.user_id, folder_id)

        if success:
            return jsonify({"success": True}), 200
        else:
            return jsonify({"error": "Note not found or unauthorized"}), 404

    except Exception as e:
        debug_log(f"Error moving note to folder: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500




@app.route("/api/notes/browse", methods=["GET"])
@supabase_auth_required
def browse_notes():
    """
    Browse public notes in the Nexus (requires .edu verification)

    Query params:
    - university: Filter by university
    - course_code: Filter by course
    - topic_tags: Comma-separated tags
    - sort: 'recent', 'popular', 'views'
    - limit: Default 50

    Returns:
    {
        "notes": [...]
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        user_id = request.user_id

        # Check Redis cache for this browse request
        university_param = request.args.get('university', '')
        sort_param = request.args.get('sort', 'recent')
        limit_param = request.args.get('limit', '50')
        offset_param = request.args.get('offset', '0')
        browse_cache_key = f"browse_api:{university_param}:{sort_param}:{limit_param}:{offset_param}"

        if redis_cache:
            try:
                cached = redis_cache.get(browse_cache_key)
                if cached:
                    debug_log(f"Browse API cache HIT")
                    return jsonify(json.loads(cached)), 200
            except:
                pass

        # Check .edu email verification requirement
        user_profile = supabase.table("user_profiles").select("edu_email_verified").eq("id", user_id).single().execute()
        if user_profile.data:
            edu_verified = user_profile.data.get('edu_email_verified', False)
            if not edu_verified:
                return jsonify({"error": "Browse requires a verified .edu email address"}), 403
        else:
            return jsonify({"error": "User profile not found"}), 404

        # Get query parameters
        university = request.args.get('university')
        course_code = request.args.get('course_code')
        topic_tags_str = request.args.get('topic_tags', '')
        topic_tags = [tag.strip() for tag in topic_tags_str.split(',') if tag.strip()] if topic_tags_str else []
        sort_by = request.args.get('sort', 'recent')
        limit = int(request.args.get('limit', 50))
        offset = int(request.args.get('offset', 0))

        # Build query -- includes pre-computed view/download counts
        query = supabase.table("notes").select(
            "id, original_filename, university, course_code, user_id, topic_tags, page_count, uploaded_at, view_count, download_count"
        ).eq("is_public", True).not_.is_("user_id", "null")

        if university:
            query = query.eq("university", university)
        if course_code:
            query = query.eq("course_code", course_code)

        # Execute query with pagination
        response = query.order("uploaded_at", desc=True).range(offset, offset + limit - 1).execute()
        notes = response.data if response.data else []

        # Filter notes by user's Nexus setting and get usernames (cached)
        filtered_notes = []
        for note in notes:
            profile = get_cached_profile(note['user_id'])
            if not profile or not profile.get('is_public', True):
                continue
            note['username'] = profile.get('username', 'Anonymous')
            filtered_notes.append(note)

        # Read pre-computed counts (updated by Celery every 30 min)
        for note in filtered_notes:
            note['view_count'] = note.get('view_count', 0) or 0
            note['usage_count'] = note.get('download_count', 0) or 0
            note.pop('download_count', None)
            note['filename'] = note.pop('original_filename', note.pop('filename', 'Unknown'))
            note.pop('user_id', None)

        notes = filtered_notes

        # Filter by tags if specified
        if topic_tags:
            notes = [n for n in notes if n.get('topic_tags') and any(tag in n['topic_tags'] for tag in topic_tags)]

        # Sort
        if sort_by == 'popular':
            notes.sort(key=lambda x: x.get('usage_count', 0), reverse=True)
        elif sort_by == 'views':
            notes.sort(key=lambda x: x.get('view_count', 0), reverse=True)
        # 'recent' is already sorted by created_at desc

        debug_log(f"Browse: Found {len(notes)} public notes")

        response_data = {"notes": notes}

        # Cache browse results for 5 minutes
        if redis_cache:
            try:
                redis_cache.setex(browse_cache_key, BROWSE_CACHE_TTL, json.dumps(response_data))
            except:
                pass

        return jsonify(response_data), 200

    except Exception as e:
        error_trace = traceback.format_exc()
        debug_log(f"❌ Browse notes error: {e}\n{error_trace}")
        # Return detailed error for debugging
        return jsonify({
            "error": str(e),
            "traceback": error_trace,
            "type": type(e).__name__
        }), 500


def cache_top_upvoted_notes():
    """
    Calculate and cache the top 10 most upvoted notes in Redis
    Called periodically or when cache expires

    Returns: list of top 10 note dicts
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        debug_log("🔄 Calculating top 10 upvoted notes...")

        # Get all ratings and count net votes per note
        ratings = supabase.table("ai_response_ratings").select("vote, cited_note_ids").execute()

        vote_counts = {}  # note_id -> net upvotes
        if ratings.data:
            for rating in ratings.data:
                cited_note_ids = rating.get('cited_note_ids', [])
                vote = rating.get('vote', 0)
                for note_id in cited_note_ids:
                    if note_id not in vote_counts:
                        vote_counts[note_id] = 0
                    vote_counts[note_id] += vote

        # Sort by net upvotes and get top 10
        sorted_notes = sorted(vote_counts.items(), key=lambda x: x[1], reverse=True)[:10]

        if not sorted_notes:
            debug_log("⚠️ No upvoted notes found, using fallback")
            # Fallback: Get 10 most recent public notes
            recent_notes = supabase.table("notes").select(
                "id, original_filename, university, course_code, user_id, uploaded_at"
            ).eq("is_public", True).not_("user_id", "is", None).order("uploaded_at", desc=True).limit(10).execute()

            if recent_notes.data:
                sorted_notes = [(note['id'], 0) for note in recent_notes.data]

        # Fetch full note data for top notes
        top_notes = []
        for note_id, net_upvotes in sorted_notes:
            note_data = supabase.table("notes").select(
                "id, original_filename, university, course_code, user_id, uploaded_at"
            ).eq("id", note_id).single().execute()

            if not note_data.data:
                continue

            note = note_data.data

            # Get username and check if Nexus is enabled
            try:
                user_profile = supabase.table("user_profiles").select("username, is_public").eq("id", note['user_id']).single().execute()
                if not user_profile.data or not user_profile.data.get('is_public', True):
                    continue  # Skip if user has Nexus disabled
                username = user_profile.data.get('username') or 'Anonymous'
            except:
                continue

            # Get content preview from first chunk
            content_preview = "Study notes and materials"
            try:
                chunks = supabase.table("note_chunks").select("chunk_text, content_summary").eq("note_id", note_id).limit(1).execute()
                if chunks.data and len(chunks.data) > 0:
                    chunk = chunks.data[0]
                    content_preview = chunk.get('content_summary') or chunk.get('chunk_text', '')
            except:
                pass

            top_notes.append({
                "note_id": note_id,
                "filename": note['original_filename'],
                "username": username,
                "university": note.get('university', 'Unknown'),
                "course_code": note.get('course_code', ''),
                "created_at": note.get('uploaded_at', ''),
                "content": content_preview,
                "upvotes": net_upvotes,
                "similarity": 0.8  # High similarity for popular notes
            })

        # Cache in Redis for 30 minutes
        if redis_cache and top_notes:
            redis_cache.setex(
                "browse:top_notes",
                BROWSE_CACHE_TTL,
                json.dumps(top_notes)
            )
            debug_log(f"✅ Cached {len(top_notes)} top upvoted notes in Redis")

        return top_notes

    except Exception as e:
        debug_log(f"❌ Error caching top notes: {e}\n{traceback.format_exc()}")
        return []


@app.route("/api/notes/top-browse", methods=["GET"])
@supabase_auth_required
def get_top_browse_notes():
    """
    Get top 10 most upvoted notes for browse page initial load
    Uses Redis cache for fast response

    Returns same format as semantic-browse
    """
    try:
        # Check .edu email verification requirement
        user_profile = supabase.table("user_profiles").select("edu_email_verified").eq("id", request.user_id).single().execute()
        if user_profile.data:
            edu_verified = user_profile.data.get('edu_email_verified', False)
            if not edu_verified:
                return jsonify({"error": "Browse requires a verified .edu email address"}), 403
        else:
            return jsonify({"error": "User profile not found"}), 404

        # Try to get from Redis cache first
        top_notes = None
        if redis_cache:
            try:
                cached_data = redis_cache.get("browse:top_notes")
                if cached_data:
                    top_notes = json.loads(cached_data)
                    debug_log(f"✅ Retrieved {len(top_notes)} top notes from Redis cache")
            except Exception as e:
                debug_log(f"⚠️ Redis cache read error: {e}")

        # If not in cache or cache miss, calculate and cache
        if not top_notes:
            debug_log("📊 Cache miss, calculating top upvoted notes...")
            top_notes = cache_top_upvoted_notes()

        if not top_notes:
            return jsonify({"notes": [], "has_more": False, "total": 0}), 200

        return jsonify({
            "notes": top_notes,
            "has_more": False,
            "total": len(top_notes)
        }), 200

    except Exception as e:
        debug_log(f"❌ Get top browse notes error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/semantic-browse", methods=["POST"])
@supabase_auth_required
def semantic_browse_notes():
    """
    Semantic search for browse page - returns note metadata + matching paragraphs

    Expects JSON:
    {
        "question": "What is photosynthesis?"
    }

    Returns:
    [
        {
            "note_id": "...",
            "filename": "Biology_Chapter_3.pdf",
            "username": "john_doe",
            "university": "MIT",
            "course_code": "BIO101",
            "created_at": "2024-01-15",
            "content": "Photosynthesis is the process...",
            "similarity": 0.89
        }
    ]
    """
    try:
        from StudyFlow.backend.supabase_client import search_notes_vector, supabase
        from StudyFlow.backend.embedding_client import generate_embedding

        # Check .edu email verification requirement
        user_profile = supabase.table("user_profiles").select("edu_email_verified").eq("id", request.user_id).single().execute()
        if user_profile.data:
            edu_verified = user_profile.data.get('edu_email_verified', False)
            if not edu_verified:
                return jsonify({"error": "Browse requires a verified .edu email address"}), 403
        else:
            return jsonify({"error": "User profile not found"}), 404

        data = request.get_json()
        if not data or not data.get('question'):
            return jsonify({"error": "Missing question"}), 400

        question = data.get('question')
        university_filter = data.get('university')  # Optional university filter
        course_code_filter = data.get('course_code')  # Optional course code filter
        professor_filter = data.get('professor')  # Optional professor filter
        username_filter = data.get('username')  # Optional uploader filter
        file_type_filter = data.get('file_type')  # Optional file type filter (pdf, docx, txt)
        sort_by = data.get('sort_by', 'relevance')  # relevance, newest, most_upvoted, most_downloaded
        offset = data.get('offset', 0)  # Pagination offset

        # Check Redis cache for recent identical searches
        filter_str = f"{university_filter or ''}:{course_code_filter or ''}:{professor_filter or ''}:{username_filter or ''}:{file_type_filter or ''}:{sort_by}"
        cache_key = f"browse:{question}:{filter_str}:{offset}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached:
                    debug_log(f"Cache HIT for browse: '{question}' (offset={offset})")
                    return jsonify(json.loads(cached)), 200
            except:
                pass  # Cache miss or error, proceed normally

        # Generate embedding for the question
        query_embedding = generate_embedding(question)
        if not query_embedding:
            return jsonify({"error": "Failed to generate query embedding"}), 500

        debug_log(f"🔍 Semantic browse search for: '{question}'" + (f" (university: {university_filter})" if university_filter else "") + f" (offset: {offset})")

        # Extract query keywords for filename search
        query_words = set(question.lower().split())
        stopwords = {'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'from', 'is', 'are', 'was', 'were', 'what', 'how', 'where', 'when', 'why', 'which', 'who'}
        query_keywords = query_words - stopwords

        # PART 1: Search using pgvector (content similarity)
        search_results = search_notes_vector(
            query_embedding=query_embedding,
            user_id=request.user_id,
            university=university_filter,
            course_code=course_code_filter,
            match_threshold=0.6,
            match_count=50
        )

        # PART 2: Search by filename (title matching)
        filename_results = []
        if query_keywords:
            try:
                # Build query to search filenames
                query_builder = supabase.table("notes").select("id, original_filename, university, course_code, user_id, uploaded_at")

                # Apply university filter if provided
                if university_filter:
                    query_builder = query_builder.eq("university", university_filter)

                # Get all public notes
                all_notes = query_builder.execute()

                # Filter by filename matching ALL query keywords (more strict)
                for note_data in all_notes.data:
                    filename_lower = note_data['original_filename'].lower()
                    # Check if ALL keywords are in the filename
                    matches = sum(1 for keyword in query_keywords if keyword in filename_lower)
                    # Only include if ALL keywords are present
                    if matches == len(query_keywords):
                        # Check if user has Nexus enabled (cached)
                        profile = get_cached_profile(note_data['user_id'])
                        if profile and profile.get('is_public', True):
                            filename_results.append({
                                'note_id': note_data['id'],
                                'note_data': note_data,
                                'username': profile.get('username', 'Anonymous'),
                                'keyword_matches': matches
                            })

                debug_log(f"📁 Filename search: Found {len(filename_results)} notes with matching titles")
            except Exception as e:
                debug_log(f"⚠️ Filename search error: {e}")
                filename_results = []

        # Format results with note metadata
        formatted_results = []
        seen_notes = set()

        # De-duplicate vector results and collect unique note IDs
        deduped_results = []
        for result in search_results:
            note_id = result['note_id']
            if note_id in seen_notes:
                continue
            seen_notes.add(note_id)
            deduped_results.append(result)

        # Batch fetch all note metadata in ONE query
        if deduped_results:
            all_note_ids = [r['note_id'] for r in deduped_results]
            notes_resp = supabase.table("notes").select(
                "id, original_filename, university, course_code, professor, file_type, user_id, uploaded_at, net_votes"
            ).in_("id", all_note_ids).execute()
            notes_map = {n['id']: n for n in (notes_resp.data or [])}
        else:
            notes_map = {}

        # Build formatted results using batch data + cached profiles
        for result in deduped_results:
            note_id = result['note_id']
            note = notes_map.get(note_id)
            if not note:
                continue

            # Skip Wikipedia articles from Nexus browse results
            if (note.get('university') or '').lower() == 'wikipedia':
                continue

            profile = get_cached_profile(note['user_id'])
            if not profile or not profile.get('is_public', True):
                continue
            username = profile.get('username', 'Anonymous')

            content = result.get('content_summary') or result.get('chunk_text', '')
            net_upvotes = note.get('net_votes', 0)

            formatted_results.append({
                "note_id": note_id,
                "filename": note['original_filename'],
                "username": username,
                "university": note.get('university', 'Unknown'),
                "course_code": note.get('course_code', ''),
                "professor": note.get('professor', ''),
                "file_type": note.get('file_type', ''),
                "created_at": note.get('uploaded_at', ''),
                "content": content,
                "similarity": round(result['similarity'], 2),
                "upvotes": net_upvotes
            })

        # Add filename search results (not already in vector results)
        # Batch fetch net_votes and content previews for all filename results
        new_filename_ids = [fr['note_id'] for fr in filename_results if fr['note_id'] not in seen_notes]
        filename_votes_map = {}
        filename_content_map = {}

        if new_filename_ids:
            # Batch net_votes
            try:
                nv_resp = supabase.table("notes").select("id, net_votes").in_("id", new_filename_ids).execute()
                for n in (nv_resp.data or []):
                    filename_votes_map[n['id']] = n.get('net_votes', 0)
            except:
                pass

            # Batch content previews
            try:
                chunks_resp = supabase.table("note_chunks").select("note_id, chunk_text, content_summary").in_("note_id", new_filename_ids).execute()
                for chunk in (chunks_resp.data or []):
                    if chunk['note_id'] not in filename_content_map:
                        filename_content_map[chunk['note_id']] = chunk.get('content_summary') or chunk.get('chunk_text', '')
            except:
                pass

        for filename_result in filename_results:
            note_id = filename_result['note_id']

            if note_id in seen_notes:
                continue
            seen_notes.add(note_id)

            note_data = filename_result['note_data']
            net_upvotes = filename_votes_map.get(note_id, 0)
            content_preview = filename_content_map.get(note_id, "Study notes and materials")

            # Add to results with low similarity (since it didn't match content)
            formatted_results.append({
                "note_id": note_id,
                "filename": note_data['original_filename'],
                "username": filename_result['username'],
                "university": note_data.get('university', 'Unknown'),
                "course_code": note_data.get('course_code', ''),
                "created_at": note_data.get('uploaded_at', ''),
                "content": content_preview,
                "similarity": 0.3,  # Low base similarity for title-only matches
                "upvotes": net_upvotes
            })

        debug_log(f"✅ Semantic browse: Found {len(formatted_results)} unique notes ({len(search_results)} from content, {len(filename_results)} from titles)")

        # Apply weighted ranking: similarity + upvote helpfulness + title match boost
        if formatted_results:

            for note in formatted_results:
                similarity = note['similarity']
                net = note.get('upvotes', 0)
                filename = note['filename'].lower()

                # Calculate helpfulness score from pre-computed net_votes
                # Normalize net votes to a -1 to 1 range (capped at +/-20)
                capped = max(-20, min(20, net))
                helpfulness_score = capped / 20.0

                # Normalize to 0-1 range
                normalized_helpfulness = (helpfulness_score + 1) / 2

                # Calculate title match score
                title_match_count = 0
                for keyword in query_keywords:
                    if keyword in filename:
                        title_match_count += 1

                # Title match boost: percentage of query keywords found in title
                if query_keywords:
                    title_match_score = title_match_count / len(query_keywords)
                else:
                    title_match_score = 0

                # Combined score: 60% similarity + 25% helpfulness + 15% title match
                combined_score = (similarity * 0.6) + (normalized_helpfulness * 0.25) + (title_match_score * 0.15)
                note['combined_score'] = round(combined_score, 4)
                note['title_match_score'] = round(title_match_score, 2)

            # Re-sort by combined score
            formatted_results.sort(key=lambda x: x.get('combined_score', 0), reverse=True)

            # Filter out low relevance results (combined score below 0.45)
            formatted_results = [note for note in formatted_results if note.get('combined_score', 0) >= 0.45]

        # Apply post-filters (professor, username, file_type)
        if professor_filter:
            pf = professor_filter.lower()
            # Need to fetch professor field for these notes
            note_ids = [n['note_id'] for n in formatted_results]
            if note_ids:
                prof_resp = supabase.table("notes").select("id, professor").in_("id", note_ids).execute()
                prof_map = {n['id']: (n.get('professor') or '').lower() for n in (prof_resp.data or [])}
                formatted_results = [n for n in formatted_results if pf in prof_map.get(n['note_id'], '')]

        if username_filter:
            uf = username_filter.lower()
            formatted_results = [n for n in formatted_results if uf in (n.get('username') or '').lower()]

        if file_type_filter:
            ft = file_type_filter.lower()
            formatted_results = [n for n in formatted_results if (n.get('filename') or '').lower().endswith('.' + ft)]

        # Apply sort override
        if sort_by == 'newest':
            formatted_results.sort(key=lambda x: x.get('created_at', ''), reverse=True)
        elif sort_by == 'most_upvoted':
            formatted_results.sort(key=lambda x: x.get('upvotes', 0), reverse=True)
        elif sort_by == 'most_downloaded':
            # Fetch download counts
            dl_note_ids = [n['note_id'] for n in formatted_results]
            if dl_note_ids:
                dl_resp = supabase.table("notes").select("id, download_count").in_("id", dl_note_ids).execute()
                dl_map = {n['id']: n.get('download_count', 0) for n in (dl_resp.data or [])}
                for n in formatted_results:
                    n['_dl'] = dl_map.get(n['note_id'], 0)
                formatted_results.sort(key=lambda x: x.get('_dl', 0), reverse=True)
                for n in formatted_results:
                    n.pop('_dl', None)
        # else: 'relevance' -- already sorted by combined_score

            debug_log(f"📊 Top 5 ranked notes (after filtering):")
            for i, note in enumerate(formatted_results[:5], 1):
                debug_log(f"  {i}. {note['filename'][:40]} - Combined: {note['combined_score']:.3f} (Sim: {note['similarity']:.2f}, Upvotes: {note['upvotes']}, Title: {note['title_match_score']:.2f})")

        # Apply pagination
        page_size = 10
        start_idx = offset
        end_idx = offset + page_size

        paginated_results = formatted_results[start_idx:end_idx]
        has_more = len(formatted_results) > end_idx

        # Clean up response - remove internal ranking fields
        for note in paginated_results:
            note.pop('combined_score', None)
            note.pop('title_match_score', None)

        response_data = {
            "notes": paginated_results,
            "has_more": has_more,
            "total": len(formatted_results)
        }

        # Track search query
        try:
            from StudyFlow.backend.search_tracker import track_search
            track_search(question, len(paginated_results), source="browse")
        except:
            pass

        # Cache result in Redis for 5 minutes
        if redis_cache:
            try:
                redis_cache.setex(cache_key, 1800, json.dumps(response_data))
            except:
                pass

        return jsonify(response_data), 200

    except Exception as e:
        debug_log(f"Semantic browse error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/topic-tags", methods=["GET"])
@supabase_auth_required
def get_topic_tags():
    """
    Get all unique topic tags from public notes

    Returns:
    {
        "tags": ["Biology", "Chemistry", ...]
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        # Get all public notes with topic tags
        response = supabase.table("notes").select("topic_tags").eq("is_public", True).execute()
        notes = response.data if response.data else []

        # Extract unique tags
        all_tags = set()
        for note in notes:
            if note.get('topic_tags'):
                for tag in note['topic_tags']:
                    all_tags.add(tag)

        # Sort alphabetically
        tags = sorted(list(all_tags))

        debug_log(f"🏷️ Found {len(tags)} unique topic tags")

        return jsonify({"tags": tags}), 200

    except Exception as e:
        debug_log(f"❌ Get topic tags error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/dmca/takedown", methods=["POST"])
def submit_takedown_request():
    """
    Submit a DMCA takedown request (no auth required for public access)

    Body:
    {
        "requestor_name": "John Doe",
        "requestor_email": "john@example.com",
        "requestor_type": "professor",
        "note_filename": "Biology_Chapter_3.pdf",
        "complaint_description": "...",
        "evidence_urls": ["https://..."],
        "good_faith": true,
        "accuracy": true
    }

    Returns:
    {
        "success": true,
        "request_id": "DMCA-20260322-001"
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from datetime import datetime

        data = request.json

        # Validate required fields
        required_fields = ['requestor_name', 'requestor_email', 'requestor_type', 'note_filename', 'complaint_description', 'good_faith', 'accuracy']
        for field in required_fields:
            if field not in data:
                return jsonify({"error": f"Missing required field: {field}"}), 400

        if not data['good_faith'] or not data['accuracy']:
            return jsonify({"error": "You must agree to the good faith and accuracy statements"}), 400

        # Generate unique request ID
        today = datetime.now().strftime("%Y%m%d")

        # Count today's requests to get next number
        count_response = supabase.table("dmca_takedown_requests").select("id", count="exact").like("request_id", f"DMCA-{today}-%").execute()
        next_num = (count_response.count if count_response.count else 0) + 1
        request_id = f"DMCA-{today}-{next_num:03d}"

        # Insert into database
        evidence_urls = data.get('evidence_urls', [])
        if isinstance(evidence_urls, str):
            evidence_urls = [url.strip() for url in evidence_urls.split('\n') if url.strip()]

        supabase.table("dmca_takedown_requests").insert({
            "request_id": request_id,
            "requestor_name": data['requestor_name'],
            "requestor_email": data['requestor_email'],
            "requestor_type": data['requestor_type'],
            "note_filename": data['note_filename'],
            "complaint_description": data['complaint_description'],
            "evidence_urls": evidence_urls,
            "copyright_claim": True,
            "status": "pending"
        }).execute()

        debug_log(f"📋 DMCA takedown request submitted: {request_id}")

        # Send email notification to admin
        send_dmca_report_notification_to_admin(
            takedown_id=request_id,
            note_filename=data['note_filename'],
            reporter_email=data['requestor_email'],
            reason=data['complaint_description'][:200]
        )

        return jsonify({
            "success": True,
            "request_id": request_id
        }), 200

    except Exception as e:
        debug_log(f"❌ DMCA takedown error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/newsletter/subscribe", methods=["POST"])
def newsletter_subscribe():
    """
    Subscribe an email to the newsletter (no auth required).

    Body: { "email": "user@example.com" }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.json
        email = data.get('email', '').strip().lower()

        if not email or '@' not in email:
            return jsonify({"error": "Valid email address required"}), 400

        # Upsert so duplicate emails don't cause errors
        supabase.table("newsletter_subscribers").upsert({
            "email": email
        }, on_conflict="email").execute()

        debug_log(f"Newsletter subscription: {email}")

        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Newsletter subscribe error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/view", methods=["POST"])
@supabase_auth_required
def track_note_view(note_id):
    """
    Track when a user views a note (for analytics)

    Returns:
    {
        "success": true
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        user_id = request.user_id

        # Get IP address
        ip_address = request.headers.get('X-Forwarded-For', request.remote_addr)
        if ',' in ip_address:
            ip_address = ip_address.split(',')[0].strip()

        # Insert view record
        supabase.table("note_views").insert({
            "user_id": user_id,
            "note_id": note_id,
            "ip_address": ip_address
        }).execute()

        debug_log(f"👁️ Note view tracked: {note_id}")

        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"❌ Track view error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/consent-log", methods=["POST"])
@supabase_auth_required
def log_consent():
    """
    Log a click-wrap consent receipt (Missouri HB 2271, SB 1324, DMCA compliance).
    Stores user_id, file_id, legal_version, timestamp, and IP for AG audit trail.
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        if not data:
            return jsonify({"error": "Missing request body"}), 400

        file_id = data.get('file_id', '')
        legal_version = data.get('legal_version', '')
        action = data.get('action', 'download_consent')

        if not file_id or not legal_version:
            return jsonify({"error": "file_id and legal_version are required"}), 400

        # Get IP address
        ip_address = request.headers.get('X-Forwarded-For', request.remote_addr)
        if ip_address and ',' in ip_address:
            ip_address = ip_address.split(',')[0].strip()

        user_agent = request.headers.get('User-Agent', '')

        supabase.table("consent_logs").insert({
            "user_id": request.user_id,
            "file_id": str(file_id),
            "legal_version": legal_version,
            "action": action,
            "ip_address": ip_address,
            "user_agent": user_agent
        }).execute()

        debug_log(f"[+] Consent logged: user={request.user_id}, file={file_id}, version={legal_version}")

        return jsonify({"status": "logged"}), 200

    except Exception as e:
        debug_log(f"[-] Consent log error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/related", methods=["GET"])
@supabase_auth_required
def get_related_notes(note_id):
    """Get up to 5 notes similar to this one using vector similarity."""
    try:
        from StudyFlow.backend.supabase_client import supabase, search_notes_vector

        # Check Redis cache
        cache_key = f"related:{note_id}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached:
                    return jsonify(json.loads(cached)), 200
            except:
                pass

        # Get the first chunk's embedding for this note
        chunk_resp = supabase.table("note_chunks").select(
            "embedding"
        ).eq("note_id", note_id).eq("chunk_index", 0).limit(1).execute()

        if not chunk_resp.data or not chunk_resp.data[0].get('embedding'):
            return jsonify({"related": []}), 200

        embedding = chunk_resp.data[0]['embedding']

        # Search for similar notes
        results = search_notes_vector(
            query_embedding=embedding,
            user_id=request.user_id,
            match_threshold=0.5,
            match_count=20
        )

        if not results:
            return jsonify({"related": []}), 200

        # Collect unique note IDs (excluding current note)
        seen = set()
        unique_note_ids = []
        for r in results:
            nid = r['note_id']
            if nid != note_id and nid not in seen:
                seen.add(nid)
                unique_note_ids.append(nid)
            if len(unique_note_ids) >= 5:
                break

        if not unique_note_ids:
            return jsonify({"related": []}), 200

        # Batch fetch note metadata
        notes_resp = supabase.table("notes").select(
            "id, original_filename, university, course_code, user_id, net_votes"
        ).in_("id", unique_note_ids).eq("is_public", True).execute()

        if not notes_resp.data:
            return jsonify({"related": []}), 200

        # Batch fetch usernames
        user_ids = list(set(n['user_id'] for n in notes_resp.data))
        profiles_resp = supabase.table("user_profiles").select("id, username, is_public").in_("id", user_ids).execute()
        profile_map = {p['id']: p for p in (profiles_resp.data or [])}

        related = []
        for n in notes_resp.data:
            profile = profile_map.get(n['user_id'])
            if not profile or not profile.get('is_public', True):
                continue
            related.append({
                "id": n['id'],
                "filename": n.get('original_filename', 'Untitled'),
                "university": n.get('university', ''),
                "course_code": n.get('course_code', ''),
                "username": profile.get('username', 'Anonymous'),
                "net_votes": n.get('net_votes', 0)
            })

        result = {"related": related[:5]}

        # Cache for 30 minutes
        if redis_cache:
            try:
                redis_cache.setex(cache_key, 1800, json.dumps(result))
            except:
                pass

        return jsonify(result), 200

    except Exception as e:
        debug_log(f"Related notes error: {e}\n{traceback.format_exc()}")
        return jsonify({"related": []}), 200


@app.route("/api/notes/<note_id>/metadata", methods=["GET"])
@supabase_auth_required
def get_note_metadata(note_id):
    """
    Get note metadata for viewing

    Returns:
    {
        "note": {
            "id": "...",
            "filename": "...",
            "file_type": "pdf|image|document|text|other",
            "university": "...",
            "course_code": "...",
            "username": "...",
            "page_count": 12
        }
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        import os

        # Get note details - check if public OR owned by current user
        note_response = supabase.table("notes").select(
            "id, original_filename, file_type, university, course_code, user_id, page_count, is_public"
        ).eq("id", note_id).execute()

        if not note_response.data:
            return jsonify({"error": "Note not found"}), 404

        note_data = note_response.data[0]

        # Check if user has access (either public or owns it)
        is_public = note_data.get('is_public', False)
        is_owner = note_data.get('user_id') == request.user_id

        if not is_public and not is_owner:
            return jsonify({"error": "Note not found or not public"}), 404

        # Get filename and file type
        filename = note_data.get('original_filename') or note_data.get('filename') or 'unknown.pdf'

        # Use stored file_type from database (more reliable than extension)
        stored_file_type = note_data.get('file_type', '').lower()

        # Map stored file_type to viewer format
        if stored_file_type == 'pdf':
            file_type = 'pdf'
        elif stored_file_type in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'webp', 'image']:
            file_type = 'image'
        elif stored_file_type in ['doc', 'docx', 'odt', 'document']:
            file_type = 'document'
        elif stored_file_type in ['txt', 'md', 'markdown', 'text']:
            file_type = 'text'
        else:
            # Fallback to extension detection if file_type not set
            ext = os.path.splitext(filename)[1].lower()
            if ext in ['.pdf']:
                file_type = 'pdf'
            elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp']:
                file_type = 'image'
            elif ext in ['.doc', '.docx', '.odt']:
                file_type = 'document'
            elif ext in ['.txt', '.md', '.markdown']:
                file_type = 'text'
            else:
                file_type = 'other'

        # Get extension for response
        ext = os.path.splitext(filename)[1].lower() if '.' in filename else ''

        # Get username
        try:
            user_profile = supabase.table("user_profiles").select("username").eq("id", note_data['user_id']).single().execute()
            username = user_profile.data['username'] if user_profile.data else 'Anonymous'
        except:
            username = 'Anonymous'

        # Build response
        response_data = {
            "id": note_data['id'],
            "filename": filename,
            "file_type": file_type,
            "extension": ext,
            "university": note_data.get('university'),
            "course_code": note_data.get('course_code'),
            "username": username,
            "page_count": note_data.get('page_count', 0),
            "is_owner": is_owner  # Include ownership status for annotation access control
        }

        return jsonify({"note": response_data}), 200

    except Exception as e:
        debug_log(f"❌ Get note metadata error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/view-file", methods=["GET"])
def view_note_file(note_id):
    """
    Get the file for viewing (PDF, image, document, etc.)
    No auth required since notes are already public - iframes can't send headers
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from flask import redirect, send_file
        import os

        # Verify note exists and is public
        note = supabase.table("notes").select("id, file_path, original_filename").eq("id", note_id).eq("is_public", True).single().execute()

        if not note.data:
            return jsonify({"error": "Note not found or not public"}), 404

        debug_log(f"View file for note {note_id}: {note.data}")

        # Get signed URL from file_path
        file_url = None
        if note.data.get('file_path'):
            try:
                # Generate signed URL (valid for 1 hour)
                signed_url_response = supabase.storage.from_("note-files").create_signed_url(
                    path=note.data['file_path'],
                    expires_in=3600  # 1 hour
                )
                file_url = signed_url_response['signedURL']
                debug_log(f"Generated signed URL for file_path: {note.data['file_path']}")
            except Exception as storage_error:
                debug_log(f"Storage error with file_path: {storage_error}")

        if file_url:
            debug_log(f"Redirecting to file: {file_url}")
            return redirect(file_url)
        else:
            debug_log(f"No file URL found for note {note_id}. Available fields: {note.data.keys()}")
            return jsonify({"error": "File URL not found", "available_fields": list(note.data.keys())}), 404

    except Exception as e:
        error_trace = traceback.format_exc()
        debug_log(f"❌ View file error: {e}\n{error_trace}")
        return jsonify({"error": str(e), "traceback": error_trace}), 500


@app.route("/api/notes/<note_id>/view-as-images", methods=["GET"])
def view_note_as_images(note_id):
    """
    Convert PDF/DOCX to images for secure viewing (prevents download/print)
    Returns JSON with base64-encoded images for each page
    Auth optional - works for public notes or user's own notes
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        import fitz  # PyMuPDF
        import base64
        from io import BytesIO

        # Check if user is authenticated
        auth_header = request.headers.get('Authorization')
        user_id = None
        if auth_header and auth_header.startswith('Bearer '):
            token = auth_header.split(' ')[1]
            try:
                user_response = supabase.auth.get_user(token)
                if user_response and user_response.user:
                    user_id = user_response.user.id
            except:
                pass

        # Fetch note - either public OR owned by authenticated user
        note_query = supabase.table("notes").select("id, file_path, original_filename, file_type, page_count, user_id, is_public").eq("id", note_id)
        note = note_query.single().execute()

        if not note.data:
            return jsonify({"error": "Note not found"}), 404

        # Check access: must be public OR owned by authenticated user
        is_public = note.data.get('is_public', False)
        note_owner_id = note.data.get('user_id')

        if not is_public and (not user_id or user_id != note_owner_id):
            return jsonify({"error": "Note not found or not public"}), 404

        # Only PDFs can be converted to images
        # Office documents (docx, xlsx, pptx) should be opened in Office editor instead
        file_type = note.data.get('file_type', '').lower()
        if file_type in ['docx', 'xlsx', 'pptx']:
            return jsonify({
                "error": "Office documents cannot be viewed as images",
                "suggestion": "Open this document in StudyFlow Office instead",
                "redirect": f"/office-editor.html?id={note_id}"
            }), 400

        if file_type != 'pdf':
            return jsonify({"error": "Only PDFs can be converted to images"}), 400

        # Download PDF from Supabase storage
        file_path = note.data.get('file_path')
        if not file_path:
            return jsonify({"error": "File path not found"}), 404

        file_data = supabase.storage.from_('note-files').download(file_path)

        # Convert PDF pages to images
        pdf_document = fitz.open(stream=file_data, filetype="pdf")
        images = []

        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            # Render page at 2x resolution for better quality
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = pix.tobytes("png")

            # Convert to base64
            img_base64 = base64.b64encode(img_data).decode('utf-8')
            images.append({
                "page": page_num + 1,
                "data": f"data:image/png;base64,{img_base64}"
            })

        pdf_document.close()

        return jsonify({
            "note_id": note_id,
            "filename": note.data.get('original_filename'),
            "total_pages": len(images),
            "images": images
        }), 200

    except Exception as e:
        error_trace = traceback.format_exc()
        debug_log(f"❌ PDF to image conversion error: {e}\n{error_trace}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/download", methods=["GET"])
@supabase_auth_required
def download_note_file(note_id):
    """
    Download a public note file with flattening and watermark.
    REQUIRES GOOD STANDING: Missouri HB 2271 compliance
    NOTE: This is a duplicate endpoint - the one at line 3474 takes precedence
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from StudyFlow.backend.pdf_flatten import flatten_pdf, flatten_image, can_flatten, convert_to_pdf
        from flask import send_file
        from datetime import date
        import io
        import uuid

        # DAILY DOWNLOAD LIMIT CHECK (3 downloads/day for free users)
        # Fetch user profile to check subscription and download count
        profile_result = supabase.table("user_profiles").select(
            "subscription_tier, subscription_status, daily_downloads_count, daily_downloads_reset_date, good_standing, permanent_bad_standing, last_verified_upload_date"
        ).eq("id", request.user_id).single().execute()
        profile = profile_result.data

        # CHECK: Good standing expires after 30 days without a quality upload (free users only)
        from datetime import datetime, timedelta
        is_scholar_check = (
            profile.get('subscription_tier') == 'pro' and
            profile.get('subscription_status') == 'active'
        )

        if not is_scholar_check and profile.get('good_standing', False):
            last_upload = profile.get('last_verified_upload_date')
            if last_upload:
                try:
                    last_upload_date = datetime.fromisoformat(last_upload.replace('Z', '+00:00'))
                    days_since_upload = (datetime.utcnow().replace(tzinfo=last_upload_date.tzinfo) - last_upload_date).days

                    if days_since_upload > 30:
                        # Revoke good standing - hasn't uploaded in 30+ days
                        supabase.table("user_profiles").update({
                            "good_standing": False
                        }).eq("id", request.user_id).execute()
                        profile['good_standing'] = False
                        debug_log(f"[!] User {request.user_id} lost good standing (no upload for {days_since_upload} days)")
                except Exception as e:
                    debug_log(f"[-] Error checking last upload date: {e}")

        # Check if Scholar's Club member (unlimited downloads)
        is_scholar = (
            profile.get('subscription_tier') == 'pro' and
            profile.get('subscription_status') == 'active'
        )

        if not is_scholar:
            # Free user - enforce 3 downloads per day limit
            today = date.today()
            reset_date = profile.get('daily_downloads_reset_date')
            download_count = profile.get('daily_downloads_count', 0)

            # Reset counter if new day
            if not reset_date or str(reset_date) != str(today):
                supabase.table("user_profiles").update({
                    "daily_downloads_count": 0,
                    "daily_downloads_reset_date": str(today)
                }).eq("id", request.user_id).execute()
                download_count = 0
                debug_log(f"[+] Reset daily download count for user {request.user_id}")

            # Check if at limit
            if download_count >= 3:
                debug_log(f"[!] User {request.user_id} hit daily download limit ({download_count}/3)")
                return jsonify({
                    "error": "Daily download limit reached",
                    "message": "You've used all 3 downloads today. Join Scholar's Club for unlimited downloads!",
                    "limit_reached": True,
                    "downloads_used": download_count,
                    "downloads_limit": 3,
                    "subscribe_url": "https://unclephilburt.github.io/studyflowwebsite/account.html"
                }), 403

        # Verify note exists (check both own notes and public notes)
        result = supabase.table("notes").select("id, file_path, original_filename, user_id").eq("id", note_id).execute()
        note_data = result.data[0] if result.data else None

        if not note_data:
            return jsonify({"error": "Note not found"}), 404

        # Check if note belongs to current user or is public
        is_own_note = note_data.get('user_id') == request.user_id

        # GOOD STANDING CHECK: Only required for downloading OTHER people's notes
        # (Reuse profile fetched earlier for download limit check)
        if not is_own_note:
            # Block if permanently banned
            if profile.get('permanent_bad_standing', False):
                return jsonify({"error": "Account suspended due to DMCA violations"}), 403

            # Check good standing (Scholar's Club members are always in good standing)
            if not is_scholar and not profile.get('good_standing', False):
                return jsonify({
                    "error": "Good Standing required to download notes",
                    "message": "Upload a verified note or subscribe to regain access",
                    "good_standing_required": True
                }), 403

        print(f"[DOWNLOAD] User {request.user_id} downloading {'own' if is_own_note else 'public'} note {note_id}", flush=True)

        file_path = note_data.get('file_path')
        original_filename = note_data.get('original_filename', 'note')

        if not file_path:
            return jsonify({"error": "File not found in storage"}), 404

        # Generate transaction code
        transaction_code = "DL-" + uuid.uuid4().hex[:8]

        # Get uploader's username for watermark
        username = "Public"
        try:
            uploader_id = note_data.get('user_id')
            if uploader_id:
                profile_result = supabase.table("user_profiles").select("username").eq("id", uploader_id).execute()
                if profile_result.data and profile_result.data[0].get("username"):
                    username = profile_result.data[0]["username"]
        except Exception:
            pass

        # Download file from Supabase Storage
        file_data = supabase.storage.from_('note-files').download(file_path)

        # Log download transaction
        ip_address = request.headers.get('X-Forwarded-For', request.remote_addr)
        if ip_address and ',' in ip_address:
            ip_address = ip_address.split(',')[0].strip()

        try:
            supabase.table("download_transactions").insert({
                "user_id": request.user_id,
                "note_id": note_id,
                "original_filename": original_filename,
                "transaction_code": transaction_code,
                "ip_address": ip_address,
                "user_agent": request.headers.get('User-Agent', '')
            }).execute()
            debug_log(f"[+] Download transaction logged: {transaction_code}")
        except Exception as tx_err:
            debug_log(f"[-] Failed to log download transaction: {tx_err}")

        # Log HB 2271 certification (Missouri legal compliance)
        try:
            supabase.table("download_certifications").insert({
                "user_id": request.user_id,
                "note_id": note_id,
                "ip_address": ip_address,
                "user_agent": request.headers.get('User-Agent')
            }).execute()
            print(f"[HB 2271] Download certification logged for user {request.user_id}", flush=True)
        except Exception as cert_err:
            debug_log(f"[-] Failed to log download certification: {cert_err}")

        # Flatten and watermark -- everything becomes a PDF
        flattenable, file_type = can_flatten(original_filename)
        name_base = original_filename.rsplit('.', 1)[0] if '.' in original_filename else original_filename

        if flattenable and file_type == 'pdf':
            file_data = flatten_pdf(file_data, username, transaction_code)
        elif flattenable and file_type == 'image':
            file_data = flatten_image(file_data, username, transaction_code)
        else:
            # Legacy file (txt, docx, etc.) -- convert to PDF first, then flatten
            pdf_data, _ = convert_to_pdf(file_data, original_filename)
            if pdf_data:
                file_data = flatten_pdf(pdf_data, username, transaction_code)

        download_name = f"{name_base}_studyflow.pdf"
        mimetype = 'application/pdf'

        # Increment download counter for free users
        if not is_scholar:
            try:
                new_count = (profile.get('daily_downloads_count', 0) or 0) + 1
                supabase.table("user_profiles").update({
                    "daily_downloads_count": new_count
                }).eq("id", request.user_id).execute()
                debug_log(f"[+] User {request.user_id} download count: {new_count}/3")
            except Exception as count_err:
                debug_log(f"[-] Failed to increment download count: {count_err}")

        return send_file(
            io.BytesIO(file_data),
            mimetype=mimetype,
            as_attachment=True,
            download_name=download_name
        )

    except Exception as e:
        debug_log(f"[-] Download file error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ SYLLABUS & CALENDAR ENDPOINTS ============

@app.route("/api/syllabus/upload", methods=["POST"])
@supabase_auth_required
def upload_syllabus():
    """
    Upload a course syllabus (PDF/DOCX) and extract calendar events using AI

    Expects multipart/form-data:
    - file: Syllabus PDF or DOCX file
    - course_name: Course name (optional)
    - course_code: Course code (optional)
    - professor_name: Professor name (optional)
    - semester: Semester (optional, e.g., "Fall 2026")

    Returns:
    {
        "syllabus_id": "uuid",
        "events": [{...extracted calendar events...}],
        "message": "Syllabus uploaded and processed"
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from StudyFlow.backend.syllabus_parser import extract_calendar_events, format_event_for_db, validate_event
        import PyPDF2
        import io

        # Get uploaded file
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400

        # Get form data
        course_name = request.form.get('course_name', 'Untitled Course')
        course_code = request.form.get('course_code')
        professor_name = request.form.get('professor_name')
        semester = request.form.get('semester')

        # Validate file type
        allowed_extensions = {'.pdf', '.docx', '.doc'}
        file_ext = os.path.splitext(file.filename)[1].lower()
        if file_ext not in allowed_extensions:
            return jsonify({"error": f"File type not supported. Please upload PDF or DOCX. Got: {file_ext}"}), 400

        # Read file content
        file_content = file.read()
        file_size = len(file_content)

        # Extract text from file with robust fallback handling
        extracted_text = ""

        if file_ext == '.pdf':
            try:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        extracted_text += page_text + "\n"

                debug_log(f"[PDF] PyPDF2 extracted {len(extracted_text)} chars")

            except Exception as e:
                debug_log(f"[PDF] Extraction error: {e}")
                import traceback
                debug_log(traceback.format_exc())
                return jsonify({
                    "error": f"Failed to process PDF file. Error: {str(e)}"
                }), 500

        elif file_ext in ['.docx', '.doc']:
            try:
                import docx
                doc = docx.Document(io.BytesIO(file_content))
                for paragraph in doc.paragraphs:
                    extracted_text += paragraph.text + "\n"
                # Also extract from tables
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            extracted_text += cell.text + "\n"
                debug_log(f"[DOCX] Extracted {len(extracted_text)} chars")
            except Exception as e:
                debug_log(f"[DOCX] Extraction error: {e}")
                return jsonify({
                    "error": f"Failed to process DOCX file. The file may be corrupted. Error: {str(e)}"
                }), 500

        # Validate we got meaningful text
        if not extracted_text.strip():
            return jsonify({
                "error": "Could not extract any text from the file. Please ensure the PDF contains text (not just images) or try uploading a DOCX file instead."
            }), 400

        debug_log(f"[VALIDATION] Extracted text length: {len(extracted_text.strip())} chars")

        # LEGAL COMPLIANCE LOGGING
        print(f"\n[LEGAL] Processing syllabus for Fair Use extraction (user: {request.user_id})", flush=True)
        print(f"[LEGAL] Extracting only factual data (dates, times, task names) per 17 USC 107", flush=True)

        # Insert syllabus record (metadata only - NO copyrighted content stored)
        syllabus_data = {
            "user_id": request.user_id,
            "course_name": course_name,
            "course_code": course_code,
            "professor_name": professor_name,
            "semester": semester,
            "file_path": None,  # SB 1324 Compliance: Do not store original file
            "original_filename": file.filename,
            "extracted_text": None,  # Copyright Compliance: Do not store copyrighted text
            "file_size": file_size
        }

        syllabus_result = supabase.table("syllabi").insert(syllabus_data).execute()
        syllabus_id = syllabus_result.data[0]['id']

        print(f"[LEGAL] Syllabus metadata saved: {syllabus_id}", flush=True)

        # Extract calendar events using AI (Fair Use - factual data only)
        print(f"\n[UPLOAD] Calling extract_calendar_events()...", flush=True)
        events = extract_calendar_events(extracted_text, course_name, course_code)
        print(f"[UPLOAD] extract_calendar_events() returned {len(events)} events", flush=True)

        # Insert events into database
        inserted_events = []
        for event in events:
            print(f"[UPLOAD] Validating event: {event.get('title', 'NO TITLE')}", flush=True)
            if validate_event(event):
                event_data = format_event_for_db(event, request.user_id, syllabus_id)
                result = supabase.table("calendar_events").insert(event_data).execute()
                if result.data:
                    inserted_events.append(result.data[0])
                    print(f"[UPLOAD] Inserted event: {event.get('title')}", flush=True)
            else:
                print(f"[UPLOAD] Event validation FAILED for: {event}", flush=True)

        print(f"\n[LEGAL] Extracted {len(inserted_events)} events (factual data only)", flush=True)
        print(f"[LEGAL] Original syllabus discarded per SB 1324 disclosure", flush=True)

        # File is NOT uploaded to storage - processing complete, file discarded
        # This satisfies: Copyright (Fair Use), Privacy (SB 1324), Academic Integrity (HB 2271)

        return jsonify({
            "success": True,
            "syllabus_id": syllabus_id,
            "events": inserted_events,
            "message": f"Syllabus uploaded and {len(inserted_events)} events extracted"
        }), 200

    except Exception as e:
        debug_log(f"❌ Syllabus upload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/canvas/import", methods=["POST"])
@supabase_auth_required
def import_canvas_calendar():
    """
    Import calendar events from Canvas iCal feed

    Expects JSON:
    {
        "feed_url": "https://canvas.school.edu/feeds/calendars/user_xxx.ics",
        "auto_sync": true
    }

    Returns:
    {
        "success": true,
        "events_count": 15,
        "message": "Imported 15 events from Canvas"
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from StudyFlow.backend.syllabus_parser import format_event_for_db, validate_event
        import requests
        from icalendar import Calendar
        from datetime import datetime

        data = request.get_json()
        feed_url = data.get('feed_url')
        auto_sync = data.get('auto_sync', False)

        if not feed_url:
            return jsonify({"error": "Missing feed_url"}), 400

        # Validate URL format
        if not feed_url.startswith('http'):
            return jsonify({"error": "Invalid feed URL"}), 400

        debug_log(f"[CANVAS] Fetching calendar feed: {feed_url}")

        # Fetch .ics file from Canvas
        try:
            response = requests.get(feed_url, timeout=30)
            response.raise_for_status()
            ics_content = response.text
        except requests.exceptions.RequestException as e:
            debug_log(f"[CANVAS] Failed to fetch feed: {e}")
            return jsonify({"error": f"Failed to fetch Canvas calendar: {str(e)}"}), 400

        # Parse iCal feed
        try:
            cal = Calendar.from_ical(ics_content)
        except Exception as e:
            debug_log(f"[CANVAS] Failed to parse iCal: {e}")
            return jsonify({"error": f"Failed to parse calendar feed: {str(e)}"}), 400

        # Extract events from iCal
        events = []
        for component in cal.walk():
            if component.name == "VEVENT":
                try:
                    # Extract event data
                    summary = str(component.get('summary', 'Untitled'))
                    description = str(component.get('description', '')) if component.get('description') else None
                    dtstart = component.get('dtstart').dt if component.get('dtstart') else None

                    # Skip if no date
                    if not dtstart:
                        continue

                    # Convert to date/time strings
                    if isinstance(dtstart, datetime):
                        due_date = dtstart.strftime('%Y-%m-%d')
                        due_time = dtstart.strftime('%H:%M:%S')
                    else:
                        due_date = dtstart.strftime('%Y-%m-%d')
                        due_time = None

                    # Detect event type from summary
                    event_type = 'other'
                    summary_lower = summary.lower()
                    if any(word in summary_lower for word in ['assignment', 'homework', 'hw']):
                        event_type = 'assignment'
                    elif any(word in summary_lower for word in ['exam', 'test', 'quiz', 'midterm', 'final']):
                        event_type = 'exam'
                    elif any(word in summary_lower for word in ['project', 'paper', 'essay', 'presentation']):
                        event_type = 'project'
                    elif any(word in summary_lower for word in ['reading', 'chapter', 'read']):
                        event_type = 'reading'

                    # Extract course name from summary or description
                    course_name = None
                    if ':' in summary:
                        course_name = summary.split(':')[0].strip()

                    event = {
                        'event_type': event_type,
                        'title': summary,
                        'description': description,
                        'due_date': due_date,
                        'due_time': due_time,
                        'course_name': course_name
                    }

                    events.append(event)
                except Exception as e:
                    debug_log(f"[CANVAS] Failed to parse event: {e}")
                    continue

        debug_log(f"[CANVAS] Parsed {len(events)} events from iCal feed")

        # Insert events into database
        inserted_events = []
        for event in events:
            if validate_event(event):
                event_data = format_event_for_db(event, request.user_id, syllabus_id=None)
                result = supabase.table("calendar_events").insert(event_data).execute()
                if result.data:
                    inserted_events.append(result.data[0])

        debug_log(f"[CANVAS] Inserted {len(inserted_events)} events")

        # Store Canvas feed URL for auto-sync if requested
        if auto_sync:
            try:
                supabase.table("user_profiles").update({
                    "canvas_feed_url": feed_url
                }).eq("id", request.user_id).execute()
                debug_log(f"[CANVAS] Saved feed URL for auto-sync")
            except Exception as e:
                debug_log(f"[CANVAS] Failed to save feed URL: {e}")

        return jsonify({
            "success": True,
            "events_count": len(inserted_events),
            "message": f"Imported {len(inserted_events)} events from Canvas"
        }), 200

    except Exception as e:
        debug_log(f"[CANVAS] Import error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

    except Exception as e:
        debug_log(f"❌ Upload syllabus error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/calendar/events", methods=["GET"])
@supabase_auth_required
def get_calendar_events():
    """
    Get user's calendar events

    Query params:
    - start_date: Filter events after this date (YYYY-MM-DD)
    - end_date: Filter events before this date (YYYY-MM-DD)
    - course_code: Filter by course code
    - completed: Filter by completion status (true/false)
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        start_date = request.args.get('start_date', '')
        end_date = request.args.get('end_date', '')
        course_code = request.args.get('course_code', '')
        completed = request.args.get('completed')

        # Check Redis cache
        cache_key = f"calendar:{request.user_id}:{start_date}:{end_date}:{course_code}:{completed}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        query = supabase.table("calendar_events").select("*").eq("user_id", request.user_id)

        if start_date:
            query = query.gte('due_date', start_date)
        if end_date:
            query = query.lte('due_date', end_date)
        if course_code:
            query = query.eq('course_code', course_code)
        if completed is not None:
            query = query.eq('completed', completed.lower() == 'true')

        query = query.order('due_date', desc=False)
        result = query.execute()

        resp = {"success": True, "events": result.data}
        if redis_cache:
            try: redis_cache.setex(cache_key, 300, json.dumps(resp))
            except: pass
        return jsonify(resp), 200

    except Exception as e:
        debug_log(f"❌ Get calendar events error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/calendar/events", methods=["POST"])
@supabase_auth_required
def create_calendar_event():
    """
    Manually create a calendar event

    Request body:
    {
        "event_type": "assignment",
        "title": "Lab Report 3",
        "description": "Write up lab 3 results",
        "due_date": "2026-04-15",
        "due_time": "23:59",
        "course_name": "Chemistry 101",
        "course_code": "CHEM101"
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from StudyFlow.backend.syllabus_parser import format_event_for_db, validate_event

        data = request.get_json()

        if not validate_event(data):
            return jsonify({"error": "Missing required fields: event_type, title, due_date"}), 400

        event_data = format_event_for_db(data, request.user_id, syllabus_id=None)
        result = supabase.table("calendar_events").insert(event_data).execute()

        return jsonify({
            "success": True,
            "event": result.data[0]
        }), 201

    except Exception as e:
        debug_log(f"❌ Create calendar event error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/calendar/events/<event_id>", methods=["PUT"])
@supabase_auth_required
def update_calendar_event(event_id):
    """
    Update a calendar event (e.g., mark as completed, edit details)

    Request body:
    {
        "completed": true,
        "title": "Updated title",
        ...
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()

        # Update event (only if it belongs to user)
        result = supabase.table("calendar_events").update(data).eq("id", event_id).eq("user_id", request.user_id).execute()

        if not result.data:
            return jsonify({"error": "Event not found or unauthorized"}), 404

        return jsonify({
            "success": True,
            "event": result.data[0]
        }), 200

    except Exception as e:
        debug_log(f"❌ Update calendar event error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/calendar/events/<event_id>", methods=["DELETE"])
@supabase_auth_required
def delete_calendar_event(event_id):
    """Delete a calendar event"""
    try:
        from StudyFlow.backend.supabase_client import supabase

        result = supabase.table("calendar_events").delete().eq("id", event_id).eq("user_id", request.user_id).execute()

        if not result.data:
            return jsonify({"error": "Event not found or unauthorized"}), 404

        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"❌ Delete calendar event error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/syllabi", methods=["GET"])
@supabase_auth_required
def get_syllabi():
    """Get user's uploaded syllabi"""
    try:
        from StudyFlow.backend.supabase_client import supabase

        result = supabase.table("syllabi").select("*").eq("user_id", request.user_id).order('created_at', desc=True).execute()

        return jsonify({
            "success": True,
            "syllabi": result.data
        }), 200

    except Exception as e:
        debug_log(f"❌ Get syllabi error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ GOOD STANDING SYSTEM ENDPOINTS ============

@app.route("/api/good-standing/status", methods=["GET"])
@supabase_auth_required
def get_good_standing_status():
    """
    Check if user is in Good Standing (Missouri HB 2271 / DMCA compliance)

    Returns:
    {
        "good_standing": true/false,
        "reason": "active_subscription" | "recent_upload" | "no_recent_uploads" | "dmca_strikes",
        "days_until_expiry": 15,  # Days until they lose good standing
        "last_verified_upload": "2026-03-20T10:30:00Z",
        "dmca_strikes": 0,
        "permanent_bad_standing": false
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from datetime import datetime, timedelta

        # Get user profile
        profile_result = supabase.table("user_profiles").select("*").eq("id", request.user_id).single().execute()
        if not profile_result.data:
            return jsonify({"error": "User profile not found"}), 404

        profile = profile_result.data

        # Check if permanently banned (3+ strikes)
        if profile.get('permanent_bad_standing', False):
            return jsonify({
                "good_standing": False,
                "reason": "dmca_strikes",
                "dmca_strikes": profile.get('dmca_strikes', 0),
                "permanent_bad_standing": True,
                "message": "Account permanently ineligible for Good Standing due to DMCA violations"
            }), 200

        # Check for active subscription (always good standing)
        # TODO: Add subscription check when Stripe integration is added
        # For now, assume no active subscription

        # Check last verified upload (30-day window)
        last_upload = profile.get('last_verified_upload_date')
        if last_upload:
            last_upload_dt = datetime.fromisoformat(last_upload.replace('Z', '+00:00'))
            days_since = (datetime.now(last_upload_dt.tzinfo) - last_upload_dt).days
            days_until_expiry = max(0, 30 - days_since)

            if days_since < 30:
                return jsonify({
                    "good_standing": True,
                    "reason": "recent_upload",
                    "days_until_expiry": days_until_expiry,
                    "last_verified_upload": last_upload,
                    "dmca_strikes": profile.get('dmca_strikes', 0),
                    "permanent_bad_standing": False
                }), 200

        # No recent uploads = bad standing
        return jsonify({
            "good_standing": False,
            "reason": "no_recent_uploads",
            "days_until_expiry": 0,
            "last_verified_upload": last_upload,
            "dmca_strikes": profile.get('dmca_strikes', 0),
            "permanent_bad_standing": False,
            "message": "Upload a verified note to regain Good Standing"
        }), 200

    except Exception as e:
        debug_log(f"Get good standing status error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/good-standing/verify-upload", methods=["POST"])
@supabase_auth_required
def verify_upload():
    """
    Verify if uploaded content is legitimate study material (not spam/blank)

    Request body:
    {
        "note_id": "uuid",
        "extracted_text": "text content from note"
    }

    Returns:
    {
        "verified": true/false,
        "confidence": 0.95,
        "rejection_reason": null | "spam" | "blank" | "invalid_format"
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        import google.generativeai as genai

        data = request.get_json()
        note_id = data.get('note_id')
        extracted_text = data.get('extracted_text', '')

        if not note_id:
            return jsonify({"error": "note_id required"}), 400

        # Quick validation: Minimum length check
        if len(extracted_text.strip()) < 50:
            # Log rejection
            supabase.table("upload_verifications").insert({
                "user_id": request.user_id,
                "note_id": note_id,
                "verification_status": "rejected",
                "rejection_reason": "blank",
                "ai_confidence": 0.0
            }).execute()

            return jsonify({
                "verified": False,
                "confidence": 0.0,
                "rejection_reason": "blank",
                "message": "Document appears to be blank or too short"
            }), 200

        # AI verification with Gemini
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        model = genai.GenerativeModel('gemini-3.1-flash-lite-preview')

        prompt = f"""Analyze this document to determine if it's legitimate study material.

DOCUMENT TEXT (first 2000 chars):
{extracted_text[:2000]}

Respond with ONLY a JSON object:
{{
    "is_study_material": true/false,
    "confidence": 0.0-1.0,
    "reason": "brief explanation",
    "category": "notes" | "textbook" | "slides" | "homework" | "spam" | "blank" | "other"
}}

Study material includes: class notes, textbook chapters, lecture slides, study guides, homework, assignments, etc.
NOT study material: grocery lists, personal emails, blank pages, random text, advertisements."""

        response = model.generate_content(prompt)
        ai_text = response.text.strip()

        try:
            from StudyFlow.backend.cost_tracker import track_ai_call
            track_ai_call("gemini", "flash-lite", "upload_verify")
        except:
            pass

        # Parse JSON response
        import json
        import re
        if ai_text.startswith("```"):
            ai_text = re.sub(r'^```json?\s*', '', ai_text)
            ai_text = re.sub(r'\s*```$', '', ai_text)

        result = json.loads(ai_text)

        is_verified = result.get('is_study_material', False)
        confidence = result.get('confidence', 0.5)
        category = result.get('category', 'other')

        # Determine rejection reason
        rejection_reason = None
        if not is_verified:
            if category == 'spam':
                rejection_reason = 'spam'
            elif category == 'blank':
                rejection_reason = 'blank'
            else:
                rejection_reason = 'invalid_format'

        # Log verification
        supabase.table("upload_verifications").insert({
            "user_id": request.user_id,
            "note_id": note_id,
            "verification_status": "verified" if is_verified else "rejected",
            "rejection_reason": rejection_reason,
            "ai_confidence": confidence
        }).execute()

        # If verified, update user's Good Standing status
        if is_verified:
            from datetime import datetime
            supabase.table("user_profiles").update({
                "good_standing": True,
                "last_verified_upload_date": datetime.utcnow().isoformat()
            }).eq("id", request.user_id).execute()

            print(f"[GOOD STANDING] User {request.user_id} verified upload, status updated", flush=True)

        return jsonify({
            "verified": is_verified,
            "confidence": confidence,
            "rejection_reason": rejection_reason,
            "message": f"Document {('verified' if is_verified else 'rejected')} as {category}"
        }), 200

    except Exception as e:
        debug_log(f"Verify upload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/download/<note_id>", methods=["POST"])
@supabase_auth_required
def download_note_with_certification(note_id):
    """
    Download a note with Missouri HB 2271 certification

    Request body:
    {
        "certified": true  # User certifies they're using for academic research/tutorial only
    }

    Returns:
    {
        "download_url": "https://...",
        "filename": "Biology_Chapter_3.pdf"
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        certified = data.get('certified', False)

        if not certified:
            return jsonify({"error": "Academic integrity certification required"}), 400

        # Check if user is in Good Standing OR has active subscription
        profile_result = supabase.table("user_profiles").select("good_standing, permanent_bad_standing").eq("id", request.user_id).single().execute()
        profile = profile_result.data

        # Block if permanently banned
        if profile.get('permanent_bad_standing', False):
            return jsonify({"error": "Account suspended due to DMCA violations"}), 403

        # Check good standing
        if not profile.get('good_standing', False):
            return jsonify({
                "error": "Good Standing required to download notes",
                "message": "Upload a verified note or subscribe to regain access",
                "good_standing_required": True
            }), 403

        # Get note details
        note_result = supabase.table("notes").select("*").eq("id", note_id).single().execute()
        if not note_result.data:
            return jsonify({"error": "Note not found"}), 404

        note = note_result.data

        # Check if note is public or owned by user
        if not note.get('is_public', False) and note.get('user_id') != request.user_id:
            return jsonify({"error": "Unauthorized access to private note"}), 403

        # Record certification (HB 2271 compliance)
        ip_address = request.headers.get('X-Forwarded-For', request.remote_addr)
        if ip_address and ',' in ip_address:
            ip_address = ip_address.split(',')[0].strip()

        supabase.table("download_certifications").insert({
            "user_id": request.user_id,
            "note_id": note_id,
            "ip_address": ip_address,
            "user_agent": request.headers.get('User-Agent')
        }).execute()

        # Notify the uploader that their note was downloaded
        try:
            note_owner_id = note.get('user_id')
            if note_owner_id and note_owner_id != request.user_id:
                # Get downloader username
                dl_profile = supabase.table("user_profiles").select("username").eq("id", request.user_id).execute()
                dl_username = dl_profile.data[0].get("username", "Someone") if dl_profile.data else "Someone"
                fname = note.get('original_filename', 'a note')
                create_notification(
                    user_id=note_owner_id,
                    notif_type="note_downloaded",
                    title="Note Downloaded",
                    message=f"@{dl_username} downloaded your note {fname}",
                    note_id=note_id,
                    actor_username=dl_username,
                )
        except Exception as notif_err:
            debug_log(f"[-] Download notification failed (non-blocking): {notif_err}")

        # Generate signed URL for download
        file_path = note.get('file_path')
        if not file_path:
            return jsonify({"error": "Note file not found"}), 404

        download_url = supabase.storage.from_('notes').create_signed_url(file_path, 3600)  # 1 hour expiry

        return jsonify({
            "download_url": download_url,
            "filename": note.get('original_filename', 'note.pdf')
        }), 200

    except Exception as e:
        debug_log(f"Download note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/dmca/report", methods=["POST"])
def submit_dmca_report():
    """
    Submit a DMCA takedown notice (does not require authentication)

    Request body:
    {
        "note_id": "uuid",
        "reporter_email": "professor@university.edu",
        "reporter_name": "Dr. John Doe",
        "reason": "This contains my copyrighted lecture slides from Fall 2025"
    }

    Returns:
    {
        "success": true,
        "message": "DMCA report submitted. Note will be reviewed within 24 hours."
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        note_id = data.get('note_id')
        reporter_email = data.get('reporter_email')
        reporter_name = data.get('reporter_name')
        reason = data.get('reason')

        if not all([note_id, reporter_email, reason]):
            return jsonify({"error": "note_id, reporter_email, and reason required"}), 400

        # Get note to find uploader and filename
        note_result = supabase.table("notes").select("user_id, original_filename").eq("id", note_id).single().execute()
        if not note_result.data:
            return jsonify({"error": "Note not found"}), 404

        uploader_id = note_result.data['user_id']
        note_filename = note_result.data.get('original_filename', 'Unknown file')

        # Insert DMCA report
        takedown_result = supabase.table("dmca_takedowns").insert({
            "note_id": note_id,
            "uploader_id": uploader_id,
            "reporter_email": reporter_email,
            "reporter_name": reporter_name,
            "takedown_reason": reason
        }).execute()

        takedown_id = takedown_result.data[0]['id'] if takedown_result.data else None

        print(f"[DMCA] Takedown report submitted for note {note_id} by {reporter_email}", flush=True)

        # Send email notification to admin
        if takedown_id:
            send_dmca_report_notification_to_admin(
                takedown_id=takedown_id,
                note_filename=note_filename,
                reporter_email=reporter_email,
                reason=reason
            )

        return jsonify({
            "success": True,
            "message": "DMCA report submitted. The note will be reviewed and removed if valid. You will receive an email confirmation within 24 hours."
        }), 200

    except Exception as e:
        debug_log(f"Submit DMCA report error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/dmca/reports", methods=["GET"])
def list_dmca_reports():
    """ADMIN ONLY: List all DMCA reports."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        reports = supabase.table("dmca_takedowns").select("*").order("created_at", desc=True).execute()

        results = []
        for r in (reports.data or []):
            # Get note info
            note_info = {}
            if r.get("note_id"):
                note_resp = supabase.table("notes").select(
                    "original_filename, university, course_code, is_public"
                ).eq("id", r["note_id"]).execute()
                if note_resp.data:
                    note_info = note_resp.data[0]

            # Get uploader username
            uploader_name = "Unknown"
            if r.get("uploader_id"):
                profile = supabase.table("user_profiles").select("username, email").eq("id", r["uploader_id"]).execute()
                if profile.data:
                    uploader_name = profile.data[0].get("username") or profile.data[0].get("email", "Unknown")

            results.append({
                "id": r["id"],
                "note_id": r.get("note_id"),
                "note_filename": note_info.get("original_filename", "Unknown"),
                "note_university": note_info.get("university"),
                "note_course": note_info.get("course_code"),
                "note_is_public": note_info.get("is_public", True),
                "uploader_id": r.get("uploader_id"),
                "uploader_name": uploader_name,
                "reporter_email": r.get("reporter_email"),
                "reporter_name": r.get("reporter_name"),
                "reason": r.get("takedown_reason"),
                "status": r.get("status", "pending"),
                "created_at": r.get("created_at")
            })

        return jsonify({"reports": results}), 200

    except Exception as e:
        debug_log(f"List DMCA reports error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/dmca/process/<takedown_id>", methods=["POST"])
def process_dmca_takedown(takedown_id):
    """
    ADMIN ONLY: Process a DMCA takedown (apply strike, remove note)

    Request body:
    {
        "admin_key": "secret_key",
        "action": "approve" | "reject"
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from datetime import datetime

        data = request.get_json()
        admin_key = data.get('admin_key')
        action = data.get('action')

        # Simple admin key check (replace with proper auth)
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        if action not in ['approve', 'reject']:
            return jsonify({"error": "Invalid action"}), 400

        # Get takedown record
        takedown_result = supabase.table("dmca_takedowns").select("*").eq("id", takedown_id).single().execute()
        if not takedown_result.data:
            return jsonify({"error": "Takedown not found"}), 404

        takedown = takedown_result.data
        note_id = takedown['note_id']
        uploader_id = takedown['uploader_id']
        reporter_email = takedown.get('reporter_email')
        reporter_name = takedown.get('reporter_name')

        # Get note filename
        note_result = supabase.table("notes").select("original_filename").eq("id", note_id).single().execute()
        note_filename = note_result.data.get('original_filename', 'Unknown file') if note_result.data else 'Unknown file'

        # Get uploader's email and username
        uploader_result = supabase.table("user_profiles").select("email, username, dmca_strikes").eq("id", uploader_id).single().execute()
        if not uploader_result.data:
            return jsonify({"error": "Uploader profile not found"}), 404

        uploader_email = uploader_result.data.get('email')
        uploader_name = uploader_result.data.get('username', 'User')
        current_strikes = uploader_result.data.get('dmca_strikes', 0)

        if action == 'approve':
            # Mark note as private (soft delete)
            supabase.table("notes").update({"is_public": False}).eq("id", note_id).execute()

            # Apply strike to uploader
            new_strikes = current_strikes + 1

            update_data = {"dmca_strikes": new_strikes}

            # Three strikes = permanent bad standing
            is_permanent_ban = new_strikes >= 3
            if is_permanent_ban:
                update_data["permanent_bad_standing"] = True
                update_data["good_standing"] = False

            supabase.table("user_profiles").update(update_data).eq("id", uploader_id).execute()

            # Mark takedown as processed
            supabase.table("dmca_takedowns").update({
                "processed_at": datetime.utcnow().isoformat(),
                "strike_applied": True
            }).eq("id", takedown_id).execute()

            print(f"[DMCA] Strike applied to user {uploader_id} ({new_strikes}/3)", flush=True)

            # Send email to user about strike
            if uploader_email:
                send_dmca_strike_notification_to_user(
                    user_email=uploader_email,
                    user_name=uploader_name,
                    strike_count=new_strikes,
                    note_filename=note_filename,
                    is_permanent_ban=is_permanent_ban
                )

            # Send confirmation to reporter
            if reporter_email:
                send_dmca_report_confirmation_to_reporter(
                    reporter_email=reporter_email,
                    reporter_name=reporter_name,
                    note_filename=note_filename,
                    action='approve'
                )

            return jsonify({
                "success": True,
                "message": f"Note removed, strike applied ({new_strikes}/3)",
                "permanent_ban": is_permanent_ban
            }), 200

        else:  # reject
            supabase.table("dmca_takedowns").update({
                "processed_at": datetime.utcnow().isoformat(),
                "strike_applied": False
            }).eq("id", takedown_id).execute()

            # Send rejection confirmation to reporter
            if reporter_email:
                send_dmca_report_confirmation_to_reporter(
                    reporter_email=reporter_email,
                    reporter_name=reporter_name,
                    note_filename=note_filename,
                    action='reject'
                )

            return jsonify({
                "success": True,
                "message": "DMCA report rejected, no action taken"
            }), 200

    except Exception as e:
        debug_log(f"Process DMCA takedown error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ SUBSCRIPTION ENDPOINTS ============

@app.route("/api/subscription/status", methods=["GET"])
@supabase_auth_required
def get_subscription_status():
    """
    Get user's current subscription status
    Returns subscription_tier and subscription_status from user_profiles
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        response = supabase.table("user_profiles").select(
            "subscription_tier, subscription_status, stripe_customer_id, stripe_subscription_id"
        ).eq("id", request.user_id).single().execute()

        if response.data:
            return jsonify(response.data), 200
        else:
            return jsonify({
                "subscription_tier": "free",
                "subscription_status": None,
                "stripe_customer_id": None,
                "stripe_subscription_id": None
            }), 200

    except Exception as e:
        debug_log(f"Get subscription status error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/downloads/daily-status", methods=["GET"])
@supabase_auth_required
def get_daily_download_status():
    """
    Get user's daily download count and limit
    Returns download count, limit, and whether user is Scholar's Club member
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        from datetime import date

        response = supabase.table("user_profiles").select(
            "subscription_tier, subscription_status, daily_downloads_count, daily_downloads_reset_date"
        ).eq("id", request.user_id).single().execute()

        if not response.data:
            return jsonify({"error": "User not found"}), 404

        profile = response.data

        # Check if Scholar's Club member (unlimited downloads)
        is_scholar = (
            profile.get('subscription_tier') == 'pro' and
            profile.get('subscription_status') == 'active'
        )

        # Reset counter if new day
        today = date.today()
        reset_date = profile.get('daily_downloads_reset_date')
        download_count = profile.get('daily_downloads_count', 0) or 0

        if not is_scholar and (not reset_date or str(reset_date) != str(today)):
            # Reset for new day
            supabase.table("user_profiles").update({
                "daily_downloads_count": 0,
                "daily_downloads_reset_date": str(today)
            }).eq("id", request.user_id).execute()
            download_count = 0

        return jsonify({
            "is_scholar": is_scholar,
            "downloads_used": download_count if not is_scholar else None,
            "downloads_limit": 3 if not is_scholar else None,
            "downloads_remaining": (3 - download_count) if not is_scholar else None,
            "reset_date": str(today) if not is_scholar else None
        }), 200

    except Exception as e:
        debug_log(f"Get daily download status error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/subscription/create-checkout", methods=["POST"])
@supabase_auth_required
def create_checkout_session():
    """
    Create a Stripe Checkout session for the Scholar's Club subscription ($4.99/mo)

    Request body:
    {
        "success_url": "https://...",
        "cancel_url": "https://..."
    }

    Returns:
    {
        "checkout_url": "https://checkout.stripe.com/..."
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        import stripe

        data = request.get_json()
        success_url = data.get('success_url')
        cancel_url = data.get('cancel_url')

        # Get user email
        user_email = request.user_email

        # Get or create Stripe customer
        response = supabase.table("user_profiles").select(
            "stripe_customer_id"
        ).eq("id", request.user_id).single().execute()

        stripe_customer_id = response.data.get('stripe_customer_id') if response.data else None

        if stripe_customer_id:
            # Retrieve existing customer
            customer = stripe.Customer.retrieve(stripe_customer_id)
        else:
            # Create new customer
            customer = stripe.Customer.create(
                email=user_email,
                metadata={'user_id': request.user_id}
            )
            stripe_customer_id = customer.id

            # Save customer ID to database
            supabase.table("user_profiles").update({
                "stripe_customer_id": stripe_customer_id
            }).eq("id", request.user_id).execute()

        # Create Stripe Checkout session for Scholar's Club ($4.99/mo)
        checkout_session = stripe.checkout.Session.create(
            customer=stripe_customer_id,
            payment_method_types=['card'],
            line_items=[{
                'price': 'price_1TEJOa9LWKaKRffVgNpExTIz',  # Scholar's Club - $4.99/month
                'quantity': 1,
            }],
            mode='subscription',
            success_url=success_url,
            cancel_url=cancel_url,
            metadata={
                'user_id': request.user_id
            }
        )

        debug_log(f"Created checkout session for user {request.user_id}: {checkout_session.id}")

        return jsonify({
            "checkout_url": checkout_session.url
        }), 200

    except Exception as e:
        debug_log(f"Create checkout session error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/subscription/customer-portal", methods=["POST"])
@supabase_auth_required
def create_customer_portal():
    """
    Create a Stripe Customer Portal session for managing subscriptions

    Request body:
    {
        "return_url": "https://..."
    }

    Returns:
    {
        "portal_url": "https://billing.stripe.com/..."
    }
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        import stripe

        data = request.get_json()
        return_url = data.get('return_url')

        # Get Stripe customer ID
        response = supabase.table("user_profiles").select(
            "stripe_customer_id"
        ).eq("id", request.user_id).single().execute()

        stripe_customer_id = response.data.get('stripe_customer_id') if response.data else None

        if not stripe_customer_id:
            return jsonify({"error": "No subscription found"}), 404

        # Create Customer Portal session
        portal_session = stripe.billing_portal.Session.create(
            customer=stripe_customer_id,
            return_url=return_url
        )

        debug_log(f"Created portal session for user {request.user_id}")

        return jsonify({
            "portal_url": portal_session.url
        }), 200

    except Exception as e:
        debug_log(f"Create customer portal error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ NOTIFICATION ENDPOINTS ============

def create_notification(user_id, notif_type, title, message, note_id=None, actor_username=None):
    """Helper to insert a notification row."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        row = {
            "user_id": user_id,
            "type": notif_type,
            "title": title,
            "message": message,
            "is_read": False,
        }
        if note_id:
            row["note_id"] = note_id
        if actor_username:
            row["actor_username"] = actor_username
        supabase.table("notifications").insert(row).execute()
    except Exception as e:
        debug_log(f"[-] Failed to create notification: {e}")


@app.route("/api/notifications", methods=["GET"])
@supabase_auth_required
def get_notifications():
    """Return the last 50 notifications for the authenticated user."""
    try:
        cache_key = f"notifications:{request.user_id}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        from StudyFlow.backend.supabase_client import supabase
        result = supabase.table("notifications") \
            .select("*") \
            .eq("user_id", request.user_id) \
            .order("created_at", desc=True) \
            .limit(50) \
            .execute()
        resp = {"notifications": result.data or []}
        if redis_cache:
            try: redis_cache.setex(cache_key, 120, json.dumps(resp))
            except: pass
        return jsonify(resp), 200
    except Exception as e:
        debug_log(f"Get notifications error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notifications/read", methods=["POST"])
@supabase_auth_required
def mark_notifications_read():
    """Mark a single notification or all notifications as read."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        data = request.get_json()

        if data.get("all"):
            supabase.table("notifications") \
                .update({"is_read": True}) \
                .eq("user_id", request.user_id) \
                .eq("is_read", False) \
                .execute()
        elif data.get("notification_id"):
            supabase.table("notifications") \
                .update({"is_read": True}) \
                .eq("id", data["notification_id"]) \
                .eq("user_id", request.user_id) \
                .execute()
        else:
            return jsonify({"error": "Provide notification_id or all:true"}), 400

        return jsonify({"success": True}), 200
    except Exception as e:
        debug_log(f"Mark notifications read error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ============ NOTE FAVORITES ENDPOINTS ============

@app.route("/api/notes/favorite", methods=["POST"])
@supabase_auth_required
def favorite_note():
    """Add a note to the user's favorites."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        data = request.get_json()
        note_id = data.get("note_id")
        if not note_id:
            return jsonify({"error": "Missing note_id"}), 400

        supabase.table("note_favorites").upsert({
            "user_id": request.user_id,
            "note_id": note_id
        }, on_conflict="user_id,note_id").execute()

        return jsonify({"success": True}), 200
    except Exception as e:
        debug_log(f"Favorite note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/favorite", methods=["DELETE"])
@supabase_auth_required
def unfavorite_note():
    """Remove a note from the user's favorites."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        data = request.get_json()
        note_id = data.get("note_id")
        if not note_id:
            return jsonify({"error": "Missing note_id"}), 400

        supabase.table("note_favorites") \
            .delete() \
            .eq("user_id", request.user_id) \
            .eq("note_id", note_id) \
            .execute()

        return jsonify({"success": True}), 200
    except Exception as e:
        debug_log(f"Unfavorite note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/favorites", methods=["GET"])
@supabase_auth_required
def get_favorite_notes():
    """Get all favorited note IDs for the user."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        result = supabase.table("note_favorites") \
            .select("note_id, created_at") \
            .eq("user_id", request.user_id) \
            .order("created_at", desc=True) \
            .execute()

        return jsonify({"favorites": result.data or []}), 200
    except Exception as e:
        debug_log(f"Get favorites error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/annotations", methods=["GET"])
@supabase_auth_required
def get_annotations(note_id):
    """Get all annotations for a note belonging to the current user."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        result = supabase.table("note_annotations") \
            .select("page_number, annotations_json") \
            .eq("user_id", request.user_id) \
            .eq("note_id", note_id) \
            .order("page_number") \
            .execute()

        annotations = []
        for row in (result.data or []):
            annotations.append({
                "page_number": row["page_number"],
                "objects": row["annotations_json"].get("objects", []) if isinstance(row["annotations_json"], dict) else []
            })

        return jsonify({"annotations": annotations}), 200
    except Exception as e:
        debug_log(f"Get annotations error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/annotations", methods=["POST"])
@supabase_auth_required
def save_annotations(note_id):
    """Upsert annotations per page for a note. Deletes pages with empty objects.
    If save_as_new=true, duplicates the note first and saves annotations to the copy."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        data = request.get_json()
        pages = data.get("annotations", [])
        save_as_new = data.get("save_as_new", False)

        if not isinstance(pages, list):
            return jsonify({"error": "annotations must be a list"}), 400

        target_note_id = note_id

        # If saving as a new note, create a new note record pointing to the SAME file
        # (annotations are stored separately, watermarks are applied during download)
        if save_as_new:
            try:
                # Fetch original note record
                original = supabase.table("notes") \
                    .select("*") \
                    .eq("id", note_id) \
                    .eq("user_id", request.user_id) \
                    .single() \
                    .execute()

                if not original.data:
                    return jsonify({"error": "Original note not found"}), 404

                orig = original.data

                # Create new note ID but reference the SAME file_path
                # This preserves the watermark (applied during download) and saves storage
                import uuid as uuid_lib
                new_note_id = str(uuid_lib.uuid4())

                # Build the new filename
                orig_filename = orig.get("original_filename", "Untitled")
                name_part = orig_filename.rsplit('.', 1)[0] if '.' in orig_filename else orig_filename
                new_filename = f"{name_part} (annotated).pdf"

                # Create new note record pointing to SAME file
                new_note_data = {
                    "id": new_note_id,
                    "user_id": request.user_id,
                    "original_filename": new_filename,
                    "file_type": orig.get("file_type", "pdf"),
                    "file_size": orig.get("file_size", 0),
                    "file_path": orig.get("file_path"),  # Same file path - no duplication
                    "page_count": orig.get("page_count", 1),
                    "processed": orig.get("processed", False),
                    "is_public": False,
                    "username": orig.get("username"),
                    "university": orig.get("university"),
                    "course_code": orig.get("course_code"),
                    "professor": orig.get("professor"),
                    "semester": orig.get("semester"),
                }
                supabase.table("notes").insert(new_note_data).execute()

                target_note_id = new_note_id
                debug_log(f"[+] Created annotated note {new_note_id} referencing original {note_id} (shared file)")

            except Exception as dup_err:
                debug_log(f"Save as new - creation error: {dup_err}\n{traceback.format_exc()}")
                return jsonify({"error": f"Failed to create annotated note: {str(dup_err)}"}), 500

        # Save annotations to the target note
        for page in pages:
            page_number = page.get("page_number")
            objects = page.get("objects", [])

            if page_number is None:
                continue

            if not objects:
                # Delete empty annotation rows
                supabase.table("note_annotations") \
                    .delete() \
                    .eq("user_id", request.user_id) \
                    .eq("note_id", target_note_id) \
                    .eq("page_number", page_number) \
                    .execute()
            else:
                # Upsert annotation data
                supabase.table("note_annotations").upsert({
                    "user_id": request.user_id,
                    "note_id": target_note_id,
                    "page_number": page_number,
                    "annotations_json": {"objects": objects},
                    "updated_at": "now()"
                }, on_conflict="user_id,note_id,page_number").execute()

        result = {"success": True}
        if save_as_new:
            result["new_note_id"] = target_note_id

        return jsonify(result), 200
    except Exception as e:
        debug_log(f"Save annotations error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/content", methods=["PUT"])
@supabase_auth_required
def update_note_content(note_id):
    """
    Update the text content of a note.
    Auth required, verifies note ownership.
    Accepts JSON body: { "content": "the new text" }
    Re-uploads text as bytes to existing file_path in Supabase Storage.
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        if not data or "content" not in data:
            return jsonify({"error": "Missing 'content' in request body"}), 400

        new_content = data["content"]

        # Get note and verify ownership
        note_response = supabase.table("notes").select(
            "id, user_id, file_path"
        ).eq("id", note_id).execute()

        if not note_response.data:
            return jsonify({"error": "Note not found"}), 404

        note_data = note_response.data[0]

        if note_data.get("user_id") != request.user_id:
            return jsonify({"error": "Not authorized to edit this note"}), 403

        file_path = note_data.get("file_path")
        if not file_path:
            return jsonify({"error": "No file path for this note"}), 404

        # Re-upload: delete old file then upload new content
        content_bytes = new_content.encode("utf-8")

        # Detect HTML content for proper content-type
        import re
        is_html = bool(re.search(r'<[a-z][\s\S]*>', new_content[:500], re.IGNORECASE))
        content_type = "text/html; charset=utf-8" if is_html else "text/plain; charset=utf-8"

        try:
            supabase.storage.from_("note-files").remove([file_path])
        except Exception as remove_err:
            debug_log(f"Warning: could not remove old file: {remove_err}")

        supabase.storage.from_("note-files").upload(
            path=file_path,
            file=content_bytes,
            file_options={"content-type": content_type}
        )

        # Update file_size in the notes table
        supabase.table("notes").update({
            "file_size": len(content_bytes)
        }).eq("id", note_id).execute()

        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Update note content error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/upload-image", methods=["POST"])
@supabase_auth_required
def upload_note_image():
    """
    Upload an image for embedding in a rich-text note.
    Accepts multipart file upload. Returns a signed URL for the image.
    """
    try:
        import uuid as _uuid
        from StudyFlow.backend.supabase_client import supabase

        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400

        file = request.files["file"]
        if not file.filename:
            return jsonify({"error": "No filename"}), 400

        # Validate image type
        allowed = {"image/png", "image/jpeg", "image/gif", "image/webp"}
        content_type = file.content_type or ""
        if content_type not in allowed:
            return jsonify({"error": "Only PNG, JPEG, GIF, and WebP images are allowed"}), 400

        # 5MB limit
        file_bytes = file.read()
        if len(file_bytes) > 5 * 1024 * 1024:
            return jsonify({"error": "Image must be under 5MB"}), 400

        # Convert to WebP for smaller file sizes (skip GIFs to preserve animation)
        from PIL import Image
        from io import BytesIO

        if content_type != "image/gif":
            img = Image.open(BytesIO(file_bytes))
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGBA")
            else:
                img = img.convert("RGB")
            buf = BytesIO()
            img.save(buf, format="WEBP", quality=85)
            file_bytes = buf.getvalue()
            content_type = "image/webp"

        ext = "gif" if content_type == "image/gif" else "webp"
        image_id = str(_uuid.uuid4())
        storage_path = f"note-images/{request.user_id}/{image_id}.{ext}"

        supabase.storage.from_("note-files").upload(
            path=storage_path,
            file=file_bytes,
            file_options={"content-type": content_type}
        )

        # Generate a long-lived signed URL (7 days)
        signed = supabase.storage.from_("note-files").create_signed_url(
            path=storage_path,
            expires_in=604800
        )
        url = signed.get("signedURL") or signed.get("signedUrl")

        return jsonify({"success": True, "url": url}), 200

    except Exception as e:
        debug_log(f"Upload note image error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/<note_id>/raw-content", methods=["GET"])
@supabase_auth_required
def get_note_raw_content(note_id):
    """
    Get the raw file content of a note (authenticated, owner only).
    Returns the file bytes directly with appropriate content-type.
    """
    try:
        from StudyFlow.backend.supabase_client import supabase

        note_response = supabase.table("notes").select(
            "id, user_id, file_path, original_filename"
        ).eq("id", note_id).execute()

        if not note_response.data:
            return jsonify({"error": "Note not found"}), 404

        note_data = note_response.data[0]

        if note_data.get("user_id") != request.user_id:
            return jsonify({"error": "Not authorized"}), 403

        file_path = note_data.get("file_path")
        if not file_path:
            return jsonify({"error": "No file path"}), 404

        try:
            file_bytes = supabase.storage.from_("note-files").download(file_path)
        except Exception:
            # File might be empty (new note) -- return empty string
            return "", 200, {"Content-Type": "text/html; charset=utf-8"}

        content_type = "text/html; charset=utf-8" if file_path.endswith(".html") else "text/plain; charset=utf-8"
        return file_bytes, 200, {"Content-Type": content_type}

    except Exception as e:
        debug_log(f"Get note raw content error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/create", methods=["POST"])
@supabase_auth_required
def create_note():
    """
    Create a new empty rich-text note.
    Accepts optional JSON body: { "title": "My Note" }
    Defaults title to "Untitled" if not provided.
    """
    try:
        import uuid as _uuid
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json(silent=True) or {}
        title = (data.get("title") or "").strip() or "Untitled"

        note_id = str(_uuid.uuid4())
        user_id = request.user_id
        file_path = f"{user_id}/{note_id}.html"

        # Upload empty HTML file to Supabase Storage
        empty_html = b""
        supabase.storage.from_("note-files").upload(
            path=file_path,
            file=empty_html,
            file_options={"content-type": "text/html; charset=utf-8"}
        )

        # Insert row into notes table
        note_row = {
            "id": note_id,
            "user_id": user_id,
            "original_filename": title + ".html",
            "file_type": "txt",
            "file_path": file_path,
            "file_size": 0,
            "is_public": False,
            "page_count": 1,
        }
        supabase.table("notes").insert(note_row).execute()

        return jsonify({"success": True, "note_id": note_id}), 201

    except Exception as e:
        debug_log(f"Create note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ──────────────────────────────────────────────
# ONLYOFFICE Private Office Endpoints
# ──────────────────────────────────────────────

ONLYOFFICE_URL = os.environ.get("ONLYOFFICE_URL", "")
ONLYOFFICE_JWT_SECRET = os.environ.get("ONLYOFFICE_JWT_SECRET", "")

OFFICE_CONTENT_TYPES = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

OFFICE_DOCUMENT_TYPES = {
    "docx": "word",
    "xlsx": "cell",
    "pptx": "slide",
}


def _create_blank_document(file_type):
    """Generate a minimal blank office document and return bytes."""
    buf = BytesIO()
    if file_type == "docx":
        from docx import Document as DocxDocument
        doc = DocxDocument()
        doc.add_paragraph("")
        doc.save(buf)
    elif file_type == "xlsx":
        from openpyxl import Workbook
        wb = Workbook()
        wb.save(buf)
    elif file_type == "pptx":
        from pptx import Presentation
        prs = Presentation()
        prs.slide_layouts[6]  # blank layout reference
        prs.slides.add_slide(prs.slide_layouts[6])
        prs.save(buf)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
    buf.seek(0)
    return buf.read()


@app.route("/api/office/create", methods=["POST"])
@supabase_auth_required
def office_create():
    """Create a new blank office document (docx/xlsx/pptx)."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        import uuid

        data = request.get_json()
        if not data or "title" not in data or "type" not in data:
            return jsonify({"error": "Missing 'title' or 'type'"}), 400

        file_type = data["type"].lower()
        if file_type not in OFFICE_CONTENT_TYPES:
            return jsonify({"error": f"Unsupported type: {file_type}. Use docx, xlsx, or pptx."}), 400

        title = data["title"].strip() or "Untitled"
        doc_id = str(uuid.uuid4())
        file_path = f"office-temp/{request.user_id}/{doc_id}.{file_type}"

        # Generate blank document
        file_bytes = _create_blank_document(file_type)

        # Upload to Supabase Storage (temporary location, not in notes table)
        supabase.storage.from_("note-files").upload(
            path=file_path,
            file=file_bytes,
            file_options={"content-type": OFFICE_CONTENT_TYPES[file_type]},
        )

        # Store in office_documents table (temporary documents, not notes)
        office_doc = {
            "id": doc_id,
            "user_id": request.user_id,
            "title": title,
            "file_path": file_path,
            "file_type": file_type,
            "file_size": len(file_bytes),
        }
        supabase.table("office_documents").insert(office_doc).execute()

        debug_log(f"[Office] Created temporary {file_type} doc '{title}' for user {request.user_id}")
        return jsonify({"success": True, "doc_id": doc_id}), 201

    except Exception as e:
        debug_log(f"Office create error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/<doc_id>/editor-config", methods=["GET"])
@supabase_auth_required
def office_editor_config(doc_id):
    """Return ONLYOFFICE editor config with signed JWT."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        import jwt

        # Get office document and verify ownership
        doc_resp = supabase.table("office_documents").select(
            "id, user_id, file_path, file_type, file_size, title"
        ).eq("id", doc_id).execute()

        if not doc_resp.data:
            return jsonify({"error": "Document not found"}), 404

        doc = doc_resp.data[0]
        is_owner = doc["user_id"] == request.user_id
        share_permission = "edit"  # owners always get edit

        if not is_owner:
            share_resp = supabase.table("office_document_shares").select(
                "permission"
            ).eq("doc_id", doc_id).eq("shared_with_user_id", request.user_id).execute()
            if not share_resp.data:
                return jsonify({"error": "Not authorized"}), 403
            share_permission = share_resp.data[0]["permission"]

        file_type = doc["file_type"]
        title = doc.get("title", "Untitled")

        # Generate 1-hour signed URL
        signed = supabase.storage.from_("note-files").create_signed_url(
            doc["file_path"], 3600
        )
        doc_url = signed.get("signedURL") or signed.get("signedUrl", "")

        # Build document key (changes when file changes to bust ONLYOFFICE cache)
        doc_key = f"{doc_id}_{doc.get('file_size', 0)}"

        callback_url = f"{BACKEND_URL}/api/office/callback?doc_id={doc_id}"

        config = {
            "document": {
                "fileType": file_type,
                "key": doc_key,
                "title": f"{title}.{file_type}",
                "url": doc_url,
                "permissions": {
                    "edit": share_permission == "edit",
                    "download": True,
                    "print": True,
                    "chat": False,
                    "comment": False,
                },
            },
            "documentType": OFFICE_DOCUMENT_TYPES.get(file_type, "word"),
            "editorConfig": {
                "callbackUrl": callback_url,
                "user": {
                    "id": request.user_id,
                    "name": request.user_email or "User",
                },
                "customization": {
                    "autosave": True,
                    "forcesave": True,
                    "compactHeader": True,
                    "customer": {
                        "name": "StudyFlow",
                        "address": "StudyFlow Suite",
                        "mail": "info@studyflowsuite.com",
                        "www": "studyflowsuite.com",
                        "info": "StudyFlow Office",
                        "logo": "https://studyflowsuite.com/icon128.png",
                        "logoDark": "https://studyflowsuite.com/icon128.png",
                    },
                    "logo": {
                        "image": "https://studyflowsuite.com/icon128.png",
                        "imageDark": "https://studyflowsuite.com/icon128.png",
                        "imageEmbedded": "https://studyflowsuite.com/icon128.png",
                        "url": "https://studyflowsuite.com",
                    },
                    "feedback": {
                        "visible": False,
                    },
                    "goback": {
                        "url": "https://studyflowsuite.com/office.html",
                        "text": "Back to Office",
                    },
                },
                "mode": "edit" if share_permission == "edit" else "view",
            },
        }

        # Add StudyFlow Assistant plugin for word documents
        if config["documentType"] == "word":
            config["editorConfig"]["plugins"] = {
                "autostart": ["asc.{7C9885A0-0001-0001-0001-000000000001}"],
                "pluginsData": ["https://studyflowsuite.com/plugins/study-assistant/config.json"]
            }

        # Sign the config JWT for ONLYOFFICE
        token = ""
        if ONLYOFFICE_JWT_SECRET:
            token = jwt.encode(config, ONLYOFFICE_JWT_SECRET, algorithm="HS256")
            config["token"] = token

        debug_log(f"[Office] Editor config for doc {doc_id}, type={file_type}")
        return jsonify({
            "config": config,
            "token": token,
            "docServerUrl": ONLYOFFICE_URL,
        }), 200

    except Exception as e:
        debug_log(f"Office editor-config error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/callback", methods=["POST"])
def office_callback():
    """
    ONLYOFFICE callback handler (server-to-server, no Supabase auth).
    Called by ONLYOFFICE when document status changes.
    Status 2 = closed after editing, 6 = force save.
    """
    try:
        from StudyFlow.backend.supabase_client import supabase
        import jwt

        doc_id = request.args.get("doc_id")
        if not doc_id:
            return jsonify({"error": 0})

        body = request.get_json(force=True)
        debug_log(f"[Office Callback] doc_id={doc_id}, status={body.get('status')}")

        # Verify ONLYOFFICE JWT if secret is configured
        if ONLYOFFICE_JWT_SECRET:
            token = body.get("token")
            if not token:
                auth_header = request.headers.get("Authorization", "")
                if auth_header.startswith("Bearer "):
                    token = auth_header[7:]
            if token:
                try:
                    jwt.decode(token, ONLYOFFICE_JWT_SECRET, algorithms=["HS256"])
                except jwt.InvalidTokenError as jwt_err:
                    debug_log(f"[Office Callback] JWT verification failed: {jwt_err}")
                    return jsonify({"error": 0})

        status = body.get("status")
        download_url = body.get("url")

        # Status 2 = closed after editing, 6 = force save
        if status in (2, 6) and download_url:
            # Download the saved file from ONLYOFFICE temp URL
            resp = requests.get(download_url, timeout=30)
            if resp.status_code == 200:
                file_bytes = resp.content

                # Get office document to find storage path
                doc_resp = supabase.table("office_documents").select(
                    "file_path"
                ).eq("id", doc_id).execute()

                if doc_resp.data:
                    file_path = doc_resp.data[0]["file_path"]
                    file_ext = file_path.rsplit(".", 1)[-1]
                    content_type = OFFICE_CONTENT_TYPES.get(file_ext, "application/octet-stream")

                    # Remove old file and upload new version
                    try:
                        supabase.storage.from_("note-files").remove([file_path])
                    except Exception as rm_err:
                        debug_log(f"[Office Callback] Remove old file warning: {rm_err}")

                    supabase.storage.from_("note-files").upload(
                        path=file_path,
                        file=file_bytes,
                        file_options={"content-type": content_type},
                    )

                    # Update file size in office_documents table
                    supabase.table("office_documents").update({
                        "file_size": len(file_bytes),
                        "updated_at": "now()"
                    }).eq("id", doc_id).execute()

                    debug_log(f"[Office Callback] Saved {len(file_bytes)} bytes for doc {doc_id}")
            else:
                debug_log(f"[Office Callback] Download failed: HTTP {resp.status_code}")

        return jsonify({"error": 0})

    except Exception as e:
        debug_log(f"Office callback error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": 0})


@app.route("/api/office/list", methods=["GET"])
@supabase_auth_required
def office_list():
    """List all office documents for the authenticated user."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        resp = supabase.table("office_documents").select(
            "id, title, file_type, file_size, created_at"
        ).eq("user_id", request.user_id).order("created_at", desc=True).execute()

        docs = []
        for row in resp.data:
            docs.append({
                "id": row["id"],
                "title": row.get("title", "Untitled"),
                "type": row["file_type"],
                "size": row.get("file_size", 0),
                "created_at": row.get("created_at", ""),
            })

        return jsonify({"documents": docs}), 200

    except Exception as e:
        debug_log(f"Office list error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/<doc_id>", methods=["DELETE"])
@supabase_auth_required
def office_delete(doc_id):
    """Delete an office document."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        doc_resp = supabase.table("office_documents").select(
            "id, user_id, file_path"
        ).eq("id", doc_id).execute()

        if not doc_resp.data:
            return jsonify({"error": "Document not found"}), 404

        doc = doc_resp.data[0]
        if doc["user_id"] != request.user_id:
            return jsonify({"error": "Not authorized"}), 403

        # Remove from storage
        try:
            supabase.storage.from_("note-files").remove([doc["file_path"]])
        except Exception as rm_err:
            debug_log(f"[Office] Storage remove warning: {rm_err}")

        # Delete row
        supabase.table("office_documents").delete().eq("id", doc_id).execute()

        debug_log(f"[Office] Deleted doc {doc_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Office delete error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/<doc_id>/rename", methods=["PATCH"])
@supabase_auth_required
def office_rename(doc_id):
    """Rename an office document."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        if not data or "title" not in data:
            return jsonify({"error": "Missing 'title'"}), 400

        doc_resp = supabase.table("office_documents").select(
            "id, user_id"
        ).eq("id", doc_id).execute()

        if not doc_resp.data:
            return jsonify({"error": "Document not found"}), 404

        if doc_resp.data[0]["user_id"] != request.user_id:
            return jsonify({"error": "Not authorized"}), 403

        supabase.table("office_documents").update({
            "title": data["title"].strip()
        }).eq("id", doc_id).execute()

        debug_log(f"[Office] Renamed doc {doc_id} to '{data['title']}'")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Office rename error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ──────────────────────────────────────────────
# Office Document Sharing Endpoints
# ──────────────────────────────────────────────

@app.route("/api/office/search-users", methods=["GET"])
@supabase_auth_required
def office_search_users():
    """Search users by username for sharing documents."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        q = request.args.get("q", "").strip()
        if len(q) < 2:
            return jsonify({"users": []}), 200

        # Escape SQL wildcards
        safe_q = q.replace("%", "").replace("_", "")
        result = supabase.table("user_profiles").select(
            "username"
        ).ilike("username", f"%{safe_q}%").neq("id", request.user_id).not_.is_("username", "null").limit(10).execute()

        users = [{"username": u["username"]} for u in (result.data or [])]
        return jsonify({"users": users}), 200

    except Exception as e:
        debug_log(f"Office search-users error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/<doc_id>/share", methods=["POST"])
@supabase_auth_required
def office_share(doc_id):
    """Share a document with a user by username. Owner only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        data = request.get_json()
        if not data or not data.get("username"):
            return jsonify({"error": "Missing username"}), 400

        username = data["username"].strip()
        permission = data.get("permission", "view")
        if permission not in ("view", "edit"):
            return jsonify({"error": "Permission must be 'view' or 'edit'"}), 400

        # Verify document ownership
        doc_resp = supabase.table("office_documents").select(
            "id, user_id"
        ).eq("id", doc_id).execute()

        if not doc_resp.data:
            return jsonify({"error": "Document not found"}), 404
        if doc_resp.data[0]["user_id"] != request.user_id:
            return jsonify({"error": "Only the document owner can share"}), 403

        # Look up target user
        user_resp = supabase.table("user_profiles").select(
            "id, username"
        ).eq("username", username).execute()

        if not user_resp.data:
            return jsonify({"error": "User not found"}), 404

        target_user_id = user_resp.data[0]["id"]
        if target_user_id == request.user_id:
            return jsonify({"error": "Cannot share with yourself"}), 400

        # Upsert share record
        supabase.table("office_document_shares").upsert({
            "doc_id": doc_id,
            "owner_id": request.user_id,
            "shared_with_user_id": target_user_id,
            "permission": permission
        }, on_conflict="doc_id,shared_with_user_id").execute()

        debug_log(f"[Office] Shared doc {doc_id} with @{username} ({permission})")

        # Notify the target user
        try:
            doc_title = doc_resp.data[0].get("title", "Untitled")
            # Get owner username
            owner_profile = supabase.table("user_profiles").select("username").eq("id", request.user_id).execute()
            owner_name = owner_profile.data[0]["username"] if owner_profile.data else "Someone"

            create_notification(
                user_id=target_user_id,
                notif_type="doc_shared",
                title="Document Shared",
                message=f"@{owner_name} shared '{doc_title}' with you ({permission})",
            )
        except Exception:
            pass  # Non-blocking

        return jsonify({"success": True, "shared_with": {"username": username, "permission": permission}}), 200

    except Exception as e:
        debug_log(f"Office share error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/<doc_id>/shares", methods=["GET"])
@supabase_auth_required
def office_list_shares(doc_id):
    """List all users a document is shared with. Owner only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        # Verify ownership
        doc_resp = supabase.table("office_documents").select(
            "id, user_id"
        ).eq("id", doc_id).execute()

        if not doc_resp.data:
            return jsonify({"error": "Document not found"}), 404
        if doc_resp.data[0]["user_id"] != request.user_id:
            return jsonify({"error": "Not authorized"}), 403

        shares_resp = supabase.table("office_document_shares").select(
            "shared_with_user_id, permission, created_at"
        ).eq("doc_id", doc_id).execute()

        shares = []
        for s in (shares_resp.data or []):
            profile = supabase.table("user_profiles").select("username").eq("id", s["shared_with_user_id"]).execute()
            username = profile.data[0]["username"] if profile.data else "Unknown"
            shares.append({
                "user_id": s["shared_with_user_id"],
                "username": username,
                "permission": s["permission"],
                "created_at": s["created_at"]
            })

        return jsonify({"shares": shares}), 200

    except Exception as e:
        debug_log(f"Office list-shares error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/<doc_id>/share/<user_id>", methods=["DELETE"])
@supabase_auth_required
def office_revoke_share(doc_id, user_id):
    """Revoke a share. Owner only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        # Verify ownership
        doc_resp = supabase.table("office_documents").select(
            "id, user_id"
        ).eq("id", doc_id).execute()

        if not doc_resp.data:
            return jsonify({"error": "Document not found"}), 404
        if doc_resp.data[0]["user_id"] != request.user_id:
            return jsonify({"error": "Not authorized"}), 403

        supabase.table("office_document_shares").delete().eq(
            "doc_id", doc_id
        ).eq("shared_with_user_id", user_id).execute()

        debug_log(f"[Office] Revoked share for doc {doc_id}, user {user_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Office revoke-share error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/office/shared-with-me", methods=["GET"])
@supabase_auth_required
def office_shared_with_me():
    """List documents shared with the current user."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        shares_resp = supabase.table("office_document_shares").select(
            "doc_id, permission, owner_id"
        ).eq("shared_with_user_id", request.user_id).order("created_at", desc=True).execute()

        if not shares_resp.data:
            return jsonify({"documents": []}), 200

        documents = []
        for s in shares_resp.data:
            doc_resp = supabase.table("office_documents").select(
                "id, title, file_type, file_size, created_at"
            ).eq("id", s["doc_id"]).execute()

            if not doc_resp.data:
                continue

            doc = doc_resp.data[0]

            # Get owner username
            owner_profile = supabase.table("user_profiles").select("username").eq("id", s["owner_id"]).execute()
            owner_username = owner_profile.data[0]["username"] if owner_profile.data else "Unknown"

            documents.append({
                "id": doc["id"],
                "title": doc.get("title", "Untitled"),
                "file_type": doc["file_type"],
                "file_size": doc.get("file_size", 0),
                "created_at": doc.get("created_at"),
                "permission": s["permission"],
                "owner_username": owner_username
            })

        return jsonify({"documents": documents}), 200

    except Exception as e:
        debug_log(f"Office shared-with-me error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ──────────────────────────────────────────────
# Study Groups Endpoints
# ──────────────────────────────────────────────

@app.route("/api/groups/discover", methods=["GET"])
@supabase_auth_required
def discover_groups():
    """Get public groups the user hasn't joined yet."""
    try:
        user_id = request.user_id

        # Get groups user is already in
        my_groups = supabase.table("study_group_members").select("group_id").eq("user_id", user_id).execute()
        my_group_ids = [g["group_id"] for g in (my_groups.data or [])]

        # Get public groups
        query = supabase.table("study_groups").select(
            "id, name, description, cover_url, member_count, category, owner_id, created_at"
        ).eq("is_public", True).order("member_count", desc=True).limit(20)

        result = query.execute()
        groups = [g for g in (result.data or []) if g["id"] not in my_group_ids]

        # Get owner usernames
        owner_ids = list(set(g["owner_id"] for g in groups))
        owners_map = {}
        if owner_ids:
            try:
                owners = supabase.table("user_profiles").select("id, username").in_("id", owner_ids).execute()
                owners_map = {o["id"]: o.get("username", "Unknown") for o in (owners.data or [])}
            except:
                pass

        for g in groups:
            g["owner_username"] = owners_map.get(g["owner_id"], "Unknown")

        return jsonify({"groups": groups}), 200

    except Exception as e:
        debug_log(f"Discover groups error: {e}\n{traceback.format_exc()}")
        return jsonify({"groups": []}), 200


@app.route("/api/groups/<group_id>/feed", methods=["GET"])
@supabase_auth_required
def get_group_feed(group_id):
    """Get posts from a specific group."""
    try:
        user_id = request.user_id

        # Verify membership
        member = supabase.table("study_group_members").select("id").eq("group_id", group_id).eq("user_id", user_id).execute()
        if not member.data:
            return jsonify({"error": "Not a member of this group"}), 403

        # Get posts for this group
        posts = supabase.table("social_posts").select("*").eq("group_id", group_id).order("created_at", desc=True).limit(50).execute()

        post_list = posts.data or []

        # Enrich with note details and avatars
        if post_list:
            note_ids = [p["note_id"] for p in post_list if p.get("note_id")]
            notes_map = {}
            if note_ids:
                try:
                    notes_res = supabase.table("notes").select("id, original_filename, thumbnail_url, file_type").in_("id", note_ids).execute()
                    notes_map = {n["id"]: n for n in (notes_res.data or [])}
                except:
                    pass

            author_ids = list(set(p["user_id"] for p in post_list))
            avatars_map = {}
            try:
                av_res = supabase.table("user_profiles").select("id, avatar_url, is_verified").in_("id", author_ids).execute()
                avatars_map = {a["id"]: {"avatar_url": a.get("avatar_url"), "is_verified": a.get("is_verified", False)} for a in (av_res.data or [])}
            except:
                pass

            for post in post_list:
                author_info = avatars_map.get(post["user_id"], {})
                post["avatar_url"] = author_info.get("avatar_url")
                post["is_verified"] = author_info.get("is_verified", False)
                if post.get("note_id"):
                    post["notes"] = notes_map.get(post["note_id"])

        return jsonify({"posts": post_list}), 200

    except Exception as e:
        debug_log(f"Group feed error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/post", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def create_group_post(group_id):
    """Create a post within a group."""
    try:
        user_id = request.user_id

        # Verify membership
        member = supabase.table("study_group_members").select("id").eq("group_id", group_id).eq("user_id", user_id).execute()
        if not member.data:
            return jsonify({"error": "Not a member of this group"}), 403

        data = request.json
        text_content = (data.get("text_content") or "").strip()
        note_id = data.get("note_id")

        if not text_content and not note_id:
            return jsonify({"error": "Post must have text or a note"}), 400

        profile = supabase.table("user_profiles").select("username").eq("id", user_id).execute()
        username = profile.data[0].get("username", "Anonymous") if profile.data else "Anonymous"

        post_type = "note" if note_id else "text"
        post_data = {
            "user_id": user_id,
            "username": username,
            "post_type": post_type,
            "text_content": text_content if text_content else None,
            "note_id": note_id,
            "group_id": group_id
        }

        result = supabase.table("social_posts").insert(post_data).execute()

        if not result.data:
            return jsonify({"error": "Failed to create post"}), 500

        return jsonify({"message": "Posted to group", "post": result.data[0]}), 201

    except Exception as e:
        debug_log(f"Group post error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/join", methods=["POST"])
@supabase_auth_required
def join_public_group(group_id):
    """Join a public group directly."""
    try:
        user_id = request.user_id

        # Check group is public
        group = supabase.table("study_groups").select("id, is_public, name").eq("id", group_id).execute()
        if not group.data:
            return jsonify({"error": "Group not found"}), 404
        if not group.data[0].get("is_public", True):
            return jsonify({"error": "This group is private. You need an invite."}), 403

        # Join
        try:
            supabase.table("study_group_members").insert({
                "group_id": group_id,
                "user_id": user_id,
                "role": "member"
            }).execute()
        except Exception as e:
            if "unique" in str(e).lower() or "duplicate" in str(e).lower():
                return jsonify({"error": "Already a member"}), 409
            raise

        # Update member count
        try:
            members = supabase.table("study_group_members").select("id", count="exact").eq("group_id", group_id).execute()
            supabase.table("study_groups").update({"member_count": members.count or 1}).eq("id", group_id).execute()
        except:
            pass

        return jsonify({"success": True, "group_name": group.data[0]["name"]}), 200

    except Exception as e:
        debug_log(f"Join group error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups", methods=["POST"])
@supabase_auth_required
def create_group():
    """Create a new study group."""
    try:
        from StudyFlow.backend.supabase_client import supabase
        import secrets

        data = request.get_json()
        if not data or not data.get("name"):
            return jsonify({"error": "Missing group name"}), 400

        name = data["name"].strip()[:60]
        invite_code = secrets.token_urlsafe(8)
        group_id = str(uuid.uuid4())

        supabase.table("study_groups").insert({
            "id": group_id,
            "name": name,
            "owner_id": request.user_id,
            "invite_code": invite_code
        }).execute()

        # Add owner as member
        supabase.table("study_group_members").insert({
            "group_id": group_id,
            "user_id": request.user_id,
            "role": "owner"
        }).execute()

        debug_log(f"[Groups] Created group '{name}' ({group_id})")
        return jsonify({"success": True, "group": {"id": group_id, "name": name, "invite_code": invite_code}}), 200

    except Exception as e:
        debug_log(f"Create group error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups", methods=["GET"])
@supabase_auth_required
def list_groups():
    """List all groups the user is a member of."""
    try:
        cache_key = f"user_groups:{request.user_id}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        from StudyFlow.backend.supabase_client import supabase

        memberships = supabase.table("study_group_members").select(
            "group_id, role"
        ).eq("user_id", request.user_id).execute()

        if not memberships.data:
            return jsonify({"groups": []}), 200

        groups = []
        for m in memberships.data:
            g = supabase.table("study_groups").select("*").eq("id", m["group_id"]).execute()
            if g.data:
                group = g.data[0]
                # Get member count
                members = supabase.table("study_group_members").select("id", count="exact").eq("group_id", m["group_id"]).execute()
                group["member_count"] = members.count if members.count else 0
                group["role"] = m["role"]
                groups.append(group)

        resp = {"groups": groups}
        if redis_cache:
            try: redis_cache.setex(cache_key, 300, json.dumps(resp))
            except: pass
        return jsonify(resp), 200

    except Exception as e:
        debug_log(f"List groups error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>", methods=["GET"])
@supabase_auth_required
def get_group(group_id):
    """Get group details including members and notes."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        # Verify membership
        member_check = supabase.table("study_group_members").select("role").eq(
            "group_id", group_id
        ).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        group = supabase.table("study_groups").select("*").eq("id", group_id).execute()
        if not group.data:
            return jsonify({"error": "Group not found"}), 404

        # Get members with usernames
        members_resp = supabase.table("study_group_members").select("user_id, role, joined_at").eq("group_id", group_id).execute()
        members = []
        for m in (members_resp.data or []):
            profile = supabase.table("user_profiles").select("username").eq("id", m["user_id"]).execute()
            members.append({
                "user_id": m["user_id"],
                "username": profile.data[0]["username"] if profile.data else "Unknown",
                "role": m["role"],
                "joined_at": m["joined_at"]
            })

        # Get shared notes (direct uploads)
        notes_resp = supabase.table("study_group_notes").select(
            "id, added_by, added_at, filename, file_path, file_size, file_type"
        ).eq("group_id", group_id).order("added_at", desc=True).execute()
        notes = []
        for n in (notes_resp.data or []):
            profile = supabase.table("user_profiles").select("username").eq("id", n["added_by"]).execute()
            notes.append({
                "id": n["id"],
                "filename": n.get("filename", "Unknown"),
                "file_size": n.get("file_size", 0),
                "file_type": n.get("file_type"),
                "added_by": profile.data[0]["username"] if profile.data else "Unknown",
                "added_at": n["added_at"]
            })

        result = group.data[0]
        result["members"] = members
        result["notes"] = notes
        result["role"] = member_check.data[0]["role"]

        return jsonify({"group": result}), 200

    except Exception as e:
        debug_log(f"Get group error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>", methods=["DELETE"])
@supabase_auth_required
def delete_group(group_id):
    """Delete a study group. Owner only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        group = supabase.table("study_groups").select("owner_id").eq("id", group_id).execute()
        if not group.data:
            return jsonify({"error": "Group not found"}), 404
        if group.data[0]["owner_id"] != request.user_id:
            return jsonify({"error": "Only the owner can delete this group"}), 403

        supabase.table("study_groups").delete().eq("id", group_id).execute()
        debug_log(f"[Groups] Deleted group {group_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Delete group error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/invite", methods=["POST"])
@supabase_auth_required
def invite_to_group(group_id):
    """Send a pending invite to a user. They must accept via notifications."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        # Verify caller is a member
        member_check = supabase.table("study_group_members").select("role").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        data = request.get_json()
        username = data.get("username", "").strip()
        if not username:
            return jsonify({"error": "Missing username"}), 400

        user_resp = supabase.table("user_profiles").select("id, username").eq("username", username).execute()
        if not user_resp.data:
            return jsonify({"error": "User not found"}), 404

        target_id = user_resp.data[0]["id"]
        if target_id == request.user_id:
            return jsonify({"error": "You are already in this group"}), 400

        # Check not already a member
        existing = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", target_id).execute()
        if existing.data:
            return jsonify({"error": "User is already a member"}), 400

        # Check for existing pending invite
        existing_invite = supabase.table("study_group_invites").select("id, status").eq(
            "group_id", group_id).eq("invited_user_id", target_id).execute()
        if existing_invite.data:
            if existing_invite.data[0]["status"] == "pending":
                return jsonify({"error": "Invite already pending"}), 400
            # Update declined invite back to pending
            supabase.table("study_group_invites").update({
                "status": "pending",
                "invited_by": request.user_id
            }).eq("id", existing_invite.data[0]["id"]).execute()
        else:
            supabase.table("study_group_invites").insert({
                "group_id": group_id,
                "invited_by": request.user_id,
                "invited_user_id": target_id,
                "status": "pending"
            }).execute()

        # Send notification
        group_resp = supabase.table("study_groups").select("name").eq("id", group_id).execute()
        group_name = group_resp.data[0]["name"] if group_resp.data else "a study group"
        inviter_profile = supabase.table("user_profiles").select("username").eq("id", request.user_id).execute()
        inviter_name = inviter_profile.data[0]["username"] if inviter_profile.data else "Someone"

        create_notification(
            user_id=target_id,
            notif_type="group_invite",
            title="Study Group Invite",
            message=f"@{inviter_name} invited you to '{group_name}'",
        )

        debug_log(f"[Groups] Invited @{username} to group {group_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Invite to group error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/invites", methods=["GET"])
@supabase_auth_required
def get_pending_invites():
    """Get all pending group invites for the current user."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        invites_resp = supabase.table("study_group_invites").select(
            "id, group_id, invited_by, created_at"
        ).eq("invited_user_id", request.user_id).eq("status", "pending").order("created_at", desc=True).execute()

        invites = []
        for inv in (invites_resp.data or []):
            group = supabase.table("study_groups").select("name").eq("id", inv["group_id"]).execute()
            inviter = supabase.table("user_profiles").select("username").eq("id", inv["invited_by"]).execute()
            invites.append({
                "id": inv["id"],
                "group_id": inv["group_id"],
                "group_name": group.data[0]["name"] if group.data else "Unknown",
                "invited_by": inviter.data[0]["username"] if inviter.data else "Unknown",
                "created_at": inv["created_at"]
            })

        return jsonify({"invites": invites}), 200

    except Exception as e:
        debug_log(f"Get invites error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/invites/<invite_id>/accept", methods=["POST"])
@supabase_auth_required
def accept_invite(invite_id):
    """Accept a group invite."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        invite = supabase.table("study_group_invites").select(
            "id, group_id, invited_user_id, status"
        ).eq("id", invite_id).execute()

        if not invite.data:
            return jsonify({"error": "Invite not found"}), 404
        if invite.data[0]["invited_user_id"] != request.user_id:
            return jsonify({"error": "Not your invite"}), 403
        if invite.data[0]["status"] != "pending":
            return jsonify({"error": "Invite already responded to"}), 400

        group_id = invite.data[0]["group_id"]

        # Add to group
        supabase.table("study_group_members").insert({
            "group_id": group_id,
            "user_id": request.user_id,
            "role": "member"
        }).execute()

        # Update invite status
        supabase.table("study_group_invites").update({
            "status": "accepted"
        }).eq("id", invite_id).execute()

        debug_log(f"[Groups] User {request.user_id} accepted invite to group {group_id}")
        return jsonify({"success": True, "group_id": group_id}), 200

    except Exception as e:
        debug_log(f"Accept invite error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/invites/<invite_id>/decline", methods=["POST"])
@supabase_auth_required
def decline_invite(invite_id):
    """Decline a group invite."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        invite = supabase.table("study_group_invites").select(
            "id, invited_user_id, status"
        ).eq("id", invite_id).execute()

        if not invite.data:
            return jsonify({"error": "Invite not found"}), 404
        if invite.data[0]["invited_user_id"] != request.user_id:
            return jsonify({"error": "Not your invite"}), 403

        supabase.table("study_group_invites").update({
            "status": "declined"
        }).eq("id", invite_id).execute()

        debug_log(f"[Groups] User {request.user_id} declined invite {invite_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Decline invite error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/join/<invite_code>", methods=["POST"])
@supabase_auth_required
def join_group(invite_code):
    """Join a group via invite code."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        group = supabase.table("study_groups").select("id, name").eq("invite_code", invite_code).execute()
        if not group.data:
            return jsonify({"error": "Invalid invite code"}), 404

        group_id = group.data[0]["id"]

        existing = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if existing.data:
            return jsonify({"error": "Already a member", "group_id": group_id}), 400

        supabase.table("study_group_members").insert({
            "group_id": group_id,
            "user_id": request.user_id,
            "role": "member"
        }).execute()

        debug_log(f"[Groups] User {request.user_id} joined group {group_id} via invite")
        return jsonify({"success": True, "group_id": group_id, "group_name": group.data[0]["name"]}), 200

    except Exception as e:
        debug_log(f"Join group error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/members/<user_id>", methods=["DELETE"])
@supabase_auth_required
def remove_group_member(group_id, user_id):
    """Remove a member from a group. Owner or self only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        group = supabase.table("study_groups").select("owner_id").eq("id", group_id).execute()
        if not group.data:
            return jsonify({"error": "Group not found"}), 404

        is_owner = group.data[0]["owner_id"] == request.user_id
        is_self = user_id == request.user_id

        if not is_owner and not is_self:
            return jsonify({"error": "Not authorized"}), 403

        if is_owner and is_self:
            return jsonify({"error": "Owner cannot leave. Delete the group instead."}), 400

        supabase.table("study_group_members").delete().eq(
            "group_id", group_id).eq("user_id", user_id).execute()

        debug_log(f"[Groups] Removed user {user_id} from group {group_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Remove member error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/notes/upload", methods=["POST"])
@supabase_auth_required
def upload_group_note(group_id):
    """Upload a file directly to the group. Does not appear in personal notes or Nexus."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        member_check = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        if "file" not in request.files:
            return jsonify({"error": "No file provided"}), 400

        file = request.files["file"]
        if not file.filename:
            return jsonify({"error": "Empty filename"}), 400

        filename = file.filename
        file_bytes = file.read()
        file_size = len(file_bytes)

        # Determine file type
        ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "bin"
        content_type = file.content_type or "application/octet-stream"

        # Upload to Supabase storage under group-notes path
        storage_path = f"group-notes/{group_id}/{uuid.uuid4()}_{filename}"
        supabase.storage.from_("note-files").upload(
            storage_path, file_bytes,
            file_options={"content-type": content_type}
        )

        # Save record
        supabase.table("study_group_notes").insert({
            "group_id": group_id,
            "added_by": request.user_id,
            "filename": filename,
            "file_path": storage_path,
            "file_size": file_size,
            "file_type": ext
        }).execute()

        debug_log(f"[Groups] Uploaded '{filename}' to group {group_id}")
        return jsonify({"success": True, "filename": filename}), 200

    except Exception as e:
        debug_log(f"Upload group note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/notes/<note_id>", methods=["DELETE"])
@supabase_auth_required
def remove_group_note(group_id, note_id):
    """Remove a note from the group. Adder or owner only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        group = supabase.table("study_groups").select("owner_id").eq("id", group_id).execute()
        if not group.data:
            return jsonify({"error": "Group not found"}), 404

        note_entry = supabase.table("study_group_notes").select("added_by, file_path").eq(
            "group_id", group_id).eq("id", note_id).execute()
        if not note_entry.data:
            return jsonify({"error": "Note not in group"}), 404

        is_owner = group.data[0]["owner_id"] == request.user_id
        is_adder = note_entry.data[0]["added_by"] == request.user_id

        if not is_owner and not is_adder:
            return jsonify({"error": "Not authorized"}), 403

        # Delete file from storage if it has a file_path
        file_path = note_entry.data[0].get("file_path")
        if file_path:
            try:
                supabase.storage.from_("note-files").remove([file_path])
            except Exception:
                pass

        supabase.table("study_group_notes").delete().eq(
            "group_id", group_id).eq("id", note_id).execute()

        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Remove group note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/notes/<note_id>/download", methods=["GET"])
@supabase_auth_required
def download_group_note(group_id, note_id):
    """Get a signed download URL for a group note."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        member_check = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        note = supabase.table("study_group_notes").select("file_path, filename").eq(
            "group_id", group_id).eq("id", note_id).execute()
        if not note.data or not note.data[0].get("file_path"):
            return jsonify({"error": "Note not found"}), 404

        signed = supabase.storage.from_("note-files").create_signed_url(
            note.data[0]["file_path"], 3600
        )
        url = signed.get("signedURL") or signed.get("signedUrl", "")

        return jsonify({"url": url, "filename": note.data[0]["filename"]}), 200

    except Exception as e:
        debug_log(f"Download group note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/messages/<message_id>/react", methods=["POST"])
@supabase_auth_required
def react_to_message(group_id, message_id):
    """Toggle a reaction on a group chat message."""
    try:
        data = request.get_json()
        emoji = data.get("emoji")
        if not emoji:
            return jsonify({"error": "emoji required"}), 400

        # Get current reactions
        msg_resp = supabase.table("study_group_messages").select("reactions").eq("id", message_id).eq("group_id", group_id).single().execute()
        if not msg_resp.data:
            return jsonify({"error": "Message not found"}), 404

        reactions = msg_resp.data.get("reactions") or {}
        user_id = request.user_id

        # Toggle: add or remove user from this emoji
        if emoji in reactions:
            users = reactions[emoji]
            if user_id in users:
                users.remove(user_id)
                if not users:
                    del reactions[emoji]
            else:
                users.append(user_id)
        else:
            reactions[emoji] = [user_id]

        # Save
        supabase.table("study_group_messages").update({"reactions": reactions}).eq("id", message_id).execute()

        return jsonify({"success": True, "reactions": reactions}), 200

    except Exception as e:
        debug_log(f"React error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/chat", methods=["POST"])
@supabase_auth_required
def group_chat(group_id):
    """Send a message in the group chat. Users only, no AI."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        # Verify membership
        member_check = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        data = request.get_json()
        message = data.get("message", "").strip()
        if not message:
            return jsonify({"error": "Missing message"}), 400

        # Get sender username
        sender_profile = supabase.table("user_profiles").select("username").eq("id", request.user_id).execute()
        sender_username = sender_profile.data[0]["username"] if sender_profile.data else "User"

        # Save user message
        supabase.table("study_group_messages").insert({
            "group_id": group_id,
            "user_id": request.user_id,
            "role": "user",
            "content": message
        }).execute()

        # Invalidate cached messages
        if redis_cache:
            try: redis_cache.delete(f"group_msgs:{group_id}")
            except: pass

        return jsonify({
            "success": True,
            "sender_username": sender_username
        }), 200

    except Exception as e:
        debug_log(f"Group chat error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/messages", methods=["GET"])
@supabase_auth_required
def get_group_messages(group_id):
    """Get chat history for a group."""
    try:
        # Check Redis cache
        cache_key = f"group_msgs:{group_id}"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached: return jsonify(json.loads(cached)), 200
            except: pass

        from StudyFlow.backend.supabase_client import supabase

        member_check = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        messages_resp = supabase.table("study_group_messages").select(
            "id, user_id, role, content, sources, created_at, pinned, reactions"
        ).eq("group_id", group_id).order("created_at").limit(100).execute()

        messages = []
        username_cache = {}
        for m in (messages_resp.data or []):
            username = None
            if m["user_id"]:
                if m["user_id"] not in username_cache:
                    profile = supabase.table("user_profiles").select("username").eq("id", m["user_id"]).execute()
                    username_cache[m["user_id"]] = profile.data[0]["username"] if profile.data else "Unknown"
                username = username_cache[m["user_id"]]

            messages.append({
                "id": m["id"],
                "role": m["role"],
                "content": m["content"],
                "sources": m.get("sources", []),
                "username": username,
                "user_id": m["user_id"],
                "created_at": m["created_at"],
                "pinned": m.get("pinned", False),
                "reactions": m.get("reactions", {})
            })

        result = {"messages": messages}
        if redis_cache:
            try: redis_cache.setex(cache_key, 60, json.dumps(result))  # 1 min for chat
            except: pass
        return jsonify(result), 200

    except Exception as e:
        debug_log(f"Get group messages error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ──────────────────────────────────────────────
# Pinned Messages & Study Schedule
# ──────────────────────────────────────────────

@app.route("/api/groups/<group_id>/messages/<message_id>/pin", methods=["POST"])
@supabase_auth_required
def pin_message(group_id, message_id):
    """Pin or unpin a group message. Owner or message sender only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        member_check = supabase.table("study_group_members").select("role").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        msg = supabase.table("study_group_messages").select("id, user_id, pinned").eq(
            "id", message_id).eq("group_id", group_id).execute()
        if not msg.data:
            return jsonify({"error": "Message not found"}), 404

        is_owner = member_check.data[0]["role"] == "owner"
        is_sender = msg.data[0]["user_id"] == request.user_id
        if not is_owner and not is_sender:
            return jsonify({"error": "Not authorized"}), 403

        new_pinned = not msg.data[0].get("pinned", False)
        supabase.table("study_group_messages").update({
            "pinned": new_pinned,
            "pinned_by": request.user_id if new_pinned else None
        }).eq("id", message_id).execute()

        return jsonify({"success": True, "pinned": new_pinned}), 200

    except Exception as e:
        debug_log(f"Pin message error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/pinned", methods=["GET"])
@supabase_auth_required
def get_pinned_messages(group_id):
    """Get all pinned messages for a group."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        member_check = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        pinned = supabase.table("study_group_messages").select(
            "id, user_id, content, created_at"
        ).eq("group_id", group_id).eq("pinned", True).order("created_at", desc=True).execute()

        messages = []
        cache = {}
        for m in (pinned.data or []):
            if m["user_id"] and m["user_id"] not in cache:
                p = supabase.table("user_profiles").select("username").eq("id", m["user_id"]).execute()
                cache[m["user_id"]] = p.data[0]["username"] if p.data else "Unknown"
            messages.append({
                "id": m["id"],
                "content": m["content"],
                "username": cache.get(m["user_id"], "Unknown"),
                "created_at": m["created_at"]
            })

        return jsonify({"pinned": messages}), 200

    except Exception as e:
        debug_log(f"Get pinned error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/events", methods=["GET"])
@supabase_auth_required
def get_group_events(group_id):
    """Get all events for a group, ordered by date."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        member_check = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        events_resp = supabase.table("study_group_events").select(
            "id, title, description, event_date, event_time, created_by, created_at"
        ).eq("group_id", group_id).order("event_date").execute()

        events = []
        cache = {}
        for ev in (events_resp.data or []):
            if ev["created_by"] not in cache:
                p = supabase.table("user_profiles").select("username").eq("id", ev["created_by"]).execute()
                cache[ev["created_by"]] = p.data[0]["username"] if p.data else "Unknown"
            events.append({
                "id": ev["id"],
                "title": ev["title"],
                "description": ev.get("description"),
                "event_date": ev["event_date"],
                "event_time": ev.get("event_time"),
                "created_by": cache[ev["created_by"]],
                "created_by_id": ev["created_by"],
                "created_at": ev["created_at"]
            })

        return jsonify({"events": events}), 200

    except Exception as e:
        debug_log(f"Get events error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/events", methods=["POST"])
@supabase_auth_required
def create_group_event(group_id):
    """Create a study event."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        member_check = supabase.table("study_group_members").select("id").eq(
            "group_id", group_id).eq("user_id", request.user_id).execute()
        if not member_check.data:
            return jsonify({"error": "Not a member"}), 403

        data = request.get_json()
        title = data.get("title", "").strip()
        if not title:
            return jsonify({"error": "Missing title"}), 400
        event_date = data.get("event_date")
        if not event_date:
            return jsonify({"error": "Missing date"}), 400

        row = {
            "group_id": group_id,
            "title": title[:100],
            "event_date": event_date,
            "created_by": request.user_id
        }
        if data.get("description"):
            row["description"] = data["description"].strip()[:500]
        if data.get("event_time"):
            row["event_time"] = data["event_time"]

        supabase.table("study_group_events").insert(row).execute()

        debug_log(f"[Groups] Created event '{title}' in group {group_id}")
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Create event error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/groups/<group_id>/events/<event_id>", methods=["DELETE"])
@supabase_auth_required
def delete_group_event(group_id, event_id):
    """Delete a study event. Creator or group owner only."""
    try:
        from StudyFlow.backend.supabase_client import supabase

        group = supabase.table("study_groups").select("owner_id").eq("id", group_id).execute()
        if not group.data:
            return jsonify({"error": "Group not found"}), 404

        event = supabase.table("study_group_events").select("created_by").eq("id", event_id).eq("group_id", group_id).execute()
        if not event.data:
            return jsonify({"error": "Event not found"}), 404

        is_owner = group.data[0]["owner_id"] == request.user_id
        is_creator = event.data[0]["created_by"] == request.user_id
        if not is_owner and not is_creator:
            return jsonify({"error": "Not authorized"}), 403

        supabase.table("study_group_events").delete().eq("id", event_id).execute()
        return jsonify({"success": True}), 200

    except Exception as e:
        debug_log(f"Delete event error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# ──────────────────────────────────────────────
# Admin: University Stats (Institutional License Prep)
# ──────────────────────────────────────────────

@app.route("/admin/review-queue", methods=["GET"])
def review_queue():
    """ADMIN: Get unreviewed notes for quality control."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        # Get IDs of already-reviewed notes
        reviewed_ids = set()
        try:
            reviewed_resp = supabase.table("reviewed_notes").select("note_id").execute()
            reviewed_ids = set(r["note_id"] for r in (reviewed_resp.data or []))
        except Exception as rev_err:
            debug_log(f"[Review] reviewed_notes query failed (table may not exist): {rev_err}")
            # Continue without filtering -- show all notes

        # Get all notes (newest first)
        notes_resp = supabase.table("notes").select(
            "id, user_id, original_filename, file_size, university, course_code, file_path, uploaded_at, is_public"
        ).order("uploaded_at", desc=True).limit(200).execute()

        # Filter out already reviewed
        unreviewed = [n for n in (notes_resp.data or []) if n["id"] not in reviewed_ids][:50]

        notes = []
        uploader_cache = {}
        for n in unreviewed:
            # Get uploader info (cached)
            uid = n["user_id"]
            if uid not in uploader_cache:
                try:
                    profile = supabase.table("user_profiles").select("username, email").eq("id", uid).execute()
                    uploader_cache[uid] = profile.data[0].get("username") or profile.data[0].get("email", "Unknown") if profile.data else "Unknown"
                except Exception:
                    uploader_cache[uid] = "Unknown"
            uploader = uploader_cache[uid]

            notes.append({
                "id": n["id"],
                "filename": n.get("original_filename", "Unknown"),
                "file_size": n.get("file_size", 0),
                "university": n.get("university"),
                "course_code": n.get("course_code"),
                "uploader": uploader,
                "uploader_id": n["user_id"],
                "is_public": n.get("is_public", True),
                "created_at": n.get("uploaded_at")
            })

        return jsonify({
            "notes": notes,
            "total": len(notes),
            "debug": {
                "total_notes_in_db": len(notes_resp.data or []),
                "reviewed_count": len(reviewed_ids),
                "unreviewed_count": len(unreviewed)
            }
        }), 200

    except Exception as e:
        debug_log(f"Review queue error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/flagged-notes", methods=["GET"])
def flagged_notes():
    """ADMIN: Get notes with net -5 or worse downvotes from chat ratings."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        # Fetch all ratings
        ratings_resp = supabase.table("ai_response_ratings").select("vote, cited_note_ids").execute()
        ratings = ratings_resp.data or []

        # Aggregate votes per note
        vote_map = {}  # note_id -> {up, down}
        for rating in ratings:
            for note_id in (rating.get("cited_note_ids") or []):
                if note_id not in vote_map:
                    vote_map[note_id] = {"up": 0, "down": 0}
                if rating["vote"] == 1:
                    vote_map[note_id]["up"] += 1
                elif rating["vote"] == -1:
                    vote_map[note_id]["down"] += 1

        # Find notes with net -5 or worse
        flagged_ids = []
        flagged_votes = {}
        for note_id, votes in vote_map.items():
            net = votes["up"] - votes["down"]
            if net <= -5:
                flagged_ids.append(note_id)
                flagged_votes[note_id] = votes

        if not flagged_ids:
            return jsonify({"notes": [], "total": 0}), 200

        # Fetch note metadata
        notes_resp = supabase.table("notes").select(
            "id, original_filename, university, course_code, user_id, file_size, uploaded_at"
        ).in_("id", flagged_ids).execute()
        notes_data = notes_resp.data or []

        # Batch fetch uploader usernames
        user_ids = list(set(n["user_id"] for n in notes_data))
        uploader_cache = {}
        if user_ids:
            profiles_resp = supabase.table("user_profiles").select("id, username").in_("id", user_ids).execute()
            for p in (profiles_resp.data or []):
                uploader_cache[p["id"]] = p.get("username") or "Unknown"

        # Build response
        result = []
        for n in notes_data:
            votes = flagged_votes.get(n["id"], {"up": 0, "down": 0})
            result.append({
                "id": n["id"],
                "filename": n.get("original_filename", "Untitled"),
                "university": n.get("university"),
                "course_code": n.get("course_code"),
                "uploader": uploader_cache.get(n["user_id"], "Unknown"),
                "file_size": n.get("file_size", 0),
                "uploaded_at": n.get("uploaded_at"),
                "upvotes": votes["up"],
                "downvotes": votes["down"],
                "net": votes["up"] - votes["down"]
            })

        # Sort worst first
        result.sort(key=lambda x: x["net"])

        debug_log(f"[Admin] Flagged notes: {len(result)} notes with net -5 or worse")
        return jsonify({"notes": result, "total": len(result)}), 200

    except Exception as e:
        debug_log(f"Flagged notes error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/search-analytics", methods=["GET"])
def admin_search_analytics():
    """ADMIN: Search analytics -- top queries, zero-result queries, volume."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.search_tracker import get_search_analytics
        data = get_search_analytics(days=30)

        if not data:
            return jsonify({
                "dates": [], "daily_volume": [],
                "top_searches": [], "zero_results": [],
                "today_top": [], "today_zero": [],
                "summary": {"total_30d": 0, "total_today": 0, "total_7d": 0, "avg_daily": 0},
                "message": "No search data yet. Tracking starts from now."
            }), 200

        return jsonify(data), 200

    except Exception as e:
        debug_log(f"Admin search analytics error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/content-stats", methods=["GET"])
def admin_content_stats():
    """ADMIN: Content stats -- top notes by views, downloads, citations. Top uploaders."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        # Check Redis cache (10 min TTL)
        cache_key = "admin:content_stats"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached:
                    return jsonify(json.loads(cached)), 200
            except:
                pass

        from StudyFlow.backend.supabase_client import supabase

        # --- Most viewed notes ---
        most_viewed = []
        try:
            # Get view counts grouped by note
            views_resp = supabase.table("note_views").select("note_id").execute()
            view_counts = {}
            for v in (views_resp.data or []):
                nid = v['note_id']
                view_counts[nid] = view_counts.get(nid, 0) + 1

            # Get top 15 by views
            top_viewed_ids = sorted(view_counts.keys(), key=lambda x: -view_counts[x])[:15]
            if top_viewed_ids:
                notes_resp = supabase.table("notes").select("id, original_filename, university, course_code, user_id").in_("id", top_viewed_ids).execute()
                notes_map = {n['id']: n for n in (notes_resp.data or [])}

                # Get usernames
                user_ids = list(set(n['user_id'] for n in (notes_resp.data or [])))
                profiles_resp = supabase.table("user_profiles").select("id, username").in_("id", user_ids).execute()
                username_map = {p['id']: p.get('username', 'Anonymous') for p in (profiles_resp.data or [])}

                for nid in top_viewed_ids:
                    note = notes_map.get(nid)
                    if note:
                        most_viewed.append({
                            "id": nid,
                            "filename": note.get('original_filename', 'Untitled'),
                            "university": note.get('university', ''),
                            "course_code": note.get('course_code', ''),
                            "uploader": username_map.get(note['user_id'], 'Anonymous'),
                            "views": view_counts[nid]
                        })
        except:
            pass

        # --- Most downloaded notes ---
        most_downloaded = []
        try:
            dl_resp = supabase.table("download_transactions").select("note_id").execute()
            dl_counts = {}
            for d in (dl_resp.data or []):
                nid = d['note_id']
                dl_counts[nid] = dl_counts.get(nid, 0) + 1

            top_dl_ids = sorted(dl_counts.keys(), key=lambda x: -dl_counts[x])[:15]
            if top_dl_ids:
                notes_resp = supabase.table("notes").select("id, original_filename, university, course_code, user_id").in_("id", top_dl_ids).execute()
                notes_map = {n['id']: n for n in (notes_resp.data or [])}

                user_ids = list(set(n['user_id'] for n in (notes_resp.data or [])))
                profiles_resp = supabase.table("user_profiles").select("id, username").in_("id", user_ids).execute()
                username_map = {p['id']: p.get('username', 'Anonymous') for p in (profiles_resp.data or [])}

                for nid in top_dl_ids:
                    note = notes_map.get(nid)
                    if note:
                        most_downloaded.append({
                            "id": nid,
                            "filename": note.get('original_filename', 'Untitled'),
                            "university": note.get('university', ''),
                            "course_code": note.get('course_code', ''),
                            "uploader": username_map.get(note['user_id'], 'Anonymous'),
                            "downloads": dl_counts[nid]
                        })
        except:
            pass

        # --- Most cited in AI chat (highest net_votes) ---
        most_cited = []
        try:
            cited_resp = supabase.table("notes").select(
                "id, original_filename, university, course_code, user_id, net_votes"
            ).not_.is_("net_votes", "null").neq("net_votes", 0).order("net_votes", desc=True).limit(15).execute()

            if cited_resp.data:
                user_ids = list(set(n['user_id'] for n in cited_resp.data))
                profiles_resp = supabase.table("user_profiles").select("id, username").in_("id", user_ids).execute()
                username_map = {p['id']: p.get('username', 'Anonymous') for p in (profiles_resp.data or [])}

                for n in cited_resp.data:
                    most_cited.append({
                        "id": n['id'],
                        "filename": n.get('original_filename', 'Untitled'),
                        "university": n.get('university', ''),
                        "course_code": n.get('course_code', ''),
                        "uploader": username_map.get(n['user_id'], 'Anonymous'),
                        "net_votes": n.get('net_votes', 0)
                    })
        except:
            pass

        # --- Top uploaders ---
        top_uploaders = []
        try:
            notes_resp = supabase.table("notes").select("user_id").execute()
            upload_counts = {}
            for n in (notes_resp.data or []):
                uid = n['user_id']
                upload_counts[uid] = upload_counts.get(uid, 0) + 1

            top_uploader_ids = sorted(upload_counts.keys(), key=lambda x: -upload_counts[x])[:15]
            if top_uploader_ids:
                profiles_resp = supabase.table("user_profiles").select("id, username, university").in_("id", top_uploader_ids).execute()
                profile_map = {p['id']: p for p in (profiles_resp.data or [])}

                for uid in top_uploader_ids:
                    p = profile_map.get(uid, {})
                    top_uploaders.append({
                        "username": p.get('username', 'Anonymous'),
                        "university": p.get('university', ''),
                        "note_count": upload_counts[uid]
                    })
        except:
            pass

        # --- Summary ---
        total_views = sum(v['views'] for v in most_viewed) if most_viewed else 0
        total_downloads = sum(d['downloads'] for d in most_downloaded) if most_downloaded else 0

        result = {
            "most_viewed": most_viewed,
            "most_downloaded": most_downloaded,
            "most_cited": most_cited,
            "top_uploaders": top_uploaders,
            "summary": {
                "total_views_tracked": total_views,
                "total_downloads_tracked": total_downloads
            }
        }

        # Cache for 10 minutes
        if redis_cache:
            try:
                redis_cache.setex(cache_key, 600, json.dumps(result))
            except:
                pass

        return jsonify(result), 200

    except Exception as e:
        debug_log(f"Admin content stats error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/ai-costs", methods=["GET"])
def admin_ai_costs():
    """ADMIN: Get AI API cost data for dashboard."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.cost_tracker import get_cost_summary
        data = get_cost_summary(days=30)

        if not data:
            return jsonify({
                "dates": [], "daily_totals": [],
                "month_total": 0, "last_month_total": 0,
                "seven_day_total": 0, "today_breakdown": [],
                "today_calls": [], "provider_totals": [],
                "message": "No cost data yet. Tracking starts from now."
            }), 200

        return jsonify(data), 200

    except Exception as e:
        debug_log(f"Admin AI costs error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/announcements", methods=["GET"])
def list_announcements():
    """ADMIN: List past announcements."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        # Get announcements (notifications with type 'announcement')
        resp = supabase.table("notifications").select(
            "id, title, message, actor_username, created_at"
        ).eq("type", "announcement").order("created_at", desc=True).limit(50).execute()

        # Dedupe by title+message+timestamp (since we bulk insert one per user)
        seen = set()
        announcements = []
        for row in (resp.data or []):
            key = f"{row['title']}|{row['message']}|{row['created_at'][:16]}"
            if key not in seen:
                seen.add(key)
                announcements.append({
                    "title": row['title'],
                    "message": row['message'],
                    "audience": row.get('actor_username', 'all'),
                    "sent_at": row['created_at']
                })

        return jsonify({"announcements": announcements}), 200

    except Exception as e:
        debug_log(f"List announcements error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/announcements", methods=["POST"])
def send_announcement():
    """ADMIN: Send announcement to users (bulk insert into notifications)."""
    try:
        data = request.get_json()
        admin_key = data.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        title = data.get("title", "").strip()
        message = data.get("message", "").strip()
        audience = data.get("audience", "all")  # 'all', 'pro', or a university name

        if not title or not message:
            return jsonify({"error": "Title and message are required"}), 400

        # Get target user IDs based on audience
        if audience == "all":
            resp = supabase.table("user_profiles").select("id").execute()
        elif audience == "pro":
            resp = supabase.table("user_profiles").select("id").neq("subscription_tier", "free").execute()
        else:
            # Audience is a university name
            resp = supabase.table("user_profiles").select("id").eq("university", audience).execute()

        user_ids = [row['id'] for row in (resp.data or [])]

        if not user_ids:
            return jsonify({"error": "No users match this audience"}), 400

        # Bulk insert notifications (batch in groups of 100 for Supabase limits)
        batch_size = 100
        total_sent = 0

        for i in range(0, len(user_ids), batch_size):
            batch = user_ids[i:i + batch_size]
            rows = [{
                "user_id": uid,
                "type": "announcement",
                "title": title,
                "message": message,
                "actor_username": audience,
                "is_read": False
            } for uid in batch]

            supabase.table("notifications").insert(rows).execute()
            total_sent += len(batch)

        debug_log(f"[Admin] Sent announcement to {total_sent} users: {title}")
        return jsonify({"success": True, "sent_to": total_sent, "audience": audience}), 200

    except Exception as e:
        debug_log(f"Send announcement error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/analytics", methods=["GET"])
def admin_analytics():
    """ADMIN: Analytics overview -- signups, uploads, conversations, downloads over last 30 days."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        # Check Redis cache (10 min TTL since analytics are expensive)
        cache_key = "admin:analytics"
        if redis_cache:
            try:
                cached = redis_cache.get(cache_key)
                if cached:
                    return jsonify(json.loads(cached)), 200
            except:
                pass

        from StudyFlow.backend.supabase_client import supabase
        from datetime import datetime, timedelta

        now = datetime.utcnow()
        thirty_days_ago = (now - timedelta(days=30)).isoformat()
        seven_days_ago = (now - timedelta(days=7)).isoformat()

        # --- Summary stats ---
        total_users = 0
        try:
            resp = supabase.table("user_profiles").select("id", count="exact").execute()
            total_users = resp.count or 0
        except:
            pass

        total_notes = 0
        try:
            resp = supabase.table("notes").select("id", count="exact").execute()
            total_notes = resp.count or 0
        except:
            pass

        total_conversations = 0
        try:
            resp = supabase.table("conversations").select("id", count="exact").is_("deleted_at", "null").execute()
            total_conversations = resp.count or 0
        except:
            pass

        total_downloads = 0
        try:
            resp = supabase.table("download_transactions").select("id", count="exact").execute()
            total_downloads = resp.count or 0
        except:
            pass

        # New users last 7 days
        new_users_7d = 0
        try:
            resp = supabase.table("user_profiles").select("id", count="exact").gte("created_at", seven_days_ago).execute()
            new_users_7d = resp.count or 0
        except:
            pass

        # New notes last 7 days
        new_notes_7d = 0
        try:
            resp = supabase.table("notes").select("id", count="exact").gte("uploaded_at", seven_days_ago).execute()
            new_notes_7d = resp.count or 0
        except:
            pass

        # --- Daily time series (last 30 days) ---

        # Signups per day
        signups_daily = {}
        try:
            resp = supabase.table("user_profiles").select("created_at").gte("created_at", thirty_days_ago).execute()
            for row in (resp.data or []):
                day = row['created_at'][:10]
                signups_daily[day] = signups_daily.get(day, 0) + 1
        except:
            pass

        # Uploads per day
        uploads_daily = {}
        try:
            resp = supabase.table("notes").select("uploaded_at").gte("uploaded_at", thirty_days_ago).execute()
            for row in (resp.data or []):
                day = row['uploaded_at'][:10]
                uploads_daily[day] = uploads_daily.get(day, 0) + 1
        except:
            pass

        # Conversations per day (proxy for active users)
        conversations_daily = {}
        try:
            resp = supabase.table("conversations").select("created_at").gte("created_at", thirty_days_ago).is_("deleted_at", "null").execute()
            for row in (resp.data or []):
                day = row['created_at'][:10]
                conversations_daily[day] = conversations_daily.get(day, 0) + 1
        except:
            pass

        # Downloads per day
        downloads_daily = {}
        try:
            resp = supabase.table("download_transactions").select("created_at").gte("created_at", thirty_days_ago).execute()
            for row in (resp.data or []):
                day = row['created_at'][:10]
                downloads_daily[day] = downloads_daily.get(day, 0) + 1
        except:
            pass

        # Build date range for last 30 days
        dates = []
        for i in range(30, -1, -1):
            d = (now - timedelta(days=i)).strftime("%Y-%m-%d")
            dates.append(d)

        # --- Top universities ---
        university_counts = {}
        try:
            resp = supabase.table("user_profiles").select("university").not_.is_("university", "null").execute()
            for row in (resp.data or []):
                uni = row.get('university')
                if uni:
                    university_counts[uni] = university_counts.get(uni, 0) + 1
        except:
            pass

        top_universities = sorted(university_counts.items(), key=lambda x: x[1], reverse=True)[:10]

        result = {
            "summary": {
                "total_users": total_users,
                "total_notes": total_notes,
                "total_conversations": total_conversations,
                "total_downloads": total_downloads,
                "new_users_7d": new_users_7d,
                "new_notes_7d": new_notes_7d
            },
            "dates": dates,
            "daily": {
                "signups": [signups_daily.get(d, 0) for d in dates],
                "uploads": [uploads_daily.get(d, 0) for d in dates],
                "conversations": [conversations_daily.get(d, 0) for d in dates],
                "downloads": [downloads_daily.get(d, 0) for d in dates]
            },
            "top_universities": [{"name": u[0], "count": u[1]} for u in top_universities]
        }

        # Cache for 10 minutes
        if redis_cache:
            try:
                redis_cache.setex(cache_key, 1800, json.dumps(result))
            except:
                pass

        return jsonify(result), 200

    except Exception as e:
        debug_log(f"Admin analytics error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/users", methods=["GET"])
def admin_users():
    """ADMIN: List users with search, filters, and stats."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        search = request.args.get("search", "").strip()
        filter_by = request.args.get("filter", "")  # suspended, banned, frozen
        sort_by = request.args.get("sort", "recent")  # recent, name, notes
        offset = int(request.args.get("offset", 0))
        limit = int(request.args.get("limit", 50))

        # Fetch user profiles
        query = supabase.table("user_profiles").select(
            "id, email, username, full_name, university, subscription_tier, "
            "good_standing, dmca_strikes, permanent_bad_standing, "
            "account_frozen, frozen_reason, age_verified, edu_email_verified, "
            "is_public, created_at"
        )

        # Apply filters
        if filter_by == "suspended":
            query = query.eq("account_frozen", True)
        elif filter_by == "banned":
            query = query.eq("permanent_bad_standing", True)
        elif filter_by == "bad_standing":
            query = query.eq("good_standing", False)

        # Sort
        if sort_by == "name":
            query = query.order("username", desc=False)
        else:
            query = query.order("created_at", desc=True)

        profiles_resp = query.range(offset, offset + limit - 1).execute()
        profiles = profiles_resp.data or []

        # Apply text search client-side (supabase text search is limited)
        if search:
            search_lower = search.lower()
            profiles = [p for p in profiles if
                search_lower in (p.get("username") or "").lower() or
                search_lower in (p.get("email") or "").lower() or
                search_lower in (p.get("full_name") or "").lower() or
                search_lower in (p.get("university") or "").lower()
            ]

        # Batch get note counts for all users
        user_ids = [p["id"] for p in profiles]
        note_counts = {}
        if user_ids:
            for uid in user_ids:
                try:
                    count_resp = supabase.table("notes").select("id", count="exact").eq("user_id", uid).execute()
                    note_counts[uid] = count_resp.count if count_resp.count else 0
                except:
                    note_counts[uid] = 0

        # Build response
        users = []
        for p in profiles:
            uid = p["id"]
            status = "active"
            if p.get("permanent_bad_standing"):
                status = "banned"
            elif p.get("account_frozen"):
                status = "suspended"
            elif not p.get("good_standing", True):
                status = "bad_standing"

            users.append({
                "id": uid,
                "email": p.get("email", ""),
                "username": p.get("username", ""),
                "full_name": p.get("full_name", ""),
                "university": p.get("university", ""),
                "subscription": p.get("subscription_tier", "free"),
                "status": status,
                "dmca_strikes": p.get("dmca_strikes", 0),
                "edu_verified": p.get("edu_email_verified", False),
                "age_verified": p.get("age_verified", False),
                "nexus_enabled": p.get("is_public", True),
                "note_count": note_counts.get(uid, 0),
                "created_at": p.get("created_at", ""),
                "frozen_reason": p.get("frozen_reason", "")
            })

        # Sort by note count if requested
        if sort_by == "notes":
            users.sort(key=lambda u: u["note_count"], reverse=True)

        debug_log(f"[Admin] Users list: {len(users)} users returned")
        return jsonify({"users": users, "total": len(users)}), 200

    except Exception as e:
        debug_log(f"Admin users error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/users/<user_id>", methods=["GET"])
def admin_user_detail(user_id):
    """ADMIN: Get detailed user profile with activity stats."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        # Get full profile
        profile_resp = supabase.table("user_profiles").select("*").eq("id", user_id).single().execute()
        if not profile_resp.data:
            return jsonify({"error": "User not found"}), 404

        p = profile_resp.data

        # Get note count and list
        notes_resp = supabase.table("notes").select(
            "id, original_filename, university, course_code, is_public, uploaded_at"
        ).eq("user_id", user_id).order("uploaded_at", desc=True).limit(20).execute()
        notes = notes_resp.data or []

        # Get DMCA history (RLS blocks anon reads, so wrap safely)
        dmca_history = []
        try:
            dmca_resp = supabase.table("dmca_takedowns").select(
                "id, note_id, reason, created_at"
            ).eq("uploader_id", user_id).order("created_at", desc=True).execute()
            dmca_history = dmca_resp.data or []
        except Exception as dmca_err:
            debug_log(f"DMCA history fetch skipped: {dmca_err}")
            dmca_history = []

        # Get total note count
        total_notes_resp = supabase.table("notes").select("id", count="exact").eq("user_id", user_id).execute()
        total_notes = total_notes_resp.count if total_notes_resp.count else 0

        # Get conversation count
        try:
            convos_resp = supabase.table("conversations").select("id", count="exact").eq("user_id", user_id).is_("deleted_at", "null").execute()
            total_conversations = convos_resp.count if convos_resp.count else 0
        except:
            total_conversations = 0

        # Determine status
        status = "active"
        if p.get("permanent_bad_standing"):
            status = "banned"
        elif p.get("account_frozen"):
            status = "suspended"
        elif not p.get("good_standing", True):
            status = "bad_standing"

        user = {
            "id": p["id"],
            "email": p.get("email", ""),
            "username": p.get("username", ""),
            "full_name": p.get("full_name", ""),
            "university": p.get("university", ""),
            "subscription": p.get("subscription_tier", "free"),
            "status": status,
            "dmca_strikes": p.get("dmca_strikes", 0),
            "edu_verified": p.get("edu_email_verified", False),
            "age_verified": p.get("age_verified", False),
            "nexus_enabled": p.get("is_public", True),
            "good_standing": p.get("good_standing", True),
            "account_frozen": p.get("account_frozen", False),
            "frozen_reason": p.get("frozen_reason", ""),
            "frozen_at": p.get("frozen_at", ""),
            "permanent_bad_standing": p.get("permanent_bad_standing", False),
            "pages_uploaded_this_month": p.get("pages_uploaded_this_month", 0),
            "created_at": p.get("created_at", ""),
            "stats": {
                "total_notes": total_notes,
                "total_conversations": total_conversations
            },
            "recent_notes": [{
                "id": n["id"],
                "filename": n.get("original_filename", "Untitled"),
                "university": n.get("university", ""),
                "course_code": n.get("course_code", ""),
                "is_public": n.get("is_public", True),
                "uploaded_at": n.get("uploaded_at", "")
            } for n in notes],
            "dmca_history": dmca_history
        }

        return jsonify(user), 200

    except Exception as e:
        debug_log(f"Admin user detail error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/users/<user_id>/suspend", methods=["POST"])
def admin_suspend_user(user_id):
    """ADMIN: Suspend (freeze) a user account."""
    try:
        data = request.get_json()
        admin_key = data.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase
        from datetime import datetime

        reason = data.get("reason", "Admin suspension")

        supabase.table("user_profiles").update({
            "account_frozen": True,
            "frozen_at": datetime.utcnow().isoformat(),
            "frozen_reason": reason
        }).eq("id", user_id).execute()

        # Notify user
        create_notification(
            user_id=user_id,
            notif_type="account_suspended",
            title="Account Suspended",
            message=f"Your account has been suspended: {reason}"
        )

        debug_log(f"[Admin] Suspended user {user_id}: {reason}")
        return jsonify({"success": True, "action": "suspended"}), 200

    except Exception as e:
        debug_log(f"Admin suspend error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/users/<user_id>/unsuspend", methods=["POST"])
def admin_unsuspend_user(user_id):
    """ADMIN: Unsuspend (unfreeze) a user account."""
    try:
        data = request.get_json()
        admin_key = data.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        supabase.table("user_profiles").update({
            "account_frozen": False,
            "frozen_at": None,
            "frozen_reason": None
        }).eq("id", user_id).execute()

        create_notification(
            user_id=user_id,
            notif_type="account_restored",
            title="Account Restored",
            message="Your account has been restored. Welcome back."
        )

        debug_log(f"[Admin] Unsuspended user {user_id}")
        return jsonify({"success": True, "action": "unsuspended"}), 200

    except Exception as e:
        debug_log(f"Admin unsuspend error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/users/<user_id>/ban", methods=["POST"])
def admin_ban_user(user_id):
    """ADMIN: Permanently ban a user (sets 3 DMCA strikes + permanent bad standing)."""
    try:
        data = request.get_json()
        admin_key = data.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase
        from datetime import datetime

        reason = data.get("reason", "Permanent ban by admin")

        supabase.table("user_profiles").update({
            "dmca_strikes": 3,
            "permanent_bad_standing": True,
            "good_standing": False,
            "account_frozen": True,
            "frozen_at": datetime.utcnow().isoformat(),
            "frozen_reason": reason
        }).eq("id", user_id).execute()

        # Make all their notes private
        supabase.table("notes").update({
            "is_public": False
        }).eq("user_id", user_id).execute()

        create_notification(
            user_id=user_id,
            notif_type="account_banned",
            title="Account Banned",
            message=f"Your account has been permanently banned: {reason}"
        )

        debug_log(f"[Admin] Banned user {user_id}: {reason}")
        return jsonify({"success": True, "action": "banned"}), 200

    except Exception as e:
        debug_log(f"Admin ban error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/review/<note_id>", methods=["POST"])
def review_note(note_id):
    """ADMIN: Approve or reject a note."""
    try:
        data = request.get_json()
        admin_key = data.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        action = data.get("action")
        if action not in ("approve", "reject"):
            return jsonify({"error": "Action must be 'approve' or 'reject'"}), 400

        from StudyFlow.backend.supabase_client import supabase

        # Get note info
        note_resp = supabase.table("notes").select("user_id, original_filename, file_path").eq("id", note_id).execute()
        if not note_resp.data:
            return jsonify({"error": "Note not found"}), 404

        note = note_resp.data[0]

        if action == "approve":
            supabase.table("reviewed_notes").upsert({
                "note_id": note_id,
                "action": "approved"
            }).execute()
            debug_log(f"[Review] Approved note {note_id}")
        else:
            reason = data.get("reason", "Does not meet quality standards")
            filename = note.get("original_filename", "your note")

            # Check if note has a social post (will be cascade-deleted)
            had_social_post = False
            try:
                social_check = supabase.table("social_posts").select("id").eq("note_id", note_id).execute()
                had_social_post = bool(social_check.data)
            except Exception:
                pass

            # Delete file from storage
            if note.get("file_path"):
                try:
                    supabase.storage.from_("note-files").remove([note["file_path"]])
                except Exception:
                    pass

            # Delete note chunks
            try:
                supabase.table("note_chunks").delete().eq("note_id", note_id).execute()
            except Exception:
                pass

            # Delete note record (cascade-deletes social post + comments/votes)
            supabase.table("notes").delete().eq("id", note_id).execute()

            # Notify the user
            message = f"Your note '{filename}' was removed: {reason}"
            if had_social_post:
                message += " The associated social post has also been removed."

            create_notification(
                user_id=note["user_id"],
                notif_type="note_removed",
                title="Note Removed",
                message=message,
            )

            debug_log(f"[Review] Rejected and deleted note {note_id} (social post: {had_social_post}): {reason}")

        return jsonify({"success": True, "action": action}), 200

    except Exception as e:
        debug_log(f"Review note error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/classification-queue", methods=["GET"])
def admin_classification_queue():
    """Admin: get recent uploads with their AI classification for review."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        limit = int(request.args.get("limit", 50))

        result = supabase.table("notes").select(
            "id, user_id, original_filename, file_type, file_size, page_count, "
            "is_public, ai_classification, uploaded_at, username"
        ).order("uploaded_at", desc=True).limit(limit).execute()

        return jsonify({"notes": result.data or []}), 200

    except Exception as e:
        debug_log(f"Classification queue error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/classification/<note_id>", methods=["POST"])
def admin_override_classification(note_id):
    """Admin: override AI classification for a note."""
    try:
        data = request.get_json()
        admin_key = data.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        action = data.get("action")
        if action not in ("make_public", "make_private"):
            return jsonify({"error": "Action must be make_public or make_private"}), 400

        is_public = action == "make_public"
        classification = "educational (admin)" if is_public else "personal (admin)"

        supabase.table("notes").update({
            "is_public": is_public,
            "ai_classification": classification
        }).eq("id", note_id).execute()

        debug_log(f"[Admin Override] Note {note_id}: {classification}")
        return jsonify({"success": True, "is_public": is_public}), 200

    except Exception as e:
        debug_log(f"Classification override error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/note-preview/<note_id>", methods=["GET"])
def admin_note_preview(note_id):
    """ADMIN: Get a signed URL for previewing a note."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        note = supabase.table("notes").select("file_path, original_filename").eq("id", note_id).execute()
        if not note.data or not note.data[0].get("file_path"):
            return jsonify({"error": "Note not found"}), 404

        signed = supabase.storage.from_("note-files").create_signed_url(note.data[0]["file_path"], 3600)
        url = signed.get("signedURL") or signed.get("signedUrl", "")

        return jsonify({"url": url, "filename": note.data[0].get("original_filename")}), 200

    except Exception as e:
        debug_log(f"Admin note preview error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/university-stats", methods=["GET"])
def university_stats():
    """Per-university metrics for institutional license sales pitches."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        from StudyFlow.backend.supabase_client import supabase

        # Get all verified users grouped by university
        profiles = supabase.table("user_profiles").select(
            "id, university, edu_email_verified"
        ).eq("edu_email_verified", True).not_.is_("university", "null").execute()

        uni_users = {}
        all_user_ids = set()
        for p in (profiles.data or []):
            uni = p["university"]
            if uni not in uni_users:
                uni_users[uni] = []
            uni_users[uni].append(p["id"])
            all_user_ids.add(p["id"])

        # Get note counts per university
        notes = supabase.table("notes").select(
            "user_id, university, course_code"
        ).not_.is_("university", "null").execute()

        uni_notes = {}
        uni_courses = {}
        for n in (notes.data or []):
            uni = n.get("university", "Unknown")
            if uni not in uni_notes:
                uni_notes[uni] = 0
                uni_courses[uni] = set()
            uni_notes[uni] += 1
            if n.get("course_code"):
                uni_courses[uni].add(n["course_code"])

        # Get AI query counts from provenance logs
        uni_queries = {}
        try:
            for uni, user_ids in uni_users.items():
                if not user_ids:
                    continue
                count_resp = supabase.table("ai_provenance_logs").select(
                    "id", count="exact"
                ).in_("user_id", user_ids[:50]).execute()
                uni_queries[uni] = count_resp.count or 0
        except Exception:
            pass  # Table might not exist yet

        # Get study group counts per university
        uni_groups = {}
        try:
            groups = supabase.table("study_groups").select(
                "institution"
            ).not_.is_("institution", "null").execute()
            for g in (groups.data or []):
                inst = g["institution"]
                uni_groups[inst] = uni_groups.get(inst, 0) + 1
        except Exception:
            pass

        # Build response
        universities = []
        all_unis = set(list(uni_users.keys()) + list(uni_notes.keys()))
        for uni in sorted(all_unis):
            universities.append({
                "university": uni,
                "verified_users": len(uni_users.get(uni, [])),
                "notes_uploaded": uni_notes.get(uni, 0),
                "courses_covered": len(uni_courses.get(uni, set())),
                "top_courses": sorted(list(uni_courses.get(uni, set())))[:10],
                "ai_queries": uni_queries.get(uni, 0),
                "study_groups": uni_groups.get(uni, 0)
            })

        # Totals
        totals = {
            "total_universities": len(all_unis),
            "total_verified_users": sum(len(v) for v in uni_users.values()),
            "total_notes": sum(uni_notes.values()),
            "total_courses": sum(len(v) for v in uni_courses.values()),
            "total_ai_queries": sum(uni_queries.values()),
            "total_study_groups": sum(uni_groups.values())
        }

        return jsonify({"totals": totals, "universities": universities}), 200

    except Exception as e:
        debug_log(f"University stats error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/notes/generate-outline", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def generate_outline():
    """
    Generate a detailed paper outline from the user's own notes.
    Does NOT write prose -- provides structured headers, key points,
    quotes, and connections the student should write about.

    Expects JSON:
    {
        "topic": "optional paper topic or prompt",
        "format": "mla" | "apa" | "chicago",
        "note_ids": [] (optional, defaults to all user notes)
    }
    """
    try:
        data = request.get_json() or {}
        topic = data.get('topic', '').strip()
        doc_format = data.get('format', 'mla').lower()
        note_ids = data.get('note_ids', [])

        user_id = request.user_id

        # Fetch user's note chunks (only their own, not Nexus)
        if note_ids:
            chunks_resp = supabase.table("note_chunks").select(
                "chunk_text, content_summary, note_id"
            ).eq("user_id", user_id).in_("note_id", note_ids).order("chunk_index").execute()
        else:
            chunks_resp = supabase.table("note_chunks").select(
                "chunk_text, content_summary, note_id"
            ).eq("user_id", user_id).order("created_at").execute()

        chunks = chunks_resp.data or []

        if not chunks:
            return jsonify({"error": "No notes found. Upload some notes first."}), 400

        # Get note filenames for source attribution
        note_id_set = list(set(c['note_id'] for c in chunks))
        notes_resp = supabase.table("notes").select(
            "id, original_filename, course_code, professor, university"
        ).in_("id", note_id_set).execute()
        notes_map = {n['id']: n for n in (notes_resp.data or [])}

        # Build context from chunks with source labels
        note_context = ""
        for chunk in chunks[:100]:  # Cap at 100 chunks to fit context window
            note = notes_map.get(chunk['note_id'], {})
            filename = note.get('original_filename', 'Unknown')
            course = note.get('course_code', '')
            source_label = filename
            if course:
                source_label += f" ({course})"
            text = chunk.get('content_summary') or chunk.get('chunk_text', '')
            if text:
                note_context += f"\n[Source: {source_label}]\n{text}\n"

        # Format instructions
        format_instructions = {
            'mla': 'MLA format: 12pt Times New Roman, double-spaced, 1-inch margins, last name and page number in header, Works Cited page format',
            'apa': 'APA format: 12pt Times New Roman, double-spaced, running head, title page with author and institution, References page format',
            'chicago': 'Chicago format: 12pt Times New Roman, double-spaced, title page, footnotes/endnotes, Bibliography page format'
        }.get(doc_format, 'MLA format')

        topic_instruction = ""
        if topic:
            topic_instruction = f"The student's paper topic or assignment prompt is: \"{topic}\"\n\nFocus the outline on this topic, using only information from their notes that is relevant."
        else:
            topic_instruction = "The student has not specified a topic. Identify the main themes and topics across all their notes and create an outline that covers the most important material."

        prompt = f"""You are an academic outline generator. Organize the student's notes into a structured paper outline.

RULES:
- You are NOT writing their paper. You are organizing their existing notes.
- Every point MUST come from the student's notes below. Do NOT add outside information.
- Do NOT write full paragraphs. Use bullet points and short phrases.
- Include direct quotes from their notes (in quotation marks) worth expanding on.
- Suggest connections between topics.
- Attribute every point to its source note.

{topic_instruction}

STUDENT'S NOTES:
{note_context}

Respond with ONLY valid JSON (no markdown, no code fences) in this exact structure:
{{
  "suggested_titles": ["Title Option 1", "Title Option 2", "Title Option 3"],
  "thesis_direction": "Your notes suggest you could argue...",
  "sections": [
    {{
      "heading": "Section Title",
      "level": 1,
      "main_idea": "One sentence framing what this section covers",
      "key_points": [
        {{"text": "Specific fact or concept from notes", "source": "filename.pdf"}},
        {{"text": "Another key point", "source": "filename.pdf"}}
      ],
      "quotes": [
        {{"text": "Direct quote worth expanding on", "source": "filename.pdf"}}
      ],
      "connections": ["How this relates to other sections"],
      "subsections": [
        {{
          "heading": "Subsection Title",
          "level": 2,
          "key_points": [
            {{"text": "Point", "source": "filename.pdf"}}
          ]
        }}
      ]
    }}
  ],
  "conclusion_direction": "What to summarize and reflect on (NOT a written conclusion)"
}}

Include 4-8 main sections with subsections where appropriate. Make key_points detailed -- 3-6 per section."""

        # Call Gemini
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        model = genai.GenerativeModel('gemini-3-flash-preview')

        response = model.generate_content(prompt)
        outline_text = response.text

        # Try to parse as JSON, fall back to raw text
        outline_json = None
        try:
            # Strip markdown code fences if present
            clean = outline_text.strip()
            if clean.startswith('```'):
                clean = clean.split('\n', 1)[1] if '\n' in clean else clean[3:]
            if clean.endswith('```'):
                clean = clean[:-3]
            clean = clean.strip()
            if clean.startswith('json'):
                clean = clean[4:].strip()
            outline_json = json.loads(clean)
        except:
            outline_json = None

        # Track cost
        try:
            from StudyFlow.backend.cost_tracker import track_ai_call
            track_ai_call("gemini", "gemini-2.5-flash", "outline_gen", tokens_estimate=len(note_context) // 4)
        except:
            pass

        # Log for compliance
        try:
            from StudyFlow.backend.conversational_noteflow import log_ai_response
            log_ai_response(
                user_id=user_id,
                question=f"Generate outline: {topic or 'all notes'}",
                response_text=outline_text[:500],
                model_used="gemini-2.5-flash",
                sources_used=[{"note_id": nid, "filename": notes_map.get(nid, {}).get('original_filename', '')} for nid in note_id_set[:10]]
            )
        except:
            pass

        return jsonify({
            "outline": outline_text,
            "outline_json": outline_json,
            "format": doc_format,
            "notes_used": len(note_id_set),
            "chunks_analyzed": len(chunks)
        }), 200

    except Exception as e:
        debug_log(f"Outline generation error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/reembed", methods=["POST"])
def trigger_reembed():
    """ADMIN: Trigger re-embedding of all note chunks with Gemini."""
    admin_key = request.args.get("key", "")
    if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
        return jsonify({"error": "Unauthorized"}), 403
    try:
        from StudyFlow.backend.tasks import reembed_all_chunks
        reembed_all_chunks.delay()
        return jsonify({"status": "Re-embedding task started"}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/admin/force-reembed", methods=["POST"])
def force_reembed_all():
    """
    ADMIN: Force re-embedding of ALL note chunks (not just NULL ones).

    This endpoint:
    1. Nullifies all existing embeddings in batches (clears old OpenAI embeddings)
    2. Triggers batch re-embedding with Gemini

    Use this to fix embedding model mismatch issues.
    """
    admin_key = request.args.get("key", "")
    if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
        return jsonify({"error": "Unauthorized"}), 403

    try:
        # Step 1: Count total chunks
        total_result = supabase.table("note_chunks").select("id", count="exact").execute()
        total_chunks = total_result.count

        debug_log(f"[FORCE REEMBED] Total chunks: {total_chunks}")

        # Step 2: Nullify embeddings in batches of 500 (using simple pagination)
        batch_size = 500
        nullified_count = 0

        debug_log(f"[FORCE REEMBED] Nullifying embeddings in batches of {batch_size}...")

        # Process in batches using offset
        offset = 0
        while offset < total_chunks:
            # Get batch of chunk IDs
            batch_result = supabase.table("note_chunks").select("id").range(offset, offset + batch_size - 1).execute()
            chunk_ids = [c["id"] for c in (batch_result.data or [])]

            if not chunk_ids:
                break

            # Nullify this batch
            supabase.table("note_chunks").update({"embedding": None}).in_("id", chunk_ids).execute()
            nullified_count += len(chunk_ids)

            debug_log(f"[FORCE REEMBED] Nullified {nullified_count}/{total_chunks} embeddings...")

            offset += batch_size

        debug_log(f"[FORCE REEMBED] All embeddings nullified ({nullified_count} total)")

        # Step 3: Trigger re-embedding task
        from StudyFlow.backend.tasks import reembed_all_chunks
        reembed_all_chunks.delay()

        debug_log(f"[FORCE REEMBED] Re-embedding task started")

        return jsonify({
            "status": "Force re-embedding started",
            "nullified_chunks": nullified_count,
            "message": "All embeddings nullified. Re-embedding task queued."
        }), 200

    except Exception as e:
        debug_log(f"[FORCE REEMBED] ERROR: {e}")
        import traceback
        debug_log(traceback.format_exc())
        return jsonify({"error": str(e)}), 500


@app.route("/admin/reembed-status", methods=["GET"])
def reembed_status():
    """ADMIN: Check re-embedding progress."""
    admin_key = request.args.get("key", "")
    if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
        return jsonify({"error": "Unauthorized"}), 403
    try:
        # Count total chunks vs chunks with embeddings
        total_result = supabase.table("note_chunks").select("id", count="exact").execute()
        total_chunks = total_result.count

        embedded_result = supabase.table("note_chunks").select("id", count="exact").not_.is_("embedding", "null").execute()
        embedded_chunks = embedded_result.count

        null_chunks = total_chunks - embedded_chunks
        percent_complete = round((embedded_chunks / total_chunks * 100) if total_chunks > 0 else 0, 1)

        return jsonify({
            "total_chunks": total_chunks,
            "embedded": embedded_chunks,
            "remaining": null_chunks,
            "percent_complete": percent_complete,
            "status": "Complete" if null_chunks == 0 else "In Progress"
        }), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/admin/flush-cache", methods=["POST"])
def flush_user_cache():
    """ADMIN: Flush all Redis caches for a user."""
    admin_key = request.args.get("key", "")
    if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
        return jsonify({"error": "Unauthorized"}), 403
    user_id = request.args.get("user_id", "")
    if not user_id:
        return jsonify({"error": "Missing user_id"}), 400
    try:
        deleted = []
        if redis_cache:
            for prefix in ["canvas_preload:", "convos_list:", "user_notes:", "profile:", "notifications:", "user_groups:", "dashboard:", "calendar:"]:
                key = prefix + user_id
                if redis_cache.delete(key):
                    deleted.append(key)
        return jsonify({"flushed": deleted}), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/admin/import-wikipedia", methods=["POST"])
def admin_import_wikipedia():
    """ADMIN: Import Wikipedia articles by subject or title."""
    admin_key = request.args.get("key", "")
    if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
        return jsonify({"error": "Unauthorized"}), 403
    try:
        data = request.get_json() or {}
        subject = data.get('subject')  # e.g. "Biology", "US History"
        article = data.get('article')  # single article title

        from scripts.wikipedia_importer import import_article, import_batch, ALL_SUBJECTS

        if article:
            result = import_article(article, course_code=data.get('course_code', 'General Knowledge'))
            return jsonify(result), 200
        elif subject and subject in ALL_SUBJECTS:
            # Run in background via Celery
            from StudyFlow.backend.tasks import celery_app
            @celery_app.task(name="wiki_import_" + subject.replace(" ", "_"))
            def do_import():
                return import_batch(ALL_SUBJECTS[subject], course_code=subject)
            do_import.delay()
            return jsonify({"status": f"Importing {len(ALL_SUBJECTS[subject])} {subject} articles in background"}), 200
        else:
            return jsonify({
                "available_subjects": list(ALL_SUBJECTS.keys()),
                "total_articles": sum(len(a) for a in ALL_SUBJECTS.values()),
                "usage": "POST with {subject: 'Biology'} or {article: 'Photosynthesis'}"
            }), 200
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/api/search", methods=["GET"])
@supabase_auth_required
def global_search():
    """Global search across notes, conversations, groups, and calendar events."""
    try:
        q = request.args.get("q", "").strip()
        if not q or len(q) < 2:
            return jsonify({"results": {}}), 200

        user_id = request.user_id
        query = f"%{q}%"
        results = {}

        # Search notes (user's own)
        try:
            notes_resp = supabase.table("notes").select(
                "id, original_filename, university, course_code, uploaded_at"
            ).eq("user_id", user_id).ilike("original_filename", query).limit(5).execute()
            if notes_resp.data:
                results["notes"] = [{
                    "id": n["id"],
                    "title": n.get("original_filename", ""),
                    "subtitle": f"{n.get('course_code', '')} {n.get('university', '')}".strip(),
                    "url": f"notes.html",
                    "type": "note"
                } for n in notes_resp.data]
        except: pass

        # Search conversations
        try:
            convos_resp = supabase.table("conversations").select(
                "id, title, updated_at"
            ).eq("user_id", user_id).eq("source", "chat").is_("deleted_at", "null").ilike("title", query).limit(5).execute()
            if convos_resp.data:
                results["chats"] = [{
                    "id": c["id"],
                    "title": c.get("title", "New Chat"),
                    "subtitle": c.get("updated_at", "")[:10] if c.get("updated_at") else "",
                    "url": f"chat.html?c={c['id']}",
                    "type": "chat"
                } for c in convos_resp.data]
        except: pass

        # Search study groups
        try:
            # Get user's group IDs first
            memberships = supabase.table("study_group_members").select("group_id").eq("user_id", user_id).execute()
            group_ids = [m["group_id"] for m in (memberships.data or [])]
            if group_ids:
                groups_resp = supabase.table("study_groups").select(
                    "id, name, description"
                ).in_("id", group_ids).ilike("name", query).limit(5).execute()
                if groups_resp.data:
                    results["groups"] = [{
                        "id": g["id"],
                        "title": g.get("name", ""),
                        "subtitle": g.get("description", "")[:50] if g.get("description") else "",
                        "url": f"study-groups.html?id={g['id']}",
                        "type": "group"
                    } for g in groups_resp.data]
        except: pass

        # Search calendar events
        try:
            events_resp = supabase.table("calendar_events").select(
                "id, title, due_date, event_type"
            ).eq("user_id", user_id).ilike("title", query).limit(5).execute()
            if events_resp.data:
                results["events"] = [{
                    "id": e["id"],
                    "title": e.get("title", ""),
                    "subtitle": e.get("due_date", "")[:10] if e.get("due_date") else "",
                    "url": "calendar.html",
                    "type": "event"
                } for e in events_resp.data]
        except: pass

        # Search Nexus (public notes, excluding Wikipedia)
        try:
            nexus_resp = supabase.table("notes").select(
                "id, original_filename, university, course_code"
            ).eq("is_public", True).neq("university", "Wikipedia").ilike("original_filename", query).limit(5).execute()
            if nexus_resp.data:
                # Filter out user's own notes (already in notes results)
                nexus_notes = [n for n in nexus_resp.data if n.get("id") not in [r["id"] for r in results.get("notes", [])]]
                if nexus_notes:
                    results["nexus"] = [{
                        "id": n["id"],
                        "title": n.get("original_filename", ""),
                        "subtitle": f"{n.get('university', '')} - {n.get('course_code', '')}".strip(' -'),
                        "url": f"note-viewer.html?id={n['id']}",
                        "type": "nexus"
                    } for n in nexus_notes[:5]]
        except: pass

        return jsonify({"results": results, "query": q}), 200

    except Exception as e:
        debug_log(f"Search error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500




# ============= FEATURE REQUESTS API =============

@app.route("/api/feature-requests", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_feature_requests():
    """Get all feature requests with vote status for current user"""
    try:
        user_id = request.user_id

        # Get filter and sort parameters
        status_filter = request.args.get("status", "all")
        sort_by = request.args.get("sort", "popular")

        # Build base query
        query = supabase.table("feature_requests").select(
            "id, user_id, username, title, description, status, vote_count, "
            "admin_response, completed_at, created_at, updated_at"
        )

        # Apply status filter
        if status_filter != "all":
            query = query.eq("status", status_filter)

        # Apply sorting
        if sort_by == "popular":
            query = query.order("vote_count", desc=True).order("created_at", desc=True)
        elif sort_by == "recent":
            query = query.order("created_at", desc=True)
        elif sort_by == "status":
            query = query.order("status").order("created_at", desc=True)

        # Execute query
        response = query.execute()
        requests_data = response.data or []

        # Get user's votes
        votes_response = supabase.table("feature_request_votes").select(
            "request_id"
        ).eq("user_id", user_id).execute()

        user_voted_ids = {v["request_id"] for v in (votes_response.data or [])}

        # Add flags to each request
        for req in requests_data:
            req["has_voted"] = req["id"] in user_voted_ids
            req["is_owner"] = req["user_id"] == user_id

        return jsonify({"requests": requests_data, "total": len(requests_data)}), 200

    except Exception as e:
        debug_log(f"Get feature requests error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/feature-requests", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def create_feature_request():
    """Submit a new feature request"""
    try:
        user_id = request.user_id
        data = request.json

        title = data.get("title", "").strip()
        description = data.get("description", "").strip()

        # Validation
        if not title or len(title) < 5:
            return jsonify({"error": "Title must be at least 5 characters"}), 400
        if len(title) > 200:
            return jsonify({"error": "Title must be less than 200 characters"}), 400
        if not description or len(description) < 10:
            return jsonify({"error": "Description must be at least 10 characters"}), 400
        if len(description) > 2000:
            return jsonify({"error": "Description must be less than 2000 characters"}), 400

        # Get username
        profile = get_cached_profile(user_id)
        username = profile.get("username", "Anonymous") if profile else "Anonymous"

        # Create feature request
        response = supabase.table("feature_requests").insert({
            "user_id": user_id,
            "username": username,
            "title": title,
            "description": description,
            "status": "pending",
            "vote_count": 0
        }).execute()

        if not response.data:
            return jsonify({"error": "Failed to create feature request"}), 500

        created_request = response.data[0]
        created_request["has_voted"] = False
        created_request["is_owner"] = True

        debug_log(f"Feature request created: {created_request['id']} by user {user_id}")

        return jsonify({
            "message": "Feature request submitted successfully",
            "request": created_request
        }), 201

    except Exception as e:
        debug_log(f"Create feature request error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/feature-requests/<request_id>/vote", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def toggle_vote_feature_request(request_id):
    """Toggle vote for a feature request"""
    try:
        user_id = request.user_id

        # Check if user already voted
        existing_vote = supabase.table("feature_request_votes").select(
            "id"
        ).eq("user_id", user_id).eq("request_id", request_id).execute()

        if existing_vote.data:
            # Remove vote
            vote_id = existing_vote.data[0]["id"]
            supabase.table("feature_request_votes").delete().eq("id", vote_id).execute()
            action = "removed"
        else:
            # Add vote
            supabase.table("feature_request_votes").insert({
                "user_id": user_id,
                "request_id": request_id
            }).execute()
            action = "added"

        # Get updated request
        updated_request = supabase.table("feature_requests").select(
            "id, vote_count"
        ).eq("id", request_id).single().execute()

        debug_log(f"Vote {action} for request {request_id} by user {user_id}")

        return jsonify({
            "message": f"Vote {action}",
            "vote_count": updated_request.data.get("vote_count", 0) if updated_request.data else 0,
            "has_voted": action == "added"
        }), 200

    except Exception as e:
        debug_log(f"Toggle vote error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/feature-requests/<request_id>", methods=["PATCH"])
@supabase_auth_required
@account_not_frozen
def update_feature_request(request_id):
    """Update a feature request (users can update their own)"""
    try:
        user_id = request.user_id
        data = request.json

        # Get existing request
        existing = supabase.table("feature_requests").select(
            "user_id"
        ).eq("id", request_id).single().execute()

        if not existing.data:
            return jsonify({"error": "Feature request not found"}), 404

        # Check ownership
        if existing.data["user_id"] != user_id:
            return jsonify({"error": "You can only update your own feature requests"}), 403

        # Users can only update title and description
        update_data = {}
        if "title" in data:
            title = data["title"].strip()
            if len(title) < 5 or len(title) > 200:
                return jsonify({"error": "Title must be 5-200 characters"}), 400
            update_data["title"] = title

        if "description" in data:
            description = data["description"].strip()
            if len(description) < 10 or len(description) > 2000:
                return jsonify({"error": "Description must be 10-2000 characters"}), 400
            update_data["description"] = description

        if not update_data:
            return jsonify({"error": "No valid fields to update"}), 400

        # Update
        response = supabase.table("feature_requests").update(update_data).eq(
            "id", request_id
        ).execute()

        debug_log(f"Feature request {request_id} updated by user {user_id}")

        return jsonify({
            "message": "Feature request updated",
            "request": response.data[0] if response.data else {}
        }), 200

    except Exception as e:
        debug_log(f"Update feature request error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/admin/feature-requests/<request_id>", methods=["PATCH"])
@supabase_auth_required
@account_not_frozen
def admin_update_feature_request(request_id):
    """Admin endpoint to update feature request status and response"""
    try:
        user_id = request.user_id
        data = request.json

        # TODO: Add proper admin role check here
        # For now, this endpoint works but should be restricted to admins only

        update_data = {}

        if "status" in data:
            valid_statuses = ["pending", "reviewing", "in-progress", "completed", "rejected"]
            if data["status"] not in valid_statuses:
                return jsonify({"error": f"Invalid status. Must be one of: {', '.join(valid_statuses)}"}), 400
            update_data["status"] = data["status"]

        if "admin_response" in data:
            update_data["admin_response"] = data["admin_response"]

        if not update_data:
            return jsonify({"error": "No fields to update"}), 400

        # Update
        response = supabase.table("feature_requests").update(update_data).eq(
            "id", request_id
        ).execute()

        if not response.data:
            return jsonify({"error": "Feature request not found"}), 404

        debug_log(f"Admin updated feature request {request_id}: {update_data}")

        return jsonify({
            "message": "Feature request updated by admin",
            "request": response.data[0]
        }), 200

    except Exception as e:
        debug_log(f"Admin update error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/feature-requests/<request_id>", methods=["DELETE"])
@supabase_auth_required
@account_not_frozen
def delete_feature_request(request_id):
    """Delete a feature request (owner only)"""
    try:
        user_id = request.user_id

        # Get existing request
        existing = supabase.table("feature_requests").select(
            "user_id"
        ).eq("id", request_id).single().execute()

        if not existing.data:
            return jsonify({"error": "Feature request not found"}), 404

        # Check ownership
        if existing.data["user_id"] != user_id:
            return jsonify({"error": "You can only delete your own feature requests"}), 403

        # Delete (votes cascade)
        supabase.table("feature_requests").delete().eq("id", request_id).execute()

        debug_log(f"Feature request {request_id} deleted by user {user_id}")

        return jsonify({"message": "Feature request deleted"}), 200

    except Exception as e:
        debug_log(f"Delete feature request error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

# =================================================




# ============= SOCIAL MEDIA API =============
# Instagram/Reddit hybrid for sharing and discovering notes

# --- POSTS & FEED ---

## ====== POST REPORTING ======

@app.route("/api/social/posts/<post_id>/report", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def report_post(post_id):
    """Report a social post for moderation."""
    try:
        data = request.get_json() or {}
        reason = data.get("reason", "").strip()
        details = data.get("details", "").strip()

        if not reason:
            return jsonify({"error": "Reason is required"}), 400

        valid_reasons = ["spam", "harassment", "inappropriate", "misinformation", "copyright", "other"]
        if reason not in valid_reasons:
            return jsonify({"error": f"Reason must be one of: {', '.join(valid_reasons)}"}), 400

        # Check post exists
        post = supabase.table("social_posts").select("id, user_id, note_id, post_type").eq("id", post_id).execute()
        if not post.data:
            return jsonify({"error": "Post not found"}), 404

        post_data = post.data[0]

        # Can't report own post
        if post_data["user_id"] == request.user_id:
            return jsonify({"error": "You cannot report your own post"}), 400

        # Copyright reports on note posts -> route to DMCA system
        if reason == "copyright" and post_data.get("note_id"):
            try:
                # Get reporter info
                reporter = supabase.table("user_profiles").select("email, full_name").eq("id", request.user_id).execute()
                reporter_email = reporter.data[0].get("email", "unknown") if reporter.data else "unknown"
                reporter_name = reporter.data[0].get("full_name", "") if reporter.data else ""

                # Get note filename
                note = supabase.table("notes").select("original_filename, user_id").eq("id", post_data["note_id"]).execute()
                note_filename = note.data[0].get("original_filename", "Unknown") if note.data else "Unknown"

                # Create DMCA takedown record
                supabase.table("dmca_takedown_requests").insert({
                    "note_id": post_data["note_id"],
                    "reporter_email": reporter_email,
                    "reporter_name": reporter_name,
                    "reason": details or "Copyright violation reported via social feed",
                    "status": "pending",
                    "note_owner_id": note.data[0]["user_id"] if note.data else None,
                    "note_filename": note_filename
                }).execute()

                # Also create a regular report for the admin social page
                try:
                    supabase.table("post_reports").insert({
                        "post_id": post_id,
                        "reporter_id": request.user_id,
                        "reason": reason,
                        "details": (details or "") + " [Routed to DMCA system]"
                    }).execute()
                except:
                    pass

                debug_log(f"Copyright report on social post {post_id} routed to DMCA for note {post_data['note_id']}")
                return jsonify({"success": True, "message": "Copyright report submitted and routed to our DMCA process. We'll review it within 24 hours."}), 201

            except Exception as dmca_err:
                debug_log(f"DMCA routing failed, falling back to regular report: {dmca_err}")
                # Fall through to regular report

        # Insert regular report (unique constraint prevents duplicate reports)
        try:
            supabase.table("post_reports").insert({
                "post_id": post_id,
                "reporter_id": request.user_id,
                "reason": reason,
                "details": details or None
            }).execute()
        except Exception as e:
            if "duplicate" in str(e).lower() or "unique" in str(e).lower():
                return jsonify({"error": "You have already reported this post"}), 409
            raise

        return jsonify({"success": True, "message": "Report submitted. We'll review it shortly."}), 201

    except Exception as e:
        debug_log(f"Report post error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


## ====== ADMIN SOCIAL MODERATION ======

@app.route("/admin/social/posts", methods=["GET"])
def admin_get_social_posts():
    """Admin: get all social posts for moderation."""
    try:
        admin_key = request.args.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        status_filter = request.args.get("status", "all")  # all, approved, denied
        page = int(request.args.get("page", 0))
        limit = int(request.args.get("limit", 50))

        query = supabase.table("social_posts").select(
            "id, user_id, username, post_type, text_content, note_id, "
            "upvote_count, downvote_count, score, comment_count, "
            "created_at, moderation_status"
        ).order("created_at", desc=True).range(page * limit, (page + 1) * limit - 1)

        if status_filter == "approved":
            query = query.eq("moderation_status", "approved")
        elif status_filter == "denied":
            query = query.eq("moderation_status", "denied")
        elif status_filter == "pending":
            query = query.is_("moderation_status", "null")

        result = query.execute()
        posts = result.data or []

        # Enrich with note filenames
        note_ids = [p['note_id'] for p in posts if p.get('note_id')]
        note_names = {}
        if note_ids:
            try:
                notes_res = supabase.table("notes").select("id, original_filename").in_("id", note_ids).execute()
                note_names = {n['id']: n.get('original_filename', 'Unknown') for n in (notes_res.data or [])}
            except:
                pass

        # Enrich with report data
        post_ids = [p['id'] for p in posts]
        reports_map = {}
        if post_ids:
            try:
                if len(post_ids) == 1:
                    reports_res = supabase.table("post_reports").select("post_id, reason, details, reporter_id, created_at, status").eq("post_id", post_ids[0]).execute()
                else:
                    reports_res = supabase.table("post_reports").select("post_id, reason, details, reporter_id, created_at, status").in_("post_id", post_ids).execute()
                for r in (reports_res.data or []):
                    pid = r['post_id']
                    if pid not in reports_map:
                        reports_map[pid] = []
                    reports_map[pid].append(r)
            except:
                pass

        for p in posts:
            if p.get('note_id'):
                p['note_filename'] = note_names.get(p['note_id'], 'Unknown')
            if not p.get('moderation_status'):
                p['moderation_status'] = 'approved'
            p['reports'] = reports_map.get(p['id'], [])
            p['report_count'] = len(p['reports'])

        # Sort: reported posts first, then by date
        posts.sort(key=lambda p: (-p['report_count'], p['created_at']), reverse=False)
        # Actually: reported first (descending by report count), then newest first
        posts.sort(key=lambda p: (-p['report_count'], ''), reverse=False)

        return jsonify({"posts": posts, "count": len(posts)}), 200

    except Exception as e:
        debug_log(f"Admin social posts error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/admin/social/posts/<post_id>", methods=["POST"])
def admin_moderate_post(post_id):
    """Admin: approve or deny a social post."""
    try:
        data = request.get_json()
        admin_key = data.get("key", "")
        if admin_key != os.getenv("ADMIN_KEY", "change_me_in_production"):
            return jsonify({"error": "Unauthorized"}), 403

        action = data.get("action")
        if action not in ("approve", "deny"):
            return jsonify({"error": "Action must be 'approve' or 'deny'"}), 400

        reason = data.get("reason", "")

        if action == "deny":
            # Get post info for notification
            post_res = supabase.table("social_posts").select("user_id, username, post_type, text_content, note_id").eq("id", post_id).execute()
            if not post_res.data:
                return jsonify({"error": "Post not found"}), 404

            post = post_res.data[0]

            # Delete the post
            supabase.table("social_posts").delete().eq("id", post_id).execute()

            # Notify the user
            post_desc = post.get('text_content', '')[:50] if post.get('text_content') else 'your post'
            msg = f"Your social post was removed by a moderator."
            if reason:
                msg += f" Reason: {reason}"

            try:
                create_notification(
                    user_id=post["user_id"],
                    notif_type="post_removed",
                    title="Post Removed",
                    message=msg,
                )
            except:
                pass

            debug_log(f"[Moderation] Denied and deleted post {post_id}: {reason}")
        else:
            # Mark as explicitly approved
            supabase.table("social_posts").update({
                "moderation_status": "approved"
            }).eq("id", post_id).execute()
            debug_log(f"[Moderation] Approved post {post_id}")

        return jsonify({"success": True, "action": action}), 200

    except Exception as e:
        debug_log(f"Admin moderate post error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/suggested-users", methods=["GET"])
@supabase_auth_required
def get_suggested_users():
    """Get suggested users to follow based on same university and activity."""
    try:
        user_id = request.user_id
        limit = int(request.args.get("limit", 5))

        # Get current user's university
        user_profile = supabase.table("user_profiles").select("university, edu_email").eq("id", user_id).execute()
        user_uni = user_profile.data[0].get("university") if user_profile.data else None

        # Get who the user already follows
        following_res = supabase.table("user_followers").select("following_id").eq("follower_id", user_id).execute()
        following_ids = set(f["following_id"] for f in (following_res.data or []))
        following_ids.add(user_id)  # Exclude self

        suggestions = []

        # Strategy 1: Users from same university
        if user_uni:
            try:
                uni_users = supabase.table("user_profiles").select(
                    "id, username, display_name, avatar_url, bio"
                ).eq("university", user_uni).limit(20).execute()

                for u in (uni_users.data or []):
                    if u["id"] not in following_ids and u.get("username"):
                        suggestions.append({
                            "user_id": u["id"],
                            "username": u["username"],
                            "display_name": u.get("display_name") or u["username"],
                            "avatar_url": u.get("avatar_url"),
                            "bio": (u.get("bio") or "")[:80],
                            "reason": user_uni
                        })
            except:
                pass

        # Strategy 2: Active posters (most posts recently)
        if len(suggestions) < limit:
            try:
                active = supabase.table("social_posts").select(
                    "user_id, username"
                ).order("created_at", desc=True).limit(50).execute()

                seen = set(s["user_id"] for s in suggestions)
                for p in (active.data or []):
                    if p["user_id"] not in following_ids and p["user_id"] not in seen and p.get("username"):
                        # Get profile
                        prof = supabase.table("user_profiles").select(
                            "id, username, display_name, avatar_url, bio"
                        ).eq("id", p["user_id"]).execute()
                        if prof.data:
                            u = prof.data[0]
                            suggestions.append({
                                "user_id": u["id"],
                                "username": u.get("username") or p["username"],
                                "display_name": u.get("display_name") or u.get("username") or p["username"],
                                "avatar_url": u.get("avatar_url"),
                                "bio": (u.get("bio") or "")[:80],
                                "reason": "Active poster"
                            })
                            seen.add(p["user_id"])
                        if len(suggestions) >= limit:
                            break
            except:
                pass

        return jsonify({"suggestions": suggestions[:limit]}), 200

    except Exception as e:
        debug_log(f"Suggested users error: {e}\n{traceback.format_exc()}")
        return jsonify({"suggestions": []}), 200


@app.route("/api/social/feed", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_social_feed():
    """Get personalized feed from followed users"""
    try:
        user_id = request.user_id
        page = int(request.args.get("page", 1))
        per_page = 20
        offset = (page - 1) * per_page

        # First, get list of users this user follows
        following_response = supabase.table("user_followers").select("following_id").eq(
            "follower_id", user_id
        ).execute()

        followed_user_ids = [f["following_id"] for f in (following_response.data or [])]
        # Add current user to see their own posts
        followed_user_ids.append(user_id)

        # Get posts from followed users + own posts
        if len(followed_user_ids) == 1:
            # Single ID - use .eq() instead of .in_()
            response = supabase.table("social_posts").select("*").eq(
                "user_id", followed_user_ids[0]
            ).order("created_at", desc=True).limit(per_page).offset(offset).execute()
        elif len(followed_user_ids) > 1:
            # Multiple IDs - use .in_()
            response = supabase.table("social_posts").select("*").in_(
                "user_id", followed_user_ids
            ).order("created_at", desc=True).limit(per_page).offset(offset).execute()
        else:
            response = None

        posts = response.data if response else []

        # Fetch note details separately for posts that have notes
        if posts:
            note_ids = [p["note_id"] for p in posts if p.get("note_id")]
            notes_map = {}
            if note_ids:
                try:
                    # Use .eq() for single ID, .in_() for multiple to avoid 400 error
                    if len(note_ids) == 1:
                        print(f"[FEED DEBUG] Fetching single note ID: {note_ids[0]}")
                        print(f"[FEED DEBUG] Using service role client: {type(supabase)}")
                        notes_response = supabase.table("notes").select(
                            "id, original_filename, thumbnail_url, file_type"
                        ).eq("id", note_ids[0]).execute()
                        print(f"[FEED DEBUG] Response data: {notes_response.data}")
                        print(f"[FEED DEBUG] Response count: {notes_response.count}")
                    else:
                        print(f"[FEED DEBUG] Fetching multiple notes: {len(note_ids)} notes")
                        notes_response = supabase.table("notes").select(
                            "id, original_filename, thumbnail_url, file_type"
                        ).in_("id", note_ids).execute()
                    notes_map = {n["id"]: n for n in (notes_response.data or [])}
                    print(f"[FEED DEBUG] Successfully mapped {len(notes_map)} notes")
                except Exception as e:
                    print(f"[FEED ERROR] Exception type: {type(e)}")
                    print(f"[FEED ERROR] Exception message: {str(e)}")
                    print(f"[FEED ERROR] Full traceback:")
                    import traceback
                    traceback.print_exc()
                    notes_map = {}

            # Add user interaction flags and note details
            post_ids = [p["id"] for p in posts]

            # Use .eq() for single ID, .in_() for multiple to avoid 400 error
            if len(post_ids) == 1:
                votes = supabase.table("post_votes").select("post_id, vote_type").eq(
                    "user_id", user_id
                ).eq("post_id", post_ids[0]).execute()
            else:
                votes = supabase.table("post_votes").select("post_id, vote_type").eq(
                    "user_id", user_id
                ).in_("post_id", post_ids).execute()
            votes_map = {v["post_id"]: v["vote_type"] for v in (votes.data or [])}

            if len(post_ids) == 1:
                bookmarks = supabase.table("post_bookmarks").select("post_id").eq(
                    "user_id", user_id
                ).eq("post_id", post_ids[0]).execute()
            else:
                bookmarks = supabase.table("post_bookmarks").select("post_id").eq(
                    "user_id", user_id
                ).in_("post_id", post_ids).execute()
            bookmark_ids = {b["post_id"] for b in (bookmarks.data or [])}

            # Fetch avatar URLs for post authors
            author_ids = list(set(p["user_id"] for p in posts))
            avatars_map = {}
            try:
                if len(author_ids) == 1:
                    av_res = supabase.table("user_profiles").select("id, avatar_url, is_verified").eq("id", author_ids[0]).execute()
                else:
                    av_res = supabase.table("user_profiles").select("id, avatar_url, is_verified").in_("id", author_ids).execute()
                avatars_map = {a["id"]: {"avatar_url": a.get("avatar_url"), "is_verified": a.get("is_verified", False)} for a in (av_res.data or [])}
            except:
                pass

            # Fetch latest 2 comments per post for inline preview
            comments_map = {}
            try:
                for pid in post_ids[:20]:  # limit to avoid too many queries
                    cres = supabase.table("post_comments").select(
                        "id, user_id, username, text_content, created_at"
                    ).eq("post_id", pid).order("created_at", desc=True).limit(2).execute()
                    if cres.data:
                        comments_map[pid] = list(reversed(cres.data))
            except:
                pass

            for post in posts:
                post["user_vote_type"] = votes_map.get(post["id"])
                post["is_bookmarked"] = post["id"] in bookmark_ids
                author_info = avatars_map.get(post["user_id"], {})
                post["avatar_url"] = author_info.get("avatar_url") if isinstance(author_info, dict) else author_info
                post["is_verified"] = author_info.get("is_verified", False) if isinstance(author_info, dict) else False
                post["recent_comments"] = comments_map.get(post["id"], [])
                if post.get("note_id"):
                    post["notes"] = notes_map.get(post["note_id"])

        return jsonify({"posts": posts, "page": page}), 200

    except Exception as e:
        debug_log(f"Get feed error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/posts/trending", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_trending_posts():
    """Get trending posts by score"""
    try:
        user_id = request.user_id
        page = int(request.args.get("page", 1))
        per_page = 20
        offset = (page - 1) * per_page

        # Get posts without join to avoid NULL note_id issues
        response = supabase.table("social_posts").select(
            "*"
        ).order("score", desc=True).order("created_at", desc=True).limit(per_page).offset(offset).execute()

        posts = response.data or []

        # Fetch note details separately for posts that have notes
        if posts:
            note_ids = [p["note_id"] for p in posts if p.get("note_id")]
            notes_map = {}
            if note_ids:
                try:
                    # Use .eq() for single ID, .in_() for multiple to avoid 400 error
                    if len(note_ids) == 1:
                        notes_response = supabase.table("notes").select(
                            "id, original_filename, thumbnail_url, file_type"
                        ).eq("id", note_ids[0]).execute()
                    else:
                        notes_response = supabase.table("notes").select(
                            "id, original_filename, thumbnail_url, file_type"
                        ).in_("id", note_ids).execute()
                    notes_map = {n["id"]: n for n in (notes_response.data or [])}
                except Exception as e:
                    debug_log(f"Error fetching notes for feed: {e}")
                    notes_map = {}

            post_ids = [p["id"] for p in posts]

            # Use .eq() for single ID, .in_() for multiple to avoid 400 error
            if len(post_ids) == 1:
                votes = supabase.table("post_votes").select("post_id, vote_type").eq(
                    "user_id", user_id
                ).eq("post_id", post_ids[0]).execute()
            else:
                votes = supabase.table("post_votes").select("post_id, vote_type").eq(
                    "user_id", user_id
                ).in_("post_id", post_ids).execute()
            votes_map = {v["post_id"]: v["vote_type"] for v in (votes.data or [])}

            if len(post_ids) == 1:
                bookmarks = supabase.table("post_bookmarks").select("post_id").eq(
                    "user_id", user_id
                ).eq("post_id", post_ids[0]).execute()
            else:
                bookmarks = supabase.table("post_bookmarks").select("post_id").eq(
                    "user_id", user_id
                ).in_("post_id", post_ids).execute()
            bookmark_ids = {b["post_id"] for b in (bookmarks.data or [])}

            # Fetch avatar URLs for post authors
            author_ids = list(set(p["user_id"] for p in posts))
            avatars_map = {}
            try:
                if len(author_ids) == 1:
                    av_res = supabase.table("user_profiles").select("id, avatar_url, is_verified").eq("id", author_ids[0]).execute()
                else:
                    av_res = supabase.table("user_profiles").select("id, avatar_url, is_verified").in_("id", author_ids).execute()
                avatars_map = {a["id"]: {"avatar_url": a.get("avatar_url"), "is_verified": a.get("is_verified", False)} for a in (av_res.data or [])}
            except:
                pass

            for post in posts:
                post["user_vote_type"] = votes_map.get(post["id"])
                post["is_bookmarked"] = post["id"] in bookmark_ids
                author_info = avatars_map.get(post["user_id"], {})
                post["avatar_url"] = author_info.get("avatar_url") if isinstance(author_info, dict) else author_info
                post["is_verified"] = author_info.get("is_verified", False) if isinstance(author_info, dict) else False
                if post.get("note_id"):
                    post["notes"] = notes_map.get(post["note_id"])

        return jsonify({"posts": posts, "page": page}), 200

    except Exception as e:
        debug_log(f"Get trending error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/posts", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def create_social_post():
    """Create a new social post"""
    try:
        user_id = request.user_id
        data = request.json

        post_type = data.get("post_type")
        note_id = data.get("note_id")
        text_content = (data.get("text_content") or "").strip()
        group_id = data.get("group_id")

        if post_type not in ["note", "text", "group_invite"]:
            return jsonify({"error": "Invalid post type"}), 400

        profile = get_cached_profile(user_id)
        username = profile.get("username", "Anonymous") if profile else "Anonymous"

        post_data = {
            "user_id": user_id,
            "username": username,
            "post_type": post_type,
            "text_content": text_content if text_content else None
        }

        if post_type == "note":
            post_data["note_id"] = note_id
        elif post_type == "group_invite":
            post_data["group_id"] = group_id

        response = supabase.table("social_posts").insert(post_data).execute()

        if not response.data:
            return jsonify({"error": "Failed to create post"}), 500

        return jsonify({"message": "Post created", "post": response.data[0]}), 201

    except Exception as e:
        debug_log(f"Create post error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/upload-note", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def upload_note_from_social():
    """Upload a note directly from social pages with optional folder assignment"""
    try:
        from StudyFlow.backend.supabase_client import (
            check_page_limit, upload_file_to_storage, create_note_record,
            increment_page_count, log_upload, get_user_profile
        )
        from StudyFlow.backend.tasks import process_note_async
        import hashlib
        import uuid

        # Check if file was uploaded
        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400

        # Get optional folder_id and content type from form
        folder_id = request.form.get('folder_id')  # Can be null
        content_type = request.form.get('content_type', 'note')  # note, photo, resource

        # Get user profile
        user_profile = get_user_profile(request.user_id)
        username = user_profile.get('username') if user_profile else None

        # Get user's IP address
        ip_address = request.headers.get('X-Forwarded-For', request.remote_addr)
        if ip_address and ',' in ip_address:
            ip_address = ip_address.split(',')[0].strip()

        # Get file info
        original_filename = file.filename
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)

        # Determine file type
        file_ext = original_filename.rsplit('.', 1)[1].lower() if '.' in original_filename else ''
        allowed_extensions = ['pdf', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']
        if file_ext not in allowed_extensions:
            return jsonify({"error": "Unsupported file type. Allowed: PDF, images"}), 400

        # Read file content
        file_content = file.read()
        file_hash = hashlib.sha256(file_content).hexdigest()

        # Log upload
        log_upload(
            user_id=request.user_id,
            file_name=original_filename,
            file_hash=file_hash,
            file_size=file_size,
            shared_with_nexus=True,  # Updated after AI classification
            ip_address=ip_address
        )

        # Extract text for OCR
        import google.generativeai as genai
        genai.configure(api_key=os.getenv("GEMINI_API_KEY"))
        model = genai.GenerativeModel('gemini-3.1-flash-lite-preview')

        if file_ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
            import PIL.Image
            import io
            image = PIL.Image.open(io.BytesIO(file_content))
            prompt = "Extract ALL text from this image. Return only the extracted text."
            response = model.generate_content([prompt, image])
            ocr_text = response.text.strip()
            page_count = 1
        elif file_ext == 'pdf':
            import PyPDF2
            import io
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
            page_count = len(pdf_reader.pages)
            ocr_text = ""
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    ocr_text += page_text + "\n\n"
            if not ocr_text.strip():
                ocr_text = f"[PDF document: {page_count} pages]"
        else:
            ocr_text = ""
            page_count = 1

        # Check page limit
        allowed, message = check_page_limit(request.user_id, page_count)
        if not allowed:
            return jsonify({"error": message, "limit_exceeded": True}), 403

        # Convert to PDF if image
        if file_ext != 'pdf':
            from StudyFlow.backend.pdf_flatten import convert_to_pdf
            pdf_data, pdf_filename = convert_to_pdf(file_content, original_filename)
            if pdf_data is None:
                return jsonify({"error": "Failed to convert to PDF"}), 500
            file_content = pdf_data
            original_filename = pdf_filename
            file_ext = 'pdf'
            file_size = len(file_content)

        # Upload to Supabase Storage
        unique_filename = f"{request.user_id}/{uuid.uuid4()}_{original_filename}"
        content_type = 'application/pdf'
        file_url = upload_file_to_storage(file_content, unique_filename, content_type)

        if not file_url:
            return jsonify({"error": "Failed to upload file to storage"}), 500

        # Generate thumbnail from first page using PyMuPDF (no poppler dependency needed)
        thumbnail_url = None
        try:
            import io
            import fitz  # PyMuPDF
            from PIL import Image

            print(f"[THUMBNAIL] Starting thumbnail generation for {original_filename}")

            # Open PDF with PyMuPDF
            pdf_doc = fitz.open(stream=file_content, filetype="pdf")
            if len(pdf_doc) > 0:
                # Render first page to pixmap (image)
                page = pdf_doc[0]
                # Render at 200 DPI (2.78x zoom since PDF is 72 DPI) for high quality
                mat = fitz.Matrix(2.78, 2.78)
                pix = page.get_pixmap(matrix=mat, alpha=False)

                # Convert to PIL Image
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)

                # Resize to larger thumbnail (1000x1000 max) for crisp display
                img.thumbnail((1000, 1000), Image.Resampling.LANCZOS)

                print(f"[THUMBNAIL] Rendered page: {pix.width}x{pix.height} -> {img.width}x{img.height}")

                # Save as JPEG
                thumb_io = io.BytesIO()
                img.save(thumb_io, format='JPEG', quality=85)
                thumb_data = thumb_io.getvalue()

                # Upload thumbnail
                thumb_filename = f"thumbnails/{request.user_id}/{uuid.uuid4()}_thumb.jpg"
                thumbnail_url = upload_file_to_storage(thumb_data, thumb_filename, 'image/jpeg')

                pdf_doc.close()
                print(f"[THUMBNAIL] Generated thumbnail: {thumbnail_url}")
            else:
                print(f"[THUMBNAIL] PDF has no pages!")
        except Exception as e:
            print(f"[THUMBNAIL] Thumbnail generation failed (non-fatal): {e}")
            import traceback
            print(f"[THUMBNAIL] Traceback: {traceback.format_exc()}")
            # Continue without thumbnail

        # Create note record with optional folder_id
        note = create_note_record(
            user_id=request.user_id,
            filename=original_filename,
            file_type=file_ext,
            file_size=file_size,
            file_path=unique_filename,
            page_count=page_count,
            course_metadata={},
            username=username,
            thumbnail_url=thumbnail_url
        )

        if not note:
            return jsonify({"error": "Failed to create note record"}), 500

        note_id = note['id']

        # Use Gemini to classify if this is educational content (Nexus-worthy)
        is_nexus = False
        try:
            text_sample = ocr_text[:1500] if ocr_text else ""
            classify_prompt = f"""Classify this uploaded file. Is it educational study material suitable for a university note-sharing platform?

Filename: {original_filename}
File type: {file_ext}
Content type selected by user: {content_type}
Text extracted from file (first 1500 chars):
{text_sample}

Respond with ONLY one word: "educational" or "personal"

Educational = lecture notes, study guides, textbook excerpts, homework, lab reports, class materials, handwritten notes, flashcards, outlines, summaries, essays, research papers
Personal = selfies, memes, screenshots of social media, personal photos, random images, non-academic content"""

            classify_response = model.generate_content(classify_prompt)
            classification = classify_response.text.strip().lower()
            is_nexus = 'educational' in classification
            debug_log(f"[AI Classification] {original_filename}: {classification} -> Nexus: {is_nexus}")
        except Exception as classify_err:
            # If classification fails, fall back to content_type from user
            is_nexus = content_type in ('note', 'resource')
            debug_log(f"[AI Classification] Failed, using fallback: {classify_err}")

        # Save classification result and set visibility
        try:
            supabase.table("notes").update({
                "is_public": is_nexus,
                "ai_classification": "educational" if is_nexus else "personal"
            }).eq("id", note_id).execute()
        except:
            pass

        # Assign to folder if specified
        if folder_id:
            try:
                supabase.table("notes").update({"folder_id": folder_id}).eq("id", note_id).execute()
                debug_log(f"Note {note_id} assigned to folder {folder_id}")
            except Exception as e:
                debug_log(f"Failed to assign folder: {e}")

        # Increment page count
        increment_page_count(request.user_id, page_count)

        # Trigger background processing
        process_note_async.delay(note_id, request.user_id, ocr_text, {}, file_hash, username)

        debug_log(f"Social note uploaded: {original_filename} ({file_size} bytes, {page_count} pages)")

        return jsonify({
            "success": True,
            "note_id": note_id,
            "filename": original_filename,
            "thumbnail_url": note.get('thumbnail_url'),
            "pages": page_count
        }), 201

    except Exception as e:
        debug_log(f"Social upload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/posts/<post_id>", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_post_detail(post_id):
    """Get single post details"""
    try:
        user_id = request.user_id

        response = supabase.table("social_posts").select(
            "*, notes(original_filename, thumbnail_url, file_type, university, course_code)"
        ).eq("id", post_id).single().execute()

        if not response.data:
            return jsonify({"error": "Post not found"}), 404

        post = response.data

        vote = supabase.table("post_votes").select("vote_type").eq(
            "user_id", user_id
        ).eq("post_id", post_id).execute()
        post["user_vote_type"] = vote.data[0]["vote_type"] if vote.data else None

        bookmark = supabase.table("post_bookmarks").select("id").eq(
            "user_id", user_id
        ).eq("post_id", post_id).execute()
        post["is_bookmarked"] = bool(bookmark.data)

        return jsonify({"post": post}), 200

    except Exception as e:
        debug_log(f"Get post detail error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/posts/<post_id>/vote", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def vote_on_post(post_id):
    """Upvote or downvote a post"""
    try:
        user_id = request.user_id
        data = request.json
        vote_type = data.get("vote_type")

        if vote_type and vote_type not in ["upvote", "downvote"]:
            return jsonify({"error": "Invalid vote type"}), 400

        existing = supabase.table("post_votes").select("id, vote_type").eq(
            "user_id", user_id
        ).eq("post_id", post_id).execute()

        if vote_type is None:
            if existing.data:
                supabase.table("post_votes").delete().eq("id", existing.data[0]["id"]).execute()
                return jsonify({"message": "Vote removed"}), 200
            return jsonify({"message": "No vote to remove"}), 200

        if existing.data:
            if existing.data[0]["vote_type"] == vote_type:
                supabase.table("post_votes").delete().eq("id", existing.data[0]["id"]).execute()
                return jsonify({"message": "Vote removed"}), 200
            else:
                supabase.table("post_votes").update({"vote_type": vote_type}).eq(
                    "id", existing.data[0]["id"]
                ).execute()
                return jsonify({"message": "Vote updated", "vote_type": vote_type}), 200
        else:
            supabase.table("post_votes").insert({
                "user_id": user_id,
                "post_id": post_id,
                "vote_type": vote_type
            }).execute()
            return jsonify({"message": "Vote added", "vote_type": vote_type}), 201

    except Exception as e:
        debug_log(f"Vote error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/posts/<post_id>/bookmark", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def toggle_bookmark(post_id):
    """Save or unsave a post"""
    try:
        user_id = request.user_id

        existing = supabase.table("post_bookmarks").select("id").eq(
            "user_id", user_id
        ).eq("post_id", post_id).execute()

        if existing.data:
            supabase.table("post_bookmarks").delete().eq("id", existing.data[0]["id"]).execute()
            return jsonify({"message": "Bookmark removed", "is_bookmarked": False}), 200
        else:
            supabase.table("post_bookmarks").insert({
                "user_id": user_id,
                "post_id": post_id
            }).execute()
            return jsonify({"message": "Bookmark added", "is_bookmarked": True}), 201

    except Exception as e:
        debug_log(f"Bookmark error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/posts/<post_id>", methods=["DELETE"])
@supabase_auth_required
@account_not_frozen
def delete_post(post_id):
    """Delete own post"""
    try:
        user_id = request.user_id

        post = supabase.table("social_posts").select("user_id").eq("id", post_id).single().execute()

        if not post.data:
            return jsonify({"error": "Post not found"}), 404

        if post.data["user_id"] != user_id:
            return jsonify({"error": "You can only delete your own posts"}), 403

        supabase.table("social_posts").delete().eq("id", post_id).execute()

        return jsonify({"message": "Post deleted"}), 200

    except Exception as e:
        debug_log(f"Delete post error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/bookmarks", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_bookmarks():
    """Get user's saved posts"""
    try:
        user_id = request.user_id
        page = int(request.args.get("page", 1))
        per_page = 20
        offset = (page - 1) * per_page

        response = supabase.table("post_bookmarks").select(
            "created_at, social_posts(*, notes(original_filename, thumbnail_url, file_type))"
        ).eq("user_id", user_id).order("created_at", desc=True).limit(per_page).offset(offset).execute()

        bookmarks = response.data or []
        posts = [{"bookmark_date": b["created_at"], **b["social_posts"]} for b in bookmarks if b.get("social_posts")]

        if posts:
            post_ids = [p["id"] for p in posts]
            votes = supabase.table("post_votes").select("post_id, vote_type").eq(
                "user_id", user_id
            ).in_("post_id", post_ids).execute()
            votes_map = {v["post_id"]: v["vote_type"] for v in (votes.data or [])}

            for post in posts:
                post["user_vote_type"] = votes_map.get(post["id"])
                post["is_bookmarked"] = True

        return jsonify({"posts": posts, "page": page}), 200

    except Exception as e:
        debug_log(f"Get bookmarks error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# --- COMMENTS ---

@app.route("/api/social/posts/<post_id>/comments", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_post_comments(post_id):
    """Get all top-level comments for a post"""
    try:
        user_id = request.user_id
        sort = request.args.get("sort", "top")

        query = supabase.table("post_comments").select("*").eq("post_id", post_id).is_("parent_id", "null")

        if sort == "top":
            query = query.order("score", desc=True).order("created_at", desc=True)
        else:
            query = query.order("created_at", desc=True)

        response = query.execute()
        comments = response.data or []

        if comments:
            comment_ids = [c["id"] for c in comments]
            votes = supabase.table("comment_votes").select("comment_id, vote_type").eq(
                "user_id", user_id
            ).in_("comment_id", comment_ids).execute()
            votes_map = {v["comment_id"]: v["vote_type"] for v in (votes.data or [])}

            for comment in comments:
                comment["user_vote_type"] = votes_map.get(comment["id"])

        return jsonify({"comments": comments}), 200

    except Exception as e:
        debug_log(f"Get comments error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/comments/<comment_id>/replies", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_comment_replies(comment_id):
    """Get all replies to a comment"""
    try:
        user_id = request.user_id

        response = supabase.table("post_comments").select("*").eq(
            "parent_id", comment_id
        ).order("created_at", desc=False).execute()

        replies = response.data or []

        if replies:
            reply_ids = [r["id"] for r in replies]
            votes = supabase.table("comment_votes").select("comment_id, vote_type").eq(
                "user_id", user_id
            ).in_("comment_id", reply_ids).execute()
            votes_map = {v["comment_id"]: v["vote_type"] for v in (votes.data or [])}

            for reply in replies:
                reply["user_vote_type"] = votes_map.get(reply["id"])

        return jsonify({"replies": replies}), 200

    except Exception as e:
        debug_log(f"Get replies error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/posts/<post_id>/comments", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def create_comment(post_id):
    """Create a new comment on a post"""
    try:
        user_id = request.user_id
        data = request.json
        content = data.get("content", "").strip()
        parent_id = data.get("parent_id")

        if not content or len(content) < 1:
            return jsonify({"error": "Comment content required"}), 400
        if len(content) > 2000:
            return jsonify({"error": "Comment too long"}), 400

        profile = get_cached_profile(user_id)
        username = profile.get("username", "Anonymous") if profile else "Anonymous"

        comment_data = {
            "post_id": post_id,
            "user_id": user_id,
            "username": username,
            "content": content
        }

        if parent_id:
            parent = supabase.table("post_comments").select("id, post_id").eq(
                "id", parent_id
            ).single().execute()
            if not parent.data or parent.data["post_id"] != post_id:
                return jsonify({"error": "Invalid parent comment"}), 400
            comment_data["parent_id"] = parent_id

        response = supabase.table("post_comments").insert(comment_data).execute()

        if not response.data:
            return jsonify({"error": "Failed to create comment"}), 500

        return jsonify({"message": "Comment created", "comment": response.data[0]}), 201

    except Exception as e:
        debug_log(f"Create comment error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/comments/<comment_id>", methods=["DELETE"])
@supabase_auth_required
@account_not_frozen
def delete_comment(comment_id):
    """Delete own comment"""
    try:
        user_id = request.user_id

        comment = supabase.table("post_comments").select("user_id").eq(
            "id", comment_id
        ).single().execute()

        if not comment.data:
            return jsonify({"error": "Comment not found"}), 404
        if comment.data["user_id"] != user_id:
            return jsonify({"error": "You can only delete your own comments"}), 403

        supabase.table("post_comments").delete().eq("id", comment_id).execute()

        return jsonify({"message": "Comment deleted"}), 200

    except Exception as e:
        debug_log(f"Delete comment error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/comments/<comment_id>/vote", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def vote_on_comment(comment_id):
    """Upvote or downvote a comment"""
    try:
        user_id = request.user_id
        data = request.json
        vote_type = data.get("vote_type")

        if vote_type and vote_type not in ["upvote", "downvote"]:
            return jsonify({"error": "Invalid vote type"}), 400

        existing = supabase.table("comment_votes").select("id, vote_type").eq(
            "user_id", user_id
        ).eq("comment_id", comment_id).execute()

        if vote_type is None:
            if existing.data:
                supabase.table("comment_votes").delete().eq("id", existing.data[0]["id"]).execute()
                return jsonify({"message": "Vote removed"}), 200
            return jsonify({"message": "No vote to remove"}), 200

        if existing.data:
            if existing.data[0]["vote_type"] == vote_type:
                supabase.table("comment_votes").delete().eq("id", existing.data[0]["id"]).execute()
                return jsonify({"message": "Vote removed"}), 200
            else:
                supabase.table("comment_votes").update({"vote_type": vote_type}).eq(
                    "id", existing.data[0]["id"]
                ).execute()
                return jsonify({"message": "Vote updated", "vote_type": vote_type}), 200
        else:
            supabase.table("comment_votes").insert({
                "user_id": user_id,
                "comment_id": comment_id,
                "vote_type": vote_type
            }).execute()
            return jsonify({"message": "Vote added", "vote_type": vote_type}), 201

    except Exception as e:
        debug_log(f"Comment vote error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# --- PROFILES & FOLLOWERS ---

@app.route("/api/user/profile", methods=["GET"])
@supabase_auth_required
def get_current_user_profile():
    """Get current user's basic profile info"""
    try:
        user_id = request.user_id

        # Get user's profile
        response = supabase.table("user_profiles").select(
            "id, username, display_name, bio, avatar_url, banner_url"
        ).eq("id", user_id).single().execute()

        if not response.data:
            return jsonify({"error": "Profile not found"}), 404

        return jsonify({"profile": response.data}), 200

    except Exception as e:
        debug_log(f"Get current user profile error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/profile/<username>", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_user_profile(username):
    """Get public user profile"""
    try:
        current_user_id = request.user_id

        profile = supabase.table("user_profiles").select(
            "id, username, display_name, bio, avatar_url, banner_url, "
            "follower_count, following_count, post_count, is_public, is_verified, created_at"
        ).eq("username", username).single().execute()

        if not profile.data:
            return jsonify({"error": "User not found"}), 404

        profile_data = profile.data
        profile_user_id = profile_data["id"]

        follow = supabase.table("user_followers").select("id").eq(
            "follower_id", current_user_id
        ).eq("following_id", profile_user_id).execute()

        profile_data["is_following"] = bool(follow.data)
        profile_data["is_own_profile"] = current_user_id == profile_user_id

        # Get posts without join to avoid NULL note_id issues
        posts = supabase.table("social_posts").select(
            "id, post_type, created_at, score, comment_count, note_id"
        ).eq("user_id", profile_user_id).order("created_at", desc=True).limit(12).execute()

        recent_posts = posts.data or []

        # Fetch note details separately for posts that have notes
        if recent_posts:
            note_ids = [p["note_id"] for p in recent_posts if p.get("note_id")]
            notes_map = {}
            if note_ids:
                # Use .eq() for single ID, .in_() for multiple to avoid 400 error
                if len(note_ids) == 1:
                    notes_response = supabase.table("notes").select(
                        "id, thumbnail_url, file_type"
                    ).eq("id", note_ids[0]).execute()
                else:
                    notes_response = supabase.table("notes").select(
                        "id, thumbnail_url, file_type"
                    ).in_("id", note_ids).execute()
                notes_map = {n["id"]: n for n in (notes_response.data or [])}

            # Add note details to posts
            for post in recent_posts:
                if post.get("note_id"):
                    post["notes"] = notes_map.get(post["note_id"])

        profile_data["recent_posts"] = recent_posts

        return jsonify({"profile": profile_data}), 200

    except Exception as e:
        debug_log(f"Get profile error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/profile/<username>/posts", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_user_posts(username):
    """Get posts for a specific user with pagination"""
    try:
        user_id = request.user_id
        page = int(request.args.get("page", 1))
        per_page = 20
        offset = (page - 1) * per_page

        # Get user ID from username
        profile = supabase.table("user_profiles").select("id").eq("username", username).single().execute()

        if not profile.data:
            return jsonify({"error": "User not found"}), 404

        profile_user_id = profile.data["id"]

        # Get posts without join to avoid NULL note_id issues
        response = supabase.table("social_posts").select("*").eq(
            "user_id", profile_user_id
        ).order("created_at", desc=True).limit(per_page).offset(offset).execute()

        posts = response.data or []

        # Fetch note details separately for posts that have notes
        if posts:
            note_ids = [p["note_id"] for p in posts if p.get("note_id")]
            notes_map = {}
            if note_ids:
                try:
                    # Use .eq() for single ID, .in_() for multiple to avoid 400 error
                    if len(note_ids) == 1:
                        notes_response = supabase.table("notes").select(
                            "id, original_filename, thumbnail_url, file_type"
                        ).eq("id", note_ids[0]).execute()
                    else:
                        notes_response = supabase.table("notes").select(
                            "id, original_filename, thumbnail_url, file_type"
                        ).in_("id", note_ids).execute()
                    notes_map = {n["id"]: n for n in (notes_response.data or [])}
                except Exception as e:
                    debug_log(f"Error fetching notes for feed: {e}")
                    notes_map = {}

            # Get user interactions
            post_ids = [p["id"] for p in posts]

            # Use .eq() for single ID, .in_() for multiple to avoid 400 error
            if len(post_ids) == 1:
                votes = supabase.table("post_votes").select("post_id, vote_type").eq(
                    "user_id", user_id
                ).eq("post_id", post_ids[0]).execute()
            else:
                votes = supabase.table("post_votes").select("post_id, vote_type").eq(
                    "user_id", user_id
                ).in_("post_id", post_ids).execute()
            votes_map = {v["post_id"]: v["vote_type"] for v in (votes.data or [])}

            if len(post_ids) == 1:
                bookmarks = supabase.table("post_bookmarks").select("post_id").eq(
                    "user_id", user_id
                ).eq("post_id", post_ids[0]).execute()
            else:
                bookmarks = supabase.table("post_bookmarks").select("post_id").eq(
                    "user_id", user_id
                ).in_("post_id", post_ids).execute()
            bookmark_ids = {b["post_id"] for b in (bookmarks.data or [])}

            for post in posts:
                post["user_vote_type"] = votes_map.get(post["id"])
                post["is_bookmarked"] = post["id"] in bookmark_ids
                # Add note details if post has a note
                if post.get("note_id"):
                    post["notes"] = notes_map.get(post["note_id"])

        return jsonify({"posts": posts, "page": page}), 200

    except Exception as e:
        debug_log(f"Get user posts error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/profile/avatar", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def upload_avatar():
    """Upload and set profile avatar"""
    try:
        from StudyFlow.backend.supabase_client import upload_file_to_storage
        import uuid

        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400

        # Check file type
        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        allowed_extensions = ['jpg', 'jpeg', 'png', 'gif', 'webp']
        if file_ext not in allowed_extensions:
            return jsonify({"error": "Invalid file type. Use JPG, PNG, GIF, or WebP"}), 400

        # Check file size (5MB max for avatars)
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        if file_size > 5 * 1024 * 1024:
            return jsonify({"error": "File too large. Maximum 5MB"}), 400

        # Read file
        file_content = file.read()

        # Upload to Supabase Storage in avatars folder
        unique_filename = f"avatars/{request.user_id}/{uuid.uuid4()}.{file_ext}"
        content_type = f'image/{file_ext}' if file_ext != 'jpg' else 'image/jpeg'
        avatar_url = upload_file_to_storage(file_content, unique_filename, content_type)

        if not avatar_url:
            return jsonify({"error": "Failed to upload file"}), 500

        # Update user profile
        response = supabase.table("user_profiles").update({
            "avatar_url": avatar_url
        }).eq("id", request.user_id).execute()

        return jsonify({"avatar_url": avatar_url}), 200

    except Exception as e:
        debug_log(f"Avatar upload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/profile/banner", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def upload_banner():
    """Upload and set profile banner"""
    try:
        from StudyFlow.backend.supabase_client import upload_file_to_storage
        import uuid

        if 'file' not in request.files:
            return jsonify({"error": "No file uploaded"}), 400

        file = request.files['file']
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400

        # Check file type
        file_ext = file.filename.rsplit('.', 1)[1].lower() if '.' in file.filename else ''
        allowed_extensions = ['jpg', 'jpeg', 'png', 'gif', 'webp']
        if file_ext not in allowed_extensions:
            return jsonify({"error": "Invalid file type. Use JPG, PNG, GIF, or WebP"}), 400

        # Check file size (10MB max for banners)
        file.seek(0, os.SEEK_END)
        file_size = file.tell()
        file.seek(0)
        if file_size > 10 * 1024 * 1024:
            return jsonify({"error": "File too large. Maximum 10MB"}), 400

        # Read file
        file_content = file.read()

        # Upload to Supabase Storage in banners folder
        unique_filename = f"banners/{request.user_id}/{uuid.uuid4()}.{file_ext}"
        content_type = f'image/{file_ext}' if file_ext != 'jpg' else 'image/jpeg'
        banner_url = upload_file_to_storage(file_content, unique_filename, content_type)

        if not banner_url:
            return jsonify({"error": "Failed to upload file"}), 500

        # Update user profile
        response = supabase.table("user_profiles").update({
            "banner_url": banner_url
        }).eq("id", request.user_id).execute()

        return jsonify({"banner_url": banner_url}), 200

    except Exception as e:
        debug_log(f"Banner upload error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/profile", methods=["PATCH"])
@supabase_auth_required
@account_not_frozen
def update_own_profile():
    """Update own profile (bio, display name, avatar, banner)"""
    try:
        user_id = request.user_id
        data = request.json

        update_data = {}

        if "bio" in data:
            bio = data["bio"].strip() if data["bio"] else None
            if bio and len(bio) > 500:
                return jsonify({"error": "Bio too long"}), 400
            update_data["bio"] = bio

        if "display_name" in data:
            display_name = data["display_name"].strip() if data["display_name"] else None
            if display_name and len(display_name) > 50:
                return jsonify({"error": "Display name too long"}), 400
            update_data["display_name"] = display_name

        if "avatar_url" in data:
            update_data["avatar_url"] = data["avatar_url"]

        if "banner_url" in data:
            update_data["banner_url"] = data["banner_url"]

        if not update_data:
            return jsonify({"error": "No fields to update"}), 400

        response = supabase.table("user_profiles").update(update_data).eq("id", user_id).execute()

        return jsonify({"message": "Profile updated", "profile": response.data[0]}), 200

    except Exception as e:
        debug_log(f"Update profile error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/follow/<username>", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def follow_user(username):
    """Follow a user"""
    try:
        follower_id = request.user_id

        target = supabase.table("user_profiles").select("id").eq("username", username).single().execute()

        if not target.data:
            return jsonify({"error": "User not found"}), 404

        following_id = target.data["id"]

        if follower_id == following_id:
            return jsonify({"error": "Cannot follow yourself"}), 400

        user1_id = min(follower_id, following_id)
        user2_id = max(follower_id, following_id)

        existing = supabase.table("user_followers").select("id").eq(
            "follower_id", follower_id
        ).eq("following_id", following_id).execute()

        if existing.data:
            return jsonify({"message": "Already following", "is_following": True}), 200

        supabase.table("user_followers").insert({
            "follower_id": follower_id,
            "following_id": following_id
        }).execute()

        return jsonify({"message": "Followed successfully", "is_following": True}), 201

    except Exception as e:
        debug_log(f"Follow error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/unfollow/<username>", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def unfollow_user(username):
    """Unfollow a user"""
    try:
        follower_id = request.user_id

        target = supabase.table("user_profiles").select("id").eq("username", username).single().execute()

        if not target.data:
            return jsonify({"error": "User not found"}), 404

        following_id = target.data["id"]

        result = supabase.table("user_followers").delete().eq(
            "follower_id", follower_id
        ).eq("following_id", following_id).execute()

        if not result.data:
            return jsonify({"message": "Not following", "is_following": False}), 200

        return jsonify({"message": "Unfollowed successfully", "is_following": False}), 200

    except Exception as e:
        debug_log(f"Unfollow error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/search/users", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def search_users():
    """Search for users by username or display name"""
    try:
        query = request.args.get("q", "").strip()

        if not query or len(query) < 2:
            return jsonify({"error": "Search query too short"}), 400

        response = supabase.table("user_profiles").select(
            "id, username, display_name, avatar_url, bio, follower_count"
        ).or_(f"username.ilike.%{query}%,display_name.ilike.%{query}%").limit(20).execute()

        users = response.data or []

        return jsonify({"users": users}), 200

    except Exception as e:
        debug_log(f"Search users error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


# --- DIRECT MESSAGES ---

@app.route("/api/social/dm/conversations", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_dm_conversations():
    """Get all DM conversations for current user"""
    try:
        user_id = request.user_id

        response = supabase.table("dm_conversations").select(
            "id, user1_id, user2_id, last_message_at, created_at"
        ).or_(f"user1_id.eq.{user_id},user2_id.eq.{user_id}").order(
            "last_message_at", desc=True
        ).execute()

        conversations = response.data or []

        enriched_conversations = []
        for conv in conversations:
            other_user_id = conv["user2_id"] if conv["user1_id"] == user_id else conv["user1_id"]

            profile = supabase.table("user_profiles").select(
                "id, username, display_name, avatar_url"
            ).eq("id", other_user_id).single().execute()

            if not profile.data:
                continue

            last_msg = supabase.table("dm_messages").select(
                "id, sender_id, content, created_at, is_read"
            ).eq("conversation_id", conv["id"]).order("created_at", desc=True).limit(1).execute()

            unread = supabase.table("dm_messages").select("id", count="exact").eq(
                "conversation_id", conv["id"]
            ).eq("is_read", False).neq("sender_id", user_id).execute()

            enriched_conversations.append({
                "conversation_id": conv["id"],
                "other_user": profile.data,
                "last_message": last_msg.data[0] if last_msg.data else None,
                "unread_count": unread.count or 0,
                "updated_at": conv["last_message_at"]
            })

        return jsonify({"conversations": enriched_conversations}), 200

    except Exception as e:
        debug_log(f"Get conversations error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/dm/conversation/<username>", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_or_create_conversation(username):
    """Get or create conversation with a user"""
    try:
        user_id = request.user_id

        target = supabase.table("user_profiles").select("id").eq("username", username).single().execute()

        if not target.data:
            return jsonify({"error": "User not found"}), 404

        other_user_id = target.data["id"]

        if user_id == other_user_id:
            return jsonify({"error": "Cannot message yourself"}), 400

        user1_id = min(user_id, other_user_id)
        user2_id = max(user_id, other_user_id)

        existing = supabase.table("dm_conversations").select("id").eq(
            "user1_id", user1_id
        ).eq("user2_id", user2_id).execute()

        if existing.data:
            conversation_id = existing.data[0]["id"]
        else:
            new_conv = supabase.table("dm_conversations").insert({
                "user1_id": user1_id,
                "user2_id": user2_id
            }).execute()
            conversation_id = new_conv.data[0]["id"]

        return jsonify({"conversation_id": conversation_id}), 200

    except Exception as e:
        debug_log(f"Get/create conversation error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/dm/conversation/<conversation_id>/messages", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_dm_messages(conversation_id):
    """Get all messages in a conversation"""
    try:
        user_id = request.user_id
        page = int(request.args.get("page", 1))
        per_page = 50
        offset = (page - 1) * per_page

        conv = supabase.table("dm_conversations").select("user1_id, user2_id").eq(
            "id", conversation_id
        ).single().execute()

        if not conv.data:
            return jsonify({"error": "Conversation not found"}), 404

        if user_id not in [conv.data["user1_id"], conv.data["user2_id"]]:
            return jsonify({"error": "Not authorized"}), 403

        response = supabase.table("dm_messages").select("*").eq(
            "conversation_id", conversation_id
        ).order("created_at", desc=True).limit(per_page).offset(offset).execute()

        messages = (response.data or [])[::-1]

        if messages:
            unread_ids = [m["id"] for m in messages if not m["is_read"] and m["sender_id"] != user_id]
            if unread_ids:
                supabase.table("dm_messages").update({"is_read": True}).in_("id", unread_ids).execute()

        return jsonify({"messages": messages, "page": page}), 200

    except Exception as e:
        debug_log(f"Get messages error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/dm/conversation/<conversation_id>/send", methods=["POST"])
@supabase_auth_required
@account_not_frozen
def send_dm_message(conversation_id):
    """Send a message in a conversation"""
    try:
        user_id = request.user_id
        data = request.json
        content = data.get("content", "").strip()

        if not content or len(content) < 1:
            return jsonify({"error": "Message content required"}), 400
        if len(content) > 2000:
            return jsonify({"error": "Message too long"}), 400

        conv = supabase.table("dm_conversations").select("user1_id, user2_id").eq(
            "id", conversation_id
        ).single().execute()

        if not conv.data:
            return jsonify({"error": "Conversation not found"}), 404

        if user_id not in [conv.data["user1_id"], conv.data["user2_id"]]:
            return jsonify({"error": "Not authorized"}), 403

        response = supabase.table("dm_messages").insert({
            "conversation_id": conversation_id,
            "sender_id": user_id,
            "content": content
        }).execute()

        if not response.data:
            return jsonify({"error": "Failed to send message"}), 500

        return jsonify({"message": "Message sent", "dm": response.data[0]}), 201

    except Exception as e:
        debug_log(f"Send message error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/social/dm/unread-count", methods=["GET"])
@supabase_auth_required
@account_not_frozen
def get_unread_dm_count():
    """Get total unread message count for user"""
    try:
        user_id = request.user_id

        convs = supabase.table("dm_conversations").select("id").or_(
            f"user1_id.eq.{user_id},user2_id.eq.{user_id}"
        ).execute()

        conv_ids = [c["id"] for c in (convs.data or [])]

        if not conv_ids:
            return jsonify({"unread_count": 0}), 200

        unread = supabase.table("dm_messages").select("id", count="exact").in_(
            "conversation_id", conv_ids
        ).eq("is_read", False).neq("sender_id", user_id).execute()

        return jsonify({"unread_count": unread.count or 0}), 200

    except Exception as e:
        debug_log(f"Get unread count error: {e}\n{traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500

# =================================================


if __name__ == "__main__":
    try:
        port = int(os.environ.get("PORT", 5000))
        app.run(host="0.0.0.0", port=port)
    except Exception as e:
        debug_log(f"Server startup error: {e}\n{traceback.format_exc()}")
